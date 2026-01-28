"""
PowerPoint Proposal Generation Module
Handles slide cloning, table updates, and presentation assembly for Phase 2.
Includes support for multi-variant product consolidation (v6.13).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from typing import List, Dict, Optional, Tuple
import io
import copy
from datetime import datetime
import pandas as pd
import math
import re
import gc  # For memory optimization
from src.helpers import get_column_value, calculate_moq, round_to_nearest_fifty_cents


def apply_marketing_rounding(price, marketing_rounding_enabled):
    """Apply marketing rounding (charm pricing) if enabled."""
    if not marketing_rounding_enabled:
        return price
    if price % 10 == 0 and price > 10:
        return price - 1
    return price


def clean_price(value):
    """Convert price string to float."""
    if pd.isna(value) or value == '':
        return 0.0
    if isinstance(value, str):
        value = value.replace('$', '').replace(',', '').strip()
        try:
            return float(value)
        except:
            return 0.0
    return float(value)


def calculate_proposal_pricing(proposal_item: Dict, get_unit_price_func, marketing_rounding: bool, discount_percent: float) -> Dict:
    """
    Calculate all pricing data for a proposal item.
    Replicates the logic from app.py lines 1033-1090.

    Returns dict with: moq, moq_price_per_unit, client_price, delivery_time, customization_setup_fee, customization_per_unit
    """
    # Use settings snapshot if available (for consistency)
    fifty_cent_rounding = False  # Default value
    if 'settings_snapshot' in proposal_item:
        fifty_cent_rounding = proposal_item['settings_snapshot'].get('fifty_cent_rounding', False)
        marketing_rounding = proposal_item['settings_snapshot']['marketing_rounding']
        discount_percent = proposal_item['settings_snapshot']['discount_percent']
        print(f"DEBUG calculate_proposal_pricing: Using settings_snapshot - fifty_cent_rounding={fifty_cent_rounding}, marketing_rounding={marketing_rounding}, discount={discount_percent}")
    else:
        print(f"DEBUG calculate_proposal_pricing: No settings_snapshot found, using defaults - fifty_cent_rounding={fifty_cent_rounding}, marketing_rounding={marketing_rounding}, discount={discount_percent}")

    product_row = pd.Series(proposal_item['product_data'])
    markup_percent = proposal_item['markup_percent']

    # Calculate MOQ using standard preliminary quantity (100 units)
    preliminary_base_price, _, _ = get_unit_price_func(product_row, 100)

    if preliminary_base_price is None:
        return None

    # Estimate total per-unit price with markup
    temp_markup_multiplier = 1 + (markup_percent / 100)
    estimated_unit_price = preliminary_base_price * temp_markup_multiplier

    # Calculate MOQ (returns dict with moq, breakdown, display_text)
    moq_result = calculate_moq(estimated_unit_price, product_row)
    moq = moq_result['moq']
    if moq is None:
        moq = 5

    # Get actual base price for MOQ quantity
    moq_base_price, moq_tier_range, _ = get_unit_price_func(product_row, moq)

    if moq_base_price is None:
        return None

    # Calculate product price WITHOUT customization
    moq_product_cost = moq_base_price * moq
    moq_markup_amount = moq_product_cost * (markup_percent / 100)
    moq_product_only_total = moq_product_cost + moq_markup_amount
    moq_product_price_per_unit = moq_product_only_total / moq

    # Apply $0.50 rounding if enabled (first)
    unrounded_moq_price = moq_product_price_per_unit
    if fifty_cent_rounding:
        moq_product_price_per_unit = round_to_nearest_fifty_cents(moq_product_price_per_unit, True)
        print(f"DEBUG: Applied $0.50 rounding to MOQ price: ${unrounded_moq_price:.2f} → ${moq_product_price_per_unit:.2f}")

    # Apply marketing rounding if enabled (second, after $0.50 rounding)
    if marketing_rounding:
        before_marketing = moq_product_price_per_unit
        moq_product_price_per_unit = apply_marketing_rounding(moq_product_price_per_unit, True)
        if before_marketing != moq_product_price_per_unit:
            print(f"DEBUG: Applied marketing rounding to MOQ price: ${before_marketing:.2f} → ${moq_product_price_per_unit:.2f}")

    # Calculate client price at MOQ (with discount)
    client_price = moq_product_price_per_unit
    if discount_percent > 0:
        client_price = moq_product_price_per_unit * (1 - discount_percent / 100)

    # Calculate price at quantity 100 (for comparison column)
    # Note: preliminary_base_price already has the price at qty 100
    price_at_100_cost = preliminary_base_price * 100
    price_at_100_markup = price_at_100_cost * (markup_percent / 100)
    price_at_100_total = price_at_100_cost + price_at_100_markup
    price_at_100_per_unit = price_at_100_total / 100

    # Apply $0.50 rounding if enabled (first)
    if fifty_cent_rounding:
        price_at_100_per_unit = round_to_nearest_fifty_cents(price_at_100_per_unit, True)

    # Apply marketing rounding if enabled (second, after $0.50 rounding)
    if marketing_rounding:
        price_at_100_per_unit = apply_marketing_rounding(price_at_100_per_unit, True)

    # Calculate client price at quantity 100 (with discount)
    client_price_at_100 = price_at_100_per_unit
    if discount_percent > 0:
        client_price_at_100 = price_at_100_per_unit * (1 - discount_percent / 100)

    # Get customization costs from product data
    setup_fee_raw = get_column_value(product_row, 'Client Price: Customization Setup Fee', 'Customization Setup Fee', '')
    per_unit_raw = get_column_value(product_row, 'Client Price: Customization Cost per Unit', 'Customization Cost per Unit', '')
    setup_fee = clean_price(setup_fee_raw) or 0.0
    per_unit_cost = clean_price(per_unit_raw) or 0.0

    # Get delivery time (default to 6-8 weeks)
    delivery_time = product_row.get('Lead Time', '6-8 weeks')
    if pd.isna(delivery_time) or not delivery_time:
        delivery_time = '6-8 weeks'

    return {
        'moq': moq,
        'moq_price_per_unit': moq_product_price_per_unit,
        'price_at_100': price_at_100_per_unit,
        'client_price': client_price,
        'client_price_at_100': client_price_at_100,
        'delivery_time': delivery_time,
        'customization_setup_fee': setup_fee,
        'customization_per_unit': per_unit_cost,
        'product_data': proposal_item['product_data']  # Include original product data for variant detection
    }


# ============================================================
# VARIANT DETECTION AND GROUPING (v6.13)
# ============================================================

def detect_variant_groups(confirmed_matches: Dict[str, str]) -> Tuple[Dict[str, List[str]], Dict[str, str]]:
    """
    Detect when multiple Google Sheets products match to the same PowerPoint slide.
    This indicates they are variants (sizes/flavors) that should be displayed together.

    Args:
        confirmed_matches: Dict mapping {gs_product_name: pptx_slide_name}

    Returns:
        Tuple of:
            - variant_groups: Dict {pptx_slide_name: [gs_product1, gs_product2, ...]}
            - single_products: Dict {pptx_slide_name: gs_product}

    Example:
        Input: {
            'Strawberry Jam - 4oz': 'STRAWBERRY JAM',
            'Strawberry Jam - 8oz': 'STRAWBERRY JAM',
            'Cutting Board': 'BUTCHER BLOCK'
        }
        Output: (
            {'STRAWBERRY JAM': ['Strawberry Jam - 4oz', 'Strawberry Jam - 8oz']},
            {'BUTCHER BLOCK': 'Cutting Board'}
        )
    """
    # Invert dictionary: group by PowerPoint slide name
    slide_to_products = {}

    for gs_product, pptx_slide in confirmed_matches.items():
        if pptx_slide not in slide_to_products:
            slide_to_products[pptx_slide] = []
        slide_to_products[pptx_slide].append(gs_product)

    # Separate into variant groups (2+ products) and single products
    variant_groups = {}
    single_products = {}

    for pptx_slide, products in slide_to_products.items():
        if len(products) > 1:
            # Multiple products matched to same slide → variant group
            variant_groups[pptx_slide] = products
        else:
            # Single product → normal case
            single_products[pptx_slide] = products[0]

    return variant_groups, single_products


def extract_variant_identifier(proposal_item: Dict, variant_index: int = 0) -> str:
    """
    Extract the variant identifier from product name for display in table's variant column.

    Tries multiple strategies to extract variant info:
    1. Text after last dash: "Product - 4oz" → "4oz"
    2. Size keywords anywhere in name: "Beeswax Candle Small" → "Small"
    3. Numeric + unit patterns: "Product 4oz", "Product 2 oz" → "4oz"
    4. Text in parentheses: "Product (New!)" → "New!"
    5. Fallback: "Option 1", "Option 2", etc.

    Args:
        proposal_item: Dict with product_data containing Product/Service name
        variant_index: Index of this variant (0, 1, 2...) for fallback naming

    Returns:
        str: Variant identifier (e.g., "4oz", "Small", "Cashews", "Option 1")
    """
    product_data = proposal_item.get('product_data', {})
    product_name = product_data.get('Product/Service', '')

    print(f"DEBUG extract_variant_identifier: Processing '{product_name}'")

    # Strategy 1: Try to extract variant after last dash
    if ' - ' in product_name:
        parts = product_name.split(' - ')
        variant = parts[-1].strip()
        # Clean up common patterns: remove parentheses
        variant = re.sub(r'[()]', '', variant)
        print(f"DEBUG extract_variant_identifier: Strategy 1 (dash) → '{variant}'")
        return variant

    # Strategy 2: Check for size keywords ANYWHERE in name (not just first word)
    size_keywords = ['Small', 'Medium', 'Large', 'Mini', 'Extra Large', 'XL', 'Regular',
                     'small', 'medium', 'large', 'mini', 'extra large', 'xl', 'regular']

    for keyword in size_keywords:
        if keyword.lower() in product_name.lower():
            # Return capitalized version
            result = keyword.capitalize()
            print(f"DEBUG extract_variant_identifier: Strategy 2 (size keyword) → '{result}'")
            return result

    # Strategy 3: Extract numeric + unit patterns (4oz, 2 oz, 8oz, etc.)
    unit_pattern = r'(\d+\.?\d*\s*(?:oz|ml|g|kg|lb|ct|pack|pk))'
    unit_match = re.search(unit_pattern, product_name, re.IGNORECASE)
    if unit_match:
        result = unit_match.group(1).strip()
        print(f"DEBUG extract_variant_identifier: Strategy 3 (unit pattern) → '{result}'")
        return result

    # Strategy 4: Extract text from parentheses (if present)
    paren_match = re.search(r'\((.*?)\)', product_name)
    if paren_match:
        result = paren_match.group(1)
        print(f"DEBUG extract_variant_identifier: Strategy 4 (parentheses) → '{result}'")
        return result

    # Fallback: Use "Option N" numbering
    result = f"Option {variant_index + 1}"
    print(f"DEBUG extract_variant_identifier: Fallback → '{result}'")
    return result


def check_pricing_consistency(pricing_data_list: List[Dict]) -> bool:
    """
    Check if all variants have identical pricing (both MOQ and client_price).

    This function determines whether a simple single-row table is sufficient
    or if a multi-row variant table is needed to show pricing differences.

    Args:
        pricing_data_list: List of pricing data dicts from calculate_proposal_pricing()

    Returns:
        bool: True if all variants share same MOQ and client_price, False otherwise

    Examples:
        >>> pricing_data_list = [
        ...     {'moq': 50, 'client_price': 12.00},
        ...     {'moq': 50, 'client_price': 12.00},
        ...     {'moq': 50, 'client_price': 12.00}
        ... ]
        >>> check_pricing_consistency(pricing_data_list)
        True

        >>> pricing_data_list = [
        ...     {'moq': 69, 'client_price': 14.50},
        ...     {'moq': 55, 'client_price': 18.50}
        ... ]
        >>> check_pricing_consistency(pricing_data_list)
        False
    """
    if len(pricing_data_list) <= 1:
        return True

    first_moq = pricing_data_list[0].get('moq', 0)
    first_price = pricing_data_list[0].get('client_price', 0)

    for item in pricing_data_list[1:]:
        item_moq = item.get('moq', 0)
        item_price = item.get('client_price', 0)

        # Check both MOQ and price (must be identical)
        if item_moq != first_moq or abs(item_price - first_price) > 0.01:
            return False

    return True


def update_cell_text_preserve_format(cell, new_text: str):
    """
    Update table cell text while preserving original font formatting.

    Args:
        cell: Table cell object
        new_text: New text to insert
    """
    # Access the text frame
    text_frame = cell.text_frame

    # Clear existing paragraphs but keep the first one for formatting reference
    if len(text_frame.paragraphs) > 0:
        # Get original font properties from first run (if exists)
        original_font = None
        if len(text_frame.paragraphs[0].runs) > 0:
            original_font = text_frame.paragraphs[0].runs[0].font

        # Clear all text
        text_frame.clear()

        # Add new paragraph with new text
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = new_text

        # Apply original font properties if they existed
        if original_font:
            if original_font.size:
                run.font.size = original_font.size
            if original_font.name:
                run.font.name = original_font.name
            if original_font.bold is not None:
                run.font.bold = original_font.bold
            if original_font.italic is not None:
                run.font.italic = original_font.italic
            # Handle color - some cells may not have color defined
            try:
                if original_font.color and hasattr(original_font.color, 'rgb') and original_font.color.rgb:
                    run.font.color.rgb = original_font.color.rgb
            except (AttributeError, ValueError):
                # Skip color if not defined or not accessible
                pass
    else:
        # Fallback if no paragraphs exist
        cell.text = new_text


def find_slide_by_product_name(prs: Presentation, product_name: str) -> Optional[int]:
    """
    Find slide index by matching product name in first shape.

    Args:
        prs: Presentation object
        product_name: Product name to find (from confirmed_matches)

    Returns:
        Slide index (int) or None if not found
    """
    for idx, slide in enumerate(prs.slides):
        if len(slide.shapes) >= 1:
            first_shape = slide.shapes[0]
            if hasattr(first_shape, "text") and first_shape.text.strip():
                if first_shape.text.strip() == product_name:
                    return idx
    return None


def clone_slide(prs: Presentation, slide_index: int) -> object:
    """
    Clone a slide from presentation, preserving all formatting.
    Uses XML deep copy approach to preserve all elements.

    Args:
        prs: Presentation object
        slide_index: Index of slide to clone (0-based)

    Returns:
        Cloned slide object
    """
    # Get source slide
    slide_list = list(prs.slides)
    source_slide = slide_list[slide_index]

    # Get blank layout (typically index 6, but we'll use the source slide's layout)
    source_layout = source_slide.slide_layout

    # Add new slide with same layout
    new_slide = prs.slides.add_slide(source_layout)

    # Remove all default shapes from new slide
    for shape in list(new_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # Copy all shapes from source to new slide
    for shape in source_slide.shapes:
        el = shape.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return new_slide


def copy_slide_with_images(source_prs: Presentation, target_prs: Presentation, slide_index: int) -> object:
    """
    Copy a slide from source to target presentation, preserving image relationships.
    Uses blank layout to avoid layout conflicts.

    Args:
        source_prs: Source presentation to copy from
        target_prs: Target presentation to copy to
        slide_index: Index of slide to copy from source

    Returns:
        New slide object in target presentation
    """
    source_slide = source_prs.slides[slide_index]

    # Use blank layout to avoid any layout conflicts
    blank_layout = target_prs.slide_layouts[6]
    new_slide = target_prs.slides.add_slide(blank_layout)

    # Remove default shapes
    for shape in list(new_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # Copy shapes from source slide
    for shape in source_slide.shapes:
        el = shape.element
        newel = copy.deepcopy(el)

        # Handle image relationships (blips)
        blips = newel.xpath('.//a:blip')

        for blip in blips:
            embed_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
            old_rId = blip.get(embed_attr)

            if old_rId:
                try:
                    # Get the image part from source slide
                    source_image_part = source_slide.part.related_part(old_rId)

                    # Add image to target presentation package
                    image_file = io.BytesIO(source_image_part.blob)
                    image_part = target_prs.part.package._image_parts.get_or_add_image_part(image_file)

                    # Create relationship from new slide to image
                    new_rId = new_slide.part.relate_to(image_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')

                    # Update the blip to use new relationship ID
                    blip.set(embed_attr, new_rId)
                except:
                    pass  # Skip failed images

        # Add the shape to new slide
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return new_slide


def update_pricing_table(slide: object, proposal_items, variant_mode: bool = False, discount_percent: float = 0.0, discount_type: str = None) -> bool:
    """
    Update pricing table in slide with proposal data.
    Supports both single-product and multi-variant tables.
    Handles table formats: 2x3, 2x4, 3x4, and multi-row variant tables (v6.13).

    Args:
        slide: Slide object with table
        proposal_items: Single Dict OR List[Dict] for variants (when variant_mode=True)
        variant_mode: If True, proposal_items is a list of variants
        discount_percent: Discount percentage applied (default: 0.0)
        discount_type: Type of discount - 'Non-profit', 'Volume Order', 'Custom', or None (default: None)

    Returns:
        True if table was updated, False if no table found

    Examples:
        Single product:
            update_pricing_table(slide, proposal_item, variant_mode=False, discount_percent=5.0, discount_type='Volume Order')

        Multi-variant (2 sizes):
            update_pricing_table(slide, [item_4oz, item_8oz], variant_mode=True, discount_percent=5.0, discount_type='Non-profit')
    """
    # Find table in slide
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if not table:
        return False

    # Convert single item to list for consistent handling
    if not variant_mode:
        proposal_items = [proposal_items]
    else:
        print(f"DEBUG update_pricing_table: variant_mode=True with {len(proposal_items)} items")

    # Detect table structure
    rows = len(table.rows)
    cols = len(table.columns)
    print(f"DEBUG update_pricing_table: Table structure {rows} rows x {cols} cols")

    # Detect customization row (last row with keywords like "Customization", "Artwork", "Branding")
    customization_row_idx = None
    for r in range(rows - 1, 0, -1):  # Start from bottom, skip header
        try:
            cell_text = table.cell(r, 0).text.lower()
            if any(keyword in cell_text for keyword in ['customization', 'artwork', 'branding', 'setup', 'sticker', 'engraving', 'no customization']):
                customization_row_idx = r
                break
        except:
            continue

    # Calculate available data rows (between header row 0 and customization row)
    if customization_row_idx:
        available_data_rows = customization_row_idx - 1  # Rows 1 to customization_row_idx-1
    else:
        available_data_rows = rows - 1  # All rows except header

    # Check if we need to add rows for all variants
    num_variants = len(proposal_items)
    if num_variants > available_data_rows:
        # Add new rows to accommodate all variants
        rows_to_add = num_variants - available_data_rows
        print(f"DEBUG: Adding {rows_to_add} rows to table for {num_variants} variants")

        import copy

        for i in range(rows_to_add):
            # Clone row 1 (first data row) using XML deep copy
            source_tr = table._tbl.tr_lst[1]  # Row 1 (first data row after header)
            new_tr = copy.deepcopy(source_tr)

            # Insert new row BEFORE customization row (if exists) or at end
            if customization_row_idx:
                # Insert before customization row
                table._tbl.insert(customization_row_idx + i, new_tr)
            else:
                # Append at end
                table._tbl.append(new_tr)

        # Update customization_row_idx (it shifted down by rows_to_add)
        if customization_row_idx:
            customization_row_idx += rows_to_add

        # Update rows count
        rows = len(table.rows)
        available_data_rows = num_variants

    # Read delivery time from PowerPoint template (preserves original template value)
    template_delivery_time = None
    if cols >= 3:
        try:
            # Read from row 1 (first data row), last column (delivery column)
            template_delivery_time = table.cell(1, cols - 1).text.strip()
            if template_delivery_time:
                print(f"DEBUG: Using template delivery time: '{template_delivery_time}'")
        except:
            pass

    # For variant mode with 4 columns, check if we can simplify the table layout
    can_simplify_table = False
    if variant_mode and cols == 4:
        # Check if price at MOQ equals price at qty 100 for the first variant
        first_variant = proposal_items[0]
        first_client_price = first_variant.get('client_price', 0)
        first_client_price_at_100 = first_variant.get('client_price_at_100', 0)

        # Can simplify if columns 2 and 3 would show identical values
        # (no tiered pricing benefit between MOQ and 100, and no discount difference)
        can_simplify_table = abs(first_client_price - first_client_price_at_100) < 0.01  # Use small epsilon for float comparison
        print(f"DEBUG: can_simplify_table={can_simplify_table} (client_price={first_client_price:.2f}, client_price_at_100={first_client_price_at_100:.2f})")

    # Update each data row with variant data
    for idx, proposal_item in enumerate(proposal_items):
        data_row_idx = idx + 1  # Row 0 is header, data starts at row 1

        # Extract variant identifier for first column
        if variant_mode:
            variant_identifier = extract_variant_identifier(proposal_item, idx)
            print(f"DEBUG update_pricing_table: Row {data_row_idx} - Variant ID: {variant_identifier}")
        else:
            # Single product: use MOQ number in first column
            variant_identifier = str(proposal_item.get('moq', 10))

        # Get pricing data
        moq = proposal_item.get('moq', 10)
        base_price = proposal_item.get('moq_price_per_unit', 0)
        price_at_100 = proposal_item.get('price_at_100', base_price)
        client_price = proposal_item.get('client_price', base_price)
        client_price_at_100 = proposal_item.get('client_price_at_100', client_price)

        # Delivery time: use template value if available, otherwise use Google Sheets data
        if template_delivery_time:
            delivery_time = template_delivery_time
        else:
            delivery_time = proposal_item.get('delivery_time', '6-8 weeks')

        # Format prices
        base_price_str = f"${base_price:.2f}"
        price_at_100_str = f"${price_at_100:.2f}"
        client_price_str = f"${client_price:.2f}"
        client_price_at_100_str = f"${client_price_at_100:.2f}"

        # Update cells based on table format
        if cols == 3:
            # Format 2x3: MOQ/Variant | Price Ea | Delivery
            update_cell_text_preserve_format(table.cell(data_row_idx, 0), variant_identifier)
            update_cell_text_preserve_format(table.cell(data_row_idx, 1), client_price_str)

            # Delivery time - always update (in variant mode, all rows show same delivery time)
            update_cell_text_preserve_format(table.cell(data_row_idx, 2), delivery_time)

            # Update header for price column (only for first variant/single product)
            if idx == 0 and not variant_mode:
                # Generate price header with discount label if applicable
                price_header = f"Price Ea\n(@ Qty {moq})"
                if discount_percent > 0 and discount_type:
                    if discount_type == 'Non-profit':
                        price_header = f"Client Price\n(5% Non-profit discount)"
                    elif discount_type == 'Volume Order':
                        price_header = f"Client Price\n(5% Volume Order discount)"
                    else:
                        price_header = f"Client Price\n({discount_percent:.1f}% discount)"

                update_cell_text_preserve_format(table.cell(0, 1), price_header)

        elif cols == 4:
            if variant_mode and can_simplify_table:
                # SIMPLIFIED LAYOUT: Variant | MOQ | Price Ea @ MOQ | Delivery
                # Use this when price at MOQ equals price at qty 100 (no tiered pricing, no discount)
                update_cell_text_preserve_format(table.cell(data_row_idx, 0), variant_identifier)
                update_cell_text_preserve_format(table.cell(data_row_idx, 1), str(moq))
                update_cell_text_preserve_format(table.cell(data_row_idx, 2), client_price_str)
                update_cell_text_preserve_format(table.cell(data_row_idx, 3), delivery_time)

                # Update headers (only once)
                if idx == 0:
                    update_cell_text_preserve_format(table.cell(0, 0), "Variant")
                    update_cell_text_preserve_format(table.cell(0, 1), "MOQ")
                    update_cell_text_preserve_format(table.cell(0, 2), "Price Ea @ MOQ")
                    update_cell_text_preserve_format(table.cell(0, 3), "Delivery")

            elif variant_mode and not can_simplify_table:
                # FULL LAYOUT: Variant | Price @ MOQ | Price @ Qty 100 | Delivery
                # Use this when prices differ (tiered pricing OR discount applied)
                update_cell_text_preserve_format(table.cell(data_row_idx, 0), variant_identifier)
                update_cell_text_preserve_format(table.cell(data_row_idx, 1), client_price_str)
                update_cell_text_preserve_format(table.cell(data_row_idx, 2), client_price_at_100_str)
                update_cell_text_preserve_format(table.cell(data_row_idx, 3), delivery_time)

                # Update headers (only once)
                if idx == 0:
                    # Generate price @ 100 header with discount label if applicable
                    price_100_header = "Price Ea\n(@ Qty 100)"
                    if discount_percent > 0 and discount_type:
                        if discount_type == 'Non-profit':
                            price_100_header = f"Client Price\n(5% Non-profit discount)"
                        elif discount_type == 'Volume Order':
                            price_100_header = f"Client Price\n(5% Volume Order discount)"
                        else:
                            price_100_header = f"Client Price\n({discount_percent:.1f}% discount)"

                    update_cell_text_preserve_format(table.cell(0, 0), "Variant")
                    update_cell_text_preserve_format(table.cell(0, 1), "Price Ea @ MOQ")
                    update_cell_text_preserve_format(table.cell(0, 2), price_100_header)
                    update_cell_text_preserve_format(table.cell(0, 3), "Delivery")

            else:
                # SINGLE PRODUCT MODE: MOQ | Price @ MOQ | Price @ Qty 100 | Delivery
                update_cell_text_preserve_format(table.cell(data_row_idx, 0), str(moq))
                update_cell_text_preserve_format(table.cell(data_row_idx, 1), client_price_str)
                update_cell_text_preserve_format(table.cell(data_row_idx, 2), client_price_at_100_str)
                update_cell_text_preserve_format(table.cell(data_row_idx, 3), delivery_time)

                # Update headers (only once)
                if idx == 0:
                    # Generate price @ 100 header with discount label if applicable
                    price_100_header = f"Price Ea\n(@ Qty 100)"
                    if discount_percent > 0 and discount_type:
                        if discount_type == 'Non-profit':
                            price_100_header = f"Client Price\n(5% Non-profit discount)"
                        elif discount_type == 'Volume Order':
                            price_100_header = f"Client Price\n(5% Volume Order discount)"
                        else:
                            price_100_header = f"Client Price\n({discount_percent:.1f}% discount)"

                    update_cell_text_preserve_format(table.cell(0, 0), "MOQ")
                    update_cell_text_preserve_format(table.cell(0, 1), f"Price Ea\n(@ Qty {moq})")
                    update_cell_text_preserve_format(table.cell(0, 2), price_100_header)
                    update_cell_text_preserve_format(table.cell(0, 3), "Delivery")

    # Update customization row if it exists (preserve at bottom, don't overwrite)
    # Only update if we have customization data from the first variant
    if customization_row_idx and not variant_mode:
        setup_fee = proposal_items[0].get('customization_setup_fee', 0)
        per_unit_cost = proposal_items[0].get('customization_per_unit', 0)

        if setup_fee > 0 or per_unit_cost > 0:
            # Format customization text
            if per_unit_cost > 0:
                customization_text = f"Artwork set-up: ${setup_fee:.2f} / Engraving per piece: From ${per_unit_cost:.2f}"
            else:
                customization_text = f"Artwork set-up: ${setup_fee:.2f}"

            update_cell_text_preserve_format(table.cell(customization_row_idx, 0), customization_text)

    return True


def add_cover_slide(prs: Presentation, client_name: str, date_str: str) -> object:
    """
    Add cover slide to beginning of presentation.

    Args:
        prs: Presentation object
        client_name: Company name from order details
        date_str: Formatted date string

    Returns:
        Cover slide object
    """
    # Use title slide layout (typically index 0)
    title_layout = prs.slide_layouts[0]

    # Add slide at beginning
    slide = prs.slides.add_slide(title_layout)

    # Set title and subtitle
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = f"Product Proposal for {client_name}"
    subtitle.text = f"Peace by Piece International\n{date_str}"

    # Move to beginning
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(0, slides[-1])

    return slide


def create_proposal_presentation(
    template_path: str,
    confirmed_matches: Dict[str, str],
    proposal_products: List[Dict],
    get_unit_price_func,
    marketing_rounding: bool = False,
    discount_percent: float = 0.0,
    discount_type: str = None
) -> Presentation:
    """
    Create new presentation with only confirmed product slides.
    Strategy: Load template, keep only needed slides, update pricing tables.

    Args:
        template_path: Path to "November All Slides.pptx"
        confirmed_matches: Dict of {gs_product_name: pptx_product_name}
        proposal_products: List of proposal items from session state
        get_unit_price_func: Function to calculate unit prices
        marketing_rounding: Whether to apply marketing rounding
        discount_percent: Discount percentage to apply

    Returns:
        New Presentation object with selected and updated slides
    """
    # Load template (we'll modify this directly)
    prs = Presentation(template_path)

    # Find slide indices to keep and their proposal data
    slides_to_keep = {}  # {slide_idx: proposal_item}

    for gs_name, pptx_name in confirmed_matches.items():
        # Find slide in template
        slide_idx = find_slide_by_product_name(prs, pptx_name)

        if slide_idx is not None:
            # Find matching proposal item
            proposal_item = next(
                (item for item in proposal_products
                 if item['product_data']['Product/Service'] == gs_name),
                None
            )

            slides_to_keep[slide_idx] = proposal_item

    # Remove slides we don't need (work backwards to preserve indices)
    slide_list = list(prs.slides)
    for idx in range(len(slide_list) - 1, -1, -1):
        if idx not in slides_to_keep:
            # Remove this slide
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]

    # Force garbage collection after slide removal (memory optimization)
    gc.collect()

    # Update pricing tables in remaining slides
    # Note: After deletion, slide indices change, so we need to match by name
    for slide in prs.slides:
        # Get product name from first shape
        if len(slide.shapes) >= 1:
            first_shape = slide.shapes[0]
            if hasattr(first_shape, "text") and first_shape.text.strip():
                product_name = first_shape.text.strip()

                # Find the proposal item for this slide
                for gs_name, pptx_name in confirmed_matches.items():
                    if pptx_name == product_name:
                        # Find proposal item
                        proposal_item = next(
                            (item for item in proposal_products
                             if item['product_data']['Product/Service'] == gs_name),
                            None
                        )

                        if proposal_item:
                            # Calculate pricing data
                            pricing_data = calculate_proposal_pricing(
                                proposal_item,
                                get_unit_price_func,
                                marketing_rounding,
                                discount_percent
                            )

                            if pricing_data:
                                update_pricing_table(slide, pricing_data, discount_percent=discount_percent, discount_type=discount_type)
                        break

    return prs


def create_proposal_presentation_with_impact(
    template_path: str,
    confirmed_matches: Dict[str, str],
    proposal_products: List[Dict],
    impact_slides_by_partner: Dict[str, Dict],
    get_unit_price_func,
    marketing_rounding: bool = False,
    discount_percent: float = 0.0,
    discount_type: str = None
) -> Presentation:
    """
    Create presentation with product AND impact slides in proper order.

    Assembly Order:
    1. Product slides (grouped by partner)
    2. Impact slides (after each partner's products)
    3. Outro slides (placeholder for future)

    Args:
        template_path: Path to PowerPoint template
        confirmed_matches: Dict of {gs_product_name: pptx_product_name}
        proposal_products: List of proposal items
        impact_slides_by_partner: Dict of {partner: {"slide_index": X, "slide_title": Y}}
        get_unit_price_func: Function to calculate unit prices
        marketing_rounding: Whether to apply marketing rounding
        discount_percent: Discount percentage to apply

    Returns:
        New Presentation with selected and ordered slides
    """
    # Load template
    prs = Presentation(template_path)

    # Build slide selection strategy
    # Group products by partner
    products_by_partner = {}

    for gs_name, pptx_name in confirmed_matches.items():
        # Find proposal item
        proposal_item = next(
            (item for item in proposal_products
             if item['product_data']['Product/Service'] == gs_name),
            None
        )

        if proposal_item:
            partner = proposal_item['product_data']['Partner']
            if partner not in products_by_partner:
                products_by_partner[partner] = []

            # Find slide index
            slide_idx = find_slide_by_product_name(prs, pptx_name)
            if slide_idx is not None:
                products_by_partner[partner].append({
                    'slide_idx': slide_idx,
                    'product_name': pptx_name,
                    'proposal_item': proposal_item,
                    'gs_name': gs_name
                })

    # Build ordered list of slides to keep
    slides_to_keep_ordered = []

    for partner in sorted(products_by_partner.keys()):
        # Add all product slides for this partner
        for product_entry in products_by_partner[partner]:
            slides_to_keep_ordered.append({
                'type': 'product',
                'index': product_entry['slide_idx'],
                'data': product_entry
            })

        # Add impact slide for this partner (if selected)
        if partner in impact_slides_by_partner:
            impact_selection = impact_slides_by_partner[partner]
            slides_to_keep_ordered.append({
                'type': 'impact',
                'index': impact_selection['slide_index'],
                'data': {
                    'partner': partner,
                    'slide_index': impact_selection['slide_index'],
                    'slide_title': impact_selection['slide_title']
                }
            })

    # Extract just the slide indices we need (in order)
    slide_indices_to_keep = set(entry['index'] for entry in slides_to_keep_ordered)

    # Remove slides we don't need (work backwards)
    slide_list = list(prs.slides)
    for idx in range(len(slide_list) - 1, -1, -1):
        if idx not in slide_indices_to_keep:
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]

    # Force garbage collection after slide removal (memory optimization)
    gc.collect()

    # Update pricing tables in product slides
    for slide in prs.slides:
        if len(slide.shapes) >= 1:
            first_shape = slide.shapes[0]
            if hasattr(first_shape, "text") and first_shape.text.strip():
                product_name = first_shape.text.strip()

                # Find if this is a product slide that needs updating
                for gs_name, pptx_name in confirmed_matches.items():
                    if pptx_name == product_name:
                        proposal_item = next(
                            (item for item in proposal_products
                             if item['product_data']['Product/Service'] == gs_name),
                            None
                        )

                        if proposal_item:
                            pricing_data = calculate_proposal_pricing(
                                proposal_item,
                                get_unit_price_func,
                                marketing_rounding,
                                discount_percent
                            )

                            if pricing_data:
                                update_pricing_table(slide, pricing_data, discount_percent=discount_percent, discount_type=discount_type)
                        break

    # Note: Outro slides now handled by create_complete_proposal_presentation()

    return prs


def download_presentation(prs: Presentation, client_name: str) -> io.BytesIO:
    """
    Convert presentation to bytes and return for download.

    Args:
        prs: Presentation object
        client_name: For filename generation

    Returns:
        BytesIO object for st.download_button
    """
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ============================================================
# PHASE 2.5: COMPLETE PROPOSAL GENERATION (INTRO + PRODUCTS + IMPACT + OUTRO)
# ============================================================

def extract_slides_from_template(template_path: str, slide_indices: List[int]) -> List:
    """
    Extract specific slides from a template by their indices.

    Args:
        template_path: Path to PowerPoint template
        slide_indices: List of slide indices to extract (0-based)

    Returns:
        List of slide objects (cloned)

    Example:
        >>> intro_slides = extract_slides_from_template("intro_outro.pptx", [0, 1, 2, 3, 4, 5, 6, 7])
    """
    prs = Presentation(template_path)
    slides = []

    for idx in slide_indices:
        if 0 <= idx < len(prs.slides):
            slides.append(prs.slides[idx])

    return slides


def create_complete_proposal_presentation(
    november_template_path: str,
    intro_outro_template_path: str,
    confirmed_matches: Dict[str, str],
    proposal_products: List[Dict],
    get_unit_price_func,
    marketing_rounding: bool = False,
    discount_percent: float = 0.0,
    impact_slide_overrides: Optional[Dict[str, Dict]] = None,
    variant_groups: Optional[Dict[str, List[str]]] = None,
    variant_grouping_prefs: Optional[Dict[str, str]] = None,
    fifty_cent_rounding: bool = False,
    discount_type: str = None
) -> Presentation:
    """
    Create complete presentation with intro, products, impact, and outro slides.
    Supports multi-variant product consolidation (v6.13).

    Slide Assembly Order:
    1. Intro slides (1-8) from Intro_Outro_Slides_PbP_Proposals.pptx
    2. Product slides (customized) from November All Slides.pptx
    3. Impact slides (one per partner, auto-selected from reference table) from November All Slides.pptx
    4. Outro slides (9-12) from Intro_Outro_Slides_PbP_Proposals.pptx

    Args:
        november_template_path: Path to November All Slides.pptx
        intro_outro_template_path: Path to Intro_Outro_Slides_PbP_Proposals.pptx
        confirmed_matches: Dict of {gs_product_name: pptx_product_name}
        proposal_products: List of proposal items
        get_unit_price_func: Function to calculate unit prices
        marketing_rounding: Whether to apply marketing rounding
        discount_percent: Discount percentage to apply
        discount_type: Type of discount - 'Non-profit', 'Volume Order', 'Custom', or None
        impact_slide_overrides: Optional dict of {partner: {"slide_index": X, "slide_title": Y}}
                                If provided, overrides reference table for that partner
        variant_groups: Optional dict of {pptx_slide_name: [gs_product1, gs_product2, ...]}
                       Identifies which products are variants of same slide
        variant_grouping_prefs: Optional dict of {pptx_slide_name: user_choice_string}
                               User preferences for how to handle each variant group

    Returns:
        Complete Presentation with all slides assembled

    Example:
        >>> prs = create_complete_proposal_presentation(
        ...     "templates/November All Slides.pptx",
        ...     "templates/Intro_Outro_Slides_PbP_Proposals.pptx",
        ...     {"Product A": "PRODUCT A"},
        ...     proposal_products,
        ...     get_unit_price_func
        ... )
    """
    # Import reference table
    from src.slide_matcher import PARTNER_IMPACT_SLIDES, extract_unique_partners

    # Load November template (we'll modify this)
    prs_november = Presentation(november_template_path)

    # Build list of slides to keep from November template
    # Format: [(slide_index, type, data), ...]
    slides_to_keep = []

    # Step 1: Find product slides
    for gs_name, pptx_name in confirmed_matches.items():
        slide_idx = find_slide_by_product_name(prs_november, pptx_name)

        if slide_idx is not None:
            # Find matching proposal item
            proposal_item = next(
                (item for item in proposal_products
                 if item['product_data']['Product/Service'] == gs_name),
                None
            )

            if proposal_item:
                slides_to_keep.append({
                    'index': slide_idx,
                    'type': 'product',
                    'data': {
                        'gs_name': gs_name,
                        'pptx_name': pptx_name,
                        'proposal_item': proposal_item
                    }
                })

    # Step 2: Find impact slides (one per partner)
    unique_partners = extract_unique_partners(proposal_products)

    for partner in unique_partners:
        # Check for override first
        if impact_slide_overrides and partner in impact_slide_overrides:
            impact_info = impact_slide_overrides[partner]
        else:
            # Use reference table
            impact_info = PARTNER_IMPACT_SLIDES.get(partner)

        if impact_info and impact_info.get('slide_index') is not None:
            slides_to_keep.append({
                'index': impact_info['slide_index'],
                'type': 'impact',
                'data': {
                    'partner': partner,
                    'slide_title': impact_info['slide_title']
                }
            })

    # Step 3: Remove slides we don't need from November template (work backwards)
    slide_indices_to_keep = set(entry['index'] for entry in slides_to_keep)
    slide_list = list(prs_november.slides)

    for idx in range(len(slide_list) - 1, -1, -1):
        if idx not in slide_indices_to_keep:
            rId = prs_november.slides._sldIdLst[idx].rId
            prs_november.part.drop_rel(rId)
            del prs_november.slides._sldIdLst[idx]

    # Force garbage collection after slide removal (memory optimization)
    gc.collect()

    # Step 4: Update pricing tables in product slides (with variant support)
    for slide in prs_november.slides:
        if len(slide.shapes) >= 1:
            first_shape = slide.shapes[0]
            if hasattr(first_shape, "text") and first_shape.text.strip():
                slide_title = first_shape.text.strip()

                # Check if this slide has variant groups
                is_variant_group = variant_groups and slide_title in variant_groups

                if is_variant_group:
                    # Get user preference for this variant group
                    user_choice = variant_grouping_prefs.get(slide_title, "") if variant_grouping_prefs else ""
                    print(f"DEBUG: Variant group detected for '{slide_title}' with {len(variant_groups[slide_title])} products")
                    print(f"DEBUG: User choice: '{user_choice}'")

                    if "single row" in user_choice.lower():
                        # Display simple single-row table (consistent pricing)
                        # Use first variant as representative since all have same pricing
                        variant_products = variant_groups[slide_title]
                        gs_name = variant_products[0]

                        proposal_item = next(
                            (item for item in proposal_products
                             if item['product_data']['Product/Service'] == gs_name),
                            None
                        )

                        if proposal_item:
                            pricing_data = calculate_proposal_pricing(
                                proposal_item,
                                get_unit_price_func,
                                marketing_rounding,
                                discount_percent
                            )

                            if pricing_data:
                                # Update table with single row (variant_mode=False for simple table)
                                update_pricing_table(slide, pricing_data, variant_mode=False, discount_percent=discount_percent, discount_type=discount_type)
                                print(f"DEBUG: Created single-row table for {slide_title} (consistent pricing)")

                    elif "all variants" in user_choice.lower() or "together" in user_choice.lower():
                        # Display variants together - collect all variants for this slide
                        variant_products = variant_groups[slide_title]
                        pricing_data_list = []

                        for gs_name in variant_products:
                            proposal_item = next(
                                (item for item in proposal_products
                                 if item['product_data']['Product/Service'] == gs_name),
                                None
                            )

                            if proposal_item:
                                pricing_data = calculate_proposal_pricing(
                                    proposal_item,
                                    get_unit_price_func,
                                    marketing_rounding,
                                    discount_percent
                                )

                                if pricing_data:
                                    pricing_data_list.append(pricing_data)

                        # Update table with all variants (variant_mode=True)
                        if pricing_data_list:
                            update_pricing_table(slide, pricing_data_list, variant_mode=True, discount_percent=discount_percent, discount_type=discount_type)
                            print(f"DEBUG: Created multi-row table for {slide_title} with {len(pricing_data_list)} variants")

                    elif "skip" in user_choice.lower():
                        # Skip - slide should have been removed already, but just in case
                        pass

                    # If "separate" was chosen, each variant gets its own slide
                    # This would require slide duplication, which is handled differently
                    # For now, we'll treat "separate" like single products (first variant wins)

                else:
                    # Single product - normal behavior
                    for gs_name, pptx_name in confirmed_matches.items():
                        if pptx_name == slide_title:
                            proposal_item = next(
                                (item for item in proposal_products
                                 if item['product_data']['Product/Service'] == gs_name),
                                None
                            )

                            if proposal_item:
                                pricing_data = calculate_proposal_pricing(
                                    proposal_item,
                                    get_unit_price_func,
                                    marketing_rounding,
                                    discount_percent
                                )

                                if pricing_data:
                                    update_pricing_table(slide, pricing_data, variant_mode=False, discount_percent=discount_percent, discount_type=discount_type)
                            break

    # Step 5: Add intro/outro slides to END in their original order
    # This keeps intro slides together, making it obvious they need to be moved to beginning
    prs_intro_outro = Presentation(intro_outro_template_path)

    # Add intro slides first (slides 1-8, indices 0-7)
    intro_slide_count = min(8, len(prs_intro_outro.slides))
    for idx in range(intro_slide_count):
        copy_slide_with_images(prs_intro_outro, prs_november, idx)

    # Add outro slides after intro (slides 9-12, indices 8-11)
    outro_start_idx = 8
    outro_end_idx = min(12, len(prs_intro_outro.slides))
    for idx in range(outro_start_idx, outro_end_idx):
        copy_slide_with_images(prs_intro_outro, prs_november, idx)

    # Final structure: Products → Impacts → Intro (at end) → Outro (at end)
    # User moves intro slides from end to beginning, leaves outro at end

    return prs_november
