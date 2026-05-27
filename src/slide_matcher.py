"""
Slide Matcher Module
Intelligent product name matching system with fuzzy matching fallback.

Logic:
1. Try exact match (case-insensitive)
2. If no exact match, use fuzzy matching with confidence threshold
3. Return match result with confidence score for user review

Dependencies:
- thefuzz: Fuzzy string matching library
- rapidfuzz: Fast string matching (auto-installed with thefuzz)
"""

from thefuzz import fuzz, process
from typing import Optional, Dict, List, Tuple
import re
from .match_manager import (
    load_manual_matches,
    get_manual_match
)
from pptx import Presentation


def find_best_match_multi_scorer(
    query: str,
    choices: List[str],
    limit: int = 3
) -> Tuple[str, int, str, List[Tuple[str, int]]]:
    """
    Use multiple fuzzy matching algorithms and return best result.

    Tests 3 different scoring algorithms and returns the match with
    the highest confidence score across all methods.

    Args:
        query: Product name from Google Sheets
        choices: List of product names from PowerPoint
        limit: Number of alternatives to return

    Returns:
        Tuple of (best_match_name, best_score, method_used, alternatives)

    Example:
        >>> find_best_match_multi_scorer("Laptop Sleeve", ["UPCYCLED LAPTOP SLEEVE", "BRIEFCASE"])
        ("UPCYCLED LAPTOP SLEEVE", 85, "token_set_ratio", [("BRIEFCASE", 20)])
    """
    scorers = [
        ('token_sort_ratio', fuzz.token_sort_ratio),
        ('token_set_ratio', fuzz.token_set_ratio),
        ('partial_ratio', fuzz.partial_ratio),
    ]

    best_overall_match = None
    best_overall_score = 0
    best_method = None
    all_alternatives = {}  # {name: best_score_seen}

    # Try each scoring method
    for method_name, scorer_func in scorers:
        matches = process.extract(
            query,
            choices,
            scorer=scorer_func,
            limit=limit
        )

        if matches:
            # Check if this method found a better match
            top_match_name, top_match_score = matches[0]
            if top_match_score > best_overall_score:
                best_overall_match = top_match_name
                best_overall_score = top_match_score
                best_method = method_name

            # Track all alternatives (keep highest score for each name)
            for match_name, match_score in matches:
                if match_name not in all_alternatives or match_score > all_alternatives[match_name]:
                    all_alternatives[match_name] = match_score

    # Build alternatives list (exclude best match, sort by score, limit to 'limit' items)
    alternatives = [
        (name, score) for name, score in all_alternatives.items()
        if name != best_overall_match
    ]
    alternatives.sort(key=lambda x: x[1], reverse=True)
    alternatives = alternatives[:limit]

    return best_overall_match, best_overall_score, best_method, alternatives


# Category keywords for boosting match confidence
CATEGORY_KEYWORDS = {
    'bags': ['BAG', 'BACKPACK', 'TOTE', 'BRIEFCASE', 'SLEEVE', 'POUCH'],
    'cutting_boards': ['CUTTING', 'BOARD', 'BUTCHER', 'BLOCK'],
    'candles': ['CANDLE', 'HOLDER', 'VOTIVE', 'TEA LIGHT'],
    'trivets': ['TRIVET'],
    'coasters': ['COASTER'],
    'bowls': ['BOWL'],
    'trays': ['TRAY'],
    'jewelry': ['BRACELET', 'NECKLACE', 'EARRING', 'JEWELRY'],
    'home_decor': ['WALL ART', 'PICTURE FRAME', 'VASE', 'PLANTER'],
}

# Keywords that suggest customization/variant (not product categories)
VARIANT_KEYWORDS = ['NOIR', 'MOF', 'LARGE', 'SMALL', 'SET', 'ENFOLD']

# Manual product mappings for known mismatches
# Format: {normalized_gs_name: exact_pptx_name}
# These override fuzzy matching with 100% confidence
MANUAL_PRODUCT_MAPPINGS = {
    'UPCYCLED EXECUTIVE URBAN BRIEFCASE': 'UPCYCLED EXECUTIVE URBAN BRIEFCASE',
    'UPCYCLED LAPTOP SLEEVE': 'UPCYCLED LAPTOP SLEEVE',
    'UPCYCLED DAY TRIPPER BACKPACK': 'UPCYCLED DAY TRIPPER BACKPACK',
    'BUTCHER BLOCK': 'BUTCHER BLOCK',
    'CANDLE HOLDERS': 'MINIMALIST CANDLE HOLDERS – Set of 3',
    # Add more known mappings as discovered during testing
}


def boost_score_if_same_category(
    query: str,
    match: str,
    base_score: int,
    boost_amount: int = 15
) -> int:
    """
    Boost confidence score if query and match share category keywords.

    Args:
        query: Product name from Google Sheets
        match: Product name from PowerPoint
        base_score: Original confidence score (0-100)
        boost_amount: Points to add if same category (default: 15)

    Returns:
        Boosted score (capped at 100)

    Example:
        >>> boost_score_if_same_category("Cutting Board - Large", "BUTCHER BLOCK", 75)
        90  # Both have 'CUTTING'/'BUTCHER' + 'BOARD'/'BLOCK' keywords
    """
    query_upper = query.upper()
    match_upper = match.upper()

    # Check if both products share keywords from the same category
    for category, keywords in CATEGORY_KEYWORDS.items():
        query_has_keyword = any(kw in query_upper for kw in keywords)
        match_has_keyword = any(kw in match_upper for kw in keywords)

        if query_has_keyword and match_has_keyword:
            # Both products are in the same category - boost confidence
            boosted_score = min(base_score + boost_amount, 100)
            return boosted_score

    # No shared category - return original score
    return base_score


def normalize_product_name(product_name: str) -> str:
    """
    Normalize product name by removing common variant suffixes.

    Strips patterns like:
    - (Noir), (Enfold), (any text in parentheses)
    - -MOF, - Large, - Small, - Set of 3
    - Extra whitespace and dashes

    Args:
        product_name: Original product name

    Returns:
        Normalized product name

    Example:
        >>> normalize_product_name("Upcycled Laptop Sleeve (Enfold)-MOF")
        "UPCYCLED LAPTOP SLEEVE"
        >>> normalize_product_name("Butcher Block - Large")
        "BUTCHER BLOCK"
        >>> normalize_product_name("Candle Holders - Set of 3")
        "CANDLE HOLDERS"
    """
    # Convert to uppercase for consistency
    normalized = product_name.upper()

    # Remove text in parentheses: (Noir), (Enfold), etc.
    normalized = re.sub(r'\([^)]*\)', '', normalized)

    # Remove common variant suffixes
    variant_patterns = [
        r'-\s*MOF',                    # -MOF
        r'-\s*LARGE',                  # - Large
        r'-\s*SMALL',                  # - Small
        r'[-–]\s*SET OF \d+',          # - Set of 3, – Set of 3
        r'-\s*ENFOLD',                 # -ENFOLD (redundant but explicit)
        r'-\s*NOIR',                   # -NOIR (redundant but explicit)
    ]

    for pattern in variant_patterns:
        normalized = re.sub(pattern, '', normalized, flags=re.IGNORECASE)

    # Clean up extra whitespace and dashes
    normalized = re.sub(r'\s+', ' ', normalized)  # Multiple spaces → single space
    normalized = re.sub(r'\s*-\s*$', '', normalized)  # Trailing dash
    normalized = re.sub(r'^\s*-\s*', '', normalized)  # Leading dash
    normalized = normalized.strip()

    return normalized


class SlideMatchResult:
    """
    Represents the result of matching a product name to a PowerPoint slide.

    Attributes:
        gs_product_name: Product name from Google Sheets
        pptx_product_name: Matched product name from PowerPoint (None if no match)
        match_type: 'exact', 'fuzzy', or 'none'
        confidence: Confidence score (0-100)
        alternatives: List of alternative matches with scores [(name, score), ...]
        match_source: 'manual' or 'automatic' (whether match came from manual override)
        match_method: Description of matching method used (e.g., "Manual override", "Token sort ratio (85%)")
    """

    def __init__(self, gs_product_name: str, pptx_product_name: Optional[str],
                 match_type: str, confidence: int, alternatives: List[Tuple[str, int]] = None,
                 match_source: str = 'automatic', match_method: str = ''):
        self.gs_product_name = gs_product_name
        self.pptx_product_name = pptx_product_name
        self.match_type = match_type
        self.confidence = confidence
        self.alternatives = alternatives or []
        self.match_source = match_source
        self.match_method = match_method

    def is_usable(self, min_confidence: int = 70) -> bool:
        """Check if match is good enough to use automatically."""
        return self.match_type != 'none' and self.confidence >= min_confidence

    def __repr__(self):
        if self.match_source == 'confirmed':
            source_indicator = "[Confirmed]"
        elif self.match_source == 'manual':
            source_indicator = "[Manual]"
        else:
            source_indicator = "[Auto]"
        return f"SlideMatchResult({source_indicator} {self.match_type}, {self.confidence}%, {self.gs_product_name} → {self.pptx_product_name})"


class SlideMatcher:
    """
    Intelligent slide matcher with exact and fuzzy matching capabilities.

    Usage:
        matcher = SlideMatcher(pptx_product_names)
        result = matcher.find_match("Upcycled Laptop Sleeve (Enfold)")
        if result.is_usable():
            print(f"Found match: {result.pptx_product_name}")
    """

    def __init__(self, pptx_product_names: List[str]):
        """
        Initialize matcher with list of product names from PowerPoint.

        Args:
            pptx_product_names: List of product names extracted from PowerPoint slides
        """
        self.pptx_product_names = pptx_product_names

        # Create uppercase mapping for exact matching
        self.pptx_upper_map = {name.upper(): name for name in pptx_product_names}

    def find_match(self, gs_product_name: str, num_alternatives: int = 3, dataset: str = 'demo', template_name: str = '') -> SlideMatchResult:
        """
        Find best matching PowerPoint slide for a Google Sheets product name.

        Args:
            gs_product_name: Product name from Google Sheets
            num_alternatives: Number of alternative matches to return
            dataset: Current dataset ('demo' or 'real') for confirmed match lookup
            template_name: Template name for confirmed match lookup

        Returns:
            SlideMatchResult with match information

        Matching Logic (IMPROVED with Confirmed Match Memory):
            0. Check confirmed matches from Google Sheets FIRST (100% confidence, highest priority)
            1. Check manual matches from JSON (100% confidence, takes precedence)
            2. Normalize product name (strip variants)
            3. Check manual product mappings
            4. Try exact match on normalized name
            5. Multi-scorer fuzzy matching (3 algorithms)
            6. Apply keyword category boosting (+15%)
            7. Return best match with confidence score
        """
        # Step 0: Check confirmed matches from Google Sheets FIRST (highest priority)
        try:
            from .match_memory import get_confirmed_match
            confirmed_match = get_confirmed_match(gs_product_name, dataset, template_name=template_name)

            if confirmed_match:
                # Confirmed match found - return with 100% confidence
                return SlideMatchResult(
                    gs_product_name=gs_product_name,
                    pptx_product_name=confirmed_match['slide_title'],
                    match_type='exact',
                    confidence=100,
                    alternatives=[],
                    match_source='confirmed',
                    match_method=f"Previously confirmed ({confirmed_match['match_type']})"
                )
        except Exception as e:
            # Log error if lookup fails
            print(f"Error looking up confirmed match for '{gs_product_name}': {str(e)}")

        # Step 1: Check manual matches from JSON (legacy system)
        manual_match_data = get_manual_match(gs_product_name, match_category="product")
        if manual_match_data:
            # Manual match found - return with 100% confidence
            return SlideMatchResult(
                gs_product_name=gs_product_name,
                pptx_product_name=manual_match_data['slide_title'],
                match_type='exact',
                confidence=100,
                alternatives=[],
                match_source='manual',
                match_method='Manual override'
            )

        # Step 2: Normalize product name
        normalized_gs_name = normalize_product_name(gs_product_name)

        # Step 3: Check manual product mappings (legacy hardcoded mappings)
        if normalized_gs_name in MANUAL_PRODUCT_MAPPINGS:
            manual_match = MANUAL_PRODUCT_MAPPINGS[normalized_gs_name]
            return SlideMatchResult(
                gs_product_name=gs_product_name,
                pptx_product_name=manual_match,
                match_type='exact',
                confidence=100,
                alternatives=[],
                match_source='automatic',
                match_method='Hardcoded mapping'
            )

        # Step 4: Try exact match on normalized name (case-insensitive)
        if normalized_gs_name in self.pptx_upper_map:
            exact_match = self.pptx_upper_map[normalized_gs_name]
            return SlideMatchResult(
                gs_product_name=gs_product_name,
                pptx_product_name=exact_match,
                match_type='exact',
                confidence=100,
                alternatives=[],
                match_source='automatic',
                match_method='Exact match'
            )

        # Step 5: Multi-scorer fuzzy matching
        # Use normalized name for better matching
        best_match_name, best_match_score, method_used, alternatives = find_best_match_multi_scorer(
            normalized_gs_name,
            self.pptx_product_names,
            limit=num_alternatives
        )

        if best_match_name is None:
            return SlideMatchResult(
                gs_product_name=gs_product_name,
                pptx_product_name=None,
                match_type='none',
                confidence=0,
                alternatives=[],
                match_source='automatic',
                match_method='No match found'
            )

        # Step 6: Apply keyword category boosting
        boosted_score = boost_score_if_same_category(
            normalized_gs_name,
            best_match_name,
            best_match_score
        )

        # Format method description
        method_description = f"{method_used.replace('_', ' ').title()} ({boosted_score}%)"
        if boosted_score > best_match_score:
            method_description += f" [boosted from {best_match_score}%]"

        return SlideMatchResult(
            gs_product_name=gs_product_name,
            pptx_product_name=best_match_name,
            match_type='fuzzy',
            confidence=boosted_score,
            alternatives=alternatives,
            match_source='automatic',
            match_method=method_description
        )

    def batch_match(self, gs_product_names: List[str], dataset: str = 'demo', template_name: str = '') -> List[SlideMatchResult]:
        """
        Match multiple products at once.

        Args:
            gs_product_names: List of product names from Google Sheets
            dataset: Current dataset ('demo' or 'real') for confirmed match lookup
            template_name: Template name for confirmed match lookup

        Returns:
            List of SlideMatchResult objects
        """
        return [self.find_match(name, dataset=dataset, template_name=template_name) for name in gs_product_names]

    def get_match_summary(self, results: List[SlideMatchResult], min_confidence: int = 70) -> Dict:
        """
        Get summary statistics for a batch of match results.

        Args:
            results: List of SlideMatchResult objects
            min_confidence: Minimum confidence threshold for counting usable matches

        Returns:
            Dictionary with summary statistics
        """
        exact_matches = sum(1 for r in results if r.match_type == 'exact')
        fuzzy_matches = sum(1 for r in results if r.match_type == 'fuzzy' and r.confidence >= min_confidence)
        poor_matches = sum(1 for r in results if r.match_type == 'fuzzy' and r.confidence < min_confidence)
        no_matches = sum(1 for r in results if r.match_type == 'none')

        return {
            'total': len(results),
            'exact': exact_matches,
            'fuzzy': fuzzy_matches,
            'poor': poor_matches,
            'none': no_matches,
            'usable': exact_matches + fuzzy_matches,
            'usable_pct': (exact_matches + fuzzy_matches) / len(results) * 100 if results else 0
        }


# Confidence threshold constants
CONFIDENCE_EXCELLENT = 90  # Green: Auto-use with high confidence
CONFIDENCE_GOOD = 70       # Yellow: Suggest, ask for confirmation
CONFIDENCE_POOR = 50       # Red: Show but warn user


def format_match_for_display(result: SlideMatchResult) -> str:
    """
    Format a match result for display to user.

    Returns:
        Human-readable string describing the match
    """
    if result.match_type == 'exact':
        return f"✓ Exact match: {result.pptx_product_name}"
    elif result.match_type == 'fuzzy':
        if result.confidence >= CONFIDENCE_EXCELLENT:
            icon = "✓"
            quality = "Excellent"
        elif result.confidence >= CONFIDENCE_GOOD:
            icon = "~"
            quality = "Good"
        else:
            icon = "?"
            quality = "Uncertain"

        return f"{icon} {quality} match ({result.confidence}%): {result.pptx_product_name}"
    else:
        return f"✗ No good match found (best guess: {result.pptx_product_name}, {result.confidence}%)"


def extract_unique_partners(proposal_products: List[Dict]) -> List[str]:
    """
    Extract unique partners from proposal products.

    Args:
        proposal_products: List of proposal products with product_data

    Returns:
        List of unique partner names (sorted alphabetically)

    Example:
        >>> products = [
        ...     {"product_data": {"Partner": "Partner X"}},
        ...     {"product_data": {"Partner": "Partner X"}},
        ...     {"product_data": {"Partner": "Partner Y"}}
        ... ]
        >>> extract_unique_partners(products)
        ['Partner X', 'Partner Y']
    """
    partners = set()
    for product in proposal_products:
        partner = product['product_data'].get('Partner', '')
        if partner:
            partners.add(partner)

    return sorted(list(partners))


def find_all_impact_slides(template_path: str) -> List[Dict]:
    """
    Find all slides in the template with text matching "* - Your Impact".

    Searches ALL text in slides (including text boxes), not just title placeholder,
    since impact slides may have titles in text boxes rather than title placeholders.

    Args:
        template_path: Path to PowerPoint template

    Returns:
        List of impact slide info:
        [
            {
                "slide_index": 156,
                "slide_title": "Hand Stitched Bags – Your Impact"
            },
            ...
        ]

    Example:
        >>> slides = find_all_impact_slides("template.pptx")
        >>> slides[0]['slide_title']
        'Hand Stitched Bags – Your Impact'
    """
    impact_slides = []

    try:
        prs = Presentation(template_path)

        for slide_idx, slide in enumerate(prs.slides):
            # Search ALL text in slide (not just title placeholder)
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    text = shape.text.strip()

                    # Check if text contains "Your Impact" (case-insensitive)
                    # Handles both "- Your Impact" and "— Your Impact" (em dash)
                    if "your impact" in text.lower():
                        impact_slides.append({
                            "slide_index": slide_idx,
                            "slide_title": text
                        })
                        break  # Only add each slide once

    except Exception as e:
        print(f"Error loading PowerPoint template: {e}")

    # Sort alphabetically by title
    impact_slides.sort(key=lambda x: x['slide_title'])

    return impact_slides


def resolve_impact_slides(
    partners: List[str],
    sheet_mapping: Dict[str, str],
    template_slides: List[Dict]
) -> Dict[str, Dict]:
    """
    Resolve impact slides by matching Google Sheet titles against current template slides.

    Args:
        partners: List of partner names in the proposal
        sheet_mapping: Dict from load_impact_slide_mapping(), e.g. {"GOEX": "Apparel -- Your Impact"}
        template_slides: List from find_all_impact_slides(), e.g. [{"slide_index": 50, "slide_title": "..."}]

    Returns:
        Dict mapping partner to resolved slide info:
        {
            "GOEX": {"slide_index": 50, "slide_title": "Apparel -- Your Impact"}
        }
        Partners with no match are omitted.
    """
    result = {}

    for partner in partners:
        expected_title = sheet_mapping.get(partner)
        if not expected_title:
            continue

        def normalize_dashes(text):
            return text.replace("\u2014", "-").replace("\u2013", "-").replace("\u2012", "-")

        normalized_expected = normalize_dashes(expected_title).lower().strip()

        # Try exact match first (after dash normalization)
        matched_slide = None
        for slide in template_slides:
            normalized_slide = normalize_dashes(slide['slide_title']).lower().strip()
            if normalized_slide == normalized_expected:
                matched_slide = slide
                break

        # Fall back to fuzzy match on title text (handles minor wording changes)
        if not matched_slide:
            best_score = 0
            for slide in template_slides:
                normalized_slide = normalize_dashes(slide['slide_title']).lower().strip()
                score = fuzz.ratio(normalized_expected, normalized_slide)
                if score > best_score:
                    best_score = score
                    matched_slide = slide if score >= 80 else None

        if matched_slide:
            result[partner] = {
                'slide_title': matched_slide['slide_title'],
                'slide_index': matched_slide['slide_index']
            }

    return result


