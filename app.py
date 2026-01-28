"""
Peace by Piece International - Order Management System
4-tab workflow: Proposals → Client Order Forms → Order & Client Info → Execution & Accounting
Version: 6.13 (Multi-Variant Product Consolidation in PowerPoint)
"""

# MEMORY OPTIMIZATION TOGGLE
# Set to False to disable memory optimization and use full caching (if issues arise)
# Disabled since upgrading to Render Standard tier (2GB RAM) - caching improves UX
USE_MEMORY_OPTIMIZATION = False

import streamlit as st
import gc  # For memory optimization
import streamlit.components.v1 as components
import gspread
from google.oauth2.service_account import Credentials
import pandas as pd
from datetime import datetime, date
import time
import copy  # For deep copying client info during backup/restore

# Import extracted modules
from src.data_loader import load_pricing_data, DATASET_CONFIGS
from src.helpers import (
    clean_price,
    apply_marketing_rounding,
    round_to_nearest_five,
    round_to_nearest_fifty_cents,
    calculate_moq,
    calculate_credit_card_fee,
    calculate_markup_from_price,
    extract_partner_contacts,
    validate_invoice_completeness,
    parse_tier_info,
    parse_tariff_rate,
    calculate_product_tariff,
    convert_proposal_to_order,
    parse_client_order_form_html,
    get_shipping_costs,
    format_shipping_display,
    get_column_value,
    get_tariff_rate
)
from src.pricing_engine import (
    determine_tier_number,
    get_unit_price_new_system,
    get_price_for_quantity,
    calculate_pbp_msrp,
    calculate_additional_costs,
    calculate_customization_costs,
    calculate_product_quote,
    calculate_order_total
)
from src.slide_matcher import SlideMatcher
from src.proposal_manager import (
    save_proposal,
    load_all_proposals,
    load_proposal_data,
    delete_proposal
)
from src.order_manager import (
    save_order,
    load_all_orders,
    load_order_data,
    delete_order
)
from src.template_loader import (
    get_template_path,
    get_template_name,
    TEMPLATE_CONFIG
)

# ============================================================
# HELPER FUNCTIONS
# ============================================================
def format_date_display(date_obj):
    """Format a date object as MM/DD/YY for display"""
    if date_obj:
        if isinstance(date_obj, str):
            return date_obj  # Return as-is if already a string
        try:
            return date_obj.strftime("%m/%d/%y")
        except:
            return str(date_obj)
    return "Not specified"

def calculate_msrp_markup(product_data):
    """
    Calculate markup percentage required to match MSRP price.
    Returns the markup if MSRP is available and valid, otherwise returns 100.0 (default).

    Args:
        product_data (dict): Product data dictionary containing MSRP and pricing info

    Returns:
        float: Markup percentage (0 if MSRP is below cost, calculated value if valid, 100.0 if no MSRP)
    """
    # Get MSRP
    msrp_raw = get_column_value(product_data, 'Vendor Published MSRP', 'MSRP', '')
    msrp = clean_price(msrp_raw)

    if msrp and msrp > 0:
        # Get base cost at quantity 100 as reference
        base_cost, _, _ = get_unit_price_new_system(product_data, 100)

        if base_cost and base_cost > 0:
            # Calculate required markup % to reach MSRP
            # Formula: MSRP = cost * (1 + markup/100)
            # Therefore: markup = ((MSRP / cost) - 1) * 100
            required_markup = ((msrp / base_cost) - 1) * 100

            # Don't allow negative markup (selling below cost)
            return max(0.0, required_markup)

    # No valid MSRP or cost, return default 100% markup
    return 100.0


def get_default_markup(product_data):
    """
    Get default markup for a product.
    Uses PBP Standard Markup from spreadsheet if available, otherwise 100%.

    Args:
        product_data: Product row data

    Returns:
        float: Markup percentage (e.g., 100.0 for 100% markup)
    """
    # Check for PBP Standard Markup column
    pbp_markup = get_column_value(product_data, 'PBP Standard Markup', None, None)

    if pbp_markup:
        try:
            # Convert from multiplier to percentage
            # Format: 2.0 = 200% total = 100% markup
            multiplier = float(pbp_markup)
            if multiplier > 0:
                return (multiplier - 1) * 100
        except (ValueError, TypeError):
            pass

    # Default to 100% markup
    return 100.0


def compute_proposal_hash(proposal_products):
    """
    Compute a hash of the proposal products for change detection.
    Includes product IDs, quantities, markups, and key settings.
    """
    import hashlib
    import json

    # Extract key fields that matter for saving
    proposal_data = []
    for item in proposal_products:
        key_fields = {
            'product': item.get('Product/Service', ''),
            'partner': item.get('Partner', ''),
            'quantity': item.get('quantity', 1),
            'markup': item.get('markup_percentage', 100.0)
        }
        proposal_data.append(key_fields)

    # Convert to JSON string and hash
    data_str = json.dumps(proposal_data, sort_keys=True)
    return hashlib.md5(data_str.encode()).hexdigest()

def compute_order_hash(order_items, client_info):
    """
    Compute a hash of the order items and client info for change detection.
    Includes all order data and client information fields.
    """
    import hashlib
    import json

    # Extract key fields from order items
    order_data = []
    for item in order_items:
        key_fields = {
            'product': item.get('Product/Service', ''),
            'partner': item.get('Partner', ''),
            'quantity': item.get('quantity', 1),
            'markup': item.get('markup_percentage', 100.0),
            'customization': item.get('has_customization', False)
        }
        order_data.append(key_fields)

    # Include client info in hash
    combined_data = {
        'order_items': order_data,
        'client_info': {
            'company': client_info.get('company_name', ''),
            'contact': client_info.get('contact_name', ''),
            'email': client_info.get('contact_email', ''),
            'client_type': client_info.get('client_type', '')
        }
    }

    # Convert to JSON string and hash
    data_str = json.dumps(combined_data, sort_keys=True, default=str)
    return hashlib.md5(data_str.encode()).hexdigest()

def has_unsaved_proposal_changes():
    """
    Check if there are unsaved changes in the current proposal.
    Returns True if products or settings have changed since last save.
    """
    if 'proposal_products' not in st.session_state or not st.session_state.proposal_products:
        return False

    current_hash = compute_proposal_hash(st.session_state.proposal_products)
    last_saved_hash = st.session_state.get('last_saved_proposal_hash', None)

    return last_saved_hash is None or current_hash != last_saved_hash

def has_unsaved_order_changes():
    """
    Check if there are unsaved changes in the current order.
    Returns True if order or client data changed since last save.
    """
    if 'order_items' not in st.session_state or not st.session_state.order_items:
        return False

    current_hash = compute_order_hash(
        st.session_state.order_items,
        st.session_state.client_info
    )
    last_saved_hash = st.session_state.get('last_saved_order_hash', None)

    return last_saved_hash is None or current_hash != last_saved_hash

def format_time_since_save(save_type):
    """
    Format the time since last save in a user-friendly way.
    save_type: 'proposal' or 'order'
    Returns string like "Saved 2 minutes ago" or "Last saved at 3:45 PM"
    """
    from datetime import datetime, timedelta

    time_key = f'last_{save_type}_save_time'
    last_save = st.session_state.get(time_key, None)

    if last_save is None:
        return None

    now = datetime.now()
    delta = now - last_save

    # If saved in the last minute
    if delta.total_seconds() < 60:
        return "Just saved"
    # If saved in the last hour
    elif delta.total_seconds() < 3600:
        minutes = int(delta.total_seconds() / 60)
        return f"Saved {minutes} minute{'s' if minutes != 1 else ''} ago"
    # If saved today
    elif delta.days == 0:
        return f"Last saved at {last_save.strftime('%-I:%M %p')}"
    # If saved yesterday
    elif delta.days == 1:
        return f"Last saved yesterday at {last_save.strftime('%-I:%M %p')}"
    else:
        return f"Last saved {delta.days} days ago"

def update_last_save_time(save_type):
    """
    Update the last save time for a given save type.
    save_type: 'proposal' or 'order'
    """
    from datetime import datetime
    st.session_state[f'last_{save_type}_save_time'] = datetime.now()

    # Also update the hash to match current state
    if save_type == 'proposal':
        st.session_state['last_saved_proposal_hash'] = compute_proposal_hash(
            st.session_state.proposal_products
        )
    elif save_type == 'order':
        st.session_state['last_saved_order_hash'] = compute_order_hash(
            st.session_state.order_items,
            st.session_state.client_info
        )

@st.cache_data(ttl=60)  # Cache for 60 seconds
def cached_load_all_proposals(refresh_counter=0):
    """
    Cached version of load_all_proposals to reduce API calls.
    Cache expires after 60 seconds to allow for updates.
    refresh_counter: Used to force cache refresh when needed.
    """
    return load_all_proposals()

@st.cache_data(ttl=60)  # Cache for 60 seconds
def cached_load_all_orders(refresh_counter=0):
    """
    Cached version of load_all_orders to reduce API calls.
    Cache expires after 60 seconds to allow for updates.
    refresh_counter: Used to force cache refresh when needed.
    """
    return load_all_orders()

def clear_saved_data_cache():
    """
    Clear the cache for saved proposals and orders.
    Call this after saving, deleting, or modifying saved data.
    """
    # Increment counter to force cache refresh
    if 'saved_data_refresh_counter' not in st.session_state:
        st.session_state.saved_data_refresh_counter = 0
    st.session_state.saved_data_refresh_counter += 1

    # Also clear the cache directly
    cached_load_all_proposals.clear()
    cached_load_all_orders.clear()

# ============================================================
# PAGE CONFIGURATION
# ============================================================
st.set_page_config(
    page_title="PBP Order Management",
    layout="wide",
    initial_sidebar_state="auto"
)

# ============================================================
# KEEP-ALIVE PING ROUTE (for GitHub Actions)
# ============================================================
# Lightweight endpoint to keep Streamlit Cloud app awake
# Responds instantly without loading data or UI
# Handle both old and new Streamlit API for query params
try:
    # Try new API (Streamlit >= 1.30)
    query_params = st.query_params
except AttributeError:
    # Fall back to old API (Streamlit < 1.30)
    query_params = st.experimental_get_query_params()

if "ping" in query_params:
    st.write("pong")
    st.stop()

# Prevent automatic page scrolling on widget interaction
st.markdown("""
<style>
    * {
       overflow-anchor: none !important;
    }
</style>
""", unsafe_allow_html=True)

# Restore scroll position after rerun (from sessionStorage)
components.html("""
    <script>
        // Check if we have a saved scroll position
        const savedScrollPos = window.parent.sessionStorage.getItem('streamlit_scroll_position');
        if (savedScrollPos !== null) {
            // Small delay to ensure DOM is ready
            setTimeout(() => {
                window.parent.document.querySelector('section.main').scrollTop = parseInt(savedScrollPos);
                // Clear the saved position after restoring
                window.parent.sessionStorage.removeItem('streamlit_scroll_position');
            }, 100);
        }
    </script>
""", height=0)

# ============================================================
# SESSION STATE INITIALIZATION
# ============================================================
# Initialize order_items if not exists
if 'order_items' not in st.session_state:
    st.session_state.order_items = []

# Ensure all order items have edited_description field (backward compatibility)
for item in st.session_state.order_items:
    if 'edited_description' not in item:
        item['edited_description'] = ''

# Initialize edit_index (None = adding new item, number = editing existing item)
if 'edit_index' not in st.session_state:
    st.session_state.edit_index = None

# Initialize order history
if 'order_history' not in st.session_state:
    st.session_state.order_history = []

# Initialize shipping in session state
if 'order_shipping' not in st.session_state:
    st.session_state.order_shipping = 0.0
if 'partner_shipping' not in st.session_state:
    st.session_state.partner_shipping = 0.0
if 'sales_tax' not in st.session_state:
    st.session_state.sales_tax = 0.0
if 'kitting_pbp_cost' not in st.session_state:
    st.session_state.kitting_pbp_cost = 0.0
if 'kitting_client_price' not in st.session_state:
    st.session_state.kitting_client_price = 0.0

# Initialize discount settings in session state
if 'order_discount_type' not in st.session_state:
    st.session_state.order_discount_type = "none"
if 'order_discount_preset' not in st.session_state:
    st.session_state.order_discount_preset = "Non-profit Discount (5%)"
if 'order_discount_custom_desc' not in st.session_state:
    st.session_state.order_discount_custom_desc = ""
if 'order_discount_custom_value' not in st.session_state:
    st.session_state.order_discount_custom_value = 0.0

# Initialize rounding settings
if 'order_use_marketing_rounding' not in st.session_state:
    st.session_state.order_use_marketing_rounding = False

if 'order_fifty_cent_rounding' not in st.session_state:
    st.session_state.order_fifty_cent_rounding = True  # Default to checked

# Initialize order confirmed flag
if 'order_confirmed' not in st.session_state:
    st.session_state.order_confirmed = False

# Initialize order backup for persistence during edit/confirm cycles
# This preserves data entered in Tab 4 when user clicks "Edit Order"
if 'order_backup' not in st.session_state:
    st.session_state.order_backup = {}

# Initialize client information
if 'client_info' not in st.session_state:
    st.session_state.client_info = {
        'is_new_client': True,
        'company_name': '',
        'contacts': [  # NEW: Support multiple contacts
            {
                'name': '',
                'email': '',
                'phone': '',
                'role': 'Primary Contact'
            }
        ],
        'client_po': '',
        'billing_address': '',
        'shipping_type': 'Ground',  # MODIFIED: Changed to dropdown default
        'shipping_address': '',
        'payment_timeline': 'Net 30',  # MODIFIED: Changed to dropdown default
        'payment_preference': 'Check',  # MODIFIED: Changed to dropdown default
        'client_in_hands_date': None,  # NEW: Target delivery date for client
        'order_submitted_by': '',  # NEW: Person submitting order
        'order_submitted_date': datetime.now().date(),  # NEW: Auto-filled submission date
        'cost_submitted_by': '',  # NEW: Person submitting costs
        'cost_submitted_date': None  # NEW: Date costs were submitted
    }

# Migrate old contact fields to new structure if needed
if 'contact_name' in st.session_state.client_info:
    # Old structure detected - migrate to new
    old_name = st.session_state.client_info.get('contact_name', '')
    old_email = st.session_state.client_info.get('contact_email', '')

    # Create contacts array with old data
    st.session_state.client_info['contacts'] = [
        {
            'name': old_name,
            'email': old_email,
            'phone': '',
            'role': 'Primary Contact'
        }
    ]

    # Remove old fields
    if 'contact_name' in st.session_state.client_info:
        del st.session_state.client_info['contact_name']
    if 'contact_email' in st.session_state.client_info:
        del st.session_state.client_info['contact_email']

# Ensure contacts array exists and has at least one contact
if 'contacts' not in st.session_state.client_info:
    st.session_state.client_info['contacts'] = [
        {
            'name': '',
            'email': '',
            'phone': '',
            'role': 'Primary Contact'
        }
    ]
elif len(st.session_state.client_info['contacts']) == 0:
    # Ensure at least one contact exists
    st.session_state.client_info['contacts'].append({
        'name': '',
        'email': '',
        'phone': '',
        'role': 'Primary Contact'
    })

# Initialize custom payment terms
if 'custom_payment_terms' not in st.session_state:
    st.session_state.custom_payment_terms = ''

# Initialize order notes (4 categories organized by audience)
if 'order_notes' not in st.session_state:
    st.session_state.order_notes = {
        'internal_pbp_team': '',      # Internal Notes (For PBP Team)
        'internal_bookkeeping': '',   # Internal Notes (For Bookkeeping)
        'external_partners': '',      # External Notes (For Partners/POs)
        'external_clients': ''        # External Notes (For Clients/Invoices)
    }

# Initialize credit card fee settings
if 'apply_cc_fee' not in st.session_state:
    st.session_state.apply_cc_fee = False
if 'cc_fee_percent' not in st.session_state:
    st.session_state.cc_fee_percent = 3.0

# Initialize proposal-specific session state (Phase 2)
if 'proposal_products' not in st.session_state:
    st.session_state.proposal_products = []

if 'proposal_marketing_rounding' not in st.session_state:
    st.session_state.proposal_marketing_rounding = False

if 'proposal_fifty_cent_rounding' not in st.session_state:
    st.session_state.proposal_fifty_cent_rounding = True  # Default to checked

if 'proposal_use_msrp' not in st.session_state:
    st.session_state.proposal_use_msrp = True  # Default to checked (use MSRP)

# Initialize order MSRP preference
if 'order_use_msrp' not in st.session_state:
    st.session_state.order_use_msrp = True  # Default to checked (use MSRP)

if 'configuring_product' not in st.session_state:
    st.session_state.configuring_product = None

if 'editing_proposal_index' not in st.session_state:
    st.session_state.editing_proposal_index = None

if 'proposal_filters' not in st.session_state:
    st.session_state.proposal_filters = {
        'min_price': 0.0,
        'client_budget': None,
        'partners': [],
        'countries': []
    }

if 'proposal_kitting_pricing' not in st.session_state:
    st.session_state.proposal_kitting_pricing = """• Impact Cards (about maker communities): No Charge
• Custom Message Cards: $65 set up/$1.50 per card
• Insertion of a card you provide: No Charge
• Kitting for domestic shipments: $10.50/box
• Kitting for international shipments: $19.00/box
• Label making & shipping coordination: $2/box

*kitting includes shipping logistics but does not include actual shipping charges
*international kitting includes customs documentation, etc."""

if 'proposal_terms' not in st.session_state:
    # Load from config file if exists
    try:
        with open('config/terms_conditions.txt', 'r') as f:
            st.session_state.proposal_terms = f.read()
    except:
        st.session_state.proposal_terms = """• Drop shipping and kitting are available
• For shipments to California zip codes, sales tax will be added
• Payment terms
• 50% to initiate a custom order
• 50% upon shipment
• Customs charges for gifts shipped internationally will be billed to client, which may take up to 120 days post shipment
• Due to the current uncertainty around US-imposed tariffs, Peace by Piece will bill client for any charges, which may take up to 120 days post shipment
• No cancellations will be accepted after any custom order has been initiated
• Gifts returned to Peace by Piece due to incorrect recipient information will incur a $20 fee plus any additional returned shipping charges from the carrier. If a new address is not provided within 30 days of the gift's return, the gift will be shipped back to the client for distribution"""

# Initialize form customization fields
if 'form_customizations' not in st.session_state:
    st.session_state.form_customizations = {
        'form_instructions': """1. Copy & paste the entire form into Docs, Word, or directly into your email reply (the format should copy along with the text)
2. Click in the gray areas to type your answers
3. For multiple choice questions, delete the options you DON'T want and keep the one you DO want
4. When finished, select all (Ctrl+A or Cmd+A), copy, and paste into your email reply
5. Fields marked with * are required""",
        'dropshipping_instructions': "For drop shipping orders, please provide: destination addresses, quantities per location, and any special delivery instructions",
        'dropshipping_placeholder': "[Input dropshipping info here]",
        'shipping_address_placeholder': "[Type full shipping address here, or N/A if drop shipping]",
        'billing_address_placeholder': '[Type billing address here, or "Same as shipping"]',
        'customization_placeholder': "[Describe any customization, logo placement, colors, etc.]",
        'impact_card_options': """Peace by Piece Impact Card
Custom Impact Card
Custom Message Card
Send us your own card""",
        'payment_options': """ACH
Check
Credit Card (3% processing fee applies)"""
    }

# Legacy compatibility - maintain sync between old and new dropshipping fields
if 'dropshipping_notes' not in st.session_state:
    st.session_state.dropshipping_notes = st.session_state.form_customizations['dropshipping_instructions']
else:
    # If dropshipping_notes exists but differs from form_customizations, sync them
    if st.session_state.dropshipping_notes != st.session_state.form_customizations['dropshipping_instructions']:
        st.session_state.dropshipping_notes = st.session_state.form_customizations['dropshipping_instructions']

if 'using_proposal_data' not in st.session_state:
    st.session_state.using_proposal_data = False

# ============================================================
# HEADER
# ============================================================
st.title("Peace by Piece Order Management System")
st.markdown("**Welcome to the PBP Order Management System** — Manage the complete order lifecycle:")
st.markdown("")
st.markdown("**→ Tab 1: Proposal Generator** - Browse products, create proposals for prospective clients")
st.markdown("**→ Tab 2: Client Order Form Generator** - Create professional order forms to send to clients")
st.markdown("**→ Tab 3: Order & Client Info** - Build orders, collect client details")
st.markdown("**→ Tab 4: Execution & Accounting** - Generate invoices and purchase orders")
st.divider()

# Initialize save tracking
if 'last_saved_proposal_hash' not in st.session_state:
    st.session_state.last_saved_proposal_hash = None
if 'last_saved_order_hash' not in st.session_state:
    st.session_state.last_saved_order_hash = None
if 'last_proposal_save_time' not in st.session_state:
    st.session_state.last_proposal_save_time = None
if 'last_order_save_time' not in st.session_state:
    st.session_state.last_order_save_time = None

# Initialize dataset selection (always use 'real')
if 'selected_dataset' not in st.session_state:
    st.session_state.selected_dataset = 'real'

# ============================================================
# DATA LOADING (before sidebar to ensure data is available)
# ============================================================
# Load data if not already loaded or if dataset changed
need_reload = (
    'df_template' not in st.session_state or
    st.session_state.get('loaded_dataset') != st.session_state.selected_dataset
)

if need_reload:
    df_template, df_metadata, df_partner_info = load_pricing_data(st.session_state.selected_dataset)
    st.session_state.df_template = df_template
    st.session_state.df_metadata = df_metadata
    st.session_state.df_partner_info = df_partner_info
    st.session_state.data_loaded_at = datetime.now()
    st.session_state.loaded_dataset = st.session_state.selected_dataset

    # Extract and store partner contacts
    st.session_state.partner_contacts = extract_partner_contacts(df_partner_info)

# ============================================================
# SIDEBAR - APP INFORMATION
# ============================================================
with st.sidebar:
    st.markdown("## Instructions & Tools")

    # Section 0: Data Source Information
    st.markdown("### Data Source")

    # Display consolidated data source information
    # Check if data is loaded (from current or previous run)
    if 'df_template' in st.session_state and st.session_state.df_template is not None:
        # Calculate from current data in session state
        unique_products = len(st.session_state.df_template)
        unique_partners = len(st.session_state.df_template['Partner'].unique())
        active_dataset_name = DATASET_CONFIGS[st.session_state.selected_dataset]['name']

        st.info(f"**{active_dataset_name}**")
        st.caption(f"{unique_products} products from {unique_partners} partners")
    else:
        st.info("Loading data...")

    # Session Status
    proposal_count = len(st.session_state.get('proposal_items', []))
    order_count = len(st.session_state.get('order_items', []))

    st.caption(f"**Proposal:** {proposal_count} items | **Order:** {order_count} items")

    st.markdown("---")

    # Saved Work Management
    st.markdown("### Saved Work")

    # Add refresh button if data seems stale
    if st.button("Refresh Saved Items", key="refresh_saved_work", help="Click to refresh saved proposals and orders", use_container_width=True):
        clear_saved_data_cache()
        st.rerun()

    st.caption("Saved items refresh automatically every 60 seconds")

    # Load saved proposals and orders (cached to reduce API calls)
    refresh_counter = st.session_state.get('saved_data_refresh_counter', 0)
    saved_proposals = cached_load_all_proposals(refresh_counter)
    saved_orders = cached_load_all_orders(refresh_counter)

    # Handle potential None returns from rate limiting
    if saved_proposals is None:
        saved_proposals = []
    if saved_orders is None:
        saved_orders = []

    # Check for unsaved proposal changes
    proposal_unsaved = has_unsaved_proposal_changes()
    proposal_save_status = format_time_since_save('proposal')

    # Build proposal header with indicators
    proposal_header = f"**Saved Proposals ({len(saved_proposals)})**"
    if proposal_unsaved:
        proposal_header += " - Unsaved changes"

    # Saved Proposals subsection
    with st.expander(proposal_header, expanded=False):
        # Show save status if available
        if proposal_save_status:
            st.caption(f"{proposal_save_status}")

        if len(saved_proposals) == 0:
            st.info("No saved proposals yet")
        else:
            # Add manage mode checkbox
            manage_proposals = st.checkbox("Manage proposals (delete)", key="manage_proposals_mode")

            for proposal in saved_proposals[:10]:  # Show max 10 most recent
                if manage_proposals:
                    # Delete mode - show delete buttons
                    col1, col2, col3 = st.columns([3, 1, 1])
                    with col1:
                        st.caption(f"{proposal['name']}")
                        st.caption(f"   {proposal['created_date'][:10]}")
                    with col2:
                        if st.button("Load", key=f"load_prop_{proposal['proposal_id']}", use_container_width=True):
                            success, proposal_data, dataset = load_proposal_data(proposal['proposal_id'])
                            if success:
                                if dataset != st.session_state.selected_dataset:
                                    st.warning(f"Dataset mismatch: {dataset} → {st.session_state.selected_dataset}")

                                st.session_state.proposal_products = proposal_data.get('proposal_products', [])
                                st.session_state.proposal_marketing_rounding = proposal_data.get('proposal_marketing_rounding', False)
                                st.session_state.proposal_use_msrp = proposal_data.get('proposal_use_msrp', True)
                                st.session_state.proposal_discount_type = proposal_data.get('proposal_discount_type', None)
                                st.session_state.proposal_discount_percent = proposal_data.get('proposal_discount_percent', 0.0)
                                st.session_state.proposal_client_budget = proposal_data.get('proposal_client_budget', 0.0)
                                
                                # Track loaded proposal info
                                st.session_state.loaded_proposal_name = proposal['name']
                                st.session_state.loaded_proposal_date = proposal.get('created_date', 'Unknown')
                                st.session_state.loaded_proposal_creator = proposal.get('created_by', 'Unknown')

                                st.success(f"Loaded: {proposal['name']}")
                                st.rerun()
                    with col3:
                        if st.button("Delete", key=f"del_prop_{proposal['proposal_id']}", use_container_width=True, help="Delete this proposal"):
                            # Store the proposal to delete in session state for confirmation
                            st.session_state.delete_proposal_confirm = proposal
                else:
                    # Normal mode - just load button
                    col1, col2 = st.columns([3, 1])
                    with col1:
                        st.caption(f"{proposal['name']}")
                        st.caption(f"   {proposal['created_date'][:10]}")
                    with col2:
                        if st.button("Load", key=f"load_prop_{proposal['proposal_id']}", use_container_width=True):
                            success, proposal_data, dataset = load_proposal_data(proposal['proposal_id'])
                            if success:
                                if dataset != st.session_state.selected_dataset:
                                    st.warning(f"Dataset mismatch: {dataset} → {st.session_state.selected_dataset}")

                                st.session_state.proposal_products = proposal_data.get('proposal_products', [])
                                st.session_state.proposal_marketing_rounding = proposal_data.get('proposal_marketing_rounding', False)
                                st.session_state.proposal_use_msrp = proposal_data.get('proposal_use_msrp', True)
                                st.session_state.proposal_discount_type = proposal_data.get('proposal_discount_type', None)
                                st.session_state.proposal_discount_percent = proposal_data.get('proposal_discount_percent', 0.0)
                                st.session_state.proposal_client_budget = proposal_data.get('proposal_client_budget', 0.0)
                                
                                # Track loaded proposal info
                                st.session_state.loaded_proposal_name = proposal['name']
                                st.session_state.loaded_proposal_date = proposal.get('created_date', 'Unknown')
                                st.session_state.loaded_proposal_creator = proposal.get('created_by', 'Unknown')

                                st.success(f"Loaded: {proposal['name']}")
                                st.rerun()

            # Handle delete confirmation
            if 'delete_proposal_confirm' in st.session_state:
                prop = st.session_state.delete_proposal_confirm
                st.warning(f"Delete '{prop['name']}'?")
                col1, col2 = st.columns(2)
                with col1:
                    if st.button("Yes, Delete", key="confirm_del_prop", type="primary"):
                        success, message = delete_proposal(prop['proposal_id'])
                        if success:
                            clear_saved_data_cache()  # Clear cache after deletion
                            st.success("Deleted successfully")
                            del st.session_state.delete_proposal_confirm
                            time.sleep(0.5)
                            st.rerun()
                        else:
                            st.error(message)
                with col2:
                    if st.button("Cancel", key="cancel_del_prop"):
                        del st.session_state.delete_proposal_confirm
                        st.rerun()

            if len(saved_proposals) > 10:
                st.caption(f"...and {len(saved_proposals) - 10} more.")

    # Check for unsaved order changes
    order_unsaved = has_unsaved_order_changes()
    order_save_status = format_time_since_save('order')

    # Build order header with indicators
    order_header = f"**Saved Orders ({len(saved_orders)})**"
    if order_unsaved:
        order_header += " - Unsaved changes"

    # Saved Orders subsection
    with st.expander(order_header, expanded=False):
        # Show save status if available
        if order_save_status:
            st.caption(f"{order_save_status}")

        if len(saved_orders) == 0:
            st.info("No saved orders yet")
        else:
            # Add manage mode checkbox
            manage_orders = st.checkbox("Manage orders (delete)", key="manage_orders_mode")

            for order in saved_orders[:10]:  # Show max 10 most recent
                if manage_orders:
                    # Delete mode - show delete buttons
                    col1, col2, col3 = st.columns([3, 1, 1])
                    with col1:
                        st.caption(f"{order['name']}")
                        st.caption(f"   {order['created_date'][:10]}")
                    with col2:
                        if st.button("Load", key=f"load_ord_{order['order_id']}", use_container_width=True):
                            success, order_data, dataset = load_order_data(order['order_id'])
                            if success:
                                if dataset != st.session_state.selected_dataset:
                                    st.warning(f"Dataset mismatch: {dataset} → {st.session_state.selected_dataset}")

                                st.session_state.order_items = order_data.get('order_items', [])

                                # Migration: Add kitting fields to old orders that don't have them
                                for item in st.session_state.order_items:
                                    if 'include_kitting' not in item:
                                        item['include_kitting'] = False
                                        item['kitting_pbp_cost'] = 0.0
                                        item['kitting_client_price'] = 0.0
                                        item['kitting_description'] = ''

                                st.session_state.order_shipping = order_data.get('order_shipping', 0.0)
                                st.session_state.partner_shipping = order_data.get('partner_shipping', 0.0)
                                st.session_state.sales_tax = order_data.get('sales_tax', 0.0)
                                st.session_state.kitting_pbp_cost = order_data.get('kitting_pbp_cost', 0.0)
                                st.session_state.kitting_client_price = order_data.get('kitting_client_price', 0.0)
                                st.session_state.order_discount_type = order_data.get('order_discount_type', 'none')
                                st.session_state.order_discount_preset = order_data.get('order_discount_preset', 'NGO Discount (5%)')
                                st.session_state.order_discount_custom_desc = order_data.get('order_discount_custom_desc', '')
                                st.session_state.order_discount_custom_value = order_data.get('order_discount_custom_value', 0.0)
                                st.session_state.order_use_marketing_rounding = order_data.get('order_use_marketing_rounding', False)
                                st.session_state.apply_cc_fee = order_data.get('apply_cc_fee', False)
                                st.session_state.cc_fee_percent = order_data.get('cc_fee_percent', 3.0)
                                st.session_state.client_info = order_data.get('client_info', st.session_state.client_info)
                                # Handle old order_notes structures - clean replacement
                                loaded_notes = order_data.get('order_notes', {})
                                if 'internal_pbp_team' in loaded_notes:
                                    # New 4-category structure - load normally
                                    st.session_state.order_notes = loaded_notes
                                else:
                                    # Old structure detected - discard and initialize fresh
                                    st.session_state.order_notes = {
                                        'internal_pbp_team': '',
                                        'internal_bookkeeping': '',
                                        'external_partners': '',
                                        'external_clients': ''
                                    }
                                st.session_state.order_confirmed = order_data.get('order_confirmed', False)

                                st.success(f"Loaded: {order['name']}")
                                st.rerun()
                    with col3:
                        if st.button("Delete", key=f"del_ord_{order['order_id']}", use_container_width=True, help="Delete this order"):
                            # Store the order to delete in session state for confirmation
                            st.session_state.delete_order_confirm = order
                else:
                    # Normal mode - just load button
                    col1, col2 = st.columns([3, 1])
                    with col1:
                        st.caption(f"{order['name']}")
                        st.caption(f"   {order['created_date'][:10]}")
                    with col2:
                        if st.button("Load", key=f"load_ord_{order['order_id']}", use_container_width=True):
                            success, order_data, dataset = load_order_data(order['order_id'])
                            if success:
                                if dataset != st.session_state.selected_dataset:
                                    st.warning(f"Dataset mismatch: {dataset} → {st.session_state.selected_dataset}")

                                st.session_state.order_items = order_data.get('order_items', [])

                                # Migration: Add kitting fields to old orders that don't have them
                                for item in st.session_state.order_items:
                                    if 'include_kitting' not in item:
                                        item['include_kitting'] = False
                                        item['kitting_pbp_cost'] = 0.0
                                        item['kitting_client_price'] = 0.0
                                        item['kitting_description'] = ''

                                st.session_state.order_shipping = order_data.get('order_shipping', 0.0)
                                st.session_state.partner_shipping = order_data.get('partner_shipping', 0.0)
                                st.session_state.sales_tax = order_data.get('sales_tax', 0.0)
                                st.session_state.kitting_pbp_cost = order_data.get('kitting_pbp_cost', 0.0)
                                st.session_state.kitting_client_price = order_data.get('kitting_client_price', 0.0)
                                st.session_state.order_discount_type = order_data.get('order_discount_type', 'none')
                                st.session_state.order_discount_preset = order_data.get('order_discount_preset', 'NGO Discount (5%)')
                                st.session_state.order_discount_custom_desc = order_data.get('order_discount_custom_desc', '')
                                st.session_state.order_discount_custom_value = order_data.get('order_discount_custom_value', 0.0)
                                st.session_state.order_use_marketing_rounding = order_data.get('order_use_marketing_rounding', False)
                                st.session_state.apply_cc_fee = order_data.get('apply_cc_fee', False)
                                st.session_state.cc_fee_percent = order_data.get('cc_fee_percent', 3.0)
                                st.session_state.client_info = order_data.get('client_info', st.session_state.client_info)
                                # Handle old order_notes structures - clean replacement
                                loaded_notes = order_data.get('order_notes', {})
                                if 'internal_pbp_team' in loaded_notes:
                                    # New 4-category structure - load normally
                                    st.session_state.order_notes = loaded_notes
                                else:
                                    # Old structure detected - discard and initialize fresh
                                    st.session_state.order_notes = {
                                        'internal_pbp_team': '',
                                        'internal_bookkeeping': '',
                                        'external_partners': '',
                                        'external_clients': ''
                                    }
                                st.session_state.order_confirmed = order_data.get('order_confirmed', False)

                                st.success(f"Loaded: {order['name']}")
                                st.rerun()

            # Handle delete confirmation
            if 'delete_order_confirm' in st.session_state:
                ord = st.session_state.delete_order_confirm
                st.warning(f"Delete '{ord['name']}'?")
                col1, col2 = st.columns(2)
                with col1:
                    if st.button("Yes, Delete", key="confirm_del_ord", type="primary"):
                        success, message = delete_order(ord['order_id'])
                        if success:
                            clear_saved_data_cache()  # Clear cache after deletion
                            st.success("Deleted successfully")
                            del st.session_state.delete_order_confirm
                            time.sleep(0.5)
                            st.rerun()
                        else:
                            st.error(message)
                with col2:
                    if st.button("Cancel", key="cancel_del_ord"):
                        del st.session_state.delete_order_confirm
                        st.rerun()

            if len(saved_orders) > 10:
                st.caption(f"...and {len(saved_orders) - 10} more.")

    st.markdown("---")

    # Section 3: Clear Current Session Button (renamed from Clear All Data)
    st.markdown("### Session Management")

    # Show current work status
    proposal_count = len(st.session_state.proposal_products)
    order_count = len(st.session_state.order_items)

    if proposal_count > 0 or order_count > 0:
        st.caption(f"Current work: {proposal_count} proposal items, {order_count} order items")

    if st.button("Reset Current Session", type="secondary", use_container_width=True):
        st.session_state.confirm_clear = True

    if st.session_state.get('confirm_clear', False):
        st.warning(f"""
        **Start a fresh session?**

        This will clear your current working session:
        - Current proposal ({proposal_count} items)
        - Current order ({order_count} items)
        - Current client info

        Your {len(saved_proposals)} saved proposals and {len(saved_orders)} saved orders will remain safe.
        """)
        col1, col2 = st.columns(2)
        with col1:
            if st.button("Yes, Reset", type="primary", use_container_width=True):
                # Clear all session state data
                st.session_state.proposal_products = []
                st.session_state.order_items = []
                # Clear proposal tracking info
                st.session_state.loaded_proposal_name = None
                st.session_state.loaded_proposal_date = None
                st.session_state.loaded_proposal_creator = None
                st.session_state.client_info = {
                    'is_new_client': True,
                    'company_name': '',
                    'contact_name': '',
                    'contact_email': '',
                    'client_po': '',
                    'billing_address': '',
                    'shipping_type': 'Ground',
                    'shipping_address': '',
                    'payment_timeline': 'Net 30',
                    'payment_preference': 'Check',
                    'client_in_hands_date': None,
                    'order_submitted_by': '',
                    'order_submitted_date': datetime.now().date(),
                    'cost_submitted_by': '',
                    'cost_submitted_date': None
                }
                st.session_state.order_notes = {
                    'internal_pbp_team': '',
                    'internal_bookkeeping': '',
                    'external_partners': '',
                    'external_clients': ''
                }
                st.session_state.order_shipping = 0.0
                st.session_state.partner_shipping = 0.0
                st.session_state.order_discount_type = "none"
                st.session_state.order_history = []
                st.session_state.confirm_clear = False

                # Show success message before rerun
                st.success(f"Session reset successfully! Your {len(saved_proposals)} saved proposals and {len(saved_orders)} saved orders are still available above.")
                time.sleep(1)  # Brief pause to show message
                st.rerun()
        with col2:
            if st.button("Cancel", use_container_width=True):
                st.session_state.confirm_clear = False
                st.rerun()

    st.markdown("---")

    # Section 3: Instructions
    with st.expander("How to Use This App", expanded=False):
        st.markdown("""
        **4-Tab Workflow:**

        **Tab 1: Proposal Generator** (for prospective clients)
        - Browse product catalog with filters
        - Add products to proposal
        - Configure quantity, markup, MSRP comparison
        - Generate MOQ-based proposal tables
        - Generate PowerPoint presentations

        **Tab 2: Client Order Form Generator**
        - Create professional HTML order forms
        - Pre-fill client information
        - Send to clients for completion

        **Tab 3: Order & Client Info** (main workflow)
        1. Import completed client forms or proposal products
        2. Enter client information (company, contact, payment)
        3. Select partner and product from dropdowns
        4. Set quantity, markup, and customization options
        5. Add to order (repeat for multiple products)
        6. Configure shipping, discounts, custom items
        7. Add order notes (partner notes, accounting notes)
        8. Review order summary

        **Tab 4: Execution & Accounting** (final step)
        - View order summary and validation warnings
        - Generate invoice and purchase order
        - Download CSV for bookkeeping
        - Export to accounting (coming soon)

        **Tips:**
        - Start with Tab 1 for proposals to prospective clients
        - Use Tab 2 to generate client order forms
        - Use Tab 3 for building and managing actual orders
        - Tab 4 requires completed order in Tab 3
        """)

    st.markdown("---")

    # Section 4: Recent Orders
    st.markdown("### Recent Orders")
    if len(st.session_state.order_history) == 0:
        st.caption("No recent orders this session")
    else:
        # Show last 5 orders, most recent first
        for idx, order in enumerate(reversed(st.session_state.order_history[-5:])):
            with st.container():
                timestamp_str = order['timestamp'].strftime('%I:%M %p')
                product_preview = ', '.join(order['product_names'][:2])
                if len(order['product_names']) > 2:
                    product_preview += f" +{len(order['product_names'])-2} more"

                st.caption(f"{timestamp_str} - {product_preview}")
                st.caption(f"${order['total_quote']:.2f} ({order['total_units']} units)")

                col1, col2 = st.columns(2)
                with col1:
                    if st.button("Load", key=f"load_order_{idx}", use_container_width=True):
                        # Reload this order
                        st.session_state.order_items = order['order_items'].copy()
                        st.session_state.order_shipping = order['shipping']
                        st.rerun()
                with col2:
                    if st.button("Delete", key=f"delete_order_{idx}", use_container_width=True):
                        # Remove from history
                        actual_idx = len(st.session_state.order_history) - 1 - idx
                        st.session_state.order_history.pop(actual_idx)
                        st.rerun()

                if idx < min(4, len(st.session_state.order_history) - 1):
                    st.markdown("---")

    st.markdown("---")

    # Section 5: Data Status
    st.markdown("### Data Status")
    if 'data_loaded_at' in st.session_state:
        load_time = st.session_state.data_loaded_at
        time_ago = datetime.now() - load_time

        if time_ago.seconds < 60:
            time_str = "Just now"
        elif time_ago.seconds < 3600:
            time_str = f"{time_ago.seconds // 60} min ago"
        else:
            time_str = load_time.strftime('%I:%M %p')

        st.caption(f"Last updated: {time_str}")

        if st.button("Refresh Data", use_container_width=True):
            # Clear cached data and reload from selected dataset
            df_template, df_metadata, df_partner_info = load_pricing_data(st.session_state.selected_dataset)
            st.session_state.df_template = df_template
            st.session_state.df_metadata = df_metadata
            st.session_state.df_partner_info = df_partner_info
            st.session_state.data_loaded_at = datetime.now()
            st.rerun()
    else:
        st.caption("Data status: Unknown")

    st.markdown("---")

    # Section 4: Download Options
    st.markdown("### Download Options")

    # Download current order as CSV
    if len(st.session_state.order_items) > 0:
        import io
        import csv

        # Build CSV content
        output = io.StringIO()
        writer = csv.writer(output)

        # Header
        writer.writerow(["Product", "Quantity", "Per Unit", "Total"])

        # Order items
        for item in st.session_state.order_items:
            writer.writerow([
                item['product_name'],
                item['quantity'],
                f"${item['total_per_unit']:.2f}",
                f"${item['product_total']:.2f}"
            ])

        # Add totals
        products_subtotal = sum(item['product_total'] for item in st.session_state.order_items)
        writer.writerow(["Shipping", "", "", f"${st.session_state.order_shipping:.2f}"])

        # Add per-product tariff lines
        for item in st.session_state.order_items:
            tariff_amount = item.get('tariff_amount', 0)
            if tariff_amount > 0:
                country = item.get('country_of_origin', 'Unknown')
                tariff_rate = item.get('tariff_rate_percent', 0)
                writer.writerow([f"Tariff: {item['product_name']} ({tariff_rate}% - {country})", "", "", f"${tariff_amount:.2f}"])

        # Calculate total tariff
        total_tariff = sum(item.get('tariff_amount', 0) for item in st.session_state.order_items)
        total_quote = products_subtotal + st.session_state.order_shipping + total_tariff
        writer.writerow(["TOTAL", "", "", f"${total_quote:.2f}"])

        csv_content = output.getvalue()

        st.download_button(
            label="Download Order (CSV)",
            data=csv_content,
            file_name=f"order_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
            mime="text/csv",
            use_container_width=True
        )
    else:
        st.caption("Add products to download order")

    # Download master pricing data
    if 'df_template' in st.session_state:
        csv_pricing = st.session_state.df_template.to_csv(index=False)

        st.download_button(
            label="Download Pricing Data (CSV)",
            data=csv_pricing,
            file_name=f"pricing_data_{datetime.now().strftime('%Y%m%d')}.csv",
            mime="text/csv",
            use_container_width=True
        )

# ============================================================
# MATCH REVIEW UI FUNCTION (PHASE 1, DAY 2)
# ============================================================
def show_match_review_ui(match_results, pptx_product_names, pptx_name_to_index=None):
    """
    Display match review UI for user confirmation of fuzzy matches.

    Args:
        match_results: List of SlideMatchResult objects from SlideMatcher
        pptx_product_names: List of all PowerPoint product names (for alternatives)
        pptx_name_to_index: Dict mapping slide names to indices (for saving confirmations)

    Returns:
        Dict of confirmed matches {gs_product_name: pptx_product_name} or None if cancelled
    """
    # Import match memory functions
    from src.match_memory import save_confirmed_match

    # Get current dataset for saving confirmations
    current_dataset = st.session_state.get('selected_dataset', 'real')
    st.markdown("---")
    st.subheader("Step 1. Review Product Matches")

    # Separate matches by type
    exact_matches = [r for r in match_results if r.match_type == 'exact']
    fuzzy_matches = [r for r in match_results if r.match_type == 'fuzzy' and r.confidence >= 70]
    poor_matches = [r for r in match_results if r.match_type == 'fuzzy' and r.confidence < 70]
    no_matches = [r for r in match_results if r.match_type == 'none']

    # Count items by status for better summary
    already_matched = len(exact_matches)  # Includes exact, previously confirmed, and manual overrides
    needs_review = len(fuzzy_matches) + len(poor_matches) + len(no_matches)

    total_products = already_matched + needs_review

    # Summary at top
    st.markdown(f"""
    **Match Summary:**
    - {already_matched} ready to use (already matched from exact or previous sessions)
    - {needs_review} need your confirmation (auto-matches and no matches)

    Total: {total_products} products
    """)

    if needs_review > 0:
        st.caption("Items are sorted by confidence (lowest first) to help you focus on matches that need the most attention.")

    if len(fuzzy_matches) == 0 and len(exact_matches) == 0:
        st.warning("No usable matches found. Cannot generate PowerPoint presentation.")
        return None

    # Initialize confirmation state if not exists
    if 'match_confirmations' not in st.session_state:
        st.session_state.match_confirmations = {}

    # ============================================================
    # UNIFIED MATCH TABLE - All Products
    # ============================================================
    st.markdown("### Product → Slide Matches")
    st.caption("Review and confirm matches. Items needing review are shown first.")

    # Combine all match types into single list
    all_matches = []

    # Add exact matches
    for result in exact_matches:
        all_matches.append({
            'result': result,
            'type': 'exact',
            'needs_review': False
        })

    # Add fuzzy matches
    for result in fuzzy_matches:
        all_matches.append({
            'result': result,
            'type': 'fuzzy',
            'needs_review': True
        })

    # Add poor/no matches
    if poor_matches or no_matches:
        all_poor = poor_matches + no_matches
        for result in all_poor:
            all_matches.append({
                'result': result,
                'type': 'poor',
                'needs_review': True
            })

    # Sort by review status (needs review first), then by confidence (lowest confidence first within needs_review)
    # This puts the most questionable matches at the top for user attention
    all_matches.sort(key=lambda x: (
        not x['needs_review'],  # Needs review first
        x['result'].confidence if x['needs_review'] else 999,  # Within needs_review: lowest confidence first (ascending)
        x['result'].gs_product_name  # Alphabetical as tiebreaker
    ))

    if all_matches:
        # Add JavaScript to capture scroll position
        components.html("""
            <script>
                const observer = new MutationObserver(function(mutations) {
                    const buttons = window.parent.document.querySelectorAll('button');
                    buttons.forEach(button => {
                        const btnText = button.textContent;
                        if (!button.dataset.scrollCaptureAttached &&
                            (btnText.includes('Confirm') ||
                             btnText.includes('Alt') ||
                             btnText.includes('Skip') ||
                             btnText.includes('Change') ||
                             btnText.includes('Search') ||
                             btnText.includes('Use'))) {
                            button.addEventListener('click', function() {
                                const scrollPos = window.parent.document.querySelector('section.main').scrollTop;
                                window.parent.sessionStorage.setItem('streamlit_scroll_position', scrollPos);
                            });
                            button.dataset.scrollCaptureAttached = 'true';
                        }
                    });
                });
                observer.observe(window.parent.document.body, { childList: true, subtree: true });
                setTimeout(function() {
                    const buttons = window.parent.document.querySelectorAll('button');
                    buttons.forEach(button => {
                        const btnText = button.textContent;
                        if (!button.dataset.scrollCaptureAttached &&
                            (btnText.includes('Confirm') ||
                             btnText.includes('Alt') ||
                             btnText.includes('Skip') ||
                             btnText.includes('Change') ||
                             btnText.includes('Search') ||
                             btnText.includes('Use'))) {
                            button.addEventListener('click', function() {
                                const scrollPos = window.parent.document.querySelector('section.main').scrollTop;
                                window.parent.sessionStorage.setItem('streamlit_scroll_position', scrollPos);
                            });
                            button.dataset.scrollCaptureAttached = 'true';
                        }
                    });
                }, 100);
            </script>
        """, height=0)

        # Table header
        hcol1, hcol2, hcol3, hcol4, hcol5, hcol6 = st.columns([2.2, 2.2, 1.8, 1.0, 1.0, 1.5])
        with hcol1:
            st.markdown("**Product**")
        with hcol2:
            st.markdown("**Matched Slide**")
        with hcol3:
            st.markdown("**Source**")
        with hcol4:
            st.markdown("**Conf%**")
        with hcol5:
            st.markdown("**Status**")
        with hcol6:
            st.markdown("**Actions**")

        st.divider()

        # Display each match in unified table
        for idx, match_item in enumerate(all_matches):
            result = match_item['result']
            match_type = match_item['type']

            # Create unique key
            match_key = f"unified_match_{idx}"

            # Check confirmation status
            confirmation = st.session_state.match_confirmations.get(result.gs_product_name, {})
            is_confirmed = confirmation.get('confirmed', False)
            is_skipped = confirmation.get('skipped', False)
            show_alternatives = confirmation.get('show_alternatives', False)
            show_search = confirmation.get('show_search', False)

            # Determine source and if it needs review
            if match_type == 'exact':
                if result.match_source == 'confirmed':
                    # Show which dataset this was confirmed in
                    source = f"Previously Confirmed"
                    needs_review = False
                elif result.match_source == 'manual':
                    source = "Manual Override"
                    needs_review = False
                else:
                    source = "Exact Match"
                    needs_review = False
            elif match_type == 'fuzzy':
                if is_confirmed:
                    source = "Auto-match"
                    needs_review = False
                elif is_skipped:
                    source = "Auto-match"
                    needs_review = False
                else:
                    source = "Auto-match"
                    needs_review = True
            else:  # poor match
                if is_confirmed:
                    source = "Manual Search"
                    needs_review = False
                elif is_skipped:
                    source = "No match found"
                    needs_review = False
                else:
                    source = "No match found"
                    needs_review = True

            col1, col2, col3, col4, col5, col6 = st.columns([2.2, 2.2, 1.8, 1.0, 1.0, 1.5])

            with col1:
                st.markdown(f"{result.gs_product_name}")

            with col2:
                if is_confirmed:
                    st.markdown(f"{confirmation['pptx_name']}")
                elif is_skipped:
                    st.markdown("(Skipped)")
                elif match_type == 'poor' or (match_type == 'fuzzy' and result.confidence < 70):
                    st.markdown("(No match)")
                else:
                    st.markdown(f"{result.pptx_product_name}")

            with col3:
                st.markdown(source)

            with col4:
                # Show confidence for fuzzy and poor matches
                if match_type == 'fuzzy' or match_type == 'poor':
                    if result.confidence >= 90:
                        st.markdown(f":green[{result.confidence}%]")
                    elif result.confidence >= 70:
                        st.markdown(f":orange[{result.confidence}%]")
                    elif result.confidence > 0:
                        st.markdown(f":red[{result.confidence}%]")
                    else:
                        st.markdown("—")
                else:
                    st.markdown("—")

            with col5:
                if needs_review and not is_confirmed and not is_skipped:
                    st.markdown("Review")
                else:
                    st.markdown("Done")

            with col6:
                if needs_review and not is_confirmed and not is_skipped:
                    # Items needing confirmation: Show "Confirm" and "Change" buttons
                    btn_col1, btn_col2 = st.columns(2)
                    with btn_col1:
                        if st.button("Confirm", key=f"{match_key}_confirm", help="Confirm this match", use_container_width=True):
                            # Set confirmation in session state FIRST
                            st.session_state.match_confirmations[result.gs_product_name] = {
                                'confirmed': True,
                                'pptx_name': result.pptx_product_name
                            }

                            # Save confirmation to Google Sheets and wait for completion
                            save_success = False
                            try:
                                if pptx_name_to_index and result.pptx_product_name in pptx_name_to_index:
                                    slide_index = pptx_name_to_index[result.pptx_product_name]
                                    save_success, save_message = save_confirmed_match(
                                        product_name=result.gs_product_name,
                                        slide_index=slide_index,
                                        slide_title=result.pptx_product_name,
                                        dataset=current_dataset,
                                        match_type='fuzzy_confirmed' if match_type == 'fuzzy' else 'poor_confirmed',
                                        confidence=result.confidence
                                    )
                                    if not save_success:
                                        st.warning(f"Match confirmed locally but couldn't save to Google Sheets: {save_message}")
                            except Exception as save_error:
                                # Don't block on save error - confirmation is already in session state
                                st.warning(f"Match confirmed locally but couldn't save to Google Sheets: {save_error}")

                            # Only rerun after save completes (successful or not - session state is set either way)
                            st.rerun()
                    with btn_col2:
                        if st.button("Change", key=f"{match_key}_change", help="Choose a different match", use_container_width=True):
                            if match_type == 'poor':
                                # For poor matches, show search
                                st.session_state.match_confirmations[result.gs_product_name] = {
                                    'confirmed': False,
                                    'show_search': True
                                }
                            else:
                                # For fuzzy matches, show alternatives
                                st.session_state.match_confirmations[result.gs_product_name] = {
                                    'confirmed': False,
                                    'show_alternatives': True
                                }
                            st.rerun()
                else:
                    # Items already ready to use: Only show "Change" button
                    if st.button("Change", key=f"{match_key}_change", use_container_width=True):
                        if match_type == 'poor':
                            # For poor matches, show search
                            st.session_state.match_confirmations[result.gs_product_name] = {
                                'confirmed': False,
                                'show_search': True
                            }
                        else:
                            # For exact/fuzzy matches, show alternatives
                            st.session_state.match_confirmations[result.gs_product_name] = {
                                'confirmed': False,
                                'skipped': False,
                                'show_alternatives': True
                            }
                        st.rerun()

            # Show alternatives inline if requested
            if show_alternatives and match_type != 'poor':
                st.markdown("")
                st.markdown(f"**Alternatives for {result.gs_product_name}:**")

                if result.alternatives:
                    st.caption("Select an alternative or search below:")

                    # Show top 3 alternatives
                    for alt_idx, (alt_name, alt_score) in enumerate(result.alternatives[:3]):
                        alt_col1, alt_col2 = st.columns([4, 1])
                        with alt_col1:
                            st.markdown(f"{alt_idx + 1}. {alt_name} ({alt_score}% confidence)")
                        with alt_col2:
                            if st.button(f"Use", key=f"{match_key}_alt_{alt_idx}", use_container_width=True):
                                # Set confirmation in session state
                                st.session_state.match_confirmations[result.gs_product_name] = {
                                    'confirmed': True,
                                    'pptx_name': alt_name
                                }

                                # Save alternative selection to Google Sheets
                                try:
                                    if pptx_name_to_index and alt_name in pptx_name_to_index:
                                        slide_index = pptx_name_to_index[alt_name]
                                        save_confirmed_match(
                                            product_name=result.gs_product_name,
                                            slide_index=slide_index,
                                            slide_title=alt_name,
                                            dataset=current_dataset,
                                            match_type='alternative_selected',
                                            confidence=alt_score
                                        )
                                except Exception as save_error:
                                    # Don't block on save error - confirmation is already in session state
                                    st.warning(f"Match saved locally but couldn't save to Google Sheets: {save_error}")

                                st.rerun()

                # Add search feature
                st.markdown("**Search for different slide:**")
                search_query = st.text_input(
                    "Search slide names",
                    key=f"{match_key}_search_input",
                    placeholder="e.g., 'short sleeve tee'"
                )

                if search_query:
                    matching_slides = [
                        name for name in pptx_product_names
                        if search_query.lower() in name.lower()
                    ]

                    if matching_slides:
                        st.caption(f"Found {len(matching_slides)} matching slides:")
                        for search_idx, slide_name in enumerate(matching_slides[:10]):
                            search_col1, search_col2 = st.columns([4, 1])
                            with search_col1:
                                st.markdown(f"- {slide_name}")
                            with search_col2:
                                if st.button(f"Use", key=f"{match_key}_search_{search_idx}", use_container_width=True):
                                    # Set confirmation in session state
                                    st.session_state.match_confirmations[result.gs_product_name] = {
                                        'confirmed': True,
                                        'pptx_name': slide_name
                                    }

                                    # Save search selection to Google Sheets
                                    try:
                                        if pptx_name_to_index and slide_name in pptx_name_to_index:
                                            slide_index = pptx_name_to_index[slide_name]
                                            save_confirmed_match(
                                                product_name=result.gs_product_name,
                                                slide_index=slide_index,
                                                slide_title=slide_name,
                                                dataset=current_dataset,
                                                match_type='search_selected',
                                                confidence=0
                                            )
                                    except Exception as save_error:
                                        # Don't block on save error - confirmation is already in session state
                                        st.warning(f"Match saved locally but couldn't save to Google Sheets: {save_error}")

                                    st.rerun()

                        if len(matching_slides) > 10:
                            st.caption(f"...and {len(matching_slides) - 10} more. Refine your search.")
                    else:
                        st.info("No slides found. Try different keywords.")

                # Add Skip and Cancel buttons
                st.markdown("")
                col_skip, col_cancel = st.columns(2)
                with col_skip:
                    if st.button("Skip this product", key=f"{match_key}_skip_alt", use_container_width=True):
                        st.session_state.match_confirmations[result.gs_product_name] = {
                            'confirmed': False,
                            'skipped': True
                        }
                        st.rerun()

                with col_cancel:
                    if st.button("Cancel", key=f"{match_key}_cancel_alt", type="secondary",
                                 help="Return without making changes", use_container_width=True):
                        # Clear the show_alternatives flag to close the change interface
                        st.session_state.match_confirmations[result.gs_product_name] = {
                            'confirmed': False,
                            'show_alternatives': False
                        }
                        st.rerun()

                st.markdown("")

            # Show search inline for poor matches
            if show_search and match_type == 'poor':
                st.markdown("")
                st.markdown(f"**Search for slide for {result.gs_product_name}:**")

                search_query = st.text_input(
                    "Enter slide name to search",
                    key=f"{match_key}_poor_search_input",
                    placeholder="e.g., 'honey', 'infused', 'sampler'"
                )

                if search_query:
                    matching_slides = [
                        name for name in pptx_product_names
                        if search_query.lower() in name.lower()
                    ]

                    if matching_slides:
                        st.caption(f"Found {len(matching_slides)} matching slides:")
                        for search_idx, slide_name in enumerate(matching_slides[:15]):
                            search_col1, search_col2 = st.columns([4, 1])
                            with search_col1:
                                st.markdown(f"- {slide_name}")
                            with search_col2:
                                if st.button(f"Use", key=f"{match_key}_poor_search_{search_idx}", use_container_width=True):
                                    # Set confirmation in session state
                                    st.session_state.match_confirmations[result.gs_product_name] = {
                                        'confirmed': True,
                                        'pptx_name': slide_name
                                    }

                                    # Save poor match search selection to Google Sheets
                                    try:
                                        if pptx_name_to_index and slide_name in pptx_name_to_index:
                                            slide_index = pptx_name_to_index[slide_name]
                                            save_confirmed_match(
                                                product_name=result.gs_product_name,
                                                slide_index=slide_index,
                                                slide_title=slide_name,
                                                dataset=current_dataset,
                                                match_type='manual_search',
                                                confidence=0
                                            )
                                    except Exception as save_error:
                                        # Don't block on save error - confirmation is already in session state
                                        st.warning(f"Match saved locally but couldn't save to Google Sheets: {save_error}")

                                    st.rerun()

                        if len(matching_slides) > 15:
                            st.caption(f"...and {len(matching_slides) - 15} more. Refine your search.")
                    else:
                        st.info("No slides found. Try different keywords.")

                # Add Skip and Cancel buttons for poor matches
                st.markdown("")
                col_skip_poor, col_cancel_poor = st.columns(2)
                with col_skip_poor:
                    if st.button("Skip this product", key=f"{match_key}_skip_poor", use_container_width=True):
                        st.session_state.match_confirmations[result.gs_product_name] = {
                            'confirmed': False,
                            'skipped': True
                        }
                        st.rerun()

                with col_cancel_poor:
                    if st.button("Cancel", key=f"{match_key}_cancel_poor", type="secondary",
                                 help="Return without making changes", use_container_width=True):
                        # Clear the show_search flag to close the change interface
                        st.session_state.match_confirmations[result.gs_product_name] = {
                            'confirmed': False,
                            'show_search': False
                        }
                        st.rerun()

                st.markdown("")

        st.divider()

    else:
        st.info("No products to match.")

    st.markdown("---")

    # Check if all fuzzy matches are confirmed
    all_confirmed = True
    pending_confirmations = []

    for result in fuzzy_matches:
        confirmation = st.session_state.match_confirmations.get(result.gs_product_name, {})
        if not confirmation.get('confirmed') and not confirmation.get('skipped'):
            all_confirmed = False
            pending_confirmations.append(result.gs_product_name)

    if not all_confirmed:
        st.warning(f"Please review {len(pending_confirmations)} pending matches before generating PowerPoint.")
        with st.expander("Pending Confirmations"):
            for product in pending_confirmations:
                st.markdown(f"- {product}")
    else:
        # Ensure all pricing edits are finalized before PowerPoint generation
        if st.session_state.get('proposal_pricing_pending_finalization', False):
            # First detection - clear flag and force rerun to ensure state is persisted
            st.session_state.proposal_pricing_pending_finalization = False
            st.info("⏳ Finalizing pricing changes...")
            st.rerun()

        # Build final confirmed matches dict
        confirmed_matches = {}

        # Add exact matches (auto-confirmed)
        for result in exact_matches:
            confirmed_matches[result.gs_product_name] = result.pptx_product_name

        # Add fuzzy matches (user-confirmed)
        for result in fuzzy_matches:
            confirmation = st.session_state.match_confirmations.get(result.gs_product_name, {})
            if confirmation.get('confirmed'):
                confirmed_matches[result.gs_product_name] = confirmation['pptx_name']

        # Add poor/no matches (manually selected)
        if poor_matches or no_matches:
            all_poor = poor_matches + no_matches
            for result in all_poor:
                confirmation = st.session_state.match_confirmations.get(result.gs_product_name, {})
                if confirmation.get('confirmed'):
                    confirmed_matches[result.gs_product_name] = confirmation['pptx_name']

        # Show summary
        st.success(f"Ready to generate PowerPoint with {len(confirmed_matches)} products!")

        # ============================================================
        # VARIANT DETECTION AND GROUPING (v6.13)
        # ============================================================
        from src.pptx_generator import detect_variant_groups, check_pricing_consistency, calculate_proposal_pricing

        # Detect if multiple products match to same slide (variant groups)
        variant_groups, single_products = detect_variant_groups(confirmed_matches)

        # Initialize variant grouping preferences and pricing consistency flags if not exists
        if 'variant_grouping_prefs' not in st.session_state:
            st.session_state.variant_grouping_prefs = {}
        if 'variant_pricing_consistent' not in st.session_state:
            st.session_state.variant_pricing_consistent = {}

        # Show variant confirmation UI if variants detected
        if variant_groups:
            st.markdown("---")
            st.subheader("Multi-Variant Products Detected")
            st.info(
                "Multiple products matched to the same PowerPoint slide. "
                "These are typically size/flavor variants."
            )

            for slide_name, products in variant_groups.items():
                # Calculate pricing for all variants to check consistency
                pricing_data_list = []
                for gs_name in products:
                    proposal_item = next(
                        (item for item in st.session_state.proposal_products
                         if item['product_data']['Product/Service'] == gs_name),
                        None
                    )
                    if proposal_item:
                        pricing_data = calculate_proposal_pricing(
                            proposal_item,
                            get_unit_price_new_system,
                            st.session_state.get('proposal_marketing_rounding', False),
                            st.session_state.get('proposal_discount_percent', 0.0)
                        )
                        if pricing_data:
                            pricing_data_list.append(pricing_data)

                # Check pricing consistency
                has_consistent_pricing = check_pricing_consistency(pricing_data_list)
                st.session_state.variant_pricing_consistent[slide_name] = has_consistent_pricing

                # Expander title with pricing indicator
                if has_consistent_pricing:
                    expander_title = f"{slide_name} ({len(products)} variants) - Consistent Pricing"
                else:
                    expander_title = f"{slide_name} ({len(products)} variants) - Variable Pricing"

                with st.expander(expander_title, expanded=True):
                    st.markdown("**Products matched to this slide:**")

                    # Show each product with its price and MOQ
                    for idx, gs_name in enumerate(products):
                        if idx < len(pricing_data_list):
                            price = pricing_data_list[idx].get('client_price', 0)
                            moq = pricing_data_list[idx].get('moq', 0)
                            st.markdown(f"• {gs_name} - MOQ: {moq}, Price: ${price:.2f}")
                        else:
                            st.markdown(f"• {gs_name}")

                    # Get partner info to check if all variants are from same partner
                    partners = set()
                    for prod in products:
                        prod_item = next(
                            (item for item in st.session_state.proposal_products
                             if item['product_data']['Product/Service'] == prod),
                            None
                        )
                        if prod_item:
                            partner = prod_item['product_data'].get('Partner', '')
                            if partner:
                                partners.add(partner)

                    # Warning if multiple partners
                    if len(partners) > 1:
                        st.warning(
                            f"These products are from different partners: {', '.join(partners)}. "
                            "Creating separate slides is recommended."
                        )

                    # Conditional radio button options based on pricing consistency
                    if has_consistent_pricing:
                        options = [
                            "Display single row (all variants have same pricing)",
                            "Display all variants (show each variant in separate row)",
                            "Create separate slides (duplicate slide for each variant)",
                            "Skip these products"
                        ]
                        default_choice = 0  # Simple table by default for consistent pricing
                    else:
                        options = [
                            "Display together (recommended - fills multiple table rows)",
                            "Create separate slides (duplicate slide for each variant)",
                            "Skip these products"
                        ]
                        default_choice = 0 if len(partners) <= 1 else 1  # Default to "together" unless multiple partners

                    # Get existing preference or use default
                    existing_pref = st.session_state.variant_grouping_prefs.get(slide_name, options[default_choice])
                    if existing_pref not in options:
                        existing_pref = options[default_choice]

                    group_choice = st.radio(
                        "How should these be displayed?",
                        options=options,
                        key=f"variant_choice_{slide_name}",
                        index=options.index(existing_pref) if existing_pref in options else default_choice
                    )

                    st.session_state.variant_grouping_prefs[slide_name] = group_choice

        # ============================================================
        # IMPACT SLIDE SELECTION (SIMPLIFIED)
        # ============================================================
        st.markdown("---")
        st.subheader("Step 2: Impact Slides")

        # Import impact slide functions
        from src.slide_matcher import extract_unique_partners, PARTNER_IMPACT_SLIDES, find_all_impact_slides

        # Extract unique partners from proposal
        unique_partners = extract_unique_partners(st.session_state.proposal_products)

        if unique_partners:
            # Initialize impact selections if not exists
            if 'impact_slide_selections' not in st.session_state:
                st.session_state.impact_slide_selections = {}

            # Auto-select from reference table for all partners
            partners_with_slides = []
            partners_without_slides = []

            for partner in unique_partners:
                auto_selected = PARTNER_IMPACT_SLIDES.get(partner)
                if auto_selected:
                    # Auto-select if not already set
                    current_selection = st.session_state.impact_slide_selections.get(partner)
                    if not current_selection or current_selection.get('slide_index') is None:
                        st.session_state.impact_slide_selections[partner] = auto_selected.copy()
                    partners_with_slides.append(partner)
                else:
                    partners_without_slides.append(partner)
                    st.session_state.impact_slide_selections[partner] = {
                        'slide_title': None,
                        'slide_index': None
                    }

            # Show simple summary message
            if partners_with_slides:
                partner_list = ", ".join(partners_with_slides)
                st.success(f"Impact slides found for: {partner_list}")

            if partners_without_slides:
                partner_list = ", ".join(partners_without_slides)
                st.warning(f"No impact slides found for: {partner_list}")

            # Collapsible customization section
            with st.expander("Customize Impact Slides (Optional)", expanded=False):
                st.caption("Change or remove impact slides for specific partners")

                from pathlib import Path

                for partner in unique_partners:
                    active_selection = st.session_state.impact_slide_selections.get(partner, {})

                    st.markdown(f"**{partner}**")

                    # Load all impact slide options (from cloud or local)
                    pptx_template = get_template_path('all_slides', show_loading=False)
                    if pptx_template:
                        all_impact_slides = find_all_impact_slides(pptx_template)
                    else:
                        st.error("Could not load PowerPoint template")
                        continue

                    dropdown_options = ["None - Skip impact slide"]
                    dropdown_options.extend([f"{slide['slide_title']}" for slide in all_impact_slides])

                    # Find current selection index
                    current_title = active_selection.get('slide_title')
                    try:
                        default_index = dropdown_options.index(current_title) if current_title else 0
                    except ValueError:
                        default_index = 0

                    selected_option = st.selectbox(
                        f"Impact slide for {partner}",
                        options=dropdown_options,
                        index=default_index,
                        key=f"impact_select_{partner}",
                        label_visibility="collapsed"
                    )

                    # Update selection immediately on change
                    if selected_option == "None - Skip impact slide":
                        st.session_state.impact_slide_selections[partner] = {
                            'slide_title': None,
                            'slide_index': None
                        }
                    else:
                        selected_slide = next(
                            (slide for slide in all_impact_slides if slide['slide_title'] == selected_option),
                            None
                        )
                        if selected_slide:
                            st.session_state.impact_slide_selections[partner] = {
                                'slide_title': selected_slide['slide_title'],
                                'slide_index': selected_slide['slide_index']
                            }

                    st.markdown("")  # Spacing

        else:
            st.info("No products selected - add products to see impact slides")
            st.session_state.impact_slide_selections = {}

        st.markdown("---")

        col1, col2, col3 = st.columns([1, 1, 1])
        with col1:
            if st.button("Reset Confirmations", use_container_width=True):
                st.session_state.match_confirmations = {}
                st.session_state.impact_slide_selections = {}
                st.session_state.generated_pptx = None
                st.rerun()

        with col2:
            if st.button("Generate PowerPoint Presentation", type="primary", use_container_width=True):
                try:
                    import time
                    start_time = time.time()

                    # Memory optimization: clear template cache before generation
                    if USE_MEMORY_OPTIMIZATION:
                        from src.template_loader import clear_template_cache
                        clear_template_cache()
                        gc.collect()
                        st.info("Memory optimization enabled - using direct download mode")

                    # Import generator functions
                    from src.pptx_generator import (
                        create_complete_proposal_presentation,
                        download_presentation
                    )

                    # Load templates from cloud/local (with loading spinner)
                    st.info(f"Using template: **{get_template_name('all_slides')}**")

                    # Load templates with memory optimization if enabled
                    november_template_path = get_template_path('all_slides', show_loading=True, use_cache=not USE_MEMORY_OPTIMIZATION)
                    intro_outro_template_path = get_template_path('intro_outro', show_loading=False, use_cache=not USE_MEMORY_OPTIMIZATION)

                    # Validation checks
                    if not november_template_path:
                        st.error("PowerPoint template could not be loaded. Please check Google Drive access.")
                        return None

                    if not intro_outro_template_path:
                        st.error("Intro/Outro template could not be loaded.")
                        return None

                    if len(confirmed_matches) == 0:
                        st.warning("No products confirmed for generation. Please confirm at least one product match.")
                        return None

                    # Create presentation with progress indicator
                    progress_container = st.empty()
                    progress_container.info(f"Step 1/5: Loading templates...")

                    # Get proposal settings
                    fifty_cent_rounding = st.session_state.get('proposal_fifty_cent_rounding', False)
                    marketing_rounding = st.session_state.get('proposal_marketing_rounding', False)
                    discount_percent = st.session_state.get('proposal_discount_percent', 0.0)
                    discount_type = st.session_state.get('proposal_discount_type', None)

                    # Get selected impact slides (for overrides)
                    # Build impact_slide_overrides dict: {partner: {"slide_index": X, "slide_title": Y}}
                    impact_slide_overrides = {}
                    if 'impact_slide_selections' in st.session_state:
                        for partner, selection in st.session_state.impact_slide_selections.items():
                            if selection.get('slide_index') is not None:
                                impact_slide_overrides[partner] = selection

                    # Count total slides
                    num_products = len(confirmed_matches)
                    num_impacts = len([s for s in impact_slide_overrides.values() if s.get('slide_index') is not None])

                    # Get variant groups and preferences
                    variant_groups_for_generation = variant_groups if variant_groups else None
                    variant_prefs_for_generation = st.session_state.get('variant_grouping_prefs', None)

                    # Debug output
                    if variant_groups_for_generation:
                        print(f"DEBUG: Passing {len(variant_groups_for_generation)} variant groups to generator")
                        for slide_name, products in variant_groups_for_generation.items():
                            print(f"  - {slide_name}: {products}")
                        print(f"DEBUG: Preferences: {variant_prefs_for_generation}")
                    else:
                        print("DEBUG: No variant groups detected")

                    # Create complete presentation (products + impacts + outro only)
                    progress_container.info(f"Step 2/4: Selecting and updating {num_products} product slide(s) + {num_impacts} impact slide(s)...")
                    prs = create_complete_proposal_presentation(
                        str(november_template_path),
                        str(intro_outro_template_path),
                        confirmed_matches,
                        st.session_state.proposal_products,
                        get_unit_price_new_system,
                        marketing_rounding,
                        discount_percent,
                        impact_slide_overrides if impact_slide_overrides else None,
                        variant_groups_for_generation,
                        variant_prefs_for_generation,
                        fifty_cent_rounding,
                        discount_type
                    )

                    progress_container.info(f"Step 3/4: Adding outro slides (4 slides)...")

                    # Convert to downloadable format
                    progress_container.info("Step 4/4: Preparing download...")
                    client_name = st.session_state.order_details.get('company_name', 'Client')
                    pptx_file = download_presentation(prs, client_name)

                    # Memory optimization: force garbage collection after generation
                    if USE_MEMORY_OPTIMIZATION:
                        gc.collect()

                    # Calculate generation time
                    generation_time = time.time() - start_time

                    # Store in session state to persist across reruns
                    st.session_state.generated_pptx = pptx_file
                    st.session_state.generated_pptx_filename = f"Proposal_{client_name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx"
                    st.session_state.pptx_product_count = len(confirmed_matches)
                    st.session_state.pptx_generation_time = generation_time

                    # Clear progress indicator
                    progress_container.empty()

                except Exception as e:
                    st.error(f"PowerPoint generation failed: {str(e)}")
                    st.error("Please check that all products have valid pricing data and try again.")
                    with st.expander("View detailed error (for debugging)"):
                        st.exception(e)
                    return None

        with col3:
            if st.button("Close", use_container_width=True):
                st.session_state.show_pptx_matching = False
                st.session_state.match_confirmations = {}
                st.session_state.generated_pptx = None
                st.rerun()

        # Show download button if presentation was generated
        if 'generated_pptx' in st.session_state and st.session_state.generated_pptx is not None:
            generation_time = st.session_state.get('pptx_generation_time', 0)
            st.success(f"PowerPoint generated successfully in {generation_time:.1f} seconds!")

            # Show generation summary
            col_summary1, col_summary2 = st.columns(2)
            with col_summary1:
                st.metric("Products Included", st.session_state.pptx_product_count)
            with col_summary2:
                st.metric("File Name", st.session_state.generated_pptx_filename)

            # Show instructions for reordering intro slides
            st.info("**Final step:** In PowerPoint, move the 8 intro slides (they're grouped together after products) to the beginning. Select slides → Drag to top (~5 seconds).")

            # Warning about PowerPoint repair message
            st.warning("**Note:** PowerPoint may show a 'repair presentation' warning when opening the file. This is normal - click 'Repair' to proceed. The content will display correctly.")

            st.download_button(
                label="Download PowerPoint Presentation",
                data=st.session_state.generated_pptx,
                file_name=st.session_state.generated_pptx_filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary",
                use_container_width=True
            )

    return None

# ============================================================
# DATA LOADING (Tab 1)
# ============================================================
# Data is now loaded before sidebar, so just retrieve from session state
try:
    df_template = st.session_state.df_template
    df_metadata = st.session_state.df_metadata
    df_partner_info = st.session_state.df_partner_info

except Exception as e:
    error_msg = str(e)

    # Check if this is a real dataset structure issue
    if st.session_state.selected_dataset == 'real' and ('Data' in error_msg or 'worksheet' in error_msg.lower()):
        st.error("Real pricing dataset is not yet properly structured.")
        st.warning("The real data spreadsheet needs to have the same structure as the demo data:\n"
                   "- Sheet 1: 'Data' (pricing data)\n"
                   "- Sheet 2: 'Metadata' (field definitions)\n"
                   "- Sheet 3: 'Partner-Specific Info' (partner contacts)")
        st.info("Please switch to 'Demo Data' in the sidebar, or complete the real dataset structure first.")
        st.stop()
    else:
        st.error(f"Failed to load data: {error_msg}")
        st.stop()

# ============================================================
# TAB STRUCTURE
# ============================================================
# Use query params to preserve active tab on rerun
# Check for Streamlit version compatibility
try:
    # Try the new API first (Streamlit >= 1.30.0)
    query_params = st.query_params
    active_tab = query_params.get("tab", "0")
except AttributeError:
    # Fall back to old API
    query_params = st.experimental_get_query_params()
    active_tab = query_params.get("tab", ["0"])[0]

# Convert to integer
try:
    active_tab_index = int(active_tab)
except (ValueError, TypeError):
    active_tab_index = 0

# Create tabs
tab_list = st.tabs([
    "Proposal Generator",
    "Client Order Form Generator",
    "Order & Client Info",
    "Execution & Accounting",
    "Executive Pricing Tool"
])

tab1, tab2, tab3, tab4, tab5 = tab_list

# ============================================================
# TAB 1: PROPOSAL GENERATOR
# ============================================================
with tab1:
    st.header("Proposal Generator - Product Catalog & Pricing")
    st.caption("Browse products, configure proposals, and generate client quotes (tables & PowerPoint)")
    st.divider()

    # ============================================================
    # SECTION 1: BROWSE & FILTER PRODUCTS (Combined Sections 1+2)
    # ============================================================
    st.subheader("1. Browse & Filter Products")
    st.caption(f"{len(df_template)} total products available")

    # Search bar - prominently placed above all filters
    search_query = st.text_input(
        "Search product names",
        placeholder="Type to search product names...",
        key="product_search",
        help="Search filters products by name only. Use the filters below for partner and country."
    )

    col1, col2, col3 = st.columns(3)

    with col1:
        st.markdown("**Client Budget Range (Per Unit)**")
        
        # Create two columns for min and max budget
        budget_col1, budget_col2 = st.columns(2)
        
        with budget_col1:
            min_budget = st.number_input(
                "Min price ($)",
                min_value=0.0,
                value=st.session_state.proposal_filters.get('min_budget', 0.0),
                step=1.0,
                key="filter_min_budget",
                help="Minimum client price per unit"
            )
        
        with budget_col2:
            max_budget = st.number_input(
                "Max price ($)",
                min_value=0.0,
                value=st.session_state.proposal_filters.get('max_budget', 0.0),
                step=1.0,
                key="filter_max_budget",
                help="Maximum client price per unit"
            )
        
        # Validation warning
        if min_budget and max_budget and min_budget > 0 and max_budget > 0 and min_budget > max_budget:
            st.error("Min price cannot be greater than Max price")

    with col2:
        st.markdown("**Partner/Maker**")
        all_partners = sorted(df_template["Partner"].unique().tolist())
        selected_partners = st.multiselect(
            "Select partners (leave empty for all)",
            options=all_partners,
            default=st.session_state.proposal_filters.get('partners', []),
            key="filter_partners"
        )

    with col3:
        st.markdown("**Country of Origin (Ships From)**")
        all_countries = sorted(df_template["Country of Origin (Ships From)"].dropna().unique().tolist())
        selected_countries = st.multiselect(
            "Select countries (leave empty for all)",
            options=all_countries,
            default=st.session_state.proposal_filters.get('countries', []),
            key="filter_countries"
        )

    # Update filters in session state
    st.session_state.proposal_filters['min_budget'] = min_budget if min_budget and min_budget > 0 else 0.0
    st.session_state.proposal_filters['max_budget'] = max_budget if max_budget and max_budget > 0 else 0.0
    st.session_state.proposal_filters['partners'] = selected_partners
    st.session_state.proposal_filters['countries'] = selected_countries

    # Apply filters
    filtered_df = df_template.copy()

    if selected_partners:
        filtered_df = filtered_df[filtered_df["Partner"].isin(selected_partners)]

    if selected_countries:
        filtered_df = filtered_df[filtered_df["Country of Origin (Ships From)"].isin(selected_countries)]

    # Price filtering based on client price range (cost * 2 for 100% markup)
    # Only apply if valid range (min <= max when both are set)
    if ((min_budget and min_budget > 0) or (max_budget and max_budget > 0)) and \
       not (min_budget and max_budget and min_budget > 0 and max_budget > 0 and min_budget > max_budget):
        price_filtered_indices = []
        for idx, row in filtered_df.iterrows():
            # Get cost estimate at quantity 100
            base_cost, _, _ = get_unit_price_new_system(row, 100)
            if base_cost:
                # Calculate client price (100% markup)
                client_price = base_cost * 2
                
                # Check if price is within range
                if min_budget and min_budget > 0 and client_price < min_budget:
                    continue
                if max_budget and max_budget > 0 and client_price > max_budget:
                    continue
                    
                price_filtered_indices.append(idx)
        filtered_df = filtered_df.loc[price_filtered_indices]

    # Search filtering - only search product names
    if search_query:
        search_lower = search_query.lower()
        search_mask = filtered_df['Product/Service'].str.lower().str.contains(search_lower, na=False)
        filtered_df = filtered_df[search_mask]

    st.divider()

    # Enhanced filter results message
    if search_query and len(filtered_df) == 0:
        st.caption(f"**No products found with name containing '{search_query}'** - Try a different search term")
    elif search_query:
        st.caption(f"**Filtered Results:** {len(filtered_df)} products with names containing '{search_query}'")
    else:
        st.caption(f"**Filtered Results:** {len(filtered_df)} products match your filters")

    # Display success message if a product was just added
    if 'show_success_message' in st.session_state and st.session_state.show_success_message:
        st.toast(f"Added {st.session_state.success_product_name} to proposal!")
        st.session_state.show_success_message = False

    # Bulk add success message
    if 'show_bulk_success_message' in st.session_state and st.session_state.show_bulk_success_message:
        st.toast(st.session_state.bulk_success_message)
        st.session_state.show_bulk_success_message = False

    # ============================================================
    # BULK ACTIONS SECTION
    # ============================================================
    if len(filtered_df) > 0:
        with st.expander("Bulk Actions - Add All Products from Partner(s)", expanded=False):
            st.caption("Quickly add all products from one or more partners to your proposal")

            # Add JavaScript to capture scroll position for bulk add button
            components.html("""
                <script>
                    // Store scroll position before bulk add button click
                    const buttons = window.parent.document.querySelectorAll('button');
                    buttons.forEach(button => {
                        if (button.textContent.includes('Add') && button.textContent.includes('Products')) {
                            button.addEventListener('click', function() {
                                const scrollPos = window.parent.document.querySelector('section.main').scrollTop;
                                window.parent.sessionStorage.setItem('streamlit_scroll_position', scrollPos);
                            });
                        }
                    });
                </script>
            """, height=0)

            # "Add All Products" button at the top (for testing)
            st.markdown("**Quick Add All Products (Testing)**")

            # Count how many products would be added
            existing_products = {item['product_data']['Product/Service'] for item in st.session_state.proposal_products}
            new_products_all = [row for idx, row in filtered_df.iterrows() if row['Product/Service'] not in existing_products]
            new_count_all = len(new_products_all)
            duplicate_count_all = len(filtered_df) - new_count_all

            col_info, col_button = st.columns([3, 1])
            with col_info:
                if new_count_all > 0:
                    st.info(f"Will add **{new_count_all} new product(s)** from filtered results ({len(filtered_df)} total)")
                    if duplicate_count_all > 0:
                        st.caption(f"({duplicate_count_all} already in proposal, will be skipped)")
                else:
                    st.warning(f"All {len(filtered_df)} filtered product(s) are already in your proposal")

            with col_button:
                if new_count_all > 0:
                    if st.button(f"Add All {new_count_all}", type="secondary", use_container_width=True, key="add_all_products_button"):
                        # Add all new products to proposal with new pricing logic
                        from src.helpers import calculate_pricing_snapshot
                        added_count = 0
                        failed_products = []

                        for product_row in new_products_all:
                            try:
                                # Use new pricing calculation
                                pricing_result = calculate_pbp_msrp(product_row.to_dict(), quantity=100)

                                # Determine markup based on pricing method
                                if st.session_state.proposal_use_msrp:
                                    # Use pricing method from spreadsheet
                                    base_cost = pricing_result['calculation_details']['per_item_cost']
                                    pbp_msrp = pricing_result['pbp_msrp']

                                    if base_cost and base_cost > 0:
                                        # Calculate markup to match PBP MSRP
                                        markup = ((pbp_msrp / base_cost) - 1) * 100
                                    else:
                                        markup = 100.0
                                else:
                                    # Use default 100% markup
                                    markup = 100.0

                                # Calculate pricing snapshot to preserve pricing when importing to orders
                                pricing_snapshot = calculate_pricing_snapshot(
                                    product_row,
                                    markup,
                                    quantity=100,  # Default for MOQ calculation
                                    discount_percent=st.session_state.get('proposal_discount_percent', 0.0),
                                    marketing_rounding=st.session_state.proposal_marketing_rounding,
                                    fifty_cent_rounding=st.session_state.proposal_fifty_cent_rounding
                                )

                                proposal_item = {
                                    'product_data': product_row.to_dict(),
                                    'markup_percent': markup,
                                    'pricing_snapshot': pricing_snapshot,
                                    'pricing_method': pricing_result['method_used'],
                                    'pricing_notes': pricing_result.get('pricing_notes', ''),
                                    'manual_override': False,
                                    'validation_warning': pricing_result.get('validation_warning', None),
                                    'settings_snapshot': {
                                        'fifty_cent_rounding': st.session_state.get('proposal_fifty_cent_rounding', False),
                                        'marketing_rounding': st.session_state.get('proposal_marketing_rounding', False),
                                        'discount_percent': st.session_state.get('proposal_discount_percent', 0.0)
                                    }
                                }
                                st.session_state.proposal_products.append(proposal_item)
                                added_count += 1

                            except (KeyError, TypeError, ValueError) as e:
                                # Track products that failed due to missing pricing data
                                product_name = product_row.get('Product/Service', 'Unknown Product')
                                failed_products.append(product_name)

                        # Set success message
                        if added_count > 0:
                            st.session_state.show_bulk_success_message = True
                            st.session_state.bulk_success_message = f"Added **{added_count} filtered products** to proposal!"

                        # Show warning for failed products
                        if failed_products:
                            st.warning(f"Failed to add {len(failed_products)} product(s): Missing pricing data in master spreadsheet")
                            with st.expander("See failed products"):
                                for product_name in failed_products:
                                    st.caption(f"- {product_name}")

                        # Keep catalog expanded after bulk add
                        st.session_state.keep_catalog_expanded = True
                        st.rerun()

            st.divider()
            st.markdown("**Add by Partner**")

            # Get unique partners from filtered results
            available_partners = sorted(filtered_df["Partner"].unique().tolist())

            # Partner selection for bulk add
            bulk_partners = st.multiselect(
                "Select partner(s) to add all their products:",
                options=available_partners,
                key="bulk_add_partners",
                help="All products from selected partners will be added (filtered products only)"
            )

            if bulk_partners:
                # Count how many new products would be added
                products_to_add = filtered_df[filtered_df["Partner"].isin(bulk_partners)]

                # Get existing product names in proposal
                existing_products = {item['product_data']['Product/Service'] for item in st.session_state.proposal_products}

                # Filter out duplicates
                new_products = []
                for idx, row in products_to_add.iterrows():
                    if row['Product/Service'] not in existing_products:
                        new_products.append(row)

                new_count = len(new_products)
                duplicate_count = len(products_to_add) - new_count

                # Show counts
                col1, col2 = st.columns([3, 1])
                with col1:
                    if new_count > 0:
                        st.info(f"Will add **{new_count} new product(s)** from {', '.join(bulk_partners)}")
                        if duplicate_count > 0:
                            st.caption(f"({duplicate_count} already in proposal, will be skipped)")
                    else:
                        st.warning(f"All {len(products_to_add)} product(s) from selected partners are already in your proposal")

                with col2:
                    if new_count > 0:
                        if st.button(f"Add {new_count} Products", type="primary", use_container_width=True, key="bulk_add_button"):
                            # Add all new products to proposal with new pricing logic
                            from src.helpers import calculate_pricing_snapshot
                            added_count = 0
                            failed_products = []

                            for product_row in new_products:
                                try:
                                    # Use new pricing calculation
                                    pricing_result = calculate_pbp_msrp(product_row.to_dict(), quantity=100)

                                    # Determine markup based on pricing method
                                    if st.session_state.proposal_use_msrp:
                                        # Use pricing method from spreadsheet
                                        base_cost = pricing_result['calculation_details']['per_item_cost']
                                        pbp_msrp = pricing_result['pbp_msrp']

                                        if base_cost and base_cost > 0:
                                            # Calculate markup to match PBP MSRP
                                            markup = ((pbp_msrp / base_cost) - 1) * 100
                                        else:
                                            markup = 100.0
                                    else:
                                        # Use default 100% markup
                                        markup = 100.0

                                    # Calculate pricing snapshot to preserve pricing when importing to orders
                                    pricing_snapshot = calculate_pricing_snapshot(
                                        product_row,
                                        markup,
                                        quantity=100,  # Default for MOQ calculation
                                        discount_percent=st.session_state.get('proposal_discount_percent', 0.0),
                                        marketing_rounding=st.session_state.proposal_marketing_rounding,
                                        fifty_cent_rounding=st.session_state.proposal_fifty_cent_rounding
                                    )

                                    proposal_item = {
                                        'product_data': product_row.to_dict(),
                                        'markup_percent': markup,
                                        'pricing_snapshot': pricing_snapshot,
                                        'pricing_method': pricing_result['method_used'],
                                        'pricing_notes': pricing_result.get('pricing_notes', ''),
                                        'manual_override': False,
                                        'validation_warning': pricing_result.get('validation_warning', None),
                                        'settings_snapshot': {
                                            'fifty_cent_rounding': st.session_state.get('proposal_fifty_cent_rounding', False),
                                            'marketing_rounding': st.session_state.get('proposal_marketing_rounding', False),
                                            'discount_percent': st.session_state.get('proposal_discount_percent', 0.0)
                                        }
                                    }
                                    st.session_state.proposal_products.append(proposal_item)
                                    added_count += 1

                                except (KeyError, TypeError, ValueError) as e:
                                    # Track products that failed due to missing pricing data
                                    product_name = product_row.get('Product/Service', 'Unknown Product')
                                    failed_products.append(product_name)

                            # Set success message
                            partner_names = ', '.join(bulk_partners)
                            if added_count > 0:
                                st.session_state.show_bulk_success_message = True
                                st.session_state.bulk_success_message = f"Added **{added_count} products** from {partner_names} to proposal!"

                            # Show warning for failed products
                            if failed_products:
                                st.warning(f"Failed to add {len(failed_products)} product(s): Missing pricing data in master spreadsheet")
                                with st.expander("See failed products"):
                                    for product_name in failed_products:
                                        st.caption(f"- {product_name}")

                            # Keep catalog expanded after bulk add
                            st.session_state.keep_catalog_expanded = True
                            st.rerun()
            else:
                st.caption("Select one or more partners above to see product count")

    if len(filtered_df) == 0:
        if search_query:
            st.warning(f"No products found with name containing '{search_query}'. Try a different search term or adjust filters.")
        else:
            st.warning("No products match your filters. Try adjusting the filter criteria above.")
    else:
        # Keep catalog expanded if user just added a product, otherwise collapse after first product added
        if 'keep_catalog_expanded' in st.session_state and st.session_state.keep_catalog_expanded:
            default_expanded = True
            st.session_state.keep_catalog_expanded = False  # Reset for next time
        else:
            default_expanded = len(st.session_state.proposal_products) == 0

        with st.expander(f"Browse Products ({len(filtered_df)} available)", expanded=default_expanded):
            # Add JavaScript to capture scroll position before button clicks
            components.html("""
                <script>
                    // Store scroll position in sessionStorage before any button click
                    const buttons = window.parent.document.querySelectorAll('button');
                    buttons.forEach(button => {
                        if (button.textContent.includes('Add to Proposal')) {
                            button.addEventListener('click', function() {
                                const scrollPos = window.parent.document.querySelector('section.main').scrollTop;
                                window.parent.sessionStorage.setItem('streamlit_scroll_position', scrollPos);
                            });
                        }
                    });
                </script>
            """, height=0)

            # Table-style header
            header_col1, header_col2, header_col3, header_col4, header_col5 = st.columns([3, 1.5, 1, 1.2, 1.5])
            with header_col1:
                st.markdown("**Product Name**")
            with header_col2:
                st.markdown("**Partner**")
            with header_col3:
                st.markdown("**Cost/Unit**")
            with header_col4:
                st.markdown("**PBP MSRP**")
            with header_col5:
                st.markdown("**Actions**")

            st.divider()

            # Display filtered products in a compact table-style format
            for idx, row in filtered_df.iterrows():
                product_data = row

                # Calculate cost and client price for display using new pricing logic
                preliminary_cost, _, _ = get_unit_price_new_system(product_data, 100)

                # Use new pricing engine to get PBP MSRP
                pricing_result = calculate_pbp_msrp(product_data.to_dict(), quantity=100)
                pbp_msrp = pricing_result['pbp_msrp']

                # Calculate MOQ based on PBP MSRP (not hardcoded 100% markup)
                moq_result = calculate_moq(pbp_msrp, product_data) if pbp_msrp else None
                estimated_moq = moq_result['moq'] if moq_result else None
                moq_display_text = moq_result['display_text'] if moq_result else ""
                moq_cost, _, _ = get_unit_price_new_system(product_data, estimated_moq) if estimated_moq else (None, None, None)

                # Get PBP MSRP at MOQ quantity (for accurate pricing)
                if estimated_moq:
                    moq_pricing_result = calculate_pbp_msrp(product_data.to_dict(), quantity=estimated_moq)
                    moq_client_price = moq_pricing_result['pbp_msrp']

                    # Calculate actual markup percentage being used
                    if moq_cost and moq_cost > 0:
                        actual_markup = ((moq_client_price / moq_cost) - 1) * 100
                    else:
                        actual_markup = 0.0
                else:
                    moq_client_price = None
                    actual_markup = 0.0

                # Compact row with all essential info
                col1, col2, col3, col4, col5 = st.columns([3, 1.5, 1, 1.2, 1.5])

                with col1:
                    st.markdown(f"**{product_data['Product/Service']}**")

                with col2:
                    st.markdown(f"{product_data['Partner']}")

                with col3:
                    if moq_cost:
                        st.markdown(f"${moq_cost:.2f}")
                    else:
                        st.markdown("—")

                with col4:
                    if moq_client_price:
                        st.markdown(f"${moq_client_price:.2f}")
                    else:
                        st.markdown("—")

                with col5:
                    # Add button - adds product to proposal with new pricing logic
                    if st.button("Add to Proposal", key=f"add_{idx}", use_container_width=True, type="primary"):
                        try:
                            # Use new pricing calculation (quantity 100 for reference)
                            pricing_result = calculate_pbp_msrp(product_data.to_dict(), quantity=100)

                            # Determine markup based on pricing method
                            if st.session_state.proposal_use_msrp:
                                # Use pricing method from spreadsheet
                                per_item_cost = pricing_result['calculation_details']['per_item_cost']
                                pbp_msrp = pricing_result['pbp_msrp']

                                if per_item_cost and per_item_cost > 0:
                                    # Calculate markup to match PBP MSRP
                                    markup = ((pbp_msrp / per_item_cost) - 1) * 100
                                else:
                                    markup = 100.0
                            else:
                                # Use default 100% markup
                                markup = 100.0

                            # Calculate pricing snapshot to preserve pricing when importing to orders
                            from src.helpers import calculate_pricing_snapshot
                            pricing_snapshot = calculate_pricing_snapshot(
                                product_data,
                                markup,
                                quantity=100,  # Default for MOQ calculation
                                discount_percent=st.session_state.get('proposal_discount_percent', 0.0),
                                marketing_rounding=st.session_state.proposal_marketing_rounding,
                                fifty_cent_rounding=st.session_state.proposal_fifty_cent_rounding
                            )

                            proposal_item = {
                                'product_data': product_data.to_dict(),
                                'markup_percent': markup,
                                'pricing_snapshot': pricing_snapshot,
                                'pricing_method': pricing_result['method_used'],
                                'pricing_notes': pricing_result.get('pricing_notes', ''),
                                'manual_override': False,
                                'validation_warning': pricing_result.get('validation_warning', None),
                                'settings_snapshot': {
                                    'fifty_cent_rounding': st.session_state.get('proposal_fifty_cent_rounding', False),
                                    'marketing_rounding': st.session_state.get('proposal_marketing_rounding', False),
                                    'discount_percent': st.session_state.get('proposal_discount_percent', 0.0)
                                }
                            }
                            st.session_state.proposal_products.append(proposal_item)

                            # Set success message
                            st.session_state.show_success_message = True
                            st.session_state.success_product_name = product_data['Product/Service']

                            # Keep catalog expanded after adding product
                            st.session_state.keep_catalog_expanded = True
                            st.rerun()

                        except (KeyError, TypeError, ValueError) as e:
                            # Handle missing pricing data gracefully
                            product_name = product_data.get('Product/Service', 'Unknown Product')
                            st.error(f"Failed to add '{product_name}' to proposal: Please add pricing data to master spreadsheet.")
                            st.caption(f"Missing data: {str(e)}")
                        except Exception as e:
                            # Catch any other unexpected errors
                            product_name = product_data.get('Product/Service', 'Unknown Product')
                            st.error(f"Failed to add '{product_name}' to proposal: {str(e)}")
                            st.caption("Please contact support if this error persists.")

                # Show additional details inline (no nested expander)
                st.caption(f"Ships From: {product_data.get('Country of Origin (Ships From)', 'N/A')} | Made In: {product_data.get('Country of Origin (Made In)', 'N/A')} | Tiered: {product_data.get('Pricing Tiers (Y/N)', 'N/A')}")

                # Show pricing method and cost basis (new schema fields)
                pricing_logic = get_column_value(product_data, 'pricing_logic', None)
                cost_basis = get_column_value(product_data, 'cost_basis', 'Per Item')

                pricing_method_display = pricing_logic if pricing_logic else "Standard markup"
                st.caption(f"Pricing Method: {pricing_method_display} | Cost Basis: {cost_basis}")

                # Show MSRP if available
                msrp_raw = get_column_value(product_data, 'Vendor Published MSRP', 'MSRP', '')
                if msrp_raw and str(msrp_raw).strip() and str(msrp_raw).strip() not in ['nan', '', '0', '0.0']:
                    from src.helpers import clean_price
                    msrp_value = clean_price(msrp_raw)
                    if msrp_value and msrp_value > 0:
                        st.caption(f"Vendor MSRP: ${msrp_value:.2f}/unit")

                # Show shipping costs
                shipping_display = format_shipping_display(product_data)
                if shipping_display != "No shipping data":
                    st.caption(f"Shipping: {shipping_display}")

                # Show estimated prices at MOQ
                if moq_cost and estimated_moq and moq_client_price:
                    st.caption(f"Est. Cost & Price at {moq_display_text}: ${moq_cost:.2f}/unit cost → ${moq_client_price:.2f}/unit client price ({actual_markup:.0f}% markup)")

                # Show description if available
                desc = product_data.get("Marketing Description", "")
                if desc and str(desc).strip() and str(desc).strip() != 'nan':
                    st.caption(f"{desc}")

                st.divider()

    # ============================================================
    # SECTION 2: PROPOSAL PREVIEW & SETTINGS
    # ============================================================
    st.divider()
    st.subheader("2. Proposal Preview & Settings")

    # ============================================================
    # SAVED PROPOSALS SECTION (Always visible)
    # ============================================================
    with st.expander("Saved Proposals", expanded=False):
        st.caption("Save your current proposal or load a previously saved one")

        # Load all saved proposals (cached to reduce API calls)
        refresh_counter = st.session_state.get('saved_data_refresh_counter', 0)
        saved_proposals = cached_load_all_proposals(refresh_counter)

        # Two columns: Load/Delete on left, Save on right
        load_col, save_col = st.columns(2)

        with load_col:
            st.markdown("**Load Proposal**")

            if len(saved_proposals) == 0:
                st.info("No saved proposals yet")
            else:
                # Create dropdown options
                proposal_options = {
                    f"{p['name']} ({p['created_date'][:10]})": p['proposal_id']
                    for p in saved_proposals
                }

                selected_proposal_label = st.selectbox(
                    "Select proposal to load:",
                    options=list(proposal_options.keys()),
                    key="load_proposal_select"
                )

                if selected_proposal_label:
                    selected_proposal_id = proposal_options[selected_proposal_label]

                    # Find full proposal data
                    selected_proposal = next(p for p in saved_proposals if p['proposal_id'] == selected_proposal_id)

                    # Show preview
                    st.caption(f"**Created by:** {selected_proposal['created_by'] or 'Unknown'}")
                    st.caption(f"**Dataset:** {selected_proposal['dataset']}")

                    # Load and Delete buttons
                    load_btn_col, delete_btn_col = st.columns(2)

                    with load_btn_col:
                        if st.button("Load", key="load_proposal_btn", type="primary", use_container_width=True):
                            success, proposal_data, dataset = load_proposal_data(selected_proposal_id)

                            if success:
                                # Check if dataset matches
                                if dataset != st.session_state.selected_dataset:
                                    st.warning(f"WARNING: This proposal was created with {dataset} dataset, but you're currently using {st.session_state.selected_dataset} dataset. Loading anyway...")

                                # Load proposal data into session state
                                st.session_state.proposal_products = proposal_data.get('proposal_products', [])
                                st.session_state.proposal_marketing_rounding = proposal_data.get('proposal_marketing_rounding', False)
                                st.session_state.proposal_use_msrp = proposal_data.get('proposal_use_msrp', True)
                                st.session_state.proposal_discount_type = proposal_data.get('proposal_discount_type', None)
                                st.session_state.proposal_discount_percent = proposal_data.get('proposal_discount_percent', 0.0)
                                st.session_state.proposal_client_budget = proposal_data.get('proposal_client_budget', 0.0)

                                st.success(f"Loaded proposal: {selected_proposal['name']}")
                                st.rerun()
                            else:
                                st.error("Failed to load proposal data")

                    with delete_btn_col:
                        if st.button("Delete", key="delete_proposal_btn", use_container_width=True):
                            # Confirmation dialog using session state
                            if 'confirm_delete_proposal_id' not in st.session_state:
                                st.session_state.confirm_delete_proposal_id = selected_proposal_id
                                st.warning(f"WARNING: Are you sure you want to delete '{selected_proposal['name']}'?")

                                confirm_col1, confirm_col2 = st.columns(2)
                                with confirm_col1:
                                    if st.button("Yes, Delete", key="confirm_delete_yes", type="primary"):
                                        success, message = delete_proposal(st.session_state.confirm_delete_proposal_id)
                                        if success:
                                            clear_saved_data_cache()  # Clear cache after deletion
                                            st.success(message)
                                            del st.session_state.confirm_delete_proposal_id
                                            st.rerun()
                                        else:
                                            st.error(message)
                                with confirm_col2:
                                    if st.button("Cancel", key="confirm_delete_no"):
                                        del st.session_state.confirm_delete_proposal_id
                                        st.rerun()

        with save_col:
            st.markdown("**Save Current Proposal**")

            # Check if there are products to save
            has_products = len(st.session_state.proposal_products) > 0

            if not has_products:
                st.info("Add products to enable saving")

            proposal_name = st.text_input(
                "Proposal name:",
                key="save_proposal_name",
                placeholder="e.g., Client ABC Winter Campaign",
                disabled=not has_products
            )

            created_by = st.text_input(
                "Your name (optional):",
                key="save_proposal_creator",
                placeholder="e.g., John Smith",
                disabled=not has_products
            )

            if st.button("Save Proposal", key="save_proposal_btn", type="primary", use_container_width=True, disabled=not has_products):
                if not proposal_name or not proposal_name.strip():
                    st.error("Please enter a proposal name")
                else:
                    # Prepare proposal data
                    proposal_data = {
                        'proposal_products': st.session_state.proposal_products,
                        'proposal_marketing_rounding': st.session_state.proposal_marketing_rounding,
                        'proposal_use_msrp': st.session_state.proposal_use_msrp,
                        'proposal_discount_type': st.session_state.get('proposal_discount_type'),
                        'proposal_discount_percent': st.session_state.get('proposal_discount_percent', 0.0),
                        'proposal_client_budget': st.session_state.get('proposal_client_budget', 0.0)
                    }

                    success, message, result = save_proposal(
                        name=proposal_name.strip(),
                        created_by=created_by.strip() if created_by else "",
                        proposal_data=proposal_data,
                        dataset=st.session_state.selected_dataset
                    )

                    if success:
                        update_last_save_time('proposal')
                        clear_saved_data_cache()  # Clear cache to show new proposal
                        # Track that this proposal was saved
                        st.session_state.loaded_proposal_name = proposal_name.strip()
                        st.session_state.loaded_proposal_date = datetime.now().strftime("%Y-%m-%d %H:%M")
                        st.session_state.loaded_proposal_creator = created_by.strip() if created_by else "Current User"
                        st.success(message)
                        st.rerun()  # Rerun to clear form
                    else:
                        # Check if it's a naming conflict
                        if result:  # result contains suggested name
                            st.error(message)
                            # Offer to save with suggested name
                            if st.button(f"Save as '{result}'", key="save_with_new_name"):
                                success2, message2, _ = save_proposal(
                                    name=result,
                                    created_by=created_by.strip() if created_by else "",
                                    proposal_data=proposal_data,
                                    dataset=st.session_state.selected_dataset
                                )
                                if success2:
                                    update_last_save_time('proposal')
                                    clear_saved_data_cache()  # Clear cache to show new proposal
                                    # Track that this proposal was saved with new name
                                    st.session_state.loaded_proposal_name = result
                                    st.session_state.loaded_proposal_date = datetime.now().strftime("%Y-%m-%d %H:%M")
                                    st.session_state.loaded_proposal_creator = created_by.strip() if created_by else "Current User"
                                    st.success(message2)
                                    st.rerun()
                        else:
                            st.error(message)

    st.divider()

    # Product count and status
    if len(st.session_state.proposal_products) == 0:
        st.info("No products added to proposal yet. Add products from the catalog above.")
    else:
        st.success(f"{len(st.session_state.proposal_products)} product(s) in proposal")

        # Proposal Settings Section
        st.markdown("### Proposal Settings")

        col1, col2, col3 = st.columns(3)

        with col1:
            # Client Budget filter for volume pricing calculations
            client_budget = st.number_input(
                "Client Budget ($)",
                min_value=0.0,
                value=st.session_state.get('proposal_client_budget', 0.0),
                step=100.0,
                help="Total potential spend by client. Used to calculate volume pricing if budget allows higher quantities.",
                key="proposal_client_budget_input"
            )
            st.session_state.proposal_client_budget = client_budget

        with col2:
            # Discount options
            discount_type = st.selectbox(
                "Client Discount",
                options=["None", "Non-profit (5%)", "Volume Order (5%)", "Custom"],
                index=0 if not st.session_state.get('proposal_discount_type') else
                      (1 if st.session_state.get('proposal_discount_type') == 'Non-profit' else
                       (2 if st.session_state.get('proposal_discount_type') == 'Volume Order' else 3)),
                key="proposal_discount_type_select"
            )

            if discount_type == "Non-profit (5%)":
                st.session_state.proposal_discount_type = 'Non-profit'
                st.session_state.proposal_discount_percent = 5.0
            elif discount_type == "Volume Order (5%)":
                st.session_state.proposal_discount_type = 'Volume Order'
                st.session_state.proposal_discount_percent = 5.0
            elif discount_type == "Custom":
                st.session_state.proposal_discount_type = 'Custom'
                custom_discount = st.number_input(
                    "Custom discount %",
                    min_value=0.0,
                    max_value=100.0,
                    value=st.session_state.get('proposal_discount_percent', 0.0),
                    step=0.5,
                    key="proposal_custom_discount"
                )
                st.session_state.proposal_discount_percent = custom_discount
            else:
                st.session_state.proposal_discount_type = None
                st.session_state.proposal_discount_percent = 0.0

        with col3:
            # Use MSRP pricing checkbox (defaults to checked)
            st.session_state.proposal_use_msrp = st.checkbox(
                "Use MSRP pricing when available",
                value=st.session_state.proposal_use_msrp,
                key="proposal_use_msrp_checkbox",
                help="When enabled, products with MSRP will have markup automatically calculated to match MSRP. Products without MSRP will use 100% markup."
            )

            # $0.50 rounding (new feature, defaults to checked)
            st.session_state.proposal_fifty_cent_rounding = st.checkbox(
                "Round prices to nearest $0.50",
                value=st.session_state.proposal_fifty_cent_rounding,
                key="proposal_fifty_cent_rounding_checkbox",
                help="Rounds all prices to nearest 50 cents (e.g., $24.37 → $24.50)"
            )

            # Marketing rounding
            st.session_state.proposal_marketing_rounding = st.checkbox(
                "Apply marketing rounding (e.g., $60 → $59)",
                value=st.session_state.proposal_marketing_rounding,
                key="proposal_marketing_rounding_checkbox",
                help="Apply charm pricing to prices ending in 0 (e.g., $60 → $59). Applied after $0.50 rounding."
            )

        st.divider()

        # Product table with pricing details and editable markup
        st.markdown("### Products in Proposal")

        # Table header
        header_col1, header_col2, header_col3, header_col4, header_col5, header_col6, header_col7, header_col8 = st.columns([3, 1.2, 1.2, 1.5, 1.2, 1.2, 1.0, 0.8])
        with header_col1:
            st.markdown("**Product**")
        with header_col2:
            st.markdown("**PBP Cost**")
        with header_col3:
            st.markdown("**Vendor MSRP**")
        with header_col4:
            st.markdown("**Pricing Method**")
        with header_col5:
            st.markdown("**PBP MSRP**")
        with header_col6:
            st.markdown("**Markup %**")
        with header_col7:
            st.markdown("**Override**")
        with header_col8:
            st.markdown("**Remove**")

        st.divider()

        # Display each product in table format
        for idx, item in enumerate(st.session_state.proposal_products):
            product_data = item['product_data']

            # Calculate PBP cost and client price at quantity 100 for display
            base_cost, _, _ = get_unit_price_new_system(product_data, 100)

            # Calculate client price using new pricing logic
            manual_override = item.get('manual_override', False)
            pricing_method = item.get('pricing_method', 'Standard markup')

            if manual_override:
                # Manual override: use simple markup calculation
                if base_cost:
                    client_price = base_cost * (1 + item['markup_percent'] / 100)
                else:
                    client_price = None
            else:
                # Use pricing method from product
                pricing_result = calculate_pbp_msrp(product_data, 100)
                client_price = pricing_result['pbp_msrp'] if pricing_result else None

            # Apply rounding if client price calculated
            unrounded_price = client_price  # Store unrounded price for comparison
            if client_price:
                # Apply $0.50 rounding if enabled
                client_price = round_to_nearest_fifty_cents(
                    client_price,
                    st.session_state.proposal_fifty_cent_rounding
                )

                # Apply marketing rounding if enabled (after $0.50 rounding)
                client_price = apply_marketing_rounding(
                    client_price,
                    st.session_state.proposal_marketing_rounding
                )

            col1, col2, col3, col4, col5, col6, col7, col8 = st.columns([3, 1.2, 1.2, 1.5, 1.2, 1.2, 1.0, 0.8])

            with col1:
                # Display product name
                product_display_name = product_data['Product/Service']
                st.markdown(f"{product_display_name}")
                st.caption(f"Partner: {product_data['Partner']}")

            with col2:
                # Show PBP cost
                if base_cost:
                    st.markdown(f"${base_cost:.2f}")
                else:
                    st.markdown("—")

            with col3:
                # Show Vendor MSRP if available
                msrp_raw = get_column_value(product_data, 'Vendor Published MSRP', 'MSRP', '')
                msrp = clean_price(msrp_raw)
                if msrp and msrp > 0:
                    st.markdown(f"${msrp:.2f}")
                else:
                    st.markdown("—")

            with col4:
                # Show pricing method
                pricing_method = item.get('pricing_method', 'Standard markup')
                manual_override = item.get('manual_override', False)
                if manual_override:
                    st.markdown("Manual override")
                elif pricing_method:
                    st.markdown(f"{pricing_method}")
                else:
                    st.markdown("Standard markup")

            with col5:
                # Editable client price field (PBP MSRP)
                if base_cost and client_price:
                    # Check if we're updating from markup to prevent circular updates
                    if st.session_state.get(f'updating_from_markup_{idx}', False):
                        # Just display the calculated price, don't create input
                        st.markdown(f"${client_price:.2f}")
                        # Show rounding note if price was rounded
                        if unrounded_price and abs(client_price - unrounded_price) > 0.01:
                            st.caption(f"Rounded from ${unrounded_price:.2f}")
                        # Clear the flag
                        st.session_state[f'updating_from_markup_{idx}'] = False
                    else:
                        new_price = st.number_input(
                            f"Price for {idx}",
                            min_value=0.01,
                            value=client_price,
                            step=1.0,
                            format="%.2f",
                            key=f"price_{idx}",
                            label_visibility="collapsed"
                        )
                        # Show rounding note if price was rounded
                        if unrounded_price and abs(client_price - unrounded_price) > 0.01:
                            st.caption(f"Rounded from ${unrounded_price:.2f}")
                        # Update markup if price changed
                        if abs(new_price - client_price) > 0.01:  # Check if meaningfully different
                            # Calculate new markup from the price
                            new_markup_calc = calculate_markup_from_price(base_cost, new_price)
                            st.session_state.proposal_products[idx]['markup_percent'] = new_markup_calc
                            # Mark that pricing state has pending changes
                            st.session_state.proposal_pricing_pending_finalization = True
                            # Store current settings with this product
                            st.session_state.proposal_products[idx]['settings_snapshot'] = {
                                'fifty_cent_rounding': st.session_state.get('proposal_fifty_cent_rounding', False),
                                'marketing_rounding': st.session_state.get('proposal_marketing_rounding', False),
                                'discount_percent': st.session_state.get('proposal_discount_percent', 0.0)
                            }
                            st.rerun()
                else:
                    st.markdown("—")

            with col6:
                # Editable markup field
                new_markup = st.number_input(
                    f"Markup for {idx}",
                    min_value=-50.0,  # Allow negative markup for below-cost pricing
                    value=item['markup_percent'],
                    step=5.0,
                    key=f"markup_{idx}",
                    label_visibility="collapsed"
                )
                # Update markup if changed
                if new_markup != item['markup_percent']:
                    st.session_state.proposal_products[idx]['markup_percent'] = new_markup
                    # Set flag to prevent circular updates
                    st.session_state[f'updating_from_markup_{idx}'] = True
                    # Mark that pricing state has pending changes
                    st.session_state.proposal_pricing_pending_finalization = True
                    # Store current settings with this product
                    st.session_state.proposal_products[idx]['settings_snapshot'] = {
                        'fifty_cent_rounding': st.session_state.get('proposal_fifty_cent_rounding', False),
                        'marketing_rounding': st.session_state.get('proposal_marketing_rounding', False),
                        'discount_percent': st.session_state.get('proposal_discount_percent', 0.0)
                    }

            with col7:
                # Manual Override Checkbox
                manual_override = item.get('manual_override', False)
                override_checked = st.checkbox(
                    "Override",
                    value=manual_override,
                    key=f"override_{idx}",
                    help="Check to manually override pricing method",
                    label_visibility="collapsed"
                )
                if override_checked != manual_override:
                    st.session_state.proposal_products[idx]['manual_override'] = override_checked
                    st.rerun()

            with col8:
                if st.button("✕", key=f"remove_{idx}", help=f"Remove {product_data['Product/Service']}", use_container_width=True):
                    st.session_state.proposal_products.pop(idx)
                    st.rerun()

            st.divider()

        # ============================================================
        # PRICING NOTES DISPLAY (Expandable)
        # ============================================================
        # Collect all products with pricing notes (and no manual override)
        products_with_notes = [
            item for item in st.session_state.proposal_products
            if item.get('pricing_notes') and not item.get('manual_override', False)
        ]

        if products_with_notes:
            with st.expander(f"Pricing Information ({len(products_with_notes)} product{'s' if len(products_with_notes) != 1 else ''})", expanded=False):
                for item in products_with_notes:
                    product_display_name = item['product_data']['Product/Service']
                    st.write(f"**{product_display_name}**")
                    st.caption(item['pricing_notes'])
                    st.write("")  # Spacing

        # ============================================================
        # VALIDATION WARNINGS DISPLAY (Expandable)
        # ============================================================
        # Collect products with validation warnings (and no manual override)
        products_with_warnings = [
            item for item in st.session_state.proposal_products
            if item.get('validation_warning') and not item.get('manual_override', False)
        ]

        if products_with_warnings:
            st.warning(f"{len(products_with_warnings)} product{'s' if len(products_with_warnings) != 1 else ''} {'have' if len(products_with_warnings) != 1 else 'has'} pricing discrepancies")

            with st.expander("View Validation Details", expanded=False):
                for item in products_with_warnings:
                    product_display_name = item['product_data']['Product/Service']
                    st.write(f"**{product_display_name}**")
                    st.caption(item['validation_warning'])
                    st.write("")  # Spacing

    # ============================================================
    # SECTION 3: GENERATE PROPOSAL TABLES
    # ============================================================
    st.divider()
    st.subheader("3. Generate Proposal Tables")

    if len(st.session_state.proposal_products) == 0:
        st.caption("Add products to generate proposal tables")
    else:
        # Default to collapsed to save space, user can expand to view
        with st.expander(f"View Proposal Tables ({len(st.session_state.proposal_products)} products)", expanded=False):
            st.markdown("Each product is presented in a separate table with MOQ pricing.")
            st.markdown("")

            # Generate a separate table for each product
            for idx, item in enumerate(st.session_state.proposal_products, 1):
                st.markdown(f"### Product {idx}: {item['product_data']['Product/Service']}")

                product_row = item['product_data']

                # Calculate MOQ using a standard preliminary quantity (100 units)
                preliminary_base_price, _, _ = get_unit_price_new_system(product_row, 100)

                if preliminary_base_price is not None:
                    # Estimate total per-unit price with markup (no customization in MOQ calc)
                    temp_markup_multiplier = 1 + (item['markup_percent'] / 100)
                    estimated_unit_price = preliminary_base_price * temp_markup_multiplier

                    # Calculate MOQ (returns dict with moq, breakdown, display_text)
                    moq_result = calculate_moq(estimated_unit_price, product_row)
                    moq = moq_result['moq'] if moq_result else None
                    moq_display_text = moq_result['display_text'] if moq_result else ""
                    if moq is None:
                        moq = 5

                    # Get actual base price for MOQ quantity
                    moq_base_price, moq_tier_range, _ = get_unit_price_new_system(product_row, moq)

                    if moq_base_price is not None:
                        # Use new pricing calculation (respects manual override if set)
                        if item.get('manual_override', False):
                            # Manual override: use stored markup percentage
                            moq_product_cost = moq_base_price * moq
                            moq_markup_amount = moq_product_cost * (item['markup_percent'] / 100)
                            moq_product_only_total = moq_product_cost + moq_markup_amount
                            moq_product_price_per_unit = moq_product_only_total / moq
                        else:
                            # Use pricing method from spreadsheet
                            pricing_result = calculate_pbp_msrp(
                                product_row,
                                quantity=moq,
                                user_markup_override=None
                            )
                            moq_product_price_per_unit = pricing_result['pbp_msrp']

                        # Apply $0.50 rounding if enabled
                        moq_product_price_per_unit = round_to_nearest_fifty_cents(
                            moq_product_price_per_unit,
                            st.session_state.proposal_fifty_cent_rounding
                        )

                        # Apply marketing rounding if enabled (after $0.50 rounding)
                        moq_product_price_per_unit = apply_marketing_rounding(
                            moq_product_price_per_unit,
                            st.session_state.proposal_marketing_rounding
                        )

                        # Calculate Client Price based on discount and budget
                        client_price = moq_product_price_per_unit
                        client_price_note = ""

                        # Get client budget and discount settings
                        client_budget = st.session_state.get('proposal_client_budget', 0.0)
                        discount_percent = st.session_state.get('proposal_discount_percent', 0.0)

                        # Check if client budget allows for higher quantity (better pricing)
                        volume_pricing_applied = False
                        volume_pricing_quantity = None
                        if client_budget > 0:
                            moq_total = moq * moq_product_price_per_unit
                            if client_budget > moq_total:
                                # Calculate what quantity the client could afford at MOQ price
                                potential_quantity = int(client_budget / moq_product_price_per_unit)

                                # Get price at that higher quantity
                                budget_qty_base_price, _, _ = get_unit_price_new_system(product_row, potential_quantity)

                                if budget_qty_base_price is not None:
                                    # Use new pricing calculation (respects manual override if set)
                                    if item.get('manual_override', False):
                                        # Manual override: use stored markup percentage
                                        budget_qty_product_cost = budget_qty_base_price * potential_quantity
                                        budget_qty_markup_amount = budget_qty_product_cost * (item['markup_percent'] / 100)
                                        budget_qty_product_only_total = budget_qty_product_cost + budget_qty_markup_amount
                                        budget_qty_price_per_unit = budget_qty_product_only_total / potential_quantity
                                    else:
                                        # Use pricing method from spreadsheet
                                        pricing_result = calculate_pbp_msrp(
                                            product_row,
                                            quantity=potential_quantity,
                                            user_markup_override=None
                                        )
                                        budget_qty_price_per_unit = pricing_result['pbp_msrp']

                                    # Apply $0.50 rounding if enabled
                                    budget_qty_price_per_unit = round_to_nearest_fifty_cents(
                                        budget_qty_price_per_unit,
                                        st.session_state.proposal_fifty_cent_rounding
                                    )

                                    # Apply marketing rounding if enabled (after $0.50 rounding)
                                    budget_qty_price_per_unit = apply_marketing_rounding(
                                        budget_qty_price_per_unit,
                                        st.session_state.proposal_marketing_rounding
                                    )

                                    # Use the better price if different from MOQ price
                                    if budget_qty_price_per_unit < moq_product_price_per_unit:
                                        client_price = budget_qty_price_per_unit
                                        volume_pricing_applied = True
                                        volume_pricing_quantity = potential_quantity

                        # Apply discount to client price
                        discount_applied = False
                        if discount_percent > 0:
                            client_price = client_price * (1 - discount_percent / 100)
                            discount_applied = True

                            # Apply $0.50 rounding after discount if enabled
                            client_price = round_to_nearest_fifty_cents(
                                client_price,
                                st.session_state.proposal_fifty_cent_rounding
                            )

                            # Apply marketing rounding again after discount if enabled
                            client_price = apply_marketing_rounding(
                                client_price,
                                st.session_state.proposal_marketing_rounding
                            )

                        # Build price note for column header
                        client_price_header = "Client Price"
                        if discount_applied or volume_pricing_applied:
                            notes = []
                            if volume_pricing_applied:
                                notes.append(f"Price ea @ Qty {volume_pricing_quantity}")
                            if discount_applied:
                                discount_type = st.session_state.get('proposal_discount_type')
                                if discount_type == 'Non-profit':
                                    notes.append("5% Non-profit discount")
                                elif discount_type == 'Volume Order':
                                    notes.append("5% Volume Order discount")
                                else:
                                    notes.append(f"{discount_percent:.1f}% discount")
                            client_price_header = f"Client Price ({', '.join(notes)})"

                        # Build proposal table
                        col_moq = "MOQ"
                        col_price = f"Price Ea (@ Qty {moq})"
                        col_client_price = client_price_header
                        col_delivery = "Delivery"

                        proposal_table = pd.DataFrame([{
                            col_moq: moq,
                            col_price: f"${moq_product_price_per_unit:.2f}",
                            col_client_price: f"${client_price:.2f}",
                            col_delivery: ""
                        }])

                        st.table(proposal_table)

                        # Show MOQ breakdown
                        st.caption(moq_display_text)

                        # ALWAYS show customization costs from product data
                        # Get customization costs from the product data
                        setup_fee_raw = get_column_value(product_row, 'Client Price: Customization Setup Fee', 'Customization Setup Fee', '')
                        per_unit_raw = get_column_value(product_row, 'Client Price: Customization Cost per Unit', 'Customization Cost per Unit', '')
                        setup_fee = clean_price(setup_fee_raw) or 0.0
                        per_unit_cost = clean_price(per_unit_raw) or 0.0

                        # Display customization costs at the bottom
                        if setup_fee > 0 or per_unit_cost > 0:
                            st.caption(f"**Customization available:** Artwork set-up: \\${setup_fee:.2f} / Branding per piece: \\${per_unit_cost:.2f}")
                        else:
                            st.caption("**Customization available:** Contact for pricing")

                        # Add download button for this product's proposal table
                        proposal_csv = proposal_table.to_csv(index=False)
                        st.download_button(
                            label=f"Download Product {idx} Proposal (CSV)",
                            data=proposal_csv,
                            file_name=f"proposal_product_{idx}_{item['product_data']['Product/Service'].replace(' ', '_')}.csv",
                            mime="text/csv",
                            key=f"download_proposal_{idx}"
                        )
                    else:
                        st.warning(f"Unable to calculate MOQ pricing for {item['product_data']['Product/Service']}")
                else:
                    st.warning(f"Product data not available for {item['product_data']['Product/Service']}")

            st.markdown("")

        st.caption("Copy these tables and paste into your proposal template.")

        # Download all proposals as CSV
        st.markdown("---")
        if st.button("Download All Proposal Tables (CSV)", use_container_width=True, type="primary"):
            # Generate comprehensive CSV matching UI display
            csv_lines = []
            csv_lines.append("PEACE BY PIECE - PRODUCT PROPOSAL")
            csv_lines.append(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
            csv_lines.append("")

            for idx, item in enumerate(st.session_state.proposal_products, 1):
                product_row = item.get('product_data', {})

                csv_lines.append(f"=== PRODUCT {idx}: {product_row.get('Product/Service', 'Unknown Product')} ===")
                csv_lines.append(f"Partner: {product_row.get('Partner', 'N/A')}")
                csv_lines.append(f"Country of Origin (Made In): {product_row.get('Country of Origin (Made In)', 'N/A')}")
                csv_lines.append(f"Country of Origin (Ships From): {product_row.get('Country of Origin (Ships From)', 'N/A')}")
                csv_lines.append("")

                # Calculate MOQ using same logic as UI display
                preliminary_base_price, _, _ = get_unit_price_new_system(product_row, 100)

                if preliminary_base_price is not None:
                    # Estimate total per-unit price with markup (no customization in MOQ calc)
                    temp_markup_multiplier = 1 + (item['markup_percent'] / 100)
                    estimated_unit_price = preliminary_base_price * temp_markup_multiplier

                    # Calculate MOQ (returns dict with moq, breakdown, display_text)
                    moq_result = calculate_moq(estimated_unit_price, product_row)
                    moq = moq_result['moq'] if moq_result else None
                    if moq is None:
                        moq = 5

                    # Get actual base price for MOQ quantity
                    moq_base_price, moq_tier_range, _ = get_unit_price_new_system(product_row, moq)

                    if moq_base_price is not None:
                        # Use new pricing calculation (respects manual override if set)
                        if item.get('manual_override', False):
                            # Manual override: use stored markup percentage
                            moq_product_cost = moq_base_price * moq
                            moq_markup_amount = moq_product_cost * (item['markup_percent'] / 100)
                            moq_product_only_total = moq_product_cost + moq_markup_amount
                            moq_product_price_per_unit = moq_product_only_total / moq
                        else:
                            # Use pricing method from spreadsheet
                            pricing_result = calculate_pbp_msrp(
                                product_row,
                                quantity=moq,
                                user_markup_override=None
                            )
                            moq_product_price_per_unit = pricing_result['pbp_msrp']

                        # Apply $0.50 rounding if enabled
                        moq_product_price_per_unit = round_to_nearest_fifty_cents(
                            moq_product_price_per_unit,
                            st.session_state.proposal_fifty_cent_rounding
                        )

                        # Apply marketing rounding if enabled (after $0.50 rounding)
                        moq_product_price_per_unit = apply_marketing_rounding(
                            moq_product_price_per_unit,
                            st.session_state.proposal_marketing_rounding
                        )

                        # Calculate Client Price based on discount and budget
                        client_price = moq_product_price_per_unit

                        # Get client budget and discount settings
                        client_budget = st.session_state.get('proposal_client_budget', 0.0)
                        discount_percent = st.session_state.get('proposal_discount_percent', 0.0)

                        # Check if client budget allows for higher quantity (better pricing)
                        volume_pricing_applied = False
                        volume_pricing_quantity = None
                        if client_budget > 0:
                            moq_total = moq * moq_product_price_per_unit
                            if client_budget > moq_total:
                                # Calculate what quantity the client could afford at MOQ price
                                potential_quantity = int(client_budget / moq_product_price_per_unit)

                                # Get price at that higher quantity
                                budget_qty_base_price, _, _ = get_unit_price_new_system(product_row, potential_quantity)

                                if budget_qty_base_price is not None:
                                    # Use new pricing calculation (respects manual override if set)
                                    if item.get('manual_override', False):
                                        # Manual override: use stored markup percentage
                                        budget_qty_product_cost = budget_qty_base_price * potential_quantity
                                        budget_qty_markup_amount = budget_qty_product_cost * (item['markup_percent'] / 100)
                                        budget_qty_product_only_total = budget_qty_product_cost + budget_qty_markup_amount
                                        budget_qty_price_per_unit = budget_qty_product_only_total / potential_quantity
                                    else:
                                        # Use pricing method from spreadsheet
                                        pricing_result = calculate_pbp_msrp(
                                            product_row,
                                            quantity=potential_quantity,
                                            user_markup_override=None
                                        )
                                        budget_qty_price_per_unit = pricing_result['pbp_msrp']

                                    # Apply $0.50 rounding if enabled
                                    budget_qty_price_per_unit = round_to_nearest_fifty_cents(
                                        budget_qty_price_per_unit,
                                        st.session_state.proposal_fifty_cent_rounding
                                    )

                                    # Apply marketing rounding if enabled (after $0.50 rounding)
                                    budget_qty_price_per_unit = apply_marketing_rounding(
                                        budget_qty_price_per_unit,
                                        st.session_state.proposal_marketing_rounding
                                    )

                                    # Use the better price if different from MOQ price
                                    if budget_qty_price_per_unit < moq_product_price_per_unit:
                                        client_price = budget_qty_price_per_unit
                                        volume_pricing_applied = True
                                        volume_pricing_quantity = potential_quantity

                        # Apply discount to client price
                        discount_applied = False
                        if discount_percent > 0:
                            client_price = client_price * (1 - discount_percent / 100)
                            discount_applied = True

                            # Apply $0.50 rounding after discount if enabled
                            client_price = round_to_nearest_fifty_cents(
                                client_price,
                                st.session_state.proposal_fifty_cent_rounding
                            )

                            # Apply marketing rounding again after discount if enabled
                            client_price = apply_marketing_rounding(
                                client_price,
                                st.session_state.proposal_marketing_rounding
                            )

                        # Build price note for column header
                        client_price_header = "Client Price"
                        if discount_applied or volume_pricing_applied:
                            notes = []
                            if volume_pricing_applied:
                                notes.append(f"Price ea @ Qty {volume_pricing_quantity}")
                            if discount_applied:
                                discount_type = st.session_state.get('proposal_discount_type')
                                if discount_type == 'Non-profit':
                                    notes.append("5% Non-profit discount")
                                elif discount_type == 'Volume Order':
                                    notes.append("5% Volume Order discount")
                                else:
                                    notes.append(f"{discount_percent:.1f}% discount")
                            client_price_header = f"Client Price ({', '.join(notes)})"

                        # Build CSV table matching UI display
                        csv_lines.append(f"MOQ,Price Ea (@ Qty {moq}),{client_price_header},Delivery")
                        csv_lines.append(f"{moq},${moq_product_price_per_unit:.2f},${client_price:.2f},")

                        # Show MOQ calculation note
                        moq_total_value = moq * moq_product_price_per_unit
                        csv_lines.append("")
                        csv_lines.append(f"MOQ calculated based on $1,000 minimum order value (MOQ {moq} units = ${moq_total_value:.2f})")

                        # ALWAYS show customization costs from product data
                        setup_fee_raw = get_column_value(product_row, 'Client Price: Customization Setup Fee', 'Customization Setup Fee', '')
                        per_unit_raw = get_column_value(product_row, 'Client Price: Customization Cost per Unit', 'Customization Cost per Unit', '')
                        setup_fee = clean_price(setup_fee_raw) or 0.0
                        per_unit_cost = clean_price(per_unit_raw) or 0.0

                        # Display customization costs
                        if setup_fee > 0 or per_unit_cost > 0:
                            csv_lines.append(f"Customization available: Artwork set-up: ${setup_fee:.2f} / Branding per piece: ${per_unit_cost:.2f}")
                        else:
                            csv_lines.append("Customization available: Contact for pricing")
                    else:
                        csv_lines.append(f"Unable to calculate MOQ pricing for {product_row.get('Product/Service', 'Unknown Product')}")
                else:
                    csv_lines.append(f"Product data not available for {product_row.get('Product/Service', 'Unknown Product')}")

                csv_lines.append("")
                csv_lines.append("")

            # Pricing for Cards & Kitting
            csv_lines.append("=== PRICING FOR CARDS & KITTING ===")
            csv_lines.append(st.session_state.proposal_kitting_pricing)
            csv_lines.append("")
            csv_lines.append("")

            # Terms & Conditions
            csv_lines.append("=== TERMS & CONDITIONS ===")
            csv_lines.append(st.session_state.proposal_terms)

            proposal_csv = "\n".join(csv_lines)

            st.download_button(
                label="Click to Download CSV",
                data=proposal_csv,
                file_name=f"proposal_tables_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                mime="text/csv",
                key="download_all_proposals_csv",
                use_container_width=True
            )

    # ============================================================
    # LEGACY SECTIONS (HIDDEN BY DEFAULT)
    # ============================================================
    # Toggle for showing legacy sections (pricing tables, copy buttons, etc.)
    SHOW_LEGACY_SECTIONS = st.session_state.get('show_legacy_pricing_sections', False)

    if SHOW_LEGACY_SECTIONS:
        # ============================================================
        # LEGACY SECTION 5: PRICING FOR CARDS & KITTING
        # ============================================================
        st.divider()
        st.subheader("(Legacy) Pricing for Cards & Kitting")

        st.session_state.proposal_kitting_pricing = st.text_area(
            "Edit kitting pricing if needed",
            value=st.session_state.proposal_kitting_pricing,
            height=200,
            key="proposal_kitting_pricing_input"
        )

        # Add copy button for kitting pricing
        if st.button("Copy Pricing for Cards & Kitting", key="copy_kitting_pricing", use_container_width=True):
            st.code(st.session_state.proposal_kitting_pricing, language=None)
            st.info("Select the text above and copy it (Ctrl+C or Cmd+C)")

        # ============================================================
        # LEGACY SECTION 6: TERMS & CONDITIONS
        # ============================================================
        st.divider()
        st.subheader("(Legacy) Terms & Conditions")

        st.session_state.proposal_terms = st.text_area(
            "Edit terms & conditions if needed",
            value=st.session_state.proposal_terms,
            height=200,
            key="proposal_terms_input"
        )

        # Add copy button for terms & conditions
        if st.button("Copy Terms & Conditions", key="copy_terms", use_container_width=True):
            st.code(st.session_state.proposal_terms, language=None)
            st.info("Select the text above and copy it (Ctrl+C or Cmd+C)")


    # ============================================================
    # SECTION 4: POWERPOINT PROPOSAL GENERATION (PHASE 2.5 COMPLETE)
    # ============================================================
    if len(st.session_state.proposal_products) > 0:
        st.divider()
        st.subheader("4. Generate PowerPoint Proposal")
        st.caption("Automatically create a customized PowerPoint presentation with matched product slides")

        # Add JavaScript to capture scroll position before PowerPoint section button clicks
        components.html("""
            <script>
                const buttons = window.parent.document.querySelectorAll('button');
                buttons.forEach(button => {
                    const btnText = button.textContent;
                    if (btnText.includes('Save Manual Match') ||
                        btnText.includes('Delete') ||
                        btnText.includes('Review Matches & Generate PowerPoint') ||
                        btnText.includes('Yes, use this slide') ||
                        btnText.includes('Show alternatives') ||
                        btnText.includes('Skip this product') ||
                        btnText.includes('Use this') ||
                        btnText.includes('Override') ||
                        btnText.includes('Apply') ||
                        btnText.includes('Cancel') ||
                        btnText.includes('Reset Confirmations') ||
                        btnText.includes('Close')) {
                        button.addEventListener('click', function() {
                            const scrollPos = window.parent.document.querySelector('section.main').scrollTop;
                            window.parent.sessionStorage.setItem('streamlit_scroll_position', scrollPos);
                        });
                    }
                });
            </script>
        """, height=0)

        # Button to trigger matching
        if st.button("Review Matches & Generate PowerPoint", type="primary", use_container_width=True, key="trigger_pptx_matching"):
            st.session_state.show_pptx_matching = True
            st.session_state.generated_pptx = None  # Clear any previous generation
            st.session_state.match_confirmations = {}  # Clear stale confirmations from previous sessions
            st.session_state.pptx_match_results = None  # Force recalculation of matches
            st.rerun()

        # Show matching UI if triggered
        if st.session_state.get('show_pptx_matching', False):
            try:
                # Only load PowerPoint and run matching if not already cached
                if 'pptx_match_results' not in st.session_state or st.session_state.pptx_match_results is None:
                    # Load PowerPoint template from cloud/local
                    from pptx import Presentation

                    # Display template name
                    st.info(f"Using template: **{get_template_name('all_slides')}**")

                    pptx_template = get_template_path('all_slides', show_loading=True)

                    if not pptx_template:
                        st.error("PowerPoint template could not be loaded. Please check Google Drive access.")
                        st.session_state.show_pptx_matching = False
                    else:
                        # Extract product names from PowerPoint
                        with st.spinner("Analyzing PowerPoint slides..."):
                            prs = Presentation(pptx_template)

                            pptx_product_names = []
                            pptx_name_to_index = {}  # Map slide names to indices for saving confirmations
                            slide_list = list(prs.slides)

                            for slide_idx, slide in enumerate(slide_list):
                                if len(slide.shapes) >= 1:
                                    first_shape = slide.shapes[0]
                                    if hasattr(first_shape, "text") and first_shape.text.strip():
                                        product_name = first_shape.text.strip()
                                        if product_name not in pptx_product_names:
                                            pptx_product_names.append(product_name)
                                            pptx_name_to_index[product_name] = slide_idx

                        st.success(f"Loaded {len(pptx_product_names)} product slides from PowerPoint")

                        # Get proposal product names
                        gs_product_names = [item['product_data']['Product/Service'] for item in st.session_state.proposal_products]

                        # Create matcher and run matching (pass dataset for confirmed match lookup)
                        with st.spinner("Matching products to slides..."):
                            # Clear cache to ensure fresh confirmed matches are loaded
                            from src.match_memory import _load_all_matches_data
                            _load_all_matches_data.clear()

                            matcher = SlideMatcher(pptx_product_names)
                            current_dataset = st.session_state.get('selected_dataset', 'real')
                            match_results = matcher.batch_match(gs_product_names, dataset=current_dataset)

                        # Cache results in session state
                        st.session_state.pptx_match_results = match_results
                        st.session_state.pptx_product_names = pptx_product_names
                        st.session_state.pptx_name_to_index = pptx_name_to_index

                # Use cached match results
                if st.session_state.pptx_match_results is not None:
                    match_results = st.session_state.pptx_match_results
                    pptx_product_names = st.session_state.pptx_product_names
                    pptx_name_to_index = st.session_state.pptx_name_to_index

                    # Show match review UI (pass name-to-index mapping for saving confirmations)
                    confirmed_matches = show_match_review_ui(match_results, pptx_product_names, pptx_name_to_index)

                    if confirmed_matches:
                        st.session_state.show_pptx_matching = False
                        st.rerun()

            except Exception as e:
                st.error(f"Error loading PowerPoint: {str(e)}")
                st.session_state.show_pptx_matching = False

    # ============================================================
    # NEXT STEPS GUIDANCE
    # ============================================================
    st.divider()

    if len(st.session_state.proposal_products) > 0:
        st.success(f"""
        **What's Next?**

        1. Download and send the proposal to your client
        2. Move to **Tab 2: Client Order Form Generator** to create a client order form
        3. Your {len(st.session_state.proposal_products)} proposal product(s) will be available for import in Tab 2
        """)
    else:
        st.info("""
        **What's Next?**

        After adding products to your proposal, you can:
        - Download proposal tables and PowerPoint presentations
        - Send to your client for review
        - Move to **Tab 2** to generate a client order form
        """)

    # Save and Navigation buttons at bottom of Tab 1
    st.divider()
    st.markdown("### Save Your Work")

    # Show unsaved changes indicator and save status
    if has_unsaved_proposal_changes():
        st.warning("You have unsaved changes in your proposal")

    save_status = format_time_since_save('proposal')
    if save_status:
        st.caption(f"{save_status}")

    # Check if there are products to save
    has_products = len(st.session_state.proposal_products) > 0

    if has_products:
        col1, col2 = st.columns(2)
        with col1:
            proposal_name = st.text_input(
                "Proposal name:",
                key="save_proposal_name_bottom",
                placeholder="e.g., Client ABC Winter Campaign"
            )
        with col2:
            created_by = st.text_input(
                "Your name (optional):",
                key="save_proposal_creator_bottom",
                placeholder="e.g., John Smith"
            )

        if st.button("Save Proposal", type="primary", use_container_width=True, key="save_proposal_btn_bottom"):
            if not proposal_name or not proposal_name.strip():
                st.error("Please enter a proposal name")
            else:
                # Prepare proposal data
                proposal_data = {
                    'proposal_products': st.session_state.proposal_products,
                    'proposal_marketing_rounding': st.session_state.proposal_marketing_rounding,
                    'proposal_use_msrp': st.session_state.proposal_use_msrp,
                    'proposal_discount_type': st.session_state.get('proposal_discount_type'),
                    'proposal_discount_percent': st.session_state.get('proposal_discount_percent', 0.0),
                    'proposal_client_budget': st.session_state.get('proposal_client_budget', 0.0)
                }

                success, message, result = save_proposal(
                    name=proposal_name.strip(),
                    created_by=created_by.strip() if created_by else "",
                    proposal_data=proposal_data,
                    dataset=st.session_state.selected_dataset
                )

                if success:
                    st.success(f"{message} - You can find it in the sidebar under 'Saved Proposals'")
                    time.sleep(1.5)
                    st.rerun()
                else:
                    # Check if it's a naming conflict
                    if result:  # result contains suggested name
                        st.error(message)
                        if st.button(f"Save as '{result}'", key="save_with_new_name_bottom"):
                            success2, message2, _ = save_proposal(
                                name=result,
                                created_by=created_by.strip() if created_by else "",
                                proposal_data=proposal_data,
                                dataset=st.session_state.selected_dataset
                            )
                            if success2:
                                st.success(f"{message2} - You can find it in the sidebar")
                                time.sleep(1.5)
                                st.rerun()
                    else:
                        st.error(message)
    else:
        st.info("Add products to your proposal to enable saving")

    # Navigation button
    st.divider()
    col1, col2, col3 = st.columns([1, 1, 1])
    with col2:
        if st.button("Continue to Tab 2: Client Order Form", type="primary", use_container_width=True, key="tab1_to_tab2"):
            st.session_state.show_tab2_prompt = True
            st.rerun()

    # Show navigation prompt if button was clicked
    if st.session_state.get('show_tab2_prompt', False):
        st.info("Click on the **'Client Order Form Generator'** tab above to continue.")
        st.session_state.show_tab2_prompt = False

# ============================================================
# TAB 2: CLIENT ORDER FORM GENERATOR (NEW)
# ============================================================
with tab2:
    st.header("Client Order Form Generator")
    st.caption("Create professional HTML order forms to send to clients")
    st.divider()

    # ============================================================
    # SECTION 1: ORDER DETAILS
    # ============================================================
    st.subheader("1. Order Details")
    st.caption("Add client information to pre-fill the order form")

    # Initialize order details in session state if not exists
    if 'order_details' not in st.session_state:
        st.session_state.order_details = {
            'client_type': 'New',
            'company_name': '',
            'contact_name': '',
            'contact_email': ''
        }

    col1, col2 = st.columns(2)

    with col1:
        # Convert client_type to boolean for checkbox
        is_new = st.session_state.order_details.get('client_type', 'New') == 'New'
        st.session_state.order_details['is_new_client'] = st.checkbox(
            "New Client?",
            value=is_new,
            key="order_detail_is_new_client"
        )
        # Store as 'New' or 'Existing' for backward compatibility
        st.session_state.order_details['client_type'] = 'New' if st.session_state.order_details['is_new_client'] else 'Existing'

        st.session_state.order_details['company_name'] = st.text_input(
            "Company Name",
            value=st.session_state.order_details.get('company_name', ''),
            key="order_detail_company_name"
        )

    with col2:
        st.session_state.order_details['contact_name'] = st.text_input(
            "Contact Name",
            value=st.session_state.order_details.get('contact_name', ''),
            key="order_detail_contact_name"
        )

        st.session_state.order_details['contact_email'] = st.text_input(
            "Contact Email",
            value=st.session_state.order_details.get('contact_email', ''),
            key="order_detail_contact_email"
        )

    # ============================================================
    # FORM CUSTOMIZATION (OPTIONAL)
    # ============================================================
    st.divider()
    with st.expander("Customize Form Template Text (Optional)", expanded=False):
        st.caption("Edit any template text that appears in the client order form")

        # Dropdown to select which field to customize
        field_labels = {
            'form_instructions': 'How to Fill Out This Form (Instructions at top)',
            'dropshipping_instructions': 'Dropshipping Instructions',
            'dropshipping_placeholder': 'Dropshipping Information Placeholder',
            'shipping_address_placeholder': 'Shipping Address Placeholder',
            'billing_address_placeholder': 'Billing Address Placeholder',
            'customization_placeholder': 'Customization/Branding Placeholder',
            'impact_card_options': 'Impact Card Options',
            'payment_options': 'Payment Options'
        }

        selected_field = st.selectbox(
            "Select field to customize",
            options=list(field_labels.keys()),
            format_func=lambda x: field_labels[x],
            index=1,  # Default to dropshipping_instructions
            key="customization_field_selector"
        )

        # Text area for editing the selected field
        current_value = st.session_state.form_customizations.get(selected_field, '')

        # Determine height based on field type
        height = 150 if selected_field in ['form_instructions', 'impact_card_options', 'payment_options'] else 100

        new_value = st.text_area(
            f"Edit: {field_labels[selected_field]}",
            value=current_value,
            height=height,
            key=f"customize_{selected_field}",
            help=f"This text appears in the '{field_labels[selected_field]}' section of the order form"
        )

        # Update session state if value changed
        if new_value != current_value:
            st.session_state.form_customizations[selected_field] = new_value
            # Update legacy dropshipping_notes for compatibility
            if selected_field == 'dropshipping_instructions':
                st.session_state.dropshipping_notes = new_value

    # Button to confirm info is ready (visual feedback)
    st.divider()
    if st.button("Update Order Form with This Info", type="primary", use_container_width=True, key="update_order_form"):
        st.session_state.show_order_form_updated = True
        st.rerun()

    # Show success message if flag is set
    if st.session_state.get('show_order_form_updated', False):
        st.success("Order form updated! The information you entered will appear in the form below.")
        st.session_state.show_order_form_updated = False

    # ============================================================
    # SECTION 2: GENERATE GOOGLE FORM (RECOMMENDED)
    # ============================================================
    st.divider()
    st.subheader("2. Generate Google Form (Recommended)")
    st.caption("Pre-fill a Google Form with proposal products and client info, then send link to client")

    # Check if proposal products exist
    has_proposal_products = len(st.session_state.proposal_products) > 0

    if not has_proposal_products:
        st.info("**No proposal products found.** Create a proposal in Tab 1 first, then return here to generate a client form.")
    else:
        st.success(f"Found {len(st.session_state.proposal_products)} products from your proposal")

        # Import forms helper
        from src.forms_helper import generate_prefilled_form_url

        # Product selection
        st.markdown("#### Select Products for Client")
        st.caption("Choose which products from your proposal to include in the form")

        # Initialize selected products tracking
        if 'google_form_selected_products' not in st.session_state:
            st.session_state.google_form_selected_products = []

        selected_products = []

        # Show proposal products with checkboxes
        for idx, item in enumerate(st.session_state.proposal_products):
            product_name = item['product_data'].get('Product/Service', '')
            quantity = item.get('pricing_snapshot', {}).get('quantity', 100)  # Default to 100 from proposal

            col1, col2 = st.columns([4, 1])

            with col1:
                include = st.checkbox(
                    f"{product_name}",
                    value=True,  # Default checked
                    key=f"google_form_product_{idx}"
                )

            with col2:
                qty = st.number_input(
                    "Qty",
                    value=quantity,
                    min_value=1,
                    key=f"google_form_qty_{idx}",
                    label_visibility="collapsed"
                )

            if include:
                selected_products.append({
                    'name': product_name,
                    'quantity': qty,
                    'customization_notes': ''  # Client can add this in the form
                })

        if selected_products:
            st.info(f"**{len(selected_products)} product(s) selected** - these will be pre-filled in the Google Form")

            # Client info (from Section 1)
            client_info = {
                'client_type': st.session_state.order_details.get('client_type', 'New'),
                'company_name': st.session_state.order_details.get('company_name', ''),
                'contact_name': st.session_state.order_details.get('contact_name', ''),
                'contact_email': st.session_state.order_details.get('contact_email', ''),
                'contact_phone': st.session_state.order_details.get('contact_phone', '')
            }

            # Generate form button
            st.markdown("#### Generate Pre-Filled Form")

            if st.button("🔗 Generate Google Form URL", type="primary", use_container_width=True):
                # Generate pre-filled URL
                form_url = generate_prefilled_form_url(client_info, selected_products)
                st.session_state.google_form_url = form_url
                st.session_state.show_google_form_url = True

            # Show generated URL
            if st.session_state.get('show_google_form_url', False) and st.session_state.get('google_form_url'):
                st.success("Form URL generated successfully!")

                st.markdown("##### Share This URL with Your Client:")

                # Show URL in text area for easy copying
                st.text_area(
                    "Copy this URL:",
                    value=st.session_state.google_form_url,
                    height=100,
                    key="google_form_url_display"
                )

                # Open in new tab button
                st.markdown(f'<a href="{st.session_state.google_form_url}" target="_blank" style="text-decoration: none;"><button style="width: 100%; padding: 0.5rem; background-color: #0066cc; color: white; border: none; border-radius: 0.25rem; cursor: pointer;">🌐 Open Form in New Tab (Preview)</button></a>', unsafe_allow_html=True)

                st.info("""
                **What's Next?**
                1. Copy the URL above
                2. Send it to your client (email, Slack, text message, etc.)
                3. Client opens the form → sees pre-filled info → completes remaining fields → submits
                4. Go to **Tab 3 → Option A** to import the completed response
                """)
        else:
            st.warning("No products selected. Check at least one product to generate a form.")

    # ============================================================
    # SECTION 3: HTML ORDER FORM (ALTERNATIVE)
    # ============================================================
    st.divider()
    st.subheader("3. HTML Order Form (Alternative)")

    st.markdown("""
    Download the HTML form below and paste it into your email to send to clients.
    The table will look professional and clients can fill it out directly.
    """)

    # Generate HTML table
    # Get order details from session state
    client_type = st.session_state.order_details.get('client_type', 'New')
    company_name = st.session_state.order_details.get('company_name', '') or '[Type company name here]'
    contact_name = st.session_state.order_details.get('contact_name', '') or '[Type your name here]'
    contact_email = st.session_state.order_details.get('contact_email', '') or '[Type your email here]'

    html_form = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>
        body {{ font-family: Arial, sans-serif; max-width: 900px; margin: 20px auto; padding: 20px; background-color: #ffffff; }}
        h2 {{ color: #2c3e50; background-color: #ffffff; border-bottom: 3px solid #3498db; padding-bottom: 10px; }}
        .instructions-box {{ background-color: #e8f4f8; border-left: 4px solid #3498db; padding: 15px; margin: 20px 0; color: #000000; }}
        .instructions-box p {{ margin: 5px 0; color: #000000; }}
        table {{ border-collapse: collapse; width: 100%; margin: 20px 0; background-color: #ffffff; }}
        th {{ background-color: #3498db !important; color: #ffffff !important; padding: 12px; text-align: left; font-weight: bold; }}
        td {{ border: 1px solid #ddd; padding: 10px; background-color: #ffffff; color: #000000; }}
        td:first-child {{ background-color: #f8f9fa !important; color: #000000 !important; font-weight: 500; width: 35%; vertical-align: top; }}
        .section-header {{ background-color: #2c3e50 !important; color: #ffffff !important; font-weight: bold; padding: 10px; }}
        .fill-in {{ background-color: #ffffff !important; color: #7f8c8d !important; min-height: 20px; font-style: italic; }}
        .product-table {{ margin: 10px 0; }}
        .product-table td {{ background-color: #ffffff !important; color: #000000 !important; }}
        .helper-text {{ color: #7f8c8d; font-size: 0.85em; display: block; margin-top: 3px; }}
        .required {{ color: #e74c3c !important; font-weight: bold; }}
    </style>
</head>
<body>
    <h2>PEACE BY PIECE CLIENT ORDER FORM</h2>

    <div class="instructions-box">
        <p><strong>HOW TO FILL OUT THIS FORM:</strong></p>"""

    # Add customizable instructions (one paragraph per line)
    instructions = st.session_state.form_customizations.get('form_instructions', '').strip()
    for i, line in enumerate(instructions.split('\n'), 1):
        if line.strip():
            html_form += f"""
        <p>{line.strip()}</p>"""

    html_form += f"""
    </div>

    <table>
        <tr>
            <td colspan="2" class="section-header">CLIENT INFORMATION</td>
        </tr>
        <tr>
            <td>New Client? <span class="required">*</span></td>
            <td class="fill-in">
                <span style="display: inline-block; margin-right: 20px;">
                    [{' X ' if client_type == 'New' else '   '}] Yes
                </span>
                <span style="display: inline-block;">
                    [{' X ' if client_type == 'Existing' else '   '}] No (Existing)
                </span>
                <span class="helper-text">Check one box</span>
            </td>
        </tr>
        <tr>
            <td>Company Name <span class="required">*</span></td>
            <td class="fill-in">{company_name}</td>
        </tr>
        <tr>
            <td>Contact Name <span class="required">*</span></td>
            <td class="fill-in">{contact_name}</td>
        </tr>
        <tr>
            <td>Contact Email <span class="required">*</span></td>
            <td class="fill-in">{contact_email}</td>
        </tr>
    </table>

    <table>
        <tr>
            <td colspan="2" class="section-header">SHIPPING & DELIVERY</td>
        </tr>
        <tr>
            <td>Drop Shipping? <span class="required">*</span></td>
            <td class="fill-in">[Delete one: Yes / No]</td>
        </tr>
        <tr>
            <td>Shipping Address<span class="helper-text">(if single location)</span></td>
            <td class="fill-in">""" + st.session_state.form_customizations.get('shipping_address_placeholder', '[Type full shipping address here, or N/A if drop shipping]') + """</td>
        </tr>
        <tr>
            <td style="font-weight: bold;">Dropshipping Instructions</td>
            <td style="background-color: #f8f9fa !important; padding: 10px; color: #000000 !important;">""" + st.session_state.form_customizations.get('dropshipping_instructions', '').replace('\n', '<br/>') + """</td>
        </tr>
        <tr>
            <td>Dropshipping Information</td>
            <td class="fill-in">""" + st.session_state.form_customizations.get('dropshipping_placeholder', '[Input dropshipping info here]') + """</td>
        </tr>
        <tr>
            <td>Billing Address</td>
            <td class="fill-in">""" + st.session_state.form_customizations.get('billing_address_placeholder', '[Type billing address here, or "Same as shipping"]') + """</td>
        </tr>
        <tr>
            <td>Client In-Hands Date <span class="required">*</span></td>
            <td class="fill-in">[Type date in format: MM/DD/YYYY]</td>
        </tr>
    </table>

    <table>
        <tr>
            <td colspan="3" class="section-header">ORDER DETAILS</td>
        </tr>
        <tr>
            <th>Product Name</th>
            <th>Quantity</th>
            <th>Customization/Branding Details</th>
        </tr>"""

    # Add product rows - either from proposal or blank rows
    customization_placeholder = st.session_state.form_customizations.get('customization_placeholder', '[Describe any customization, logo placement, colors, etc.]')
    if len(st.session_state.proposal_products) > 0:
        for prop_item in st.session_state.proposal_products:
            product_name = prop_item.get('product_data', {}).get('Product/Service', 'Unknown Product')
            quantity = prop_item.get('quantity', '')
            # Show placeholder text if quantity is empty
            quantity_display = quantity if quantity else '<span style="color: #7f8c8d; font-style: italic;">[Input Qty]</span>'
            html_form += f"""
        <tr>
            <td class="product-table">{product_name}</td>
            <td class="product-table">{quantity_display}</td>
            <td class="product-table" style="color: #7f8c8d; font-style: italic;">{customization_placeholder}</td>
        </tr>"""
    else:
        # Add 3 blank rows if no products in proposal
        for i in range(3):
            html_form += f"""
        <tr>
            <td class="product-table" style="color: #7f8c8d; font-style: italic;">[Product name]</td>
            <td class="product-table" style="color: #7f8c8d; font-style: italic;">[Qty]</td>
            <td class="product-table" style="color: #7f8c8d; font-style: italic;">{customization_placeholder}</td>
        </tr>"""

    html_form += """
    </table>

    <table>
        <tr>
            <td colspan="2" class="section-header">IMPACT CARDS</td>
        </tr>
        <tr>
            <td>Impact Card Preference <span class="required">*</span></td>
            <td class="fill-in">[Delete all except the ONE option you want]<br/><br/>"""

    # Add customizable impact card options
    impact_options = st.session_state.form_customizations.get('impact_card_options', '')
    for line in impact_options.split('\n'):
        if line.strip():
            html_form += f"""
                {line.strip()}<br/>"""

    html_form += """
            </td>
        </tr>
    </table>

    <table>
        <tr>
            <td colspan="2" class="section-header">PAYMENT</td>
        </tr>
        <tr>
            <td>Payment Preference <span class="required">*</span></td>
            <td class="fill-in">[Delete all except the ONE option you want]<br/><br/>"""

    # Add customizable payment options
    payment_options = st.session_state.form_customizations.get('payment_options', '')
    for line in payment_options.split('\n'):
        if line.strip():
            html_form += f"""
                {line.strip()}<br/>"""

    html_form += """
            </td>
        </tr>
    </table>

</body>
</html>"""

    # Show preview in expander
    with st.expander("Preview HTML Form", expanded=False):
        st.components.v1.html(html_form, height=800, scrolling=True)

    # Download buttons
    col1, col2, col3 = st.columns(3)

    with col1:
        st.download_button(
            label="Download HTML Form",
            data=html_form,
            file_name=f"client_order_form_{datetime.now().strftime('%Y%m%d')}.html",
            mime="text/html",
            key="download_client_form_html",
            type="primary"
        )

    with col2:
        # Keep TXT version for backup
        client_form_text = """CLIENT ORDER FORM

New Client? [ ] Yes  [ ] No (Existing)
Company Name: _______________________
Contact: _______________________
Contact Email: _______________________
Drop Shipping? [ ] Y  [ ] N
Shipping address if one location: _______________________
Destination breakdown if drop shipping internationally: _______________________
Billing address: _______________________
Client In-Hands Date: _______________________

Order Details:
Product Name | Quantity | Customization/Branding Details
___________|_________|_____________________________
___________|_________|_____________________________
___________|_________|_____________________________

Impact Cards: [ ] Peace by Piece Impact Card  [ ] Custom Impact Card
              [ ] Custom Message Card  [ ] Send us their own card

Payment Preference: [ ] ACH  [ ] Check  [ ] Credit Card (3% processing fee)
"""
        st.download_button(
            label="Download TXT (Backup)",
            data=client_form_text,
            file_name="client_order_form.txt",
            mime="text/plain",
            key="download_client_form_txt"
        )

    with col3:
        # Generate CSV version
        csv_lines = []
        csv_lines.append("FIELD,VALUE")
        csv_lines.append("New Client?,")
        csv_lines.append("Company Name,")
        csv_lines.append("Contact Name,")
        csv_lines.append("Contact Email,")
        csv_lines.append("Drop Shipping?,")
        csv_lines.append("Shipping Address,")
        csv_lines.append("Destination Breakdown,")
        csv_lines.append("Billing Address,")
        csv_lines.append("Client In-Hands Date,")
        csv_lines.append("")
        csv_lines.append("ORDER DETAILS")
        csv_lines.append("Product Name,Quantity,Customization Details")

        # Add placeholder rows for each product in proposal
        if len(st.session_state.proposal_products) > 0:
            for prop_item in st.session_state.proposal_products:
                product_name = prop_item.get('product_data', {}).get('Product/Service', 'Unknown Product')
                quantity = prop_item.get('quantity', '')
                csv_lines.append(f"{product_name},{quantity},")
        else:
            # Add 3 blank rows
            for i in range(3):
                csv_lines.append(",,")

        csv_lines.append("")
        csv_lines.append("IMPACT CARDS")
        csv_lines.append("Impact Card Type,")
        csv_lines.append("")
        csv_lines.append("PAYMENT")
        csv_lines.append("Payment Preference,")

        client_form_csv = "\n".join(csv_lines)

        st.download_button(
            label="Download CSV (Backup)",
            data=client_form_csv,
            file_name=f"client_order_form_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
            mime="text/csv",
            key="download_client_form_csv"
        )

    st.info("Tip: Download the HTML form, open it in your browser, then copy the entire page and paste it into your email. It will preserve all formatting!")

    # ============================================================
    # NEXT STEPS GUIDANCE
    # ============================================================
    st.divider()

    st.info("""
    **What's Next?**

    1. Download and send the client order form to your client
    2. When your client returns the completed form, move to **Tab 3: Order & Client Info** to process the order
    3. You can import the completed HTML form directly in Tab 3
    """)

    # Navigation button at bottom of Tab 2
    st.divider()
    col1, col2, col3 = st.columns([1, 1, 1])
    with col2:
        if st.button("Continue to Tab 3: Order & Client Info", type="primary", use_container_width=True, key="tab2_to_tab3"):
            st.session_state.show_tab3_prompt = True
            st.rerun()

    # Show navigation prompt if button was clicked
    if st.session_state.get('show_tab3_prompt', False):
        st.info("Click on the **'Order & Client Info'** tab above to continue.")
        st.session_state.show_tab3_prompt = False

# ============================================================
# TAB 3: ORDER & CLIENT INFO
# ============================================================
with tab3:
    st.header("Order & Client Info - Input Order & Client Details")

    # Load data for product matching (needed for Google Form and HTML import)
    df_template, df_metadata, df_partner_info = load_pricing_data(st.session_state.selected_dataset)

    # MIGRATION: Add Phase 3 fields to old order items (backward compatibility)
    for item in st.session_state.order_items:
        if 'pricing_method' not in item:
            item['pricing_method'] = 'Standard markup'
        if 'pricing_notes' not in item:
            item['pricing_notes'] = ''
        if 'manual_override' not in item:
            item['manual_override'] = False
        if 'validation_warning' not in item:
            item['validation_warning'] = None

    # Show unsaved changes indicator and save status
    col1_status, col2_status = st.columns([1, 1])
    with col1_status:
        if has_unsaved_order_changes():
            st.warning("You have unsaved changes in your order")
    with col2_status:
        save_status = format_time_since_save('order')
        if save_status:
            st.caption(f"{save_status}")

    # Quick Save Section at top
    has_order = len(st.session_state.order_items) > 0
    if has_order:
        with st.container():
            col1, col2, col3 = st.columns([2, 2, 1])
            with col1:
                quick_order_name = st.text_input(
                    "Quick save order:",
                    key="quick_save_order_name",
                    placeholder="e.g., Client ABC Q1 2025",
                    label_visibility="collapsed"
                )
            with col2:
                quick_creator = st.text_input(
                    "Your name:",
                    key="quick_save_creator",
                    placeholder="Your name (optional)",
                    label_visibility="collapsed"
                )
            with col3:
                if st.button("Save", type="primary", use_container_width=True, key="quick_save_order_btn"):
                    if not quick_order_name or not quick_order_name.strip():
                        st.error("Please enter an order name")
                    else:
                        # Prepare order data
                        order_data = {
                            'order_items': st.session_state.order_items,
                            'order_shipping': st.session_state.order_shipping,
                            'partner_shipping': st.session_state.partner_shipping,
                            'sales_tax': st.session_state.sales_tax,
                            'kitting_pbp_cost': st.session_state.kitting_pbp_cost,
                            'kitting_client_price': st.session_state.kitting_client_price,
                            'order_discount_type': st.session_state.order_discount_type,
                            'order_discount_preset': st.session_state.order_discount_preset,
                            'order_discount_custom_desc': st.session_state.order_discount_custom_desc,
                            'order_discount_custom_value': st.session_state.order_discount_custom_value,
                            'order_use_marketing_rounding': st.session_state.order_use_marketing_rounding,
                            'apply_cc_fee': st.session_state.apply_cc_fee,
                            'cc_fee_percent': st.session_state.cc_fee_percent,
                            'client_info': st.session_state.client_info,
                            'order_notes': st.session_state.order_notes,
                            'order_confirmed': st.session_state.order_confirmed
                        }

                        success, message, result = save_order(
                            name=quick_order_name.strip(),
                            created_by=quick_creator.strip() if quick_creator else "",
                            order_data=order_data,
                            dataset=st.session_state.selected_dataset
                        )

                        if success:
                            update_last_save_time('order')
                            clear_saved_data_cache()  # Clear cache to show new order
                            st.success(f"{message} - Available in sidebar under 'Saved Orders'")
                            time.sleep(1)
                            st.rerun()
                        else:
                            if result:
                                st.error(f"{message} Try: {result}")
                            else:
                                st.error(message)

    st.divider()

    # ============================================================
    # WORKFLOW GUIDANCE
    # ============================================================
    st.subheader("Getting Started - Choose Your Workflow")

    # Determine which options are available
    has_proposal = len(st.session_state.proposal_products) > 0

    if has_proposal:
        st.markdown("""
        There are **4 ways** to build an order in this tab. Choose the option that matches your situation:

        **RECOMMENDED:** If you sent a Google Form link (from Tab 2) and client submitted response → Use **Option A** below

        **Alternative:** If you sent an HTML form (from Tab 2) and received it back completed → Use **Option B** below

        **Alternative:** If you have a proposal (from Tab 1) but no completed client form → Use **Option C** below

        **Fallback:** If starting fresh without a proposal or form → Use **Option D** below
        """)
    else:
        st.markdown("""
        There are **3 ways** to build an order in this tab. Choose the option that matches your situation:

        **RECOMMENDED:** If you sent a Google Form link (from Tab 2) and client submitted response → Use **Option A** below

        **Alternative:** If you sent an HTML form (from Tab 2) and received it back completed → Use **Option B** below

        **Alternative:** If starting fresh without a proposal or form → Use **Option C** below

        **Tip:** If you want to create a proposal first, go back to Tab 1 to build a proposal, then return here to import it.
        """)
    st.divider()

    # ============================================================
    # SAVED ORDERS SECTION (Always visible)
    # ============================================================
    with st.expander("Saved Orders", expanded=False):
        st.caption("Save your current order or load a previously saved one")

        # Load all saved orders
        refresh_counter = st.session_state.get('saved_data_refresh_counter', 0)
        saved_orders = cached_load_all_orders(refresh_counter)  # Use cached version to reduce API calls

        # Two columns: Load/Delete on left, Save on right
        load_col, save_col = st.columns(2)

        with load_col:
            st.markdown("**Load Order**")

            if len(saved_orders) == 0:
                st.info("No saved orders yet")
            else:
                # Create dropdown options
                order_options = {
                    f"{o['name']} ({o['created_date'][:10]})": o['order_id']
                    for o in saved_orders
                }

                selected_order_label = st.selectbox(
                    "Select order to load:",
                    options=list(order_options.keys()),
                    key="load_order_select"
                )

                if selected_order_label:
                    selected_order_id = order_options[selected_order_label]

                    # Find full order data
                    selected_order = next(o for o in saved_orders if o['order_id'] == selected_order_id)

                    # Show preview
                    st.caption(f"**Created by:** {selected_order['created_by'] or 'Unknown'}")
                    st.caption(f"**Dataset:** {selected_order['dataset']}")

                    # Load and Delete buttons
                    load_btn_col, delete_btn_col = st.columns(2)

                    with load_btn_col:
                        if st.button("Load", key="load_order_btn", type="primary", use_container_width=True):
                            success, order_data, dataset = load_order_data(selected_order_id)

                            if success:
                                # Check if dataset matches
                                if dataset != st.session_state.selected_dataset:
                                    st.warning(f"WARNING: This order was created with {dataset} dataset, but you're currently using {st.session_state.selected_dataset} dataset. Loading anyway...")

                                # Load order data into session state
                                st.session_state.order_items = order_data.get('order_items', [])

                                # Migration: Add kitting fields to old orders that don't have them
                                for item in st.session_state.order_items:
                                    if 'include_kitting' not in item:
                                        item['include_kitting'] = False
                                        item['kitting_pbp_cost'] = 0.0
                                        item['kitting_client_price'] = 0.0
                                        item['kitting_description'] = ''

                                st.session_state.order_shipping = order_data.get('order_shipping', 0.0)
                                st.session_state.partner_shipping = order_data.get('partner_shipping', 0.0)
                                st.session_state.sales_tax = order_data.get('sales_tax', 0.0)
                                st.session_state.kitting_pbp_cost = order_data.get('kitting_pbp_cost', 0.0)
                                st.session_state.kitting_client_price = order_data.get('kitting_client_price', 0.0)
                                st.session_state.order_discount_type = order_data.get('order_discount_type', 'none')
                                st.session_state.order_discount_preset = order_data.get('order_discount_preset', 'NGO Discount (5%)')
                                st.session_state.order_discount_custom_desc = order_data.get('order_discount_custom_desc', '')
                                st.session_state.order_discount_custom_value = order_data.get('order_discount_custom_value', 0.0)
                                st.session_state.order_use_marketing_rounding = order_data.get('order_use_marketing_rounding', False)
                                st.session_state.apply_cc_fee = order_data.get('apply_cc_fee', False)
                                st.session_state.cc_fee_percent = order_data.get('cc_fee_percent', 3.0)
                                st.session_state.client_info = order_data.get('client_info', st.session_state.client_info)
                                # Handle old order_notes structures - clean replacement
                                loaded_notes = order_data.get('order_notes', {})
                                if 'internal_pbp_team' in loaded_notes:
                                    # New 4-category structure - load normally
                                    st.session_state.order_notes = loaded_notes
                                else:
                                    # Old structure detected - discard and initialize fresh
                                    st.session_state.order_notes = {
                                        'internal_pbp_team': '',
                                        'internal_bookkeeping': '',
                                        'external_partners': '',
                                        'external_clients': ''
                                    }
                                st.session_state.order_confirmed = order_data.get('order_confirmed', False)

                                st.success(f"Loaded order: {selected_order['name']}")
                                st.rerun()
                            else:
                                st.error("Failed to load order data")

                    with delete_btn_col:
                        if st.button("Delete", key="delete_order_btn", use_container_width=True):
                            # Confirmation dialog using session state
                            if 'confirm_delete_order_id' not in st.session_state:
                                st.session_state.confirm_delete_order_id = selected_order_id
                                st.warning(f"WARNING: Are you sure you want to delete '{selected_order['name']}'?")

                                confirm_col1, confirm_col2 = st.columns(2)
                                with confirm_col1:
                                    if st.button("Yes, Delete", key="confirm_delete_order_yes", type="primary"):
                                        success, message = delete_order(st.session_state.confirm_delete_order_id)
                                        if success:
                                            clear_saved_data_cache()  # Clear cache after deletion
                                            st.success(message)
                                            del st.session_state.confirm_delete_order_id
                                            st.rerun()
                                        else:
                                            st.error(message)
                                with confirm_col2:
                                    if st.button("Cancel", key="confirm_delete_order_no"):
                                        del st.session_state.confirm_delete_order_id
                                        st.rerun()

        with save_col:
            st.markdown("**Save Current Order**")

            # Check if there are products to save
            has_products = len(st.session_state.order_items) > 0

            if not has_products:
                st.info("Add products to enable saving")

            order_name = st.text_input(
                "Order name:",
                key="save_order_name",
                placeholder="e.g., Client ABC Q1 2025 Order",
                disabled=not has_products
            )

            created_by = st.text_input(
                "Your name (optional):",
                key="save_order_creator",
                placeholder="e.g., John Smith",
                disabled=not has_products
            )

            if st.button("Save Order", key="save_order_btn", type="primary", use_container_width=True, disabled=not has_products):
                if not order_name or not order_name.strip():
                    st.error("Please enter an order name")
                else:
                    # Prepare order data
                    order_data = {
                        'order_items': st.session_state.order_items,
                        'order_shipping': st.session_state.order_shipping,
                        'partner_shipping': st.session_state.partner_shipping,
                        'sales_tax': st.session_state.sales_tax,
                        'kitting_pbp_cost': st.session_state.kitting_pbp_cost,
                        'kitting_client_price': st.session_state.kitting_client_price,
                        'order_discount_type': st.session_state.order_discount_type,
                        'order_discount_preset': st.session_state.order_discount_preset,
                        'order_discount_custom_desc': st.session_state.order_discount_custom_desc,
                        'order_discount_custom_value': st.session_state.order_discount_custom_value,
                        'order_use_marketing_rounding': st.session_state.order_use_marketing_rounding,
                        'apply_cc_fee': st.session_state.apply_cc_fee,
                        'cc_fee_percent': st.session_state.cc_fee_percent,
                        'client_info': st.session_state.client_info,
                        'order_notes': st.session_state.order_notes,
                        'order_confirmed': st.session_state.order_confirmed
                    }

                    success, message, result = save_order(
                        name=order_name.strip(),
                        created_by=created_by.strip() if created_by else "",
                        order_data=order_data,
                        dataset=st.session_state.selected_dataset
                    )

                    if success:
                        update_last_save_time('order')
                        clear_saved_data_cache()  # Clear cache to show new order
                        st.success(message)
                        st.rerun()  # Rerun to clear form
                    else:
                        # Check if it's a naming conflict
                        if result:  # result contains suggested name
                            st.error(message)
                            # Offer to save with suggested name
                            if st.button(f"Save as '{result}'", key="save_order_with_new_name"):
                                success2, message2, _ = save_order(
                                    name=result,
                                    created_by=created_by.strip() if created_by else "",
                                    order_data=order_data,
                                    dataset=st.session_state.selected_dataset
                                )
                                if success2:
                                    update_last_save_time('order')
                                    clear_saved_data_cache()  # Clear cache to show new order
                                    st.success(message2)
                                    st.rerun()
                        else:
                            st.error(message)

    st.divider()

    # ============================================================
    # OPTION A: GOOGLE FORM RESPONSE IMPORT (RECOMMENDED)
    # ============================================================
    st.header("Option A: Import from Google Form Response (RECOMMENDED)")
    st.markdown("**Use this if:** You sent a Google Form link from Tab 2 and client submitted their response")
    st.caption("View and import client responses from your Google Form submissions")

    # Import helpers
    from src.data_loader import connect_to_sheets
    from src.forms_helper import get_unimported_responses, parse_form_response, mark_response_imported, format_product_summary

    if st.button("🔄 Load Recent Form Responses", key="load_google_form_responses"):
        with st.spinner("Loading responses from Google Sheets..."):
            try:
                gc = connect_to_sheets()
                df_unimported = get_unimported_responses(gc)

                if df_unimported.empty:
                    st.info("No new responses found. All responses have been imported.")
                else:
                    st.success(f"Found **{len(df_unimported)}** unimported response(s)")
                    st.session_state.google_form_responses = df_unimported
                    st.session_state.google_form_gc = gc

            except Exception as e:
                st.error(f"Error loading responses: {e}")
                st.session_state.google_form_responses = None

    # Show responses if loaded
    if st.session_state.get('google_form_responses') is not None:
        df_responses = st.session_state.google_form_responses

        if not df_responses.empty:
            st.markdown("#### Available Responses:")

            for idx, row in df_responses.iterrows():
                # Parse response
                response_data = parse_form_response(row)

                # Create expander for each response
                company = response_data['client_info'].get('company_name', 'Unknown Company')
                timestamp = response_data['metadata'].get('timestamp', 'No timestamp')

                with st.expander(f"📋 {company} - {timestamp}", expanded=False):
                    col1, col2 = st.columns(2)

                    with col1:
                        st.markdown("**Client Information:**")
                        st.write(f"- Type: {response_data['client_info'].get('client_type', 'N/A')}")
                        st.write(f"- Company: {response_data['client_info'].get('company_name', 'N/A')}")
                        st.write(f"- Contact: {response_data['client_info'].get('contact_name', 'N/A')}")
                        st.write(f"- Email: {response_data['client_info'].get('contact_email', 'N/A')}")

                    with col2:
                        st.markdown("**Order Details:**")
                        st.write(f"- Products: {format_product_summary(response_data['products'])}")
                        st.write(f"- In-Hands: {response_data['shipping_info'].get('in_hands_date', 'N/A')}")
                        st.write(f"- Payment: {response_data['payment_info'].get('payment_preference', 'N/A')}")

                    st.markdown("**Products:**")
                    for product in response_data['products']:
                        st.write(f"- {product['name']} (Qty: {product['quantity']})")
                        if product.get('customization_notes'):
                            st.caption(f"  Custom: {product['customization_notes']}")

                    # Import button
                    if st.button(f"Import This Response", key=f"import_google_response_{idx}"):
                        # Import client info
                        client_type = response_data['client_info'].get('client_type', 'New')
                        drop_shipping = response_data['shipping_info'].get('drop_shipping', 'No')

                        # Build contacts array from form data
                        contacts = [{
                            'name': response_data['client_info'].get('contact_name', ''),
                            'email': response_data['client_info'].get('contact_email', ''),
                            'phone': response_data['client_info'].get('contact_phone', ''),
                            'role': 'Primary Contact'
                        }]

                        # Parse in-hands date string to date object (Google Forms returns YYYY-MM-DD string)
                        in_hands_date_str = response_data['shipping_info'].get('in_hands_date', None)
                        if in_hands_date_str:
                            try:
                                in_hands_date = datetime.strptime(in_hands_date_str, '%Y-%m-%d').date()
                            except (ValueError, AttributeError):
                                in_hands_date = None
                        else:
                            in_hands_date = None

                        st.session_state.client_info = {
                            'is_new_client': client_type == 'New',
                            'company_name': response_data['client_info'].get('company_name', ''),
                            'contacts': contacts,
                            'client_po': '',
                            'billing_address': response_data['shipping_info'].get('billing_address', ''),
                            'shipping_type': 'Drop Shipping' if drop_shipping == 'Yes' else 'Ground',
                            'shipping_address': response_data['shipping_info'].get('shipping_address', ''),
                            'payment_timeline': response_data['payment_info'].get('payment_preference', 'Net 30'),
                            'payment_preference': response_data['payment_info'].get('payment_method', 'Check'),
                            'client_in_hands_date': in_hands_date,
                            'order_submitted_by': '',
                            'order_submitted_date': datetime.now().date(),
                            'cost_submitted_by': '',
                            'cost_submitted_date': None
                        }

                        # Import products (using new pricing logic from Phase 1)
                        from src.pricing_engine import get_unit_price_new_system, calculate_pbp_msrp
                        from src.helpers import get_tariff_rate, calculate_product_tariff, get_shipping_costs, clean_price, get_column_value

                        products_imported = 0
                        products_skipped = []
                        max_pbp_shipping = 0.0

                        for product in response_data['products']:
                            product_name = product['name']

                            # Match product name to catalog (exact match, case-insensitive)
                            df_filtered = df_template[df_template['Product/Service'].str.lower().str.strip() == product_name.lower().strip()]

                            if not df_filtered.empty:
                                product_data = df_filtered.iloc[0].to_dict()
                                quantity = product['quantity']

                                # Check if this product was in the proposal - if so, use proposal markup
                                markup_percent = None
                                for prop_item in st.session_state.proposal_products:
                                    if prop_item['product_data'].get('Product/Service', '').lower() == product_name.lower():
                                        markup_percent = prop_item['markup_percent']
                                        break

                                # Calculate pricing using new system (Phase 3)
                                pricing_result = calculate_pbp_msrp(product_data, quantity, user_markup_override=markup_percent)

                                # Extract pricing method and notes
                                pricing_method = pricing_result['method_used']
                                pricing_notes = get_column_value(product_data, 'Pricing Notes', '', '')

                                # Generate validation warning if mismatch
                                validation_warning = None
                                if pricing_result['validation_status'] == 'mismatch':
                                    spreadsheet_val = pricing_result['spreadsheet_msrp']
                                    calculated_val = pricing_result['pbp_msrp']
                                    validation_warning = f"Price mismatch: Spreadsheet=${spreadsheet_val:.2f} | Calculated=${calculated_val:.2f}"

                                # Determine markup based on pricing method
                                if pricing_method == "Standard markup":
                                    # Use markup from calculation or default
                                    markup_percent = pricing_result['calculation_details'].get('markup_percent', 100.0)
                                else:
                                    # MSRP-based method - calculate implied markup
                                    base_cost = pricing_result['calculation_details']['per_item_cost']
                                    pbp_msrp = pricing_result['pbp_msrp']
                                    if base_cost > 0:
                                        markup_percent = ((pbp_msrp / base_cost) - 1) * 100
                                    else:
                                        markup_percent = 100.0

                                markup_multiplier = markup_percent / 100.0

                                # Get base price for quantity
                                base_price_per_unit, tier_info, tier_num = get_unit_price_new_system(product_data, quantity)

                                # Calculate costs with calculated markup
                                product_cost_subtotal = base_price_per_unit * quantity
                                markup_amount = product_cost_subtotal * markup_multiplier
                                product_total = product_cost_subtotal + markup_amount

                                # Parse tariff
                                tariff_rate_percent = get_tariff_rate(product_data, product_cost_subtotal)
                                tariff_base = product_cost_subtotal
                                tariff_amount = calculate_product_tariff(tariff_base, tariff_rate_percent)

                                # Build order item (matching HTML import structure + Phase 3 new fields)
                                order_item = {
                                    'product_name': product_data.get('Product/Service', 'Unknown Product'),
                                    'product_ref': product_data.get('Purchase Description', ''),
                                    'partner': product_data.get('Partner', 'Unknown Partner'),
                                    'product_data': product_data,
                                    'product_data_row': product_data,
                                    'is_custom': False,
                                    'quantity': quantity,
                                    'base_price': base_price_per_unit,
                                    'tier_range': tier_info if tier_info else '',
                                    'tier_column': f'T{tier_num}' if tier_num else '',
                                    'markup_percent': markup_percent,
                                    'markup_amount': markup_amount,
                                    'proposal_markup_percent': markup_percent,  # Track original markup for warning system
                                    # Phase 3: New pricing fields
                                    'pricing_method': pricing_method,
                                    'pricing_notes': pricing_notes,
                                    'validation_warning': validation_warning,
                                    'manual_override': False,
                                    'include_customization': False,
                                    'customization_description': product_data.get('Customization Info', ''),
                                    'customization_setup_fee': float(clean_price(get_column_value(product_data, 'Client Price: Customization Setup Fee', 'Customization Setup Fee', '')) or 0.0),
                                    'customization_per_unit': float(clean_price(get_column_value(product_data, 'Client Price: Customization Cost per Unit', 'Customization Cost per Unit', '')) or 0.0),
                                    'customization_setup_total': 0.0,
                                    'customization_unit_total': 0.0,
                                    'apply_custom_minimum': False,
                                    'customization_minimum_qty': 0,
                                    'product_subtotal': product_cost_subtotal,
                                    'subtotal_before_markup': product_cost_subtotal,
                                    'product_total': product_total,
                                    'total_per_unit': product_total / quantity,  # Fix: divide by quantity for per-unit price
                                    'quoted_price_per_unit': (product_cost_subtotal + markup_amount) / quantity,  # Fix: divide by quantity
                                    'tariff_info': f"{product_data.get('Country', 'N/A')} - {tariff_rate_percent}%" if tariff_rate_percent > 0 else '',
                                    'tariff_rate_percent': tariff_rate_percent,
                                    'tariff_base': tariff_base,
                                    'tariff_amount': tariff_amount,
                                    'edited_description': '',
                                    # Per-product kitting fields
                                    'include_kitting': False,
                                    'kitting_pbp_cost': 0.0,
                                    'kitting_client_price': 0.0,
                                    'kitting_description': ''
                                }

                                st.session_state.order_items.append(order_item)
                                products_imported += 1

                                # Track maximum PBP shipping cost
                                pbp_shipping, _ = get_shipping_costs(product_data)
                                max_pbp_shipping = max(max_pbp_shipping, pbp_shipping)
                            else:
                                products_skipped.append(product_name)

                        # Auto-populate partner shipping with the maximum cost found
                        if max_pbp_shipping > 0 and st.session_state.partner_shipping == 0:
                            st.session_state.partner_shipping = max_pbp_shipping

                        # Mark as imported
                        gc = st.session_state.google_form_gc
                        row_index = idx + 2  # +2 because: +1 for header, +1 for 1-indexed
                        order_id = f"IMPORTED-{datetime.now().strftime('%Y%m%d-%H%M%S')}"
                        mark_response_imported(gc, row_index, order_id)

                        # Show results
                        if products_imported > 0:
                            st.toast(f"Imported {products_imported} product(s) from {company}")

                        if products_skipped:
                            st.warning(f"Could not match {len(products_skipped)} product(s): {', '.join(products_skipped)}")

                        st.info("**Next:** Scroll down to Section 2 to configure your order")

                        # Rerun to show imported products in Section 2
                        st.rerun()

    st.divider()

    # ============================================================
    # OPTION B: HTML CLIENT ORDER FORM IMPORT (ALTERNATIVE)
    # ============================================================
    st.header("Option B: Import from HTML Order Form (Alternative)")
    st.markdown("**Use this if:** You sent an HTML form from Tab 2 and received it back completed")

    with st.expander("Upload Completed Client Order Form", expanded=False):
        st.caption("Upload an HTML order form completed by your client to auto-populate client information.")
        st.info("This will import client info, shipping, payment details, and order products from the form.")

        uploaded_file = st.file_uploader(
            "Upload Client Order Form (HTML)",
            type=['html', 'htm'],
            key="import_client_form_top",
            help="Upload the HTML form that your client filled out and returned to you"
        )

        if uploaded_file is not None:
            # Read and parse the HTML content
            content = uploaded_file.read().decode('utf-8')
            parsed_data = parse_client_order_form_html(content)

            # Show parsing errors if any
            if parsed_data['parse_errors']:
                st.warning("Parsing warnings:")
                for error in parsed_data['parse_errors']:
                    st.caption(f"- {error}")

            # Show preview of extracted data
            st.markdown("**Preview of Extracted Data:**")

            preview_data = {
                "New Client?": "Yes" if parsed_data['client_type'] == 'New' else ("No" if parsed_data['client_type'] == 'Existing' else "[Not filled]"),
                "Company Name": parsed_data['company_name'] or "[Not filled]",
                "Contact Name": parsed_data['contact_name'] or "[Not filled]",
                "Contact Email": parsed_data['contact_email'] or "[Not filled]",
                "Drop Shipping": parsed_data['drop_shipping'] or "[Not filled]",
                "Shipping Address": parsed_data['shipping_address'] or "[Not filled]",
                "Dropshipping Info": parsed_data['dropshipping_info'] or "[Not filled]",
                "Billing Address": parsed_data['billing_address'] or "[Not filled]",
                "Client In-Hands Date": parsed_data['client_in_hands_date'] or "[Not filled]",
                "Impact Card Preference": parsed_data['impact_card_preference'] or "[Not filled]",
                "Payment Preference": parsed_data['payment_preference'] or "[Not filled]"
            }

            # Display as table
            df_preview = pd.DataFrame(list(preview_data.items()), columns=["Field", "Value"])
            st.dataframe(df_preview, use_container_width=True, hide_index=True)

            # Show extracted products
            if parsed_data['products']:
                st.markdown(f"**Products Found:** {len(parsed_data['products'])} product(s)")
                products_df = pd.DataFrame(parsed_data['products'], columns=["Product Name"])
                st.dataframe(products_df, use_container_width=True, hide_index=True)
            else:
                st.caption("No products found in order form")

            # Import button
            st.divider()
            if st.button("Import Client Information", type="primary", use_container_width=True, key="import_client_data_btn_top"):
                # Apply to session state
                st.session_state.client_info['is_new_client'] = (parsed_data['client_type'] == 'New')
                st.session_state.client_info['company_name'] = parsed_data['company_name']
                st.session_state.client_info['contact_name'] = parsed_data['contact_name']
                st.session_state.client_info['contact_email'] = parsed_data['contact_email']
                st.session_state.client_info['shipping_address'] = parsed_data['shipping_address']
                st.session_state.client_info['billing_address'] = parsed_data['billing_address']

                # Set shipping type based on drop shipping answer
                # Default to 'One Location' (show shipping address) unless clearly "Yes" for drop shipping
                if parsed_data['drop_shipping'] == 'Yes':
                    st.session_state.client_info['shipping_type'] = 'Drop Shipping'
                else:
                    # For "No", empty, or any unclear answer, default to One Location
                    st.session_state.client_info['shipping_type'] = 'One Location'

                # Parse and apply date
                if parsed_data['client_in_hands_date']:
                    try:
                        from datetime import datetime
                        # Try to parse MM/DD/YYYY format
                        date_obj = datetime.strptime(parsed_data['client_in_hands_date'], '%m/%d/%Y')
                        st.session_state.client_info['client_in_hands_date'] = date_obj.date()
                    except:
                        # If parsing fails, leave empty and let user enter manually
                        st.session_state.client_info['client_in_hands_date'] = None

                # Map payment preference to our dropdown values
                payment_map = {
                    'ACH': 'ACH',
                    'Check': 'Check',
                    'Credit Card': 'Credit Card (3% processing fee applies)'
                }
                if parsed_data['payment_preference'] in payment_map:
                    st.session_state.client_info['payment_preference'] = payment_map[parsed_data['payment_preference']]

                st.success("Client information imported successfully! Review and edit the fields below as needed.")
                st.rerun()

            # ============================================================
            # PRODUCT SELECTION FROM ORDER FORM
            # ============================================================
            if parsed_data['products']:
                st.divider()
                st.markdown("**Select Products from Order Form:**")
                st.caption("Products will be matched against the catalog and added with default settings (quantity 1, 100% markup).")

                # Match products against catalog
                matched_products = []
                unmatched_products = []

                for product_name in parsed_data['products']:
                    # Try exact match first
                    exact_match = st.session_state.df_template[
                        st.session_state.df_template['Product/Service'].str.lower() == product_name.lower()
                    ]

                    if len(exact_match) > 0:
                        # Exact match found
                        product_row = exact_match.iloc[0].to_dict()
                        matched_products.append({
                            'name': product_name,
                            'match_type': 'Exact',
                            'product_data': product_row
                        })
                    else:
                        # Try partial match
                        partial_match = st.session_state.df_template[
                            st.session_state.df_template['Product/Service'].str.contains(product_name, case=False, na=False)
                        ]

                        if len(partial_match) > 0:
                            # Partial match found - use first match
                            product_row = partial_match.iloc[0].to_dict()
                            matched_products.append({
                                'name': product_name,
                                'match_type': 'Partial',
                                'catalog_name': product_row['Product/Service'],
                                'product_data': product_row
                            })
                        else:
                            # No match found
                            unmatched_products.append(product_name)

                # Show matched products with checkboxes
                if matched_products:
                    st.markdown(f"**Matched Products ({len(matched_products)}):**")

                    selected_product_indices = []

                    for idx, match in enumerate(matched_products):
                        col1, col2 = st.columns([4, 1])

                        with col1:
                            is_selected = st.checkbox(
                                f"{match['name']} ({match['match_type']} match)",
                                key=f"select_form_product_{idx}",
                                value=True  # Default to checked
                            )

                            # Show catalog name if different (partial match)
                            if match['match_type'] == 'Partial':
                                st.caption(f"Catalog: {match['catalog_name']}")

                            # Show partner
                            st.caption(f"Partner: {match['product_data'].get('Partner', 'N/A')}")

                        with col2:
                            if is_selected:
                                selected_product_indices.append(idx)

                    # Add selected button
                    if len(selected_product_indices) > 0:
                        if st.button(f"Add {len(selected_product_indices)} Selected Product(s) to Order", type="primary", use_container_width=True, key="add_form_products_btn"):
                            # Add products to order (using new pricing logic from Phase 3)
                            max_pbp_shipping = 0.0
                            for idx in selected_product_indices:
                                match = matched_products[idx]
                                product_data = match['product_data']

                                # Calculate pricing using new system (Phase 3)
                                pricing_result = calculate_pbp_msrp(product_data, quantity=1)

                                # Extract pricing method and notes
                                pricing_method = pricing_result['method_used']
                                pricing_notes = get_column_value(product_data, 'Pricing Notes', '', '')

                                # Generate validation warning if mismatch
                                validation_warning = None
                                if pricing_result['validation_status'] == 'mismatch':
                                    spreadsheet_val = pricing_result['spreadsheet_msrp']
                                    calculated_val = pricing_result['pbp_msrp']
                                    validation_warning = f"Price mismatch: Spreadsheet=${spreadsheet_val:.2f} | Calculated=${calculated_val:.2f}"

                                # Determine markup based on pricing method
                                if pricing_method == "Standard markup":
                                    markup_percent = pricing_result['calculation_details'].get('markup_percent', 100.0)
                                else:
                                    # MSRP-based method - calculate implied markup
                                    base_cost = pricing_result['calculation_details']['per_item_cost']
                                    pbp_msrp = pricing_result['pbp_msrp']
                                    if base_cost > 0:
                                        markup_percent = ((pbp_msrp / base_cost) - 1) * 100
                                    else:
                                        markup_percent = 100.0

                                # Get base price for quantity 1
                                base_price_per_unit, tier_info, tier_num = get_unit_price_new_system(product_data, 1)

                                # Calculate costs
                                product_cost_subtotal = base_price_per_unit * 1
                                markup_amount = product_cost_subtotal * (markup_percent / 100.0)
                                product_total = product_cost_subtotal + markup_amount

                                # Parse tariff
                                tariff_rate_percent = get_tariff_rate(product_data, product_cost_subtotal)
                                tariff_base = product_cost_subtotal
                                tariff_amount = calculate_product_tariff(tariff_base, tariff_rate_percent)

                                # Build order item (with Phase 3 new fields)
                                order_item = {
                                    'product_name': product_data.get('Product/Service', 'Unknown Product'),
                                    'product_ref': product_data.get('Purchase Description', ''),
                                    'partner': product_data.get('Partner', 'Unknown Partner'),
                                    'product_data': product_data,
                                    'product_data_row': product_data,
                                    'is_custom': False,
                                    'quantity': 1,
                                    'base_price': base_price_per_unit,
                                    'tier_range': tier_info if tier_info else '',
                                    'tier_column': f'T{tier_num}' if tier_num else '',
                                    'markup_percent': markup_percent,
                                    'markup_amount': markup_amount,
                                    # Phase 3: New pricing fields
                                    'pricing_method': pricing_method,
                                    'pricing_notes': pricing_notes,
                                    'validation_warning': validation_warning,
                                    'manual_override': False,
                                    'include_customization': False,
                                    'customization_description': product_data.get('Customization Info', ''),
                                    'customization_setup_fee': float(clean_price(get_column_value(product_data, 'Client Price: Customization Setup Fee', 'Customization Setup Fee', '')) or 0.0),
                                    'customization_per_unit': float(clean_price(get_column_value(product_data, 'Client Price: Customization Cost per Unit', 'Customization Cost per Unit', '')) or 0.0),
                                    'customization_setup_total': 0.0,
                                    'customization_unit_total': 0.0,
                                    'apply_custom_minimum': False,
                                    'customization_minimum_qty': 0,
                                    'product_subtotal': product_cost_subtotal,
                                    'subtotal_before_markup': product_cost_subtotal,
                                    'product_total': product_total,
                                    'total_per_unit': product_total,
                                    'quoted_price_per_unit': (product_cost_subtotal + markup_amount),
                                    'tariff_info': f"{product_data.get('Country', 'N/A')} - {tariff_rate_percent}%" if tariff_rate_percent > 0 else '',
                                    'tariff_rate_percent': tariff_rate_percent,
                                    'tariff_base': tariff_base,
                                    'tariff_amount': tariff_amount,
                                    'edited_description': '',  # Initialize with empty for user to edit later
                                    # Per-product kitting fields
                                    'include_kitting': False,
                                    'kitting_pbp_cost': 0.0,
                                    'kitting_client_price': 0.0,
                                    'kitting_description': ''
                                }

                                st.session_state.order_items.append(order_item)

                                # Track maximum PBP shipping cost
                                pbp_shipping, _ = get_shipping_costs(product_data)
                                max_pbp_shipping = max(max_pbp_shipping, pbp_shipping)

                            # Auto-populate partner shipping with the maximum cost found
                            if max_pbp_shipping > 0 and st.session_state.partner_shipping == 0:
                                st.session_state.partner_shipping = max_pbp_shipping

                            st.toast(f"Added {len(selected_product_indices)} product(s) from order form!")
                            st.rerun()
                    else:
                        st.caption("Select at least one product above to add to order.")

                # Show unmatched products
                if unmatched_products:
                    st.divider()
                    st.warning(f"**Not Found in Catalog ({len(unmatched_products)}):**")
                    for product_name in unmatched_products:
                        st.caption(f"- {product_name}")
                    st.caption("These products will not be added. You can add them manually using Option C below.")

    st.divider()

    # ============================================================
    # OPTION C: PROPOSAL PRODUCTS SELECTION (if available)
    # ============================================================
    if len(st.session_state.proposal_products) > 0:
        st.header("Option C: Import Products from Proposal (Tab 1)")
        st.markdown("**Use this if:** You created a proposal in Tab 1 but don't have a completed client form")
        
        # Display proposal source information
        proposal_info_msg = f"**{len(st.session_state.proposal_products)} product(s) available**"
        
        if 'loaded_proposal_name' in st.session_state and st.session_state.loaded_proposal_name:
            # This is a saved/loaded proposal
            proposal_info_msg += f" from saved proposal: **'{st.session_state.loaded_proposal_name}'**"
            if 'loaded_proposal_date' in st.session_state:
                proposal_info_msg += f"\n\nCreated/Saved: {st.session_state.loaded_proposal_date}"
            if 'loaded_proposal_creator' in st.session_state and st.session_state.loaded_proposal_creator != 'Unknown':
                proposal_info_msg += f" | By: {st.session_state.loaded_proposal_creator}"
        else:
            # This is an unsaved proposal from current session
            proposal_info_msg += " from **Current Session Proposal (unsaved)**"
            proposal_info_msg += "\n\nTip: Save your proposal in Tab 1 to preserve it for future use"
        
        st.info(proposal_info_msg)
        st.session_state.using_proposal_data = True

        # Import All button at top level
        col1, col2 = st.columns([1, 1])
        with col1:
            if st.button("Import All Products from Proposal", type="primary", use_container_width=True, key="import_all_proposal"):
                # Import all proposal products to order
                imported_count = 0
                max_pbp_shipping = 0.0  # Track maximum shipping cost among all products

                for prop_item in st.session_state.proposal_products:
                    order_item = convert_proposal_to_order(
                        prop_item,
                        get_unit_price_new_system,
                        calculate_product_tariff
                    )
                    st.session_state.order_items.append(order_item)
                    imported_count += 1

                    # Track maximum PBP shipping cost
                    pbp_shipping, _ = get_shipping_costs(prop_item.get('product_data', {}))
                    max_pbp_shipping = max(max_pbp_shipping, pbp_shipping)

                # Auto-populate partner shipping with the maximum cost found
                if max_pbp_shipping > 0 and st.session_state.partner_shipping == 0:
                    st.session_state.partner_shipping = max_pbp_shipping

                st.toast(f"Imported all {imported_count} product(s) from proposal!")
                st.rerun()

        with col2:
            st.caption("Or select individually below:")

        with st.expander("Select Individual Products from Proposal", expanded=False):
            st.markdown("Select specific products from your proposal to add to this order. You can edit quantities and settings after adding.")

            # Build selection checkboxes
            selected_proposal_indices = []

            for idx, prop_item in enumerate(st.session_state.proposal_products):
                product_data = prop_item.get('product_data', {})

                col1, col2 = st.columns([4, 1])

                with col1:
                    # Display product name with variant (if applicable)
                    from src.helpers import format_product_with_variant
                    product_display_name = format_product_with_variant(
                        product_data.get('Product/Service', 'Unknown Product'),
                        prop_item.get('selected_variant')
                    )

                    is_selected = st.checkbox(
                        f"{product_display_name} - {product_data.get('Partner', 'N/A')}",
                        key=f"select_proposal_{idx}"
                    )

                    # Show proposal details
                    st.caption(f"Quantity: {prop_item.get('quantity', 'N/A')} | Markup: {prop_item.get('markup_percent', 0)}%")

                    if prop_item.get('include_customization', False):
                        setup_fee = prop_item.get('customization_setup_fee', 0)
                        per_unit = prop_item.get('customization_per_unit', 0)
                        st.caption(f"Customization: ${setup_fee:.2f} setup + ${per_unit:.2f}/unit")

                with col2:
                    if is_selected:
                        selected_proposal_indices.append(idx)

            # Add selected button
            if len(selected_proposal_indices) > 0:
                if st.button(f"Add {len(selected_proposal_indices)} Selected Product(s) to Order", type="primary", use_container_width=True):
                    # Convert and add to order
                    max_pbp_shipping = 0.0
                    for idx in selected_proposal_indices:
                        prop_item = st.session_state.proposal_products[idx]
                        order_item = convert_proposal_to_order(
                            prop_item,
                            get_unit_price_new_system,
                            calculate_product_tariff
                        )
                        st.session_state.order_items.append(order_item)

                        # Track maximum PBP shipping cost
                        pbp_shipping, _ = get_shipping_costs(prop_item.get('product_data', {}))
                        max_pbp_shipping = max(max_pbp_shipping, pbp_shipping)

                    # Auto-populate partner shipping with the maximum cost found
                    if max_pbp_shipping > 0 and st.session_state.partner_shipping == 0:
                        st.session_state.partner_shipping = max_pbp_shipping

                    st.toast(f"Added {len(selected_proposal_indices)} product(s) to order!")
                    st.rerun()
            else:
                st.caption("Select at least one product above to add to order.")

        st.divider()
    else:
        st.session_state.using_proposal_data = False


    # ============================================================
    # OPTION D (or C): MANUAL PRODUCT SELECTION
    # ============================================================
    # Display success message if a product was just added
    if 'show_add_to_order_success' in st.session_state and st.session_state.show_add_to_order_success:
        st.toast(f"Added {st.session_state.add_to_order_product_name} to order")
        st.session_state.show_add_to_order_success = False

    # Adjust option label based on whether proposal exists
    option_label = "Option D" if has_proposal else "Option C"
    st.header(f"{option_label}: Manual Product Selection")
    st.markdown("**Use this if:** You're starting from scratch without a proposal or completed form")
    st.caption("Add products to your order, then configure settings for each product below")

    # Create dropdowns for filtering
    col1, col2, col3 = st.columns([2, 1, 1])

    with col1:
        col_partner, col_product = st.columns(2)

        with col_partner:
            # Partner dropdown (using "Partner" column from Template sheet)
            partners = sorted(df_template["Partner"].unique().tolist())
            selected_partner = st.selectbox("Select Partner", partners, key="add_partner_select")

        with col_product:
            # Filter products based on partner selection (using "Product/Service" column)
            available_products = df_template[df_template["Partner"] == selected_partner]["Product/Service"].unique().tolist()
            selected_product = st.selectbox("Select Product/Service", available_products, key="add_product_select")

    # Check if selected product has variants
    from src.helpers import has_variants, parse_variant_types
    selected_product_data = df_template[
        (df_template["Partner"] == selected_partner) &
        (df_template["Product/Service"] == selected_product)
    ].iloc[0] if selected_product else None

    product_has_variants = has_variants(selected_product_data) if selected_product_data is not None else False
    variant_types = parse_variant_types(selected_product_data) if product_has_variants else []

    # Variant selector (if product has variants)
    selected_variant_manual = None
    if product_has_variants and variant_types:
        col_var = st.columns([2, 1, 1])[0]  # Match layout
        with col_var:
            selected_variant_manual = st.selectbox(
                "Select Variant:",
                options=[''] + variant_types,  # Empty option first
                key="manual_variant_select"
            )
            if not selected_variant_manual:
                st.caption("Variant recommended but not required")

    with col2:
        st.write("")  # Spacing
        st.session_state.order_use_msrp = st.checkbox(
            "Use MSRP pricing",
            value=st.session_state.order_use_msrp,
            key="order_use_msrp_checkbox",
            help="When enabled, products with MSRP will have markup auto-calculated to match MSRP. Products without MSRP will use 100% markup."
        )

    with col3:
        st.write("")  # Spacing
        st.write("")  # Spacing
        if st.button("Add to Order", type="primary", use_container_width=True):
            # Get selected product details
            product_data = df_template[
                (df_template["Partner"] == selected_partner) &
                (df_template["Product/Service"] == selected_product)
            ].iloc[0]

            # Get default customization costs from spreadsheet
            default_setup_fee_raw = get_column_value(product_data, 'Client Price: Customization Setup Fee', 'Customization Setup Fee', '')
            default_per_unit_raw = get_column_value(product_data, 'Client Price: Customization Cost per Unit', 'Customization Cost per Unit', '')
            default_setup_fee = clean_price(default_setup_fee_raw) or 0.0
            default_per_unit = clean_price(default_per_unit_raw) or 0.0

            # Calculate pricing using new system (Phase 3)
            from src.pricing_engine import calculate_pbp_msrp

            # Determine if using MSRP pricing
            use_msrp = st.session_state.order_use_msrp

            # Get pricing result
            pricing_result = calculate_pbp_msrp(product_data.to_dict(), quantity=1)

            # Extract pricing method and notes
            pricing_method = pricing_result['method_used']
            pricing_notes = get_column_value(product_data, 'Pricing Notes', '', '')

            # Generate validation warning if mismatch
            validation_warning = None
            if pricing_result['validation_status'] == 'mismatch':
                spreadsheet_val = pricing_result['spreadsheet_msrp']
                calculated_val = pricing_result['pbp_msrp']
                validation_warning = f"Price mismatch: Spreadsheet=${spreadsheet_val:.2f} | Calculated=${calculated_val:.2f}"

            # Determine markup based on pricing method
            if use_msrp:
                if pricing_method == "Standard markup":
                    markup = pricing_result['calculation_details'].get('markup_percent', 100.0)
                else:
                    # MSRP-based method - calculate implied markup
                    base_cost = pricing_result['calculation_details']['per_item_cost']
                    pbp_msrp = pricing_result['pbp_msrp']
                    if base_cost > 0:
                        markup = ((pbp_msrp / base_cost) - 1) * 100
                    else:
                        markup = 100.0
            else:
                # User unchecked MSRP - use standard markup
                markup = get_default_markup(product_data.to_dict())

            # Manual override: if user unchecked MSRP, treat as manual override
            manual_override = not use_msrp

            # Add product with defaults
            new_item = {
                'product_name': product_data.get('Product/Service', 'Unknown Product'),
                'partner': product_data.get('Partner', 'Unknown Partner'),
                'product_data': product_data.to_dict(),
                'quantity': 1,
                'markup_percent': markup,
                'selected_variant': selected_variant_manual if selected_variant_manual else None,
                # Phase 3: New pricing fields
                'pricing_method': pricing_method,
                'pricing_notes': pricing_notes,
                'validation_warning': validation_warning,
                'manual_override': manual_override,
                'include_customization': False,
                'customization_setup_fee': float(default_setup_fee),
                'customization_per_unit': float(default_per_unit),
                'customization_minimum_qty': 0,
                'apply_custom_minimum': False,
                'include_tariff': False,
                'is_custom': False,
                # Per-product kitting fields
                'include_kitting': False,
                'kitting_pbp_cost': 0.0,
                'kitting_client_price': 0.0,
                'kitting_description': ''
            }

            # Calculate pricing for this item (will be recalculated when edited)
            base_price, tier_range, tier_column = get_unit_price_new_system(product_data, 1)

            if base_price:
                # Add calculated fields
                new_item.update({
                    'base_price': base_price,
                    'tier_range': tier_range,
                    'tier_column': tier_column,
                    'product_ref': product_data.get('Partner Product SKU/REF', 'N/A'),
                    'country_of_origin_made_in': product_data.get('Country of Origin (Made In)', 'N/A'),
                    'country_of_origin_ships_from': product_data.get('Country of Origin (Ships From)', 'N/A'),
                    'customization_description': product_data.get('Customization Info', ''),
                    'product_subtotal': base_price * 1,
                    'customization_setup_total': 0.0,
                    'customization_unit_total': 0.0,
                    'subtotal_before_markup': base_price * 1,
                    'markup_amount': (base_price * 1) * (markup / 100),
                    'product_total': (base_price * 1) + ((base_price * 1) * (markup / 100)),
                    'total_per_unit': ((base_price * 1) + ((base_price * 1) * (markup / 100))) / 1,
                    'tariff_rate_percent': 0.0,
                    'tariff_amount': 0.0,
                    'edited_description': ''  # Initialize with empty for user to edit later
                })

                st.session_state.order_items.append(new_item)

                # Auto-populate partner shipping if this product has shipping data
                pbp_shipping_cost, _ = get_shipping_costs(product_data.to_dict())
                if pbp_shipping_cost > 0 and st.session_state.partner_shipping == 0:
                    st.session_state.partner_shipping = pbp_shipping_cost

                # Set success message for deferred toast (include variant if applicable)
                from src.helpers import format_product_with_variant
                st.session_state.show_add_to_order_success = True
                st.session_state.add_to_order_product_name = format_product_with_variant(
                    product_data.get('Product/Service', 'product'),
                    selected_variant_manual
                )
                st.rerun()
            else:
                st.error("Could not determine pricing for this product")

    # Show product details in expander
    product_data_preview = df_template[
        (df_template["Partner"] == selected_partner) &
        (df_template["Product/Service"] == selected_product)
    ].iloc[0]

    origin_made = product_data_preview.get("Country of Origin (Made In)", "N/A")
    origin_ships = product_data_preview.get("Country of Origin (Ships From)", "N/A")
    has_tiers = product_data_preview.get("Pricing Tiers (Y/N)", "N/A")

    with st.expander("Show Product Details"):
        col1, col2 = st.columns(2)
        with col1:
            st.markdown(f"**Partner:** {product_data_preview['Partner']}")
            st.markdown(f"**Product/Service:** {product_data_preview['Product/Service']}")
        with col2:
            st.markdown(f"**Made In:** {origin_made if origin_made else 'N/A'}")
            st.markdown(f"**Ships From:** {origin_ships if origin_ships else 'N/A'}")
            st.markdown(f"**Tiered Pricing:** {has_tiers}")

        # Show product description if available
        description = product_data_preview.get("Marketing Description", "")
        if description and description.strip():
            st.markdown("---")
            st.markdown("**Marketing Description:**")
            st.write(description)

        # Show pricing tier info if applicable
        tier_info = product_data_preview.get("Pricing Tiers Info", "")
        if tier_info and tier_info.strip() and tier_info != "NA":
            st.markdown("---")
            st.markdown("**Pricing Tier Information:**")
            st.markdown("This product uses tiered pricing - the price per unit decreases as you order more. The tier ranges below show which price applies based on your order quantity.")
            st.markdown("")
            st.markdown(f"**Tier Ranges:** {tier_info}")
            st.caption("Your order quantity will automatically match to the correct tier and price.")

    st.divider()

    # ============================================================
    # OPTION D: CREATE CUSTOM PRODUCT
    # ============================================================
    option_label_d = "Option D" if has_proposal else "Option C"
    st.header(f"{option_label_d}: Create Custom Product")
    st.markdown("**Use this if:** You need to add a unique product not in the catalog")
    st.caption("Create one-off items, executive samples, or products pending catalog addition")

    with st.expander("Create Custom Product", expanded=False):
        st.caption("Enter basic info - you'll configure quantity, markup, and customization after adding")

        col1, col2, col3 = st.columns([2, 2, 1])

        with col1:
            custom_product_name = st.text_input(
                "Product Name*",
                key="custom_product_name",
                placeholder="e.g., Custom Gold Engraving"
            )

        with col2:
            # Partner selection with "Custom/Other" as first option
            partner_options = ["Custom/Other"] + sorted(df_template['Partner'].unique().tolist())
            custom_partner = st.selectbox(
                "Partner*",
                options=partner_options,
                key="custom_partner",
                help="Select partner for POC tracking, or 'Custom/Other'"
            )

        with col3:
            custom_base_cost = st.number_input(
                "Base Cost/Unit*",
                min_value=0.01,
                value=10.00,
                step=0.50,
                key="custom_base_cost",
                format="%.2f",
                help="What PBP pays per unit"
            )

        # Add button
        if st.button("Add to Order", key="add_custom_product_btn", type="primary", use_container_width=True):
            # Validation
            if not custom_product_name or len(custom_product_name.strip()) < 3:
                st.error("Product name is required (min 3 characters)")
            elif custom_base_cost <= 0:
                st.error("Base cost must be greater than $0")
            else:
                # Create custom product with DEFAULTS (just like catalog products)
                # User will configure quantity, markup, customization inline after adding

                # Create minimal product_data dict (simulates spreadsheet row)
                custom_product_data = {
                    'Product/Service': custom_product_name.strip(),
                    'Partner': custom_partner,
                    'Pricing Tiers (Y/N)': 'N',  # Always flat-rate for custom
                    'PBP Cost (No Tiers)': custom_base_cost,
                    'PBP Standard Markup': 100.0,  # Default markup
                    'Country of Origin (Made In)': '',
                    'Country of Origin (Ships From)': '',
                    'Vendor Published MSRP': 0,
                    'Customization Setup Fee': 0,
                    'Customization Cost per Unit': 0,
                    'Customization Info': '',
                    'Marketing Description': '',
                    'Tariff Estimate (%)': 0,
                    'MOQ (PBP)': '',
                    'MOV (PBP)': '',
                    'MOQ (Partner)': '',
                    'MOV (Partner)': '',
                    'Tariff Info': '',
                    'Purchase Description': '',
                    'Units per Package': 1,
                }

                # Use same structure as catalog products
                base_price = custom_base_cost  # No tiers, so base price = entered cost
                tier_range = "No Tiers"
                tier_column = "PBP Cost (No Tiers)"
                markup = 100.0  # Default markup

                # Create order item with DEFAULTS (same structure as catalog products + Phase 3 fields)
                new_item = {
                    'product_name': custom_product_name.strip(),
                    'partner': custom_partner,
                    'product_data': custom_product_data,
                    'quantity': 1,  # DEFAULT - user edits inline
                    'markup_percent': markup,  # DEFAULT - user edits inline
                    'selected_variant': None,
                    # Phase 3: New pricing fields (custom products use Standard markup)
                    'pricing_method': 'Standard markup',
                    'pricing_notes': '',
                    'validation_warning': None,
                    'manual_override': False,
                    'include_customization': False,  # DEFAULT - user enables inline
                    'customization_setup_fee': 0.0,
                    'customization_per_unit': 0.0,
                    'customization_minimum_qty': 0,
                    'apply_custom_minimum': False,
                    'include_tariff': False,
                    'is_custom_product': True,  # FLAG to distinguish from catalog and legacy custom
                    'source': 'custom',
                    # Per-product kitting fields
                    'include_kitting': False,
                    'kitting_pbp_cost': 0.0,
                    'kitting_client_price': 0.0,
                    'kitting_description': ''
                }

                # Add calculated fields (same as catalog products)
                new_item.update({
                    'base_price': base_price,
                    'tier_range': tier_range,
                    'tier_column': tier_column,
                    'product_ref': 'CUSTOM',
                    'country_of_origin_made_in': '',
                    'country_of_origin_ships_from': '',
                    'customization_description': '',
                    'product_subtotal': base_price * 1,
                    'customization_setup_total': 0.0,
                    'customization_unit_total': 0.0,
                    'subtotal_before_markup': base_price * 1,
                    'markup_amount': (base_price * 1) * (markup / 100),
                    'product_total': (base_price * 1) + ((base_price * 1) * (markup / 100)),
                    'total_per_unit': ((base_price * 1) + ((base_price * 1) * (markup / 100))) / 1,
                    'tariff_rate_percent': 0.0,
                    'tariff_amount': 0.0,
                    'edited_description': ''
                })

                # Add to order
                st.session_state.order_items.append(new_item)

                # Set success message for deferred toast (same pattern as Option C)
                st.session_state.show_add_to_order_success = True
                st.session_state.add_to_order_product_name = custom_product_name.strip()

                st.rerun()

    st.divider()

    # ============================================================
    # CURRENT ORDER (with inline editing)
    # ============================================================
    st.header("2. Current Order")

    if len(st.session_state.order_items) == 0:
        st.info("No products in order yet. Add products above to get started.")
    else:
        st.caption("Edit settings for each product below. Changes update totals in real-time.")
        st.write("")

        # Iterate through order items and display editable cards
        for idx, item in enumerate(st.session_state.order_items):
            # Check if this is OLD-STYLE custom line item (before custom product enhancement)
            if item.get('is_custom', False) and not item.get('is_custom_product', False):
                # Legacy custom line items - show simplified view
                st.write("---")
                st.subheader(f"{item['product_name']}")
                st.caption(f"Custom Line Item (Legacy)")

                col1, col2 = st.columns([4, 1])
                with col1:
                    st.write(f"**Description:** {item.get('custom_description', 'N/A')}")
                    st.write(f"**Quantity:** {item['quantity']} | **Unit Price:** ${item['total_per_unit']:.2f} | **Total:** ${item['product_total']:.2f}")
                with col2:
                    if st.button("Remove", key=f"remove_custom_{idx}", type="secondary"):
                        st.session_state.order_items.pop(idx)
                        st.rerun()
                continue

            # Regular product card (INCLUDING new custom products with is_custom_product=True)
            st.write("---")

            # Header with product name and remove button
            col_header, col_remove = st.columns([5, 1])
            with col_header:
                # Display product name with variant (if applicable)
                from src.helpers import format_product_with_variant
                product_display_name = format_product_with_variant(
                    item['product_name'],
                    item.get('selected_variant')
                )
                st.subheader(f"{product_display_name}")

                # Show custom product indicator or regular product info
                if item.get('is_custom_product', False):
                    st.caption(f"Custom Product | Partner: {item['partner']} | Base Cost: ${item['base_price']:.2f}/unit")
                else:
                    st.caption(f"Partner: {item['partner']} | Origin: {item.get('country_of_origin', 'N/A')}")

                    # Show price source indicator (pricing snapshot vs. recalculated) - only for catalog products
                    if item.get('from_proposal_snapshot'):
                        st.caption("✓ Pricing preserved from proposal (saved configuration)")
                    elif item.get('source') == 'proposal':
                        st.caption("Pricing recalculated from current spreadsheet")
            with col_remove:
                if st.button("Remove", key=f"remove_product_{idx}", type="secondary"):
                    st.session_state.order_items.pop(idx)
                    st.rerun()

            # Get product data for recalculations
            product_data = item['product_data']

            # PRICING METHOD & MANUAL OVERRIDE (Phase 3)
            pricing_method = item.get('pricing_method', 'Standard markup')
            manual_override = item.get('manual_override', False)

            # QUANTITY & PRICING SECTION
            st.markdown("##### Quantity & Pricing")

            # Calculate base price for current quantity first
            new_quantity = item['quantity']  # Default value
            base_price, tier_range, tier_column = get_unit_price_new_system(product_data, item['quantity'])

            # Table header
            header_col1, header_col2, header_col3, header_col4, header_col5, header_col6 = st.columns([1.2, 1.2, 1.2, 1.5, 1.5, 1.2])
            with header_col1:
                st.markdown("**Quantity**")
            with header_col2:
                st.markdown("**PBP Cost**")
            with header_col3:
                st.markdown("**Vendor MSRP**")
            with header_col4:
                st.markdown("**Pricing Method**")
            with header_col5:
                st.markdown("**PBP MSRP (Client Price)**")
            with header_col6:
                st.markdown("**Markup %**")

            # Data row
            col1, col2, col3, col4, col5, col6 = st.columns([1.2, 1.2, 1.2, 1.5, 1.5, 1.2])

            with col1:
                new_quantity = st.number_input(
                    f"Quantity for {idx}",
                    min_value=1,
                    value=item['quantity'],
                    step=1,
                    key=f"prod_qty_{idx}",
                    label_visibility="collapsed"
                )
                # Highlight if quantity is 1 (warning)
                if new_quantity == 1:
                    st.caption("WARNING: Qty = 1")

                # Recalculate base price if quantity changed
                if new_quantity != item['quantity']:
                    base_price, tier_range, tier_column = get_unit_price_new_system(product_data, new_quantity)

            with col2:
                # Show PBP cost
                if base_price:
                    st.markdown(f"${base_price:.2f}")
                else:
                    st.markdown("—")

            with col3:
                # Show Vendor MSRP if available
                msrp_raw = get_column_value(product_data, 'Vendor Published MSRP', 'MSRP', '')
                msrp = clean_price(msrp_raw)
                if msrp and msrp > 0:
                    st.markdown(f"${msrp:.2f}")
                else:
                    st.markdown("—")

            with col4:
                # Show pricing method
                if manual_override:
                    st.markdown("Manual override")
                elif pricing_method:
                    st.markdown(f"{pricing_method}")
                else:
                    st.markdown("Standard markup")

            with col5:
                # Calculate client price (base + markup, before customization)
                if base_price:
                    product_subtotal_calc = base_price * new_quantity

                    # Use current markup for calculation
                    current_markup = st.session_state.get(f"prod_markup_{idx}", item['markup_percent'])
                    if not st.session_state.get(f'updating_from_price_tab3_{idx}', False):
                        current_markup = item['markup_percent']  # Will be updated by markup input below

                    markup_amount_calc = product_subtotal_calc * (current_markup / 100)
                    client_price_raw = product_subtotal_calc + markup_amount_calc
                    client_price_per_unit_unrounded = client_price_raw / new_quantity

                    # Apply $0.50 rounding if enabled (default: True)
                    client_price_per_unit = round_to_nearest_fifty_cents(
                        client_price_per_unit_unrounded,
                        st.session_state.order_fifty_cent_rounding
                    )

                    # Check if we're updating from markup to prevent circular updates
                    if st.session_state.get(f'updating_from_markup_tab3_{idx}', False):
                        # Just display the calculated price, don't create input
                        st.markdown(f"${client_price_per_unit:.2f}")
                        # Show rounding note if price was rounded
                        if st.session_state.order_fifty_cent_rounding and abs(client_price_per_unit - client_price_per_unit_unrounded) > 0.01:
                            st.caption(f"Rounded from ${client_price_per_unit_unrounded:.2f}")
                        # Clear the flag
                        st.session_state[f'updating_from_markup_tab3_{idx}'] = False
                    else:
                        # Editable client price field
                        new_client_price = st.number_input(
                            f"Client Price for {idx}",
                            min_value=0.01,
                            value=client_price_per_unit,
                            step=1.0,
                            format="%.2f",
                            key=f"prod_price_{idx}",
                            label_visibility="collapsed"
                        )
                        # Show rounding note if price was rounded
                        if st.session_state.order_fifty_cent_rounding and abs(client_price_per_unit - client_price_per_unit_unrounded) > 0.01:
                            st.caption(f"Rounded from ${client_price_per_unit_unrounded:.2f}")

                        # Update markup if price changed
                        if abs(new_client_price - client_price_per_unit) > 0.01:  # Check if meaningfully different
                            # Calculate new markup from the price
                            new_markup_calc = calculate_markup_from_price(base_price, new_client_price)
                            st.session_state.order_items[idx]['markup_percent'] = new_markup_calc
                            st.session_state[f'updating_from_price_tab3_{idx}'] = True
                            st.rerun()
                else:
                    st.markdown("—")

            with col6:
                # Check if we're updating from client price to prevent circular updates
                if st.session_state.get(f'updating_from_price_tab3_{idx}', False):
                    # Just display the calculated markup, don't create input
                    st.markdown(f"{item['markup_percent']:.1f}%")
                    # Clear the flag
                    st.session_state[f'updating_from_price_tab3_{idx}'] = False
                else:
                    new_markup = st.number_input(
                        f"Markup for {idx}",
                        min_value=-50.0,  # Allow negative markup for below-cost pricing
                        value=item['markup_percent'],
                        step=5.0,
                        key=f"prod_markup_{idx}",
                        label_visibility="collapsed"
                    )
                    # Set flag if markup changed to prevent circular updates
                    if new_markup != item['markup_percent']:
                        st.session_state[f'updating_from_markup_tab3_{idx}'] = True

            st.divider()

            # Show tier info below table
            if tier_range != "No Tiers":
                st.caption(f"Using tier: {tier_range} | ${base_price:.2f}/unit")
            else:
                st.caption(f"Flat pricing: ${base_price:.2f}/unit")

            # Manual override checkbox
            new_manual_override = st.checkbox(
                "Enable Manual Price Override",
                value=manual_override,
                key=f"manual_override_tab3_{idx}",
                help="Check to manually override the pricing method and set custom markup/price"
            )
            if new_manual_override != manual_override:
                st.session_state.order_items[idx]['manual_override'] = new_manual_override
                st.rerun()

            # Show quoted price warning if this came from a proposal and price changed
            if base_price:
                quoted_price = item.get('quoted_price_per_unit', 0.0)
                if quoted_price > 0:
                    # Product came from proposal - show comparison
                    if abs(client_price_per_unit - quoted_price) > 0.01:  # Allow for rounding errors
                        # Determine WHY the price changed
                        reasons = []

                        # Check if tier changed
                        proposal_tier = item.get('proposal_tier_column', '')
                        current_tier = item.get('tier_column', '')
                        proposal_tier_range = item.get('proposal_tier_range', '')
                        current_tier_range = tier_range

                        if proposal_tier and current_tier and proposal_tier != current_tier:
                            reasons.append(f"Tier change: {proposal_tier} ({proposal_tier_range}) → {current_tier} ({current_tier_range})")

                        # Check if markup changed
                        proposal_markup = item.get('proposal_markup_percent', 0)
                        current_markup = item['markup_percent']  # Use the actual markup from the item

                        if abs(proposal_markup - current_markup) > 0.01:
                            reasons.append(f"Markup change: {proposal_markup:.0f}% → {current_markup:.0f}%")

                        # Display warning with reasons
                        if reasons:
                            reason_text = " | ".join(reasons)
                            st.warning(f"WARNING: Price changed from proposal (${quoted_price:.2f}/unit → ${client_price_per_unit:.2f}/unit)\n\nReason: {reason_text}")
                        else:
                            st.warning(f"WARNING: Price changed from proposal: Quoted price was ${quoted_price:.2f}/unit, current is ${client_price_per_unit:.2f}/unit")
                    else:
                        st.info(f"Matches quoted price: ${quoted_price:.2f}/unit")
                else:
                    # Product added manually (not from proposal)
                    st.caption("Quoted price: Unknown (product added manually)")

            # VALIDATION WARNING (Phase 3) - Only show if not manually overridden
            validation_warning = item.get('validation_warning', None)
            if validation_warning and not manual_override:
                st.warning(f"**Pricing Validation:** {validation_warning}")
                st.caption("The calculated price doesn't match the spreadsheet value. Enable manual override if this is intentional.")

            # PRICING NOTES (Phase 3) - Only show if notes exist and not manually overridden
            pricing_notes = item.get('pricing_notes', '')
            if pricing_notes and pricing_notes.strip() and not manual_override:
                with st.expander("ℹ️ Pricing Information", expanded=False):
                    st.caption(pricing_notes)

            # CUSTOMIZATION SECTION (always available)
            customization_info = product_data.get("Customization Info", "")
            st.markdown("##### Customization (Optional)")
            if customization_info and customization_info.strip():
                st.caption(f"Available: {customization_info}")
            else:
                st.caption("No customization details from spreadsheet - set custom values below if needed")

            new_include_custom = st.checkbox(
                "Include Customization",
                value=item.get('include_customization', False),
                key=f"prod_custom_{idx}"
            )

            if new_include_custom:
                    col_setup, col_perunit = st.columns(2)

                    with col_setup:
                        st.markdown("**Client Pricing:**")
                        # Always read default from product_data, use item value if user has edited it
                        default_setup_raw = get_column_value(product_data, 'Client Price: Customization Setup Fee', 'Customization Setup Fee', '')
                        default_setup = clean_price(default_setup_raw) or 0.0
                        stored_setup = item.get('customization_setup_fee', 0.0)
                        # Use stored value if it's non-zero OR if no default exists, otherwise use default
                        display_setup = stored_setup if (stored_setup > 0 or default_setup == 0) else default_setup

                        new_setup_fee = st.number_input(
                            "Setup Fee (to Client)",
                            min_value=0.0,
                            value=float(display_setup),
                            step=1.0,
                            key=f"prod_setup_{idx}"
                        )

                        st.markdown("**Partner Cost:**")
                        # Load partner cost from spreadsheet
                        default_partner_setup = clean_price(product_data.get('PBP Price: Customization Setup Fee', '')) or 0.0
                        stored_partner_setup = item.get('partner_customization_setup_fee', 0.0)
                        display_partner_setup = stored_partner_setup if (stored_partner_setup > 0 or default_partner_setup == 0) else default_partner_setup

                        new_partner_setup_fee = st.number_input(
                            "Setup Fee (from Partner)",
                            min_value=0.0,
                            value=float(display_partner_setup),
                            step=1.0,
                            key=f"prod_partner_setup_{idx}",
                            help="Cost PBP pays to partner for setup"
                        )

                    with col_perunit:
                        st.markdown("**Client Pricing:**")
                        # Always read default from product_data, use item value if user has edited it
                        default_perunit_raw = get_column_value(product_data, 'Client Price: Customization Cost per Unit', 'Customization Cost per Unit', '')
                        default_perunit = clean_price(default_perunit_raw) or 0.0
                        stored_perunit = item.get('customization_per_unit', 0.0)
                        # Use stored value if it's non-zero OR if no default exists, otherwise use default
                        display_perunit = stored_perunit if (stored_perunit > 0 or default_perunit == 0) else default_perunit

                        new_perunit_cost = st.number_input(
                            "Per-Unit Cost (to Client)",
                            min_value=0.0,
                            value=float(display_perunit),
                            step=0.1,
                            key=f"prod_perunit_{idx}"
                        )

                        st.markdown("**Partner Cost:**")
                        # Load partner cost from spreadsheet
                        default_partner_perunit = clean_price(product_data.get('PBP Price: Customization Cost per Unit', '')) or 0.0
                        stored_partner_perunit = item.get('partner_customization_per_unit', 0.0)
                        display_partner_perunit = stored_partner_perunit if (stored_partner_perunit > 0 or default_partner_perunit == 0) else default_partner_perunit

                        new_partner_perunit_cost = st.number_input(
                            "Per-Unit Cost (from Partner)",
                            min_value=0.0,
                            value=float(display_partner_perunit),
                            step=0.1,
                            key=f"prod_partner_perunit_{idx}",
                            help="Cost PBP pays to partner per unit"
                        )

                    # Customization minimum
                    new_apply_minimum = st.checkbox(
                        "Apply minimum quantity for customization",
                        value=item.get('apply_custom_minimum', False),
                        key=f"prod_apply_min_{idx}",
                        help="Charge for minimum units even if ordering fewer"
                    )

                    if new_apply_minimum:
                        # Get stored minimum or default to max(100, current_quantity)
                        stored_min = item.get('customization_minimum_qty', 0)
                        default_min = max(100, new_quantity)
                        display_min = stored_min if stored_min > 0 else default_min

                        new_custom_min_qty = st.number_input(
                            "Minimum Customization Quantity",
                            min_value=1,
                            value=display_min,
                            step=1,
                            key=f"prod_min_qty_{idx}"
                        )

                        if new_custom_min_qty > new_quantity:
                            st.info(f"Charging for {new_custom_min_qty} customization units (ordering {new_quantity} product units)")
                    else:
                        new_custom_min_qty = 0

                    # CUSTOMIZATION ADD-ONS SECTION
                    st.divider()
                    st.markdown("**Customization Add-Ons**")
                    st.caption("Add additional customization options (e.g., 2nd color, special materials, extra features)")

                    # Initialize add-ons list if not exists
                    if 'customization_addons' not in item:
                        item['customization_addons'] = []

                    # Button to add new add-on
                    if st.button("➕ Add Customization Option", key=f"add_addon_{idx}"):
                        item['customization_addons'].append({
                            'name': '',
                            'client_setup_fee': 0.0,
                            'client_per_unit_cost': 0.0,
                            'partner_setup_fee': 0.0,
                            'partner_per_unit_cost': 0.0
                        })
                        st.rerun()

                    # Display existing add-ons
                    addons_to_remove = []
                    for addon_idx, addon in enumerate(item.get('customization_addons', [])):
                        # Backward compatibility: migrate old field names to new ones
                        if 'setup_fee' in addon and 'client_setup_fee' not in addon:
                            # Assume old values were client prices
                            addon['client_setup_fee'] = addon.get('setup_fee', 0.0)
                            addon['partner_setup_fee'] = 0.0  # Default to zero
                        if 'per_unit_cost' in addon and 'client_per_unit_cost' not in addon:
                            addon['client_per_unit_cost'] = addon.get('per_unit_cost', 0.0)
                            addon['partner_per_unit_cost'] = 0.0  # Default to zero

                        with st.expander(f"Add-On {addon_idx + 1}: {addon.get('name', 'Unnamed')}", expanded=True):
                            # Option name and remove button
                            name_col, remove_col = st.columns([4, 0.5])
                            with name_col:
                                addon['name'] = st.text_input(
                                    "Option Name",
                                    value=addon.get('name', ''),
                                    placeholder="e.g., Second Color, Premium Wood",
                                    key=f"addon_name_{idx}_{addon_idx}"
                                )
                            with remove_col:
                                st.write("")  # Spacing
                                if st.button("Remove", key=f"remove_addon_{idx}_{addon_idx}"):
                                    addons_to_remove.append(addon_idx)

                            # Client Pricing Section
                            st.markdown("**Client Pricing:**")
                            col1, col2 = st.columns(2)
                            with col1:
                                addon['client_setup_fee'] = st.number_input(
                                    "Setup Fee (to Client)",
                                    min_value=0.0,
                                    value=float(addon.get('client_setup_fee', 0.0)),
                                    step=1.0,
                                    key=f"addon_client_setup_{idx}_{addon_idx}",
                                    help="One-time setup fee charged to client for this option"
                                )
                            with col2:
                                addon['client_per_unit_cost'] = st.number_input(
                                    "Per Unit Cost (to Client)",
                                    min_value=0.0,
                                    value=float(addon.get('client_per_unit_cost', 0.0)),
                                    step=0.1,
                                    key=f"addon_client_perunit_{idx}_{addon_idx}",
                                    help="Additional cost per unit charged to client"
                                )

                            # Partner Cost Section
                            st.markdown("**Partner Cost:**")
                            col3, col4 = st.columns(2)
                            with col3:
                                addon['partner_setup_fee'] = st.number_input(
                                    "Setup Fee (from Partner)",
                                    min_value=0.0,
                                    value=float(addon.get('partner_setup_fee', 0.0)),
                                    step=1.0,
                                    key=f"addon_partner_setup_{idx}_{addon_idx}",
                                    help="Cost PBP pays to partner for setup"
                                )
                            with col4:
                                addon['partner_per_unit_cost'] = st.number_input(
                                    "Per Unit Cost (from Partner)",
                                    min_value=0.0,
                                    value=float(addon.get('partner_per_unit_cost', 0.0)),
                                    step=0.1,
                                    key=f"addon_partner_perunit_{idx}_{addon_idx}",
                                    help="Cost PBP pays to partner per unit"
                                )

                    # Remove add-ons marked for deletion
                    if addons_to_remove:
                        for addon_idx in sorted(addons_to_remove, reverse=True):
                            item['customization_addons'].pop(addon_idx)
                        st.rerun()

                    # Show total add-on costs
                    if item.get('customization_addons'):
                        # Calculate totals for client and partner (with backward compatibility)
                        total_addon_client_setup = sum(
                            addon.get('client_setup_fee', addon.get('setup_fee', 0.0))
                            for addon in item['customization_addons'] if addon.get('name')
                        )
                        total_addon_client_perunit = sum(
                            addon.get('client_per_unit_cost', addon.get('per_unit_cost', 0.0))
                            for addon in item['customization_addons'] if addon.get('name')
                        )
                        total_addon_partner_setup = sum(
                            addon.get('partner_setup_fee', 0.0)
                            for addon in item['customization_addons'] if addon.get('name')
                        )
                        total_addon_partner_perunit = sum(
                            addon.get('partner_per_unit_cost', 0.0)
                            for addon in item['customization_addons'] if addon.get('name')
                        )

                        # Display with clear distinction
                        if total_addon_client_setup > 0 or total_addon_client_perunit > 0 or total_addon_partner_setup > 0 or total_addon_partner_perunit > 0:
                            col_summary1, col_summary2 = st.columns(2)
                            with col_summary1:
                                st.info(f"**Client Total:** ${total_addon_client_setup:.2f} setup + ${total_addon_client_perunit:.2f}/unit")
                            with col_summary2:
                                st.info(f"**Partner Total:** ${total_addon_partner_setup:.2f} setup + ${total_addon_partner_perunit:.2f}/unit")
            else:
                new_setup_fee = 0.0
                new_perunit_cost = 0.0
                new_partner_setup_fee = 0.0
                new_partner_perunit_cost = 0.0
                new_apply_minimum = False
                new_custom_min_qty = 0
                # Clear add-ons if customization is disabled
                item['customization_addons'] = []

            # COUNTRY & TARIFF SECTION (only for custom products)
            if item.get('is_custom_product', False):
                st.markdown("##### Country & Tariff")
                st.caption("Set country of origin for tariff estimates and tracking")

                col1, col2 = st.columns(2)
                with col1:
                    country_made = st.text_input(
                        "Made In",
                        value=item.get('country_of_origin_made_in', ''),
                        placeholder="e.g., USA, China, India, Vietnam, Mexico",
                        key=f"custom_made_in_{idx}",
                        help="Enter country where product is manufactured"
                    )

                    # Update item if changed
                    if country_made != item.get('country_of_origin_made_in', ''):
                        st.session_state.order_items[idx]['country_of_origin_made_in'] = country_made
                        # Update product_data as well
                        st.session_state.order_items[idx]['product_data']['Country of Origin (Made In)'] = country_made
                        # Auto-set ships from to same as made in if not set
                        if not item.get('country_of_origin_ships_from'):
                            st.session_state.order_items[idx]['country_of_origin_ships_from'] = country_made
                            st.session_state.order_items[idx]['product_data']['Country of Origin (Ships From)'] = country_made

                with col2:
                    country_ships = st.text_input(
                        "Ships From",
                        value=item.get('country_of_origin_ships_from', ''),
                        placeholder="e.g., USA, China, India, Vietnam, Mexico",
                        key=f"custom_ships_from_{idx}",
                        help="Enter country where product ships from (affects tariffs)"
                    )

                    # Update item if changed
                    if country_ships != item.get('country_of_origin_ships_from', ''):
                        st.session_state.order_items[idx]['country_of_origin_ships_from'] = country_ships
                        st.session_state.order_items[idx]['product_data']['Country of Origin (Ships From)'] = country_ships

                # Tariff estimate
                st.caption("Tariff estimate (optional, for reference)")
                tariff_pct = st.number_input(
                    "Tariff Estimate (%)",
                    min_value=0.0,
                    max_value=100.0,
                    value=item.get('tariff_rate_percent', 0.0),
                    step=0.5,
                    key=f"custom_tariff_{idx}",
                    format="%.1f",
                    help="Estimated tariff rate based on country and product type"
                )

                if tariff_pct != item.get('tariff_rate_percent', 0.0):
                    st.session_state.order_items[idx]['tariff_rate_percent'] = tariff_pct

            # KITTING SECTION (per-product)
            st.markdown("##### Kitting & Packaging")
            st.caption("Add costs for repackaging, special boxes, or product-specific assembly")

            # Checkbox to enable kitting for this product
            kitting_enabled = st.checkbox(
                "Include kitting costs for this product",
                value=item.get('include_kitting', False),
                key=f"prod_kitting_enabled_{idx}",
                help="Add costs for repackaging, special boxes, or product-specific assembly"
            )

            if kitting_enabled:
                col_kit1, col_kit2 = st.columns(2)

                with col_kit1:
                    kitting_pbp = st.number_input(
                        "Kitting Cost (PBP) ($)",
                        min_value=0.0,
                        value=item.get('kitting_pbp_cost', 0.0),
                        step=5.0,
                        key=f"prod_kitting_pbp_{idx}",
                        help="What PBP pays for this product's kitting"
                    )

                with col_kit2:
                    kitting_client = st.number_input(
                        "Kitting Price (Client) ($)",
                        min_value=0.0,
                        value=item.get('kitting_client_price', 0.0),
                        step=5.0,
                        key=f"prod_kitting_client_{idx}",
                        help="What client pays for this product's kitting"
                    )

                # Optional description
                kitting_desc = st.text_input(
                    "Description (Optional)",
                    value=item.get('kitting_description', ''),
                    key=f"prod_kitting_desc_{idx}",
                    placeholder="e.g., Premium gift box, Repackaging"
                )

                # Update order item
                st.session_state.order_items[idx]['include_kitting'] = True
                st.session_state.order_items[idx]['kitting_pbp_cost'] = kitting_pbp
                st.session_state.order_items[idx]['kitting_client_price'] = kitting_client
                st.session_state.order_items[idx]['kitting_description'] = kitting_desc
            else:
                # Reset kitting if unchecked
                st.session_state.order_items[idx]['include_kitting'] = False
                st.session_state.order_items[idx]['kitting_pbp_cost'] = 0.0
                st.session_state.order_items[idx]['kitting_client_price'] = 0.0
                st.session_state.order_items[idx]['kitting_description'] = ''

            # RECALCULATE PRICING
            # Get base price for new quantity
            base_price, tier_range, tier_column = get_unit_price_new_system(product_data, new_quantity)

            if base_price:
                # Calculate customization costs
                if new_include_custom:
                    effective_custom_qty = new_custom_min_qty if (new_apply_minimum and new_custom_min_qty > new_quantity) else new_quantity
                    customization_setup_total = new_setup_fee
                    customization_unit_total = new_perunit_cost * effective_custom_qty

                    # Add customization add-on costs (client pricing)
                    for addon in item.get('customization_addons', []):
                        if addon.get('name'):  # Only include add-ons with names
                            # Use new field names with backward compatibility
                            customization_setup_total += addon.get('client_setup_fee', addon.get('setup_fee', 0.0))
                            customization_unit_total += addon.get('client_per_unit_cost', addon.get('per_unit_cost', 0.0)) * effective_custom_qty
                else:
                    customization_setup_total = 0.0
                    customization_unit_total = 0.0

                # Calculate totals
                product_subtotal = base_price * new_quantity
                subtotal_before_markup = product_subtotal + customization_setup_total + customization_unit_total
                # Use the item's markup_percent which gets updated by the input fields
                markup_amount = product_subtotal * (item['markup_percent'] / 100)
                product_total = subtotal_before_markup + markup_amount
                total_per_unit = product_total / new_quantity

                # Calculate partner customization costs (for accounting)
                if new_include_custom:
                    effective_custom_qty = new_custom_min_qty if (new_apply_minimum and new_custom_min_qty > new_quantity) else new_quantity
                    partner_customization_setup_total = new_partner_setup_fee
                    partner_customization_unit_total = new_partner_perunit_cost * effective_custom_qty

                    # Add partner costs for add-ons
                    for addon in item.get('customization_addons', []):
                        if addon.get('name'):  # Only include add-ons with names
                            partner_customization_setup_total += addon.get('partner_setup_fee', 0.0)
                            partner_customization_unit_total += addon.get('partner_per_unit_cost', 0.0) * effective_custom_qty
                else:
                    partner_customization_setup_total = 0.0
                    partner_customization_unit_total = 0.0

                # Update item in session state
                # Get the current markup (either from the input field or the item's current value)
                current_markup = item['markup_percent']  # Default to current value
                # If there's a markup input field in session state, use that
                if f"prod_markup_{idx}" in st.session_state:
                    current_markup = st.session_state[f"prod_markup_{idx}"]

                st.session_state.order_items[idx].update({
                    'quantity': new_quantity,
                    'markup_percent': current_markup,
                    'include_customization': new_include_custom,
                    'customization_setup_fee': new_setup_fee,
                    'customization_per_unit': new_perunit_cost,
                    'partner_customization_setup_fee': new_partner_setup_fee,
                    'partner_customization_per_unit': new_partner_perunit_cost,
                    'partner_customization_setup_total': partner_customization_setup_total,
                    'partner_customization_unit_total': partner_customization_unit_total,
                    'apply_custom_minimum': new_apply_minimum,
                    'customization_minimum_qty': new_custom_min_qty,
                    'base_price': base_price,
                    'tier_range': tier_range,
                    'tier_column': tier_column,
                    'product_subtotal': product_subtotal,
                    'customization_setup_total': customization_setup_total,
                    'customization_unit_total': customization_unit_total,
                    'subtotal_before_markup': subtotal_before_markup,
                    'markup_amount': markup_amount,
                    'product_total': product_total,
                    'total_per_unit': total_per_unit,
                    'tariff_base': product_subtotal  # Base for tariff calculation (product cost only, excludes markup and customization)
                })

                # Recalculate tariff if rate is set
                if item.get('tariff_rate_percent', 0) > 0:
                    item['tariff_amount'] = calculate_product_tariff(
                        product_subtotal,
                        item['tariff_rate_percent']
                    )

                # PRICING BREAKDOWN DISPLAY
                st.markdown("##### Pricing Breakdown")
                st.caption(f"**Product: {item['product_name']}**")

                # Import helper function for formatting rows
                from src.helpers import format_pricing_breakdown_row

                breakdown_data = []

                # Base product row - include product name in description
                product_pbp_cost = product_subtotal
                client_per_unit = (product_subtotal + markup_amount) / new_quantity if new_quantity > 0 else 0

                # Apply $0.50 rounding to client per unit price (default: True)
                client_per_unit = round_to_nearest_fifty_cents(
                    client_per_unit,
                    st.session_state.order_fifty_cent_rounding
                )

                # Recalculate total based on rounded per-unit price
                product_client_price = client_per_unit * new_quantity

                breakdown_data.append(
                    format_pricing_breakdown_row(
                        f"Base Product: {item['product_name']}",
                        new_quantity,
                        base_price,  # PBP per unit
                        product_pbp_cost,  # PBP total
                        client_per_unit,  # Client per unit
                        product_client_price  # Client total
                    )
                )

                # Customization setup fee (if applicable)
                if customization_setup_total > 0:
                    breakdown_data.append(
                        format_pricing_breakdown_row(
                            "Customization Setup",
                            "one-time",
                            partner_customization_setup_total,  # PBP per unit (same as total for one-time)
                            partner_customization_setup_total,  # PBP total
                            customization_setup_total,  # Client per unit (same as total for one-time)
                            customization_setup_total  # Client total
                        )
                    )

                # Customization per-unit (if applicable)
                if customization_unit_total > 0:
                    effective_custom_qty = new_custom_min_qty if (new_apply_minimum and new_custom_min_qty > new_quantity) else new_quantity
                    # Calculate per-unit costs including add-ons
                    partner_per_unit_with_addons = partner_customization_unit_total / effective_custom_qty if effective_custom_qty > 0 else 0
                    client_per_unit_with_addons = customization_unit_total / effective_custom_qty if effective_custom_qty > 0 else 0
                    breakdown_data.append(
                        format_pricing_breakdown_row(
                            "Customization Per-Unit",
                            effective_custom_qty,
                            partner_per_unit_with_addons,  # PBP per unit cost (includes add-ons)
                            partner_customization_unit_total,  # PBP total
                            client_per_unit_with_addons,  # Client per unit price (includes add-ons)
                            customization_unit_total  # Client total
                        )
                    )

                # Per-product kitting (if applicable)
                kitting_pbp = item.get('kitting_pbp_cost', 0.0)
                kitting_client = item.get('kitting_client_price', 0.0)
                if item.get('include_kitting', False) and (kitting_pbp > 0 or kitting_client > 0):
                    kitting_desc = item.get('kitting_description', 'Kitting')
                    breakdown_data.append(
                        format_pricing_breakdown_row(
                            f"Kitting: {kitting_desc}",
                            "one-time",
                            kitting_pbp,  # PBP per unit (same as total for one-time)
                            kitting_pbp,  # PBP total
                            kitting_client,  # Client per unit (same as total for one-time)
                            kitting_client  # Client total
                        )
                    )

                # Create DataFrame with new column structure
                breakdown_df = pd.DataFrame(
                    breakdown_data,
                    columns=["Description", "Units", "PBP Cost (Per Unit)", "PBP Cost", "Client Price (Per Unit)", "Client Price"]
                )
                st.table(breakdown_df)

                # Show totals summary (include kitting if present)
                kitting_pbp = item.get('kitting_pbp_cost', 0.0) if item.get('include_kitting', False) else 0.0
                kitting_client = item.get('kitting_client_price', 0.0) if item.get('include_kitting', False) else 0.0
                total_pbp_cost = product_pbp_cost + partner_customization_setup_total + partner_customization_unit_total + kitting_pbp
                total_client_price = product_total + kitting_client

                st.caption(f"**Totals:** PBP Cost: ${total_pbp_cost:.2f} | Client Price: ${total_client_price:.2f} | Margin: ${total_client_price - total_pbp_cost:.2f}")

                # Store additional cost fields for order summary
                st.session_state.order_items[idx].update({
                    'customization_setup_cost': partner_customization_setup_total,  # What PBP pays for setup
                    'customization_unit_cost': partner_customization_unit_total,    # What PBP pays per unit
                })

                # Add note if minimum customization quantity is applied
                if new_include_custom and new_apply_minimum and new_custom_min_qty > new_quantity:
                    st.caption(f"Note: Customization minimum of {new_custom_min_qty} units applied (ordering {new_quantity} product units)")

        st.write("---")

        # Clear order button
        if st.button("Clear Entire Order", type="secondary"):
            st.session_state.order_items = []
            st.rerun()

        # PRICING SUMMARY (Phase 3) - Show all products with notes or warnings
        st.write("---")
        st.markdown("##### Pricing Summary")

        # Collect items with pricing notes (excluding manually overridden)
        items_with_notes = [
            item for item in st.session_state.order_items
            if item.get('pricing_notes') and item.get('pricing_notes').strip() and not item.get('manual_override', False)
        ]

        # Collect items with validation warnings (excluding manually overridden)
        items_with_warnings = [
            item for item in st.session_state.order_items
            if item.get('validation_warning') and not item.get('manual_override', False)
        ]

        # Show pricing notes if any
        if items_with_notes:
            with st.expander(f"ℹ️ Pricing Information ({len(items_with_notes)} product{'s' if len(items_with_notes) > 1 else ''})", expanded=False):
                for item in items_with_notes:
                    from src.helpers import format_product_with_variant
                    product_display = format_product_with_variant(item['product_name'], item.get('selected_variant'))
                    st.markdown(f"**{product_display}**")
                    st.caption(item['pricing_notes'])
                    st.write("")  # Spacing

        # Show validation warnings if any
        if items_with_warnings:
            st.warning(f"{len(items_with_warnings)} product{'s' if len(items_with_warnings) > 1 else ''} {'have' if len(items_with_warnings) > 1 else 'has'} pricing discrepancies")
            with st.expander("View Validation Details", expanded=False):
                for item in items_with_warnings:
                    from src.helpers import format_product_with_variant
                    product_display = format_product_with_variant(item['product_name'], item.get('selected_variant'))
                    st.markdown(f"**{product_display}**")
                    st.caption(item['validation_warning'])
                    st.write("")  # Spacing

    # ============================================================
    # ORDER SETTINGS
    # ============================================================
    st.divider()
    st.header("3. Order Settings")

    if len(st.session_state.order_items) == 0:
        st.caption("Add products to your order first, then configure order settings here.")
    else:
        # Shipping & Tariffs - Side by Side
        st.subheader("Shipping & Tariffs")

        col_shipping, col_tariff = st.columns(2)

        with col_shipping:
            st.markdown("**Client Shipping:**")
            st.session_state.order_shipping = st.number_input(
                "Shipping Price to Client ($)",
                min_value=0.0,
                value=st.session_state.order_shipping,
                step=10.0,
                key="shipping_input",
                help="Shipping cost charged to client"
            )

            st.markdown("**Partner Shipping:**")
            st.session_state.partner_shipping = st.number_input(
                "Shipping Cost from Partner ($)",
                min_value=0.0,
                value=st.session_state.partner_shipping,
                step=10.0,
                key="partner_shipping_input",
                help="Shipping cost PBP pays to partner (for Purchase Orders)"
            )

            st.markdown("**Sales Tax:**")
            st.session_state.sales_tax = st.number_input(
                "Estimated Sales Tax ($)",
                min_value=0.0,
                value=st.session_state.sales_tax,
                step=5.0,
                key="sales_tax_input",
                help="Estimated sales tax amount to be charged to client"
            )

        with col_tariff:
            # Calculate total tariff for expander label
            total_tariff = sum(item.get('tariff_amount', 0.0) for item in st.session_state.order_items)

            with st.expander(f"Tariff Configuration (Total: ${total_tariff:.2f})", expanded=False):
                st.caption("Default rates applied based on country of origin. Expand to customize per product.")
                st.markdown("""
Tariffs are import duties based on product country of origin.
Rates default to current estimates but can be adjusted as needed.
""")

                # Build editable tariff table with detailed breakdown
                tariff_table_rows = []

                for idx, item in enumerate(st.session_state.order_items):
                    # Get tariff base (product cost + markup, excludes customization)
                    tariff_base = item.get('tariff_base', 0.0)
                    tariff_base_per_unit = tariff_base / item['quantity'] if item['quantity'] > 0 else 0

                    # Display product info
                    st.markdown(f"**{idx + 1}. {item['product_name']}**")

                    col1, col2, col3 = st.columns([2, 2, 2])

                    with col1:
                        country = item.get('country_of_origin', 'N/A')
                        st.write(f"**Country:** {country if country else 'N/A'}")
                        st.write(f"**Quantity:** {item['quantity']} units")

                    with col2:
                        st.write(f"**Unit Cost:** ${tariff_base_per_unit:.2f}")
                        st.write(f"**Total Cost:** ${tariff_base:.2f}")
                        st.caption("(Product + Markup, excludes customization)")

                    with col3:
                        # Editable tariff rate
                        current_rate = item.get('tariff_rate_percent', 0.0)
                        new_rate = st.number_input(
                            "Tariff Rate (%)",
                            min_value=0.0,
                            max_value=100.0,
                            value=current_rate,
                            step=0.5,
                            key=f"tariff_rate_{idx}",
                            format="%.1f"
                        )

                        # Update if changed
                        if new_rate != current_rate:
                            item['tariff_rate_percent'] = new_rate
                            item['tariff_amount'] = calculate_product_tariff(tariff_base, new_rate)

                        tariff_amount = item.get('tariff_amount', 0.0)
                        st.write(f"**Tariff Amount:** ${tariff_amount:.2f}")
                        if tariff_base > 0 and new_rate > 0:
                            st.caption(f"${tariff_base:.2f} × {new_rate}% = ${tariff_amount:.2f}")

                    # Show tariff info if available
                    tariff_info = item.get('tariff_info', '')
                    if tariff_info and tariff_info.strip():
                        st.caption(f"ℹ️ {tariff_info}")

                    st.markdown("")  # Spacing

                st.caption("Tariff is calculated on product cost + markup (excludes customization fees and shipping)")

        # Order Adjustments - Consolidated Section
        st.divider()
        st.subheader("Order Adjustments")

        col1, col2, col3 = st.columns(3)

        with col1:
            # Discount as dropdown (like in Proposals tab)
            discount_options = ["None", "Non-profit (5%)", "Volume Order (5%)", "Custom"]
            current_discount = "None"
            if st.session_state.order_discount_type == "preset":
                # Check the preset value to determine which option
                if "Non-profit" in st.session_state.order_discount_preset:
                    current_discount = "Non-profit (5%)"
                elif "Volume Order" in st.session_state.order_discount_preset:
                    current_discount = "Volume Order (5%)"
            elif st.session_state.order_discount_type == "custom":
                current_discount = "Custom"

            discount_selection = st.selectbox(
                "Client Discount",
                options=discount_options,
                index=discount_options.index(current_discount),
                key="order_discount_select"
            )

            # Show what discount was quoted in proposal (if any)
            proposal_discount_type = st.session_state.get('proposal_discount_type')
            proposal_discount_percent = st.session_state.get('proposal_discount_percent', 0.0)

            if proposal_discount_type == 'Non-profit':
                st.caption("Discount Quoted to Client: Non-profit Discount (5%)")
            elif proposal_discount_type == 'Volume Order':
                st.caption("Discount Quoted to Client: Volume Order Discount (5%)")
            elif proposal_discount_type == 'Custom' and proposal_discount_percent > 0:
                st.caption(f"Discount Quoted to Client: Custom ({proposal_discount_percent}%)")
            else:
                st.caption("Discount Quoted to Client: None")

            # Update session state based on selection
            if discount_selection == "Non-profit (5%)":
                st.session_state.order_discount_type = "preset"
                st.session_state.order_discount_preset = "Non-profit Discount (5%)"
                st.session_state.order_discount_custom_value = 0.0
                st.session_state.order_discount_custom_desc = ""
            elif discount_selection == "Volume Order (5%)":
                st.session_state.order_discount_type = "preset"
                st.session_state.order_discount_preset = "Volume Order Discount (5%)"
                st.session_state.order_discount_custom_value = 0.0
                st.session_state.order_discount_custom_desc = ""
            elif discount_selection == "Custom":
                st.session_state.order_discount_type = "custom"
                # Show custom discount input below
            else:
                st.session_state.order_discount_type = "none"
                st.session_state.order_discount_custom_value = 0.0
                st.session_state.order_discount_custom_desc = ""

        with col2:
            st.session_state.order_fifty_cent_rounding = st.checkbox(
                "Round prices to nearest $0.50",
                value=st.session_state.order_fifty_cent_rounding,
                key="order_fifty_cent_rounding_checkbox",
                help="Rounds all prices to nearest 50 cents (e.g., $24.37 → $24.50)"
            )

            st.session_state.order_use_marketing_rounding = st.checkbox(
                "Apply marketing rounding (e.g., $60 → $59)",
                value=st.session_state.order_use_marketing_rounding,
                key="marketing_rounding_checkbox",
                help="Apply charm pricing to prices ending in 0. Applied after $0.50 rounding."
            )

        with col3:
            st.session_state.apply_cc_fee = st.checkbox(
                "Credit card fee",
                value=st.session_state.apply_cc_fee,
                key="cc_fee_checkbox",
                help="Add credit card processing fee to total (default 3%)"
            )

        # Row 2: Conditional inputs for Custom Discount and CC Fee
        if discount_selection == "Custom" or st.session_state.apply_cc_fee:
            col1_row2, col2_row2, col3_row2 = st.columns(3)

            with col1_row2:
                if discount_selection == "Custom":
                    st.session_state.order_discount_custom_value = st.number_input(
                        "Custom discount %",
                        min_value=0.0,
                        max_value=100.0,
                        value=st.session_state.order_discount_custom_value,
                        step=0.5,
                        key="order_custom_discount_input"
                    )

            with col2_row2:
                pass  # Empty column for alignment

            with col3_row2:
                if st.session_state.apply_cc_fee:
                    st.session_state.cc_fee_percent = st.number_input(
                        "CC fee %",
                        min_value=0.0,
                        max_value=10.0,
                        value=st.session_state.cc_fee_percent,
                        step=0.1,
                        key="cc_fee_percent_input",
                        help="Percentage fee charged for credit card payments"
                    )

        # Kitting & Gift Set Pricing
        st.divider()
        st.subheader("Kitting & Gift Set Pricing (Sale-wide)")
        st.caption("Add costs for gift boxes or packaging that apply to the entire order. For product-specific kitting, see product settings above.")

        col1_kitting, col2_kitting = st.columns(2)

        with col1_kitting:
            st.session_state.kitting_pbp_cost = st.number_input(
                "Kitting Cost (PBP) ($)",
                min_value=0.0,
                value=st.session_state.kitting_pbp_cost,
                step=10.0,
                key="kitting_pbp_input",
                help="What PBP pays for gift set assembly and packaging"
            )

        with col2_kitting:
            st.session_state.kitting_client_price = st.number_input(
                "Kitting Price (Client) ($)",
                min_value=0.0,
                value=st.session_state.kitting_client_price,
                step=10.0,
                key="kitting_client_input",
                help="What client pays for gift sets and custom packaging"
            )

        # Custom Line Items & Order Notes - Side by Side
        st.divider()

        # Count custom items and filled notes
        custom_item_count = sum(1 for item in st.session_state.order_items if item.get('is_custom', False))
        filled_notes_count = sum(1 for note in st.session_state.order_notes.values() if note and note.strip())

        col_custom, col_notes = st.columns(2)

        with col_custom:
            with st.expander(f"Add Custom Line Item ({custom_item_count} added)", expanded=False):
                st.caption("Add unique services or customizations not in the catalog")

                custom_name = st.text_input(
                    "Product/Service Name*",
                    key="custom_name_input",
                    placeholder="e.g., Custom Engraving Service"
                )
                custom_quantity = st.number_input(
                    "Quantity*",
                    min_value=1,
                    value=1,
                    step=1,
                    key="custom_quantity_input"
                )
                custom_description = st.text_input(
                    "Description",
                    key="custom_description_input",
                    placeholder="e.g., Laser engraving on wooden items"
                )
                custom_price = st.number_input(
                    "Total Price ($)*",
                    min_value=0.0,
                    value=0.0,
                    step=10.0,
                    key="custom_price_input",
                    help="Total price for this line item (quantity × unit price)"
                )

                if st.button("Add Custom Item to Order", type="secondary", use_container_width=True, key="add_custom_item_btn"):
                    # Validation
                    if not custom_name or custom_price <= 0:
                        st.error("Please fill in Product/Service Name and set Total Price greater than $0")
                    else:
                        # Create custom item
                        custom_item = {
                            'product_name': custom_name,
                            'product_ref': "CUSTOM",
                            'partner': "Custom",
                            'quantity': custom_quantity,
                            'markup_percent': 0.0,
                            'include_labels': False,
                            'base_price': custom_price / custom_quantity,
                            'tier_range': "N/A",
                            'tier_column': "N/A",
                            'additional_costs': {},
                            'product_subtotal': custom_price,
                            'art_setup_total': 0,
                            'label_cost_total': 0,
                            'subtotal_before_markup': custom_price,
                            'markup_amount': 0,
                            'product_total': custom_price,
                            'total_per_unit': custom_price / custom_quantity,
                            'is_custom': True,
                            'custom_description': custom_description if custom_description else "Custom line item",
                            'country_of_origin': '',
                            'tariff_rate_percent': 0.0,
                            'tariff_info': '',
                            'tariff_base': 0.0,
                            'tariff_amount': 0.0,
                            'edited_description': '',  # Empty for custom items, they use custom_description
                            # Per-product kitting fields
                            'include_kitting': False,
                            'kitting_pbp_cost': 0.0,
                            'kitting_client_price': 0.0,
                            'kitting_description': ''
                        }

                        st.session_state.order_items.append(custom_item)
                        st.toast(f"Added custom item: {custom_name}")
                        st.rerun()

        with col_notes:
            st.write(f"**Order Notes** ({filled_notes_count} filled)")
            st.caption("Add important details about this order")

    # Order Notes Section - Always visible text areas in 2x2 layout organized by audience
    st.divider()
    st.subheader("Order Notes")
    st.caption("Organize notes by audience: Internal (stays within PBP) vs External (shared with partners/clients).")

    # First row - 2 internal fields
    notes_col1, notes_col2 = st.columns(2)

    with notes_col1:
        internal_team_value = st.session_state.order_notes.get('internal_pbp_team', '')
        st.session_state.order_notes['internal_pbp_team'] = st.text_area(
            "Internal Notes (For PBP Team)",
            value=internal_team_value,
            placeholder="Rush order - coordinate with warehouse. Contact Jim for priority handling.",
            height=120,
            key="internal_pbp_team_input",
            help="Team coordination, workflow notes, internal reminders"
        )
        if internal_team_value:
            word_count = len(internal_team_value.split())
            st.caption(f"{word_count} words")

    with notes_col2:
        internal_bookkeeping_value = st.session_state.order_notes.get('internal_bookkeeping', '')
        st.session_state.order_notes['internal_bookkeeping'] = st.text_area(
            "Internal Notes (For Bookkeeping)",
            value=internal_bookkeeping_value,
            placeholder="Net 30 terms, send invoice after PO approval. Track against Q1 budget.",
            height=120,
            key="internal_bookkeeping_input",
            help="Accounting, billing, payment tracking"
        )
        if internal_bookkeeping_value:
            word_count = len(internal_bookkeeping_value.split())
            st.caption(f"{word_count} words")

    # Second row - 2 external fields
    notes_col3, notes_col4 = st.columns(2)

    with notes_col3:
        external_partners_value = st.session_state.order_notes.get('external_partners', '')
        st.session_state.order_notes['external_partners'] = st.text_area(
            "External Notes (For Partners/POs)",
            value=external_partners_value,
            placeholder="Kitting required: gift box with ribbon. Include branded tissue paper. Ship to warehouse by 3/15.",
            height=120,
            key="external_partners_input",
            help="Instructions for partners, PO details, shipping requirements"
        )
        if external_partners_value:
            word_count = len(external_partners_value.split())
            st.caption(f"{word_count} words")

    with notes_col4:
        external_clients_value = st.session_state.order_notes.get('external_clients', '')
        st.session_state.order_notes['external_clients'] = st.text_area(
            "External Notes (For Clients/Invoices)",
            value=external_clients_value,
            placeholder="Dropship directly to client. Include care instructions card with each unit.",
            height=120,
            key="external_clients_input",
            help="Client-facing information, special requests, delivery instructions"
        )
        if external_clients_value:
            word_count = len(external_clients_value.split())
            st.caption(f"{word_count} words")

    # Use session state values for calculations
    shipping = st.session_state.order_shipping
    # Calculate total tariff from all products
    tariff = sum(item.get('tariff_amount', 0.0) for item in st.session_state.order_items)

    # Calculate discount
    discount_percent = 0.0
    discount_description = ""

    if st.session_state.order_discount_type == "preset":
        # Extract percentage from preset string (e.g., "NGO Discount (5%)" -> 5.0)
        preset = st.session_state.order_discount_preset
        discount_description = preset
        # Parse percentage from string like "Non-profit Discount (5%)"
        if "(" in preset and "%" in preset:
            percent_str = preset.split("(")[1].split("%")[0]
            discount_percent = float(percent_str)

    elif st.session_state.order_discount_type == "custom":
        discount_percent = st.session_state.order_discount_custom_value
        discount_description = st.session_state.order_discount_custom_desc if st.session_state.order_discount_custom_desc else f"Custom Discount ({discount_percent}%)"

    # ============================================================
    # TOTAL ORDER CALCULATION
    # ============================================================
    st.divider()
    st.header("4. Order Summary")

    if len(st.session_state.order_items) == 0:
        st.caption("Add products to your order to see the total quote calculation.")
    else:
        # Import helper functions
        from src.helpers import calculate_split_totals

        # Calculate split totals
        split_totals = calculate_split_totals(st.session_state.order_items)

        # Apply discount ONLY to products client price (not customization)
        discount_amount = split_totals['products_only_client_price'] * (discount_percent / 100)
        products_after_discount = split_totals['products_only_client_price'] - discount_amount

        # Build subtotal: products after discount + customization (no discount on customization)
        subtotal_after_discount = products_after_discount + split_totals['customization_client_price']

        # Get sales tax amount
        sales_tax = st.session_state.sales_tax

        # Calculate total kitting: global + per-product
        global_kitting_pbp = st.session_state.kitting_pbp_cost
        global_kitting_client = st.session_state.kitting_client_price

        # Sum per-product kitting
        per_product_kitting_pbp = 0.0
        per_product_kitting_client = 0.0
        for item in st.session_state.order_items:
            if item.get('include_kitting', False):
                per_product_kitting_pbp += item.get('kitting_pbp_cost', 0.0)
                per_product_kitting_client += item.get('kitting_client_price', 0.0)

        # Total kitting = global + per-product
        total_kitting_pbp = global_kitting_pbp + per_product_kitting_pbp
        total_kitting_client = global_kitting_client + per_product_kitting_client

        # Get kitting costs (for backward compatibility with variable names below)
        kitting_pbp = total_kitting_pbp
        kitting_client = total_kitting_client

        # Calculate base total before CC fee
        total_before_cc = subtotal_after_discount + shipping + sales_tax + kitting_client + tariff

        # Calculate credit card fee (applied to total before CC fee)
        cc_fee_amount = calculate_credit_card_fee(total_before_cc, st.session_state.apply_cc_fee, st.session_state.cc_fee_percent)

        # Final total
        total_quote = total_before_cc + cc_fee_amount

        # Apply $0.50 rounding if enabled
        total_quote = round_to_nearest_fifty_cents(total_quote, st.session_state.order_fifty_cent_rounding)

        # Apply marketing rounding if enabled (after $0.50 rounding)
        total_quote = apply_marketing_rounding(total_quote, st.session_state.order_use_marketing_rounding)

        total_units = sum(item['quantity'] for item in st.session_state.order_items)

        summary_items = []

        # SECTION 1: Products with product names as headers
        st.markdown("##### Products")
        for item in st.session_state.order_items:
            # Skip custom line items for now
            if item.get('is_custom', False):
                summary_items.append([
                    item['product_name'],
                    item['quantity'],
                    f"${item['total_per_unit']:.2f}",
                    f"${item['product_total']:.2f}",
                    f"${item['total_per_unit']:.2f}",
                    f"${item['product_total']:.2f}"  # Custom items have same cost and price
                ])
                continue

            # Product header
            st.caption(f"**{item['product_name']}**")

            # Regular product: show base product with product name
            product_pbp_cost = item.get('product_subtotal', 0)
            product_client_price = product_pbp_cost + item.get('markup_amount', 0)

            # Add per-product kitting if included
            kitting_note = ""
            if item.get('include_kitting', False):
                kitting_pbp = item.get('kitting_pbp_cost', 0.0)
                kitting_client = item.get('kitting_client_price', 0.0)
                product_pbp_cost += kitting_pbp
                product_client_price += kitting_client
                if kitting_client > 0:
                    kitting_desc = item.get('kitting_description', 'kitting')
                    kitting_note = f" (includes ${kitting_client:.2f} {kitting_desc})"

            product_pbp_per_unit = product_pbp_cost / item['quantity'] if item['quantity'] > 0 else 0
            product_client_per_unit = product_client_price / item['quantity'] if item['quantity'] > 0 else 0

            summary_items.append([
                f"Base Product: {item['product_name']}{kitting_note}",
                item['quantity'],
                f"${product_pbp_per_unit:.2f}",
                f"${product_pbp_cost:.2f}",
                f"${product_client_per_unit:.2f}",
                f"${product_client_price:.2f}"
            ])

        # Products subtotal row
        summary_items.append([
            "**Products Subtotal**",
            "",
            "",
            f"**${split_totals['products_pbp_cost']:.2f}**",
            "",
            f"**${split_totals['products_client_price']:.2f}**"
        ])

        # SECTION 2: Customization (if any)
        if split_totals['customization_client_price'] > 0:
            summary_items.append(["", "", "", "", "", ""])  # Empty row for spacing
            st.markdown("##### Customization")

            for item in st.session_state.order_items:
                if item.get('is_custom', False):
                    continue

                # Show customization if present
                setup_pbp = item.get('customization_setup_cost', 0)
                setup_client = item.get('customization_setup_total', 0)
                unit_pbp = item.get('customization_unit_cost', 0)
                unit_client = item.get('customization_unit_total', 0)
                partner_per_unit = item.get('partner_customization_per_unit', 0)

                if setup_client > 0:
                    summary_items.append([
                        f"{item['product_name']} - Setup",
                        "one-time",
                        f"${setup_pbp:.2f}",  # Per unit same as total for one-time
                        f"${setup_pbp:.2f}",
                        f"${setup_client:.2f}",  # Per unit same as total for one-time
                        f"${setup_client:.2f}"
                    ])

                if unit_client > 0:
                    effective_custom_qty = item.get('customization_minimum_qty', item['quantity']) if item.get('apply_custom_minimum', False) else item['quantity']
                    # Calculate per-unit costs including add-ons
                    pbp_custom_per_unit = unit_pbp / effective_custom_qty if effective_custom_qty > 0 else 0
                    client_custom_per_unit = unit_client / effective_custom_qty if effective_custom_qty > 0 else 0

                    summary_items.append([
                        f"{item['product_name']} - Per Unit",
                        effective_custom_qty,
                        f"${pbp_custom_per_unit:.2f}",  # Includes add-ons
                        f"${unit_pbp:.2f}",
                        f"${client_custom_per_unit:.2f}",  # Includes add-ons
                        f"${unit_client:.2f}"
                    ])

            # Customization subtotal row
            summary_items.append([
                "**Customization Subtotal**",
                "",
                "",
                f"**${split_totals['customization_pbp_cost']:.2f}**",
                "",
                f"**${split_totals['customization_client_price']:.2f}**"
            ])

        # SECTION 3: Discounts, Shipping, Tariffs, Fees
        summary_items.append(["", "", "", "", "", ""])  # Empty row for spacing

        # Add discount line if applicable (applies to products only, not customization)
        if discount_percent > 0:
            summary_items.append([
                f"Discount ({discount_description})",
                "",
                "",
                "",
                f"-${discount_amount:.2f}",
                ""
            ])

        # Shipping (same for PBP and client)
        summary_items.append(["Shipping", "", "", f"${shipping:.2f}", "", f"${shipping:.2f}"])

        # Sales Tax (only affects client price)
        if sales_tax > 0:
            summary_items.append(["Sales Tax (Estimated)", "", "", "", "", f"${sales_tax:.2f}"])

        # Kitting/Gift Set Assembly (Sale-wide) - only show if global kitting > 0
        if global_kitting_pbp > 0 or global_kitting_client > 0:
            summary_items.append(["Kitting/Gift Set Assembly (Sale-wide)", "", "", f"${global_kitting_pbp:.2f}", "", f"${global_kitting_client:.2f}"])

        # Add tariff for each product (if > 0)
        for item in st.session_state.order_items:
            tariff_amount = item.get('tariff_amount', 0)
            if tariff_amount > 0:
                country = item.get('country_of_origin', 'Unknown')
                tariff_rate = item.get('tariff_rate_percent', 0)
                summary_items.append([
                    f"Tariff: {item['product_name']} ({tariff_rate}% - {country})",
                    "",
                    "",
                    f"${tariff_amount:.2f}",
                    "",
                    f"${tariff_amount:.2f}"
                ])

        # Add credit card fee if applicable (client pays)
        if st.session_state.apply_cc_fee and cc_fee_amount > 0:
            summary_items.append([
                f"Credit Card Fee ({st.session_state.cc_fee_percent}%)",
                "",
                "",
                "",
                "",
                f"${cc_fee_amount:.2f}"
            ])

        # Calculate total PBP cost (sales tax only affects client, not PBP)
        total_pbp_cost = split_totals['total_pbp_cost'] + shipping + kitting_pbp + tariff

        # Final total row
        summary_items.append([
            "**TOTAL**",
            f"**{total_units} units**",
            "",
            f"**${total_pbp_cost:.2f}**",
            "",
            f"**${total_quote:.2f}**"
        ])

        # Create and display table with new column structure
        summary_df = pd.DataFrame(summary_items, columns=["Item", "Qty", "PBP Cost (Per Unit)", "PBP Cost", "Client Price (Per Unit)", "Client Price"])
        st.table(summary_df)

        # Display total
        avg_per_unit = total_quote / total_units if total_units > 0 else 0
        st.success(f"Total Quote: ${total_quote:.2f}  ({total_units} total units @ ${avg_per_unit:.2f} avg per unit)")

        # Add download button for order summary
        summary_csv = summary_df.to_csv(index=False)
        st.download_button(
            label="Download Order Summary (CSV)",
            data=summary_csv,
            file_name=f"order_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
            mime="text/csv",
            key="download_order_summary"
        )

        # Save to history button
        if st.button("Save Quote to History", type="secondary"):
            # Create order history entry
            order_entry = {
                'timestamp': datetime.now(),
                'total_quote': total_quote,
                'total_units': total_units,
                'num_products': len(st.session_state.order_items),
                'product_names': [item['product_name'] for item in st.session_state.order_items],
                'order_items': [item.copy() for item in st.session_state.order_items],
                'shipping': shipping,
                'tariff': tariff,
                'discount_type': st.session_state.order_discount_type,
                'discount_description': discount_description,
                'discount_percent': discount_percent,
                'discount_amount': discount_amount,
                'use_marketing_rounding': st.session_state.order_use_marketing_rounding
            }
            st.session_state.order_history.append(order_entry)
            st.success("Quote saved to history!")
            st.rerun()

    # ============================================================
    # CLIENT INFORMATION UI
    # ============================================================
    st.divider()
    st.header("5. Client & Order Information")

    with st.expander("Client Details", expanded=False):
        st.markdown("Enter client information for invoices and purchase orders.")

        col1, col2 = st.columns(2)

        with col1:
            st.session_state.client_info['is_new_client'] = st.checkbox(
                "New Client?",
                value=st.session_state.client_info['is_new_client']
            )

            st.session_state.client_info['company_name'] = st.text_input(
                "Company Name",
                value=st.session_state.client_info['company_name'],
                placeholder="e.g., Acme Corp"
            )


            st.session_state.client_info['client_po'] = st.text_input(
                "Client PO Number (optional)",
                value=st.session_state.client_info['client_po'],
                placeholder="e.g., PO-2025-001"
            )

        with col2:
            st.session_state.client_info['billing_address'] = st.text_area(
                "Billing Address",
                value=st.session_state.client_info['billing_address'],
                placeholder="123 Main St\nCity, State ZIP",
                height=100
            )

            st.session_state.client_info['shipping_type'] = st.selectbox(
                "Shipping Type",
                options=['One Location', 'Drop Shipping'],
                index=0 if st.session_state.client_info['shipping_type'] == 'One Location' else 1
            )

            if st.session_state.client_info['shipping_type'] == 'One Location':
                st.session_state.client_info['shipping_address'] = st.text_area(
                    "Shipping Address",
                    value=st.session_state.client_info['shipping_address'],
                    placeholder="456 Shipping Lane\nCity, State ZIP",
                    height=100
                )
            else:
                st.session_state.client_info['shipping_address'] = ''
                st.caption("Drop shipping details to be arranged separately")

        # Contacts Management Section
        st.markdown("---")
        st.markdown("**Contacts Information**")

        # Display all contacts
        contacts = st.session_state.client_info.get('contacts', [])
        if not contacts:
            contacts = [{'name': '', 'email': '', 'phone': '', 'role': 'Primary Contact'}]
            st.session_state.client_info['contacts'] = contacts

        # Add/Remove buttons
        col_add, col_spacer = st.columns([1, 3])
        with col_add:
            if st.button("Add Another Contact", key="add_contact_tab3"):
                st.session_state.client_info['contacts'].append({
                    'name': '',
                    'email': '',
                    'phone': '',
                    'role': ''
                })
                st.rerun()

        # Display each contact
        for idx, contact in enumerate(contacts):
            contact_num = idx + 1

            # Contact header with remove button
            if len(contacts) > 1:
                col_header, col_remove = st.columns([3, 1])
                with col_header:
                    st.markdown(f"**Contact {contact_num}**")
                with col_remove:
                    if st.button(f"Remove", key=f"remove_contact_{idx}_tab3"):
                        st.session_state.client_info['contacts'].pop(idx)
                        st.rerun()
            else:
                st.markdown(f"**Contact {contact_num}**")

            # Contact fields in 2 columns
            col_contact1, col_contact2 = st.columns(2)

            with col_contact1:
                contact['name'] = st.text_input(
                    "Name",
                    value=contact.get('name', ''),
                    key=f"contact_name_{idx}_tab3",
                    placeholder="e.g., John Smith"
                )

                contact['email'] = st.text_input(
                    "Email",
                    value=contact.get('email', ''),
                    key=f"contact_email_{idx}_tab3",
                    placeholder="e.g., john@company.com"
                )

            with col_contact2:
                contact['phone'] = st.text_input(
                    "Phone",
                    value=contact.get('phone', ''),
                    key=f"contact_phone_{idx}_tab3",
                    placeholder="e.g., (555) 123-4567"
                )

                role_options = ['Primary Contact', 'Billing Contact', 'Technical Contact', 'Shipping Contact', 'Other']
                current_role = contact.get('role', 'Primary Contact')
                if current_role not in role_options and current_role:
                    role_options.append(current_role)

                contact['role'] = st.selectbox(
                    "Role",
                    options=role_options,
                    index=role_options.index(current_role) if current_role in role_options else 0,
                    key=f"contact_role_{idx}_tab3"
                )

            if idx < len(contacts) - 1:
                st.markdown("")  # Add spacing between contacts

        st.markdown("---")
        st.markdown("**Payment & Delivery Terms**")

        col3, col4 = st.columns(2)

        with col3:
            # Updated to dropdown with Net 15 and Custom option
            payment_terms_options = ['Net 30', 'Net 15', 'Net 60', 'Due on Receipt', '50% Deposit', 'Custom']
            current_timeline = st.session_state.client_info.get('payment_timeline', 'Net 30')

            # If current value is custom (not in standard options), select Custom
            if current_timeline not in payment_terms_options[:-1]:  # Exclude 'Custom' from check
                if current_timeline and current_timeline != 'Custom':
                    st.session_state.custom_payment_terms = current_timeline
                selected_index = payment_terms_options.index('Custom')
            else:
                selected_index = payment_terms_options.index(current_timeline)

            selected_payment = st.selectbox(
                "Payment Terms",
                options=payment_terms_options,
                index=selected_index
            )

            # If Custom is selected, show text input
            if selected_payment == "Custom":
                custom_terms = st.text_input(
                    "Enter Custom Payment Terms",
                    value=st.session_state.custom_payment_terms,
                    placeholder="e.g., Net 45, 2/10 Net 30, etc.",
                    key="custom_payment_input_tab3"
                )
                st.session_state.custom_payment_terms = custom_terms
                st.session_state.client_info['payment_timeline'] = custom_terms if custom_terms else "Custom"
            else:
                st.session_state.client_info['payment_timeline'] = selected_payment

        with col4:
            # Updated to match template options
            payment_method_options = ['Check', 'ACH', 'Credit Card', 'Wire Transfer']
            current_preference = st.session_state.client_info.get('payment_preference', 'Check')
            if current_preference not in payment_method_options:
                payment_method_options.append(current_preference)

            st.session_state.client_info['payment_preference'] = st.selectbox(
                "Payment Method",
                options=payment_method_options,
                index=payment_method_options.index(current_preference) if current_preference in payment_method_options else 0
            )

        col5, col6 = st.columns(2)

        with col5:
            # NEW: Ship method dropdown
            ship_method_options = ['Ground', 'Air', 'Freight', 'Other']
            current_ship_type = st.session_state.client_info.get('shipping_type', 'Ground')
            # Map old values to new
            if current_ship_type == 'One Location':
                current_ship_type = 'Ground'
            elif current_ship_type == 'Drop Shipping':
                current_ship_type = 'Other'

            if current_ship_type not in ship_method_options:
                ship_method_options.append(current_ship_type)

            # Note: We're overriding shipping_type to use ship method
            ship_method = st.selectbox(
                "Ship Method",
                options=ship_method_options,
                index=ship_method_options.index(current_ship_type) if current_ship_type in ship_method_options else 0,
                help="How products will be shipped to client"
            )

        with col6:
            # NEW: Client in-hands date
            st.session_state.client_info['client_in_hands_date'] = st.date_input(
                "Client In-Hands Date",
                value=st.session_state.client_info.get('client_in_hands_date'),
                help="Target delivery date for client to receive products"
            )

        st.markdown("---")
        st.markdown("**Order Submission Details**")

        col7, col8 = st.columns(2)

        with col7:
            st.session_state.client_info['order_submitted_by'] = st.text_input(
                "Order Submitted By",
                value=st.session_state.client_info.get('order_submitted_by', ''),
                placeholder="Your name",
                help="Person creating this order"
            )

            st.session_state.client_info['cost_submitted_by'] = st.text_input(
                "Cost Submitted By",
                value=st.session_state.client_info.get('cost_submitted_by', ''),
                placeholder="Finance contact name",
                help="Person who submitted/verified costs"
            )

        with col8:
            # Order submitted date (auto-filled, read-only display)
            order_date = st.session_state.client_info.get('order_submitted_date', datetime.now().date())
            st.date_input(
                "Order Submitted Date",
                value=order_date,
                disabled=True,
                help="Auto-filled when order created"
            )

            st.session_state.client_info['cost_submitted_date'] = st.date_input(
                "Cost Submitted Date",
                value=st.session_state.client_info.get('cost_submitted_date'),
                help="Date when costs were submitted/verified"
            )

    # ============================================================
    # CONFIRM ORDER
    # ============================================================
    st.divider()

    if not st.session_state.order_confirmed:
        st.markdown("### Ready to finalize your order?")
        st.caption("Review the order summary above and client information, then confirm to proceed to Tab 3.")

        if st.button("Confirm Order", type="primary", use_container_width=True):
            st.session_state.order_confirmed = True
            st.rerun()
    else:
        st.success("Order complete! Your order summary is ready.")
        st.info("Go to **Tab 4: Execution & Accounting** to generate Invoice & Purchase Order for this order.")

        # Optional: Add button to go back and edit
        if st.button("Edit Order", type="secondary"):
            st.session_state.order_confirmed = False
            st.rerun()

    # Save Your Work section at bottom of Tab 3
    st.divider()
    st.markdown("### Save Your Work")

    # Show unsaved changes indicator and save status
    if has_unsaved_order_changes():
        st.warning("You have unsaved changes in your order")

    save_status_bottom = format_time_since_save('order')
    if save_status_bottom:
        st.caption(f"{save_status_bottom}")

    # Check if there are products to save
    has_order_items = len(st.session_state.order_items) > 0

    if has_order_items:
        col1, col2 = st.columns(2)
        with col1:
            order_name = st.text_input(
                "Order name:",
                key="save_order_name_bottom",
                placeholder="e.g., Client ABC Q1 2025 Order"
            )
        with col2:
            created_by = st.text_input(
                "Your name (optional):",
                key="save_order_creator_bottom",
                placeholder="e.g., John Smith"
            )

        if st.button("Save Order", type="primary", use_container_width=True, key="save_order_btn_bottom"):
            if not order_name or not order_name.strip():
                st.error("Please enter an order name")
            else:
                # Prepare order data
                order_data = {
                    'order_items': st.session_state.order_items,
                    'order_shipping': st.session_state.order_shipping,
                    'partner_shipping': st.session_state.partner_shipping,
                    'sales_tax': st.session_state.sales_tax,
                    'kitting_pbp_cost': st.session_state.kitting_pbp_cost,
                    'kitting_client_price': st.session_state.kitting_client_price,
                    'order_discount_type': st.session_state.order_discount_type,
                    'order_discount_preset': st.session_state.order_discount_preset,
                    'order_discount_custom_desc': st.session_state.order_discount_custom_desc,
                    'order_discount_custom_value': st.session_state.order_discount_custom_value,
                    'order_use_marketing_rounding': st.session_state.order_use_marketing_rounding,
                    'apply_cc_fee': st.session_state.apply_cc_fee,
                    'cc_fee_percent': st.session_state.cc_fee_percent,
                    'client_info': st.session_state.client_info,
                    'order_notes': st.session_state.order_notes,
                    'order_confirmed': st.session_state.order_confirmed
                }

                success, message, result = save_order(
                    name=order_name.strip(),
                    created_by=created_by.strip() if created_by else "",
                    order_data=order_data,
                    dataset=st.session_state.selected_dataset
                )

                if success:
                    update_last_save_time('order')
                    clear_saved_data_cache()  # Clear cache to show new order
                    st.success(f"{message} - You can find it in the sidebar under 'Saved Orders'")
                    time.sleep(1.5)
                    st.rerun()
                else:
                    # Check if it's a naming conflict
                    if result:  # result contains suggested name
                        st.error(message)
                        if st.button(f"Save as '{result}'", key="save_with_new_name_order_bottom"):
                            success2, message2, _ = save_order(
                                name=result,
                                created_by=created_by.strip() if created_by else "",
                                order_data=order_data,
                                dataset=st.session_state.selected_dataset
                            )
                            if success2:
                                update_last_save_time('order')
                                clear_saved_data_cache()  # Clear cache to show new order
                                st.success(f"{message2} - You can find it in the sidebar")
                                time.sleep(1.5)
                                st.rerun()
                    else:
                        st.error(message)
    else:
        st.info("Add products to your order to enable saving")

    # Navigation button at bottom of Tab 3
    st.divider()
    col1, col2, col3 = st.columns([1, 1, 1])
    with col2:
        if st.button("Continue to Tab 4: Execution & Accounting", type="primary", use_container_width=True, key="tab3_to_tab4"):
            st.session_state.show_tab4_prompt = True
            st.rerun()

    # Show navigation prompt if button was clicked
    if st.session_state.get('show_tab4_prompt', False):
        st.info("Click on the **'Execution & Accounting'** tab above to continue.")
        st.session_state.show_tab4_prompt = False

# ============================================================
# TAB 4: EXECUTION & ACCOUNTING
# ============================================================
with tab4:
    st.header("Execution & Accounting - Invoice & Purchase Order Management")
    st.caption("Generate invoices and purchase orders for confirmed orders")
    st.divider()

    # Check if order exists in Tab 3
    if len(st.session_state.order_items) == 0:
        st.info("No order found. Please build an order in Tab 3 first.")
        st.markdown("### To create an invoice/PO:")
        st.markdown("1. Go to **Tab 3: Order & Client Info**")
        st.markdown("2. Complete Sections 1-8 (client info, products, settings, summary)")
        st.markdown("3. Return to this tab to generate Invoice/PO")
    else:
        # ============================================================
        # SECTION 1: REVIEW & EDIT ORDER INFORMATION
        # ============================================================
        st.subheader("1. Review & Edit Order Information")
        st.caption("Review and complete any missing information before generating the invoice/PO")

        client_info = st.session_state.client_info
        validation_warnings = validate_invoice_completeness(client_info, st.session_state.order_items)

        if validation_warnings:
            st.warning(f"{len(validation_warnings)} field(s) need attention. Complete the missing information below:")
        else:
            st.success("All required fields complete - ready to generate Invoice/PO")

        # Editable fields for missing information
        with st.expander("Edit Order Information", expanded=bool(validation_warnings)):
            col1, col2 = st.columns(2)

            with col1:
                st.markdown("**Client Information**")

                # New Client checkbox with callback
                def update_is_new_client():
                    st.session_state.client_info['is_new_client'] = st.session_state.tab4_is_new_client

                st.checkbox(
                    "New Client?",
                    value=st.session_state.client_info.get('is_new_client', True),
                    key="tab4_is_new_client",
                    on_change=update_is_new_client
                )

                # Use callback to update client_info to ensure persistence
                def update_company_name():
                    st.session_state.client_info['company_name'] = st.session_state.tab3_company_name

                st.text_input(
                    "Company Name",
                    value=st.session_state.client_info.get('company_name', ''),
                    key="tab3_company_name",
                    on_change=update_company_name
                )


                def update_billing_address():
                    st.session_state.client_info['billing_address'] = st.session_state.tab3_billing_address

                st.text_area(
                    "Billing Address",
                    value=st.session_state.client_info.get('billing_address', ''),
                    key="tab3_billing_address",
                    height=80,
                    on_change=update_billing_address
                )

                def update_shipping_address():
                    st.session_state.client_info['shipping_address'] = st.session_state.tab3_shipping_address

                st.text_area(
                    "Shipping Address",
                    value=st.session_state.client_info.get('shipping_address', ''),
                    key="tab3_shipping_address",
                    height=80,
                    on_change=update_shipping_address
                )

            # Contacts Management Section
            st.markdown("---")
            st.markdown("**Contacts Management**")

            # Display all contacts
            contacts = st.session_state.client_info.get('contacts', [])
            if not contacts:
                contacts = [{'name': '', 'email': '', 'phone': '', 'role': 'Primary Contact'}]
                st.session_state.client_info['contacts'] = contacts

            # Add button
            if st.button("Add Contact", key="add_contact_tab4"):
                st.session_state.client_info['contacts'].append({
                    'name': '',
                    'email': '',
                    'phone': '',
                    'role': ''
                })
                st.rerun()

            # Display each contact
            for idx, contact in enumerate(contacts):
                # Display contact header with remove button if multiple contacts
                if len(contacts) > 1:
                    contact_header_col, remove_col = st.columns([5, 1])
                    with contact_header_col:
                        st.markdown(f"**Contact {idx + 1}**")
                    with remove_col:
                        if st.button("Remove", key=f"tab4_remove_contact_{idx}"):
                            st.session_state.client_info['contacts'].pop(idx)
                            st.rerun()
                else:
                    st.markdown(f"**Contact {idx + 1}**")

                contact_col1, contact_col2 = st.columns(2)

                with contact_col1:
                    # Create callback functions for each field
                    def update_contact_name(idx=idx):
                        st.session_state.client_info['contacts'][idx]['name'] = st.session_state[f"tab4_contact_name_{idx}"]

                    def update_contact_email(idx=idx):
                        st.session_state.client_info['contacts'][idx]['email'] = st.session_state[f"tab4_contact_email_{idx}"]

                    contact['name'] = st.text_input(
                        "Name",
                        value=contact.get('name', ''),
                        key=f"tab4_contact_name_{idx}",
                        on_change=lambda idx=idx: st.session_state.client_info['contacts'][idx].update(
                            {'name': st.session_state[f"tab4_contact_name_{idx}"]}
                        )
                    )

                    contact['email'] = st.text_input(
                        "Email",
                        value=contact.get('email', ''),
                        key=f"tab4_contact_email_{idx}",
                        on_change=lambda idx=idx: st.session_state.client_info['contacts'][idx].update(
                            {'email': st.session_state[f"tab4_contact_email_{idx}"]}
                        )
                    )

                with contact_col2:
                    contact['phone'] = st.text_input(
                        "Phone",
                        value=contact.get('phone', ''),
                        key=f"tab4_contact_phone_{idx}",
                        on_change=lambda idx=idx: st.session_state.client_info['contacts'][idx].update(
                            {'phone': st.session_state[f"tab4_contact_phone_{idx}"]}
                        )
                    )

                    role_options = ['Primary Contact', 'Billing Contact', 'Technical Contact', 'Shipping Contact', 'Other']
                    current_role = contact.get('role', 'Primary Contact')
                    if current_role not in role_options and current_role:
                        role_options.append(current_role)

                    contact['role'] = st.selectbox(
                        "Role",
                        options=role_options,
                        index=role_options.index(current_role) if current_role in role_options else 0,
                        key=f"tab4_contact_role_{idx}",
                        on_change=lambda idx=idx: st.session_state.client_info['contacts'][idx].update(
                            {'role': st.session_state[f"tab4_contact_role_{idx}"]}
                        )
                    )

                # Add separator between contacts if not the last one
                if idx < len(contacts) - 1:
                    st.divider()

            with col2:
                st.markdown("**Order Details**")

                def update_in_hands_date():
                    st.session_state.client_info['client_in_hands_date'] = st.session_state.tab3_in_hands_date

                st.date_input(
                    "Client In-Hands Date",
                    value=st.session_state.client_info.get('client_in_hands_date'),
                    key="tab3_in_hands_date",
                    on_change=update_in_hands_date
                )

                def update_ship_method():
                    st.session_state.client_info['shipping_type'] = st.session_state.tab3_ship_method

                ship_method_options = ['Ground', 'Air', 'Freight', 'Other']
                current_ship = st.session_state.client_info.get('shipping_type', 'Ground')
                st.selectbox(
                    "Ship Method",
                    options=ship_method_options,
                    index=ship_method_options.index(current_ship) if current_ship in ship_method_options else 0,
                    key="tab3_ship_method",
                    on_change=update_ship_method
                )

                def update_payment_terms():
                    selected = st.session_state.tab3_payment_terms
                    if selected == "Custom":
                        # Use the custom terms if available
                        if st.session_state.custom_payment_terms:
                            st.session_state.client_info['payment_timeline'] = st.session_state.custom_payment_terms
                    else:
                        st.session_state.client_info['payment_timeline'] = selected

                payment_terms_options = ['Net 30', 'Net 15', 'Net 60', 'Due on Receipt', '50% Deposit', 'Custom']
                current_terms = st.session_state.client_info.get('payment_timeline', 'Net 30')

                # Handle custom terms display
                if current_terms not in payment_terms_options[:-1]:  # Exclude 'Custom' from check
                    if current_terms and current_terms != 'Custom':
                        st.session_state.custom_payment_terms = current_terms
                    display_index = payment_terms_options.index("Custom")
                else:
                    display_index = payment_terms_options.index(current_terms)

                selected_payment = st.selectbox(
                    "Payment Terms",
                    options=payment_terms_options,
                    index=display_index,
                    key="tab3_payment_terms",
                    on_change=update_payment_terms
                )

                # If Custom is selected, show text input
                if selected_payment == "Custom":
                    def update_custom_terms():
                        st.session_state.custom_payment_terms = st.session_state.custom_terms_edit_input
                        st.session_state.client_info['payment_timeline'] = st.session_state.custom_terms_edit_input

                    st.text_input(
                        "Custom Payment Terms",
                        value=st.session_state.custom_payment_terms,
                        placeholder="e.g., Net 45, 2/10 Net 30, etc.",
                        key="custom_terms_edit_input",
                        on_change=update_custom_terms
                    )

                def update_payment_method():
                    st.session_state.client_info['payment_preference'] = st.session_state.tab3_payment_method

                payment_method_options = ['Check', 'ACH', 'Credit Card', 'Wire Transfer']
                current_method = st.session_state.client_info.get('payment_preference', 'Check')
                st.selectbox(
                    "Payment Method",
                    options=payment_method_options,
                    index=payment_method_options.index(current_method) if current_method in payment_method_options else 0,
                    key="tab3_payment_method",
                    on_change=update_payment_method
                )

                def update_order_submitted_by():
                    st.session_state.client_info['order_submitted_by'] = st.session_state.tab3_order_submitted_by

                st.text_input(
                    "Order Submitted By",
                    value=st.session_state.client_info.get('order_submitted_by', ''),
                    key="tab3_order_submitted_by",
                    on_change=update_order_submitted_by
                )

                def update_cost_submitted_by():
                    st.session_state.client_info['cost_submitted_by'] = st.session_state.tab3_cost_submitted_by

                st.text_input(
                    "Cost Submitted By",
                    value=st.session_state.client_info.get('cost_submitted_by', ''),
                    key="tab3_cost_submitted_by",
                    on_change=update_cost_submitted_by
                )

                def update_cost_submitted_date():
                    st.session_state.client_info['cost_submitted_date'] = st.session_state.tab3_cost_submitted_date

                st.date_input(
                    "Cost Submitted Date",
                    value=st.session_state.client_info.get('cost_submitted_date'),
                    key="tab3_cost_submitted_date",
                    on_change=update_cost_submitted_date
                )

        # ============================================================
        # EDITABLE PRODUCT DESCRIPTIONS
        # ============================================================
        with st.expander("Edit Product Descriptions", expanded=False):
            st.caption("Customize product descriptions for clearer invoices and bookkeeping")

            if st.session_state.order_items:
                for idx, item in enumerate(st.session_state.order_items):
                    # Create callback for this specific item
                    def update_description(idx=idx):
                        st.session_state.order_items[idx]['edited_description'] = st.session_state[f'desc_edit_{idx}']

                    # Generate default description
                    if item.get('is_custom', False):
                        # Custom items use their custom_description
                        default_desc = item.get('custom_description', item.get('product_name', 'Custom Item'))
                    else:
                        # Regular products use product name
                        default_desc = item.get('product_name', 'Unknown Product')
                        # Add quantity info if more than 1
                        quantity = item.get('quantity', 1)
                        if quantity > 1:
                            default_desc += f" (Qty: {quantity})"

                    # Get current edited description or use default
                    current_desc = item.get('edited_description', '')
                    if not current_desc:
                        current_desc = default_desc

                    # Display text input for description
                    st.text_input(
                        f"{item.get('product_name', 'Product')} - {item.get('partner', 'Unknown Partner')}",
                        value=current_desc,
                        key=f"desc_edit_{idx}",
                        on_change=update_description,
                        placeholder=default_desc,
                        help=f"Default: {default_desc}"
                    )
            else:
                st.info("No products in order to edit descriptions for.")

        st.divider()

        # ============================================================
        # EDITABLE ORDER SETTINGS
        # ============================================================
        with st.expander("Edit Order Settings", expanded=False):
            st.caption("Make adjustments to order settings here. Changes sync to Tab 2.")

            # Shipping & Tariffs - Side by Side
            st.subheader("Shipping & Tariffs")

            col_shipping, col_tariff = st.columns(2)

            with col_shipping:
                st.markdown("**Client Shipping:**")
                st.session_state.order_shipping = st.number_input(
                    "Shipping Price to Client ($)",
                    min_value=0.0,
                    value=st.session_state.order_shipping,
                    step=10.0,
                    key="tab3_shipping_input",
                    help="Shipping cost charged to client"
                )

                st.markdown("**Partner Shipping:**")
                st.session_state.partner_shipping = st.number_input(
                    "Shipping Cost from Partner ($)",
                    min_value=0.0,
                    value=st.session_state.partner_shipping,
                    step=10.0,
                    key="tab3_partner_shipping_input",
                    help="Shipping cost PBP pays to partner (for Purchase Orders)"
                )

                st.markdown("**Sales Tax:**")
                st.session_state.sales_tax = st.number_input(
                    "Estimated Sales Tax ($)",
                    min_value=0.0,
                    value=st.session_state.sales_tax,
                    step=5.0,
                    key="tab3_sales_tax_input",
                    help="Estimated sales tax amount to be charged to client"
                )

            with col_tariff:
                # Calculate total tariff for display
                total_tariff = sum(item.get('tariff_amount', 0.0) for item in st.session_state.order_items)

                st.write(f"**Tariff Total:** ${total_tariff:.2f}")

                show_tariff_details = st.checkbox(
                    "Customize Tariff Rates",
                    key="tab3_show_tariff_details",
                    help="Default rates applied based on country of origin. Check to customize per product."
                )

                if show_tariff_details:
                    st.caption("Tariffs are import duties based on product country of origin. Rates default to current estimates but can be adjusted as needed.")

                    # Build editable tariff table with detailed breakdown
                    for idx, item in enumerate(st.session_state.order_items):
                        # Get tariff base (product cost + markup, excludes customization)
                        tariff_base = item.get('tariff_base', 0.0)
                        tariff_base_per_unit = tariff_base / item['quantity'] if item['quantity'] > 0 else 0

                        # Display product info
                        st.markdown(f"**{idx + 1}. {item['product_name']}**")

                        col1, col2, col3 = st.columns([2, 2, 2])

                        with col1:
                            country = item.get('country_of_origin', 'N/A')
                            st.write(f"**Country:** {country if country else 'N/A'}")
                            st.write(f"**Quantity:** {item['quantity']} units")

                        with col2:
                            st.write(f"**Unit Cost:** ${tariff_base_per_unit:.2f}")
                            st.write(f"**Total Cost:** ${tariff_base:.2f}")
                            st.caption("(Product + Markup, excludes customization)")

                        with col3:
                            # Editable tariff rate
                            current_rate = item.get('tariff_rate_percent', 0.0)
                            new_rate = st.number_input(
                                "Tariff Rate (%)",
                                min_value=0.0,
                                max_value=100.0,
                                value=current_rate,
                                step=0.5,
                                key=f"tab3_tariff_rate_{idx}",
                                format="%.1f"
                            )

                            # Update if changed
                            if new_rate != current_rate:
                                item['tariff_rate_percent'] = new_rate
                                item['tariff_amount'] = calculate_product_tariff(tariff_base, new_rate)

                            tariff_amount = item.get('tariff_amount', 0.0)
                            st.write(f"**Tariff Amount:** ${tariff_amount:.2f}")
                            if tariff_base > 0 and new_rate > 0:
                                st.caption(f"${tariff_base:.2f} × {new_rate}% = ${tariff_amount:.2f}")

                        # Show tariff info if available
                        tariff_info = item.get('tariff_info', '')
                        if tariff_info and tariff_info.strip():
                            st.caption(f"ℹ️ {tariff_info}")

                        st.markdown("")  # Spacing

                    st.caption("Tariff is calculated on product cost + markup (excludes customization fees and shipping)")

            # Order Adjustments - Consolidated Section
            st.divider()
            st.subheader("Order Adjustments")

            col1, col2, col3 = st.columns(3)

            with col1:
                # Discount as dropdown
                discount_options = ["None", "Non-profit (5%)", "Volume Order (5%)", "Custom"]
                current_discount = "None"
                if st.session_state.order_discount_type == "preset":
                    # Check the preset value to determine which option
                    if "Non-profit" in st.session_state.order_discount_preset:
                        current_discount = "Non-profit (5%)"
                    elif "Volume Order" in st.session_state.order_discount_preset:
                        current_discount = "Volume Order (5%)"
                elif st.session_state.order_discount_type == "custom":
                    current_discount = "Custom"

                discount_selection = st.selectbox(
                    "Client Discount",
                    options=discount_options,
                    index=discount_options.index(current_discount),
                    key="tab3_order_discount_select"
                )

                # Update session state based on selection
                if discount_selection == "Non-profit (5%)":
                    st.session_state.order_discount_type = "preset"
                    st.session_state.order_discount_preset = "Non-profit Discount (5%)"
                    st.session_state.order_discount_custom_value = 0.0
                    st.session_state.order_discount_custom_desc = ""
                elif discount_selection == "Volume Order (5%)":
                    st.session_state.order_discount_type = "preset"
                    st.session_state.order_discount_preset = "Volume Order Discount (5%)"
                    st.session_state.order_discount_custom_value = 0.0
                    st.session_state.order_discount_custom_desc = ""
                elif discount_selection == "Custom":
                    st.session_state.order_discount_type = "custom"
                else:
                    st.session_state.order_discount_type = "none"
                    st.session_state.order_discount_custom_value = 0.0
                    st.session_state.order_discount_custom_desc = ""

            with col2:
                st.session_state.order_fifty_cent_rounding = st.checkbox(
                    "Round prices to nearest $0.50",
                    value=st.session_state.order_fifty_cent_rounding,
                    key="tab4_fifty_cent_rounding_checkbox",
                    help="Rounds all prices to nearest 50 cents (e.g., $24.37 → $24.50)"
                )

                st.session_state.order_use_marketing_rounding = st.checkbox(
                    "Apply marketing rounding (e.g., $60 → $59)",
                    value=st.session_state.order_use_marketing_rounding,
                    key="tab3_marketing_rounding_checkbox",
                    help="Apply charm pricing to prices ending in 0. Applied after $0.50 rounding."
                )

            with col3:
                st.session_state.apply_cc_fee = st.checkbox(
                    "Credit card fee",
                    value=st.session_state.apply_cc_fee,
                    key="tab3_cc_fee_checkbox",
                    help="Add credit card processing fee to total (default 3%)"
                )

            # Row 2: Conditional inputs for Custom Discount and CC Fee
            if discount_selection == "Custom" or st.session_state.apply_cc_fee:
                col1_row2, col2_row2, col3_row2 = st.columns(3)

                with col1_row2:
                    if discount_selection == "Custom":
                        st.session_state.order_discount_custom_value = st.number_input(
                            "Custom discount %",
                            min_value=0.0,
                            max_value=100.0,
                            value=st.session_state.order_discount_custom_value,
                            step=0.5,
                            key="tab3_order_custom_discount_input"
                        )

                with col2_row2:
                    pass  # Empty column for alignment

                with col3_row2:
                    if st.session_state.apply_cc_fee:
                        st.session_state.cc_fee_percent = st.number_input(
                            "CC fee %",
                            min_value=0.0,
                            max_value=10.0,
                            value=st.session_state.cc_fee_percent,
                            step=0.1,
                            key="tab3_cc_fee_percent_input",
                            help="Percentage fee charged for credit card payments"
                        )

            # Kitting & Gift Set Pricing
            st.divider()
            st.subheader("Kitting & Gift Set Pricing")

            col1_kitting, col2_kitting = st.columns(2)

            with col1_kitting:
                st.session_state.kitting_pbp_cost = st.number_input(
                    "Kitting Cost (PBP) ($)",
                    min_value=0.0,
                    value=st.session_state.kitting_pbp_cost,
                    step=10.0,
                    key="tab4_kitting_pbp_input",
                    help="What PBP pays for gift set assembly and packaging"
                )

            with col2_kitting:
                st.session_state.kitting_client_price = st.number_input(
                    "Kitting Price (Client) ($)",
                    min_value=0.0,
                    value=st.session_state.kitting_client_price,
                    step=10.0,
                    key="tab4_kitting_client_input",
                    help="What client pays for gift sets and custom packaging"
                )

            # Custom Line Items & Order Notes - Side by Side
            st.divider()

            # Count custom items and filled notes
            custom_item_count = sum(1 for item in st.session_state.order_items if item.get('is_custom', False))
            filled_notes_count = sum(1 for note in st.session_state.order_notes.values() if note and note.strip())

            col_custom, col_notes = st.columns(2)

            with col_custom:
                st.write(f"**Custom Line Items** ({custom_item_count} added)")

                show_custom_item_form = st.checkbox(
                    "Add Custom Line Item",
                    key="tab3_show_custom_item_form",
                    help="Add unique services or customizations not in the catalog"
                )

                if show_custom_item_form:
                    st.caption("Add unique services or customizations not in the catalog")

                    custom_name = st.text_input(
                        "Product/Service Name*",
                        key="tab3_custom_name_input",
                        placeholder="e.g., Custom Engraving Service"
                    )
                    custom_quantity = st.number_input(
                        "Quantity*",
                        min_value=1,
                        value=1,
                        step=1,
                        key="tab3_custom_quantity_input"
                    )
                    custom_description = st.text_input(
                        "Description",
                        key="tab3_custom_description_input",
                        placeholder="e.g., Laser engraving on wooden items"
                    )
                    custom_price = st.number_input(
                        "Total Price ($)*",
                        min_value=0.0,
                        value=0.0,
                        step=10.0,
                        key="tab3_custom_price_input",
                        help="Total price for this line item (quantity × unit price)"
                    )

                    if st.button("Add Custom Item to Order", type="secondary", use_container_width=True, key="tab3_add_custom_item_btn"):
                        # Validation
                        if not custom_name or custom_price <= 0:
                            st.error("Please fill in Product/Service Name and set Total Price greater than $0")
                        else:
                            # Create custom item
                            custom_item = {
                                'product_name': custom_name,
                                'product_ref': "CUSTOM",
                                'partner': "Custom",
                                'quantity': custom_quantity,
                                'markup_percent': 0.0,
                                'include_labels': False,
                                'base_price': custom_price / custom_quantity,
                                'tier_range': "N/A",
                                'tier_column': "N/A",
                                'additional_costs': {},
                                'product_subtotal': custom_price,
                                'art_setup_total': 0,
                                'label_cost_total': 0,
                                'subtotal_before_markup': custom_price,
                                'markup_amount': 0,
                                'product_total': custom_price,
                                'total_per_unit': custom_price / custom_quantity,
                                'is_custom': True,
                                'custom_description': custom_description if custom_description else "Custom line item",
                                'country_of_origin': '',
                                'tariff_rate_percent': 0.0,
                                'tariff_info': '',
                                'tariff_base': 0.0,
                                'tariff_amount': 0.0,
                                'edited_description': ''  # Empty for custom items, they use custom_description
                            }

                            st.session_state.order_items.append(custom_item)
                            st.toast(f"Added custom item: {custom_name}")
                            st.rerun()

            with col_notes:
                st.write(f"**Order Notes** ({filled_notes_count} filled)")

                show_notes_form = st.checkbox(
                    "Add Order Notes",
                    key="tab3_show_notes_form",
                    help="Add specific details for this order"
                )

                if show_notes_form:
                    st.caption("Add specific details for this order (organized by audience)")

                    st.session_state.order_notes['internal_pbp_team'] = st.text_area(
                        "Internal Notes (For PBP Team)",
                        value=st.session_state.order_notes.get('internal_pbp_team', ''),
                        placeholder="Rush order - coordinate with warehouse. Contact Jim for priority handling.",
                        height=70,
                        key="tab3_internal_pbp_team",
                        help="Team coordination, workflow notes, internal reminders"
                    )

                    st.session_state.order_notes['internal_bookkeeping'] = st.text_area(
                        "Internal Notes (For Bookkeeping)",
                        value=st.session_state.order_notes.get('internal_bookkeeping', ''),
                        placeholder="Net 30 terms, send invoice after PO approval. Track against Q1 budget.",
                        height=70,
                        key="tab3_internal_bookkeeping",
                        help="Accounting, billing, payment tracking"
                    )

                    st.session_state.order_notes['external_partners'] = st.text_area(
                        "External Notes (For Partners/POs)",
                        value=st.session_state.order_notes.get('external_partners', ''),
                        placeholder="Kitting required: gift box with ribbon. Include branded tissue paper. Ship to warehouse by 3/15.",
                        height=70,
                        key="tab3_external_partners",
                        help="Instructions for partners, PO details, shipping requirements"
                    )

                    st.session_state.order_notes['external_clients'] = st.text_area(
                        "External Notes (For Clients/Invoices)",
                        value=st.session_state.order_notes.get('external_clients', ''),
                        placeholder="Dropship directly to client. Include care instructions card with each unit.",
                        height=70,
                        key="tab3_external_clients",
                        help="Client-facing information, special requests, delivery instructions"
                    )

        st.divider()

        # ============================================================
        # SECTION 3: INVOICE & PURCHASE ORDER GENERATION
        # ============================================================
        st.subheader("2. Generate Invoice & Purchase Order")

        st.markdown("### INVOICE AND PURCHASE ORDER REQUEST FORM")
        st.caption("Complete form matching the bookkeeper template requirements")

        # ============================================================
        # TABLE 1: CLIENT/COMPANY INFORMATION
        # ============================================================
        st.markdown("#### 1. Client/Company Information")

        client_info_data = []

        # Company/Client Name
        company_name = client_info.get('company_name', 'Not specified')
        client_info_data.append(["Company/Client Name", company_name])

        # New Client Status (as separate row with checkbox format)
        is_new = client_info.get('is_new_client', False)
        client_status_display = f"[X] Yes  [ ] No" if is_new else "[ ] Yes  [X] No"
        client_info_data.append(["New Client?", client_status_display])

        # Contact + Email (using new contacts array - show all contacts if multiple)
        contacts = client_info.get('contacts', [])
        if contacts and len(contacts) > 0:
            if len(contacts) == 1:
                # Single contact
                contact = contacts[0]
                contact_name = contact.get('name', 'Not specified')
                contact_email = contact.get('email', 'Not specified')
                contact_display = f"{contact_name} <{contact_email}>"
            else:
                # Multiple contacts - show all with roles
                contact_lines = []
                for contact in contacts:
                    name = contact.get('name', 'Not specified')
                    email = contact.get('email', 'Not specified')
                    role = contact.get('role', '')
                    if role:
                        contact_lines.append(f"{name} <{email}> ({role})")
                    else:
                        contact_lines.append(f"{name} <{email}>")
                contact_display = " | ".join(contact_lines)
        else:
            # Fallback to old fields if they exist (backward compatibility)
            contact_name = client_info.get('contact_name', 'Not specified')
            contact_email = client_info.get('contact_email', 'Not specified')
            contact_display = f"{contact_name} <{contact_email}>"
        client_info_data.append(["Contact + Email", contact_display])

        # Company Billing Address + Email (using new contacts array)
        billing_address = client_info.get('billing_address', 'Not specified')
        if contacts and len(contacts) > 0:
            billing_email = contacts[0].get('email', 'Not specified')
        else:
            # Fallback to old field if it exists (backward compatibility)
            billing_email = client_info.get('contact_email', 'Not specified')
        client_info_data.append(["Company Billing Address + Email", f"{billing_address} | {billing_email}"])

        # Company Shipping Address
        shipping_address = client_info.get('shipping_address', 'Not specified')
        shipping_type = client_info.get('shipping_type', 'Not specified')
        client_info_data.append(["Company Shipping Address", f"{shipping_address} ({shipping_type})"])

        # Client PO #
        po_number = client_info.get('client_po', 'N/A')
        client_info_data.append(["Client PO #", po_number])

        # Display as table
        client_info_df = pd.DataFrame(client_info_data, columns=["Field", "Value"])
        st.table(client_info_df)

        st.divider()

        # ============================================================
        # TABLE 2: PARTNERS + POINT OF CONTACTS
        # ============================================================
        st.markdown("#### 2. Partners + Point of Contacts")

        # Get unique partners from order items
        partners_in_order = list(set(item['partner'] for item in st.session_state.order_items if not item.get('is_custom', False)))

        partners_data = []
        if partners_in_order and hasattr(st.session_state, 'partner_contacts'):
            for partner_name in partners_in_order:
                partner_contact = st.session_state.partner_contacts.get(partner_name, {})
                poc_name = partner_contact.get('poc_name', 'Not specified')
                poc_email = partner_contact.get('poc_email', 'Not specified')
                partners_data.append({
                    "Partner": partner_name,
                    "Point of Contact (POC)": f"{poc_name} <{poc_email}>"
                })
        else:
            partners_data.append({
                "Partner": "No partners in order",
                "Point of Contact (POC)": "N/A"
            })

        partners_df = pd.DataFrame(partners_data)
        st.table(partners_df)

        st.divider()

        # ============================================================
        # TABLE 3: ORDER DETAILS (DATES, SHIPPING, PAYMENT)
        # ============================================================
        st.markdown("#### 3. Order Details")

        order_details_data = []

        # Client In-Hands Date
        client_in_hands = client_info.get('client_in_hands_date', 'Not specified')
        order_details_data.append(["Client In-Hands Date", format_date_display(client_in_hands)])

        # Ship Method
        ship_method = client_info.get('shipping_type', 'Not specified')
        order_details_data.append(["Ship Method", ship_method])

        # Payment Terms
        payment_terms = client_info.get('payment_timeline', 'Not specified')
        order_details_data.append(["Payment Terms", payment_terms])

        # Payment Method
        payment_method = client_info.get('payment_preference', 'Not specified')
        order_details_data.append(["Payment Method", payment_method])

        # Order Submitted By + Date
        order_submitted_by = client_info.get('order_submitted_by', 'Not specified')
        order_submitted_date = client_info.get('order_submitted_date', datetime.now().date())
        order_details_data.append(["Order Submitted By", f"{order_submitted_by} (Date: {format_date_display(order_submitted_date)})"])

        # Cost Submitted By + Date
        cost_submitted_by = client_info.get('cost_submitted_by', 'Not specified')
        cost_submitted_date = client_info.get('cost_submitted_date', 'Not specified')
        formatted_cost_date = format_date_display(cost_submitted_date) if cost_submitted_date else 'Not specified'
        order_details_data.append(["Cost Submitted By", f"{cost_submitted_by} (Date: {formatted_cost_date})"])

        # Display as table
        order_details_df = pd.DataFrame(order_details_data, columns=["Field", "Value"])
        st.table(order_details_df)

        st.divider()

        # ============================================================
        # DESCRIPTION HELPER FUNCTIONS
        # ============================================================
        def get_description_for_invoice(product_data, product_name):
            """
            Get product description for client-facing invoices.

            Hierarchy:
            1. Billing Description (to Client)
            2. Marketing Description (Website)
            3. Product/Service Name
            """
            # Try Billing Description first
            billing_desc = get_column_value(product_data, 'billing_description', None)
            if billing_desc and str(billing_desc).strip() and str(billing_desc).strip().lower() != 'nan':
                return str(billing_desc).strip()

            # Fallback to Marketing Description
            marketing_desc = get_column_value(product_data, 'marketing_description', None)
            if marketing_desc and str(marketing_desc).strip() and str(marketing_desc).strip().lower() != 'nan':
                return str(marketing_desc).strip()

            # Final fallback to product name
            return product_name

        def get_description_for_po(product_data, product_name):
            """
            Get product description for partner purchase orders.

            Hierarchy:
            1. Purchase Description (to Partner)
            2. Billing Description (to Client)
            3. Product/Service Name
            """
            # Try Purchase Description first
            purchase_desc = get_column_value(product_data, 'purchase_description', None)
            if purchase_desc and str(purchase_desc).strip() and str(purchase_desc).strip().lower() != 'nan':
                return str(purchase_desc).strip()

            # Fallback to Billing Description
            billing_desc = get_column_value(product_data, 'billing_description', None)
            if billing_desc and str(billing_desc).strip() and str(billing_desc).strip().lower() != 'nan':
                return str(billing_desc).strip()

            # Final fallback to product name
            return product_name

        # ============================================================
        # TABLE 4: INVOICE AND PURCHASE ORDER ITEM DETAILS
        # ============================================================
        st.markdown("#### 4. Invoice and Purchase Order Item Details")
        st.caption("""
        This cost-to-sell segment outlines our partners' cost, our sell price to client,
        and our partners' requested in-hands date. Our in-hands date for clients may be
        later than the in-hands date to Peace by Piece for kitting purposes.
        """)

        # Build line items table with per-unit and total columns
        invoice_line_items = []
        for item in st.session_state.order_items:
            # Check if custom item
            if item.get('is_custom', False):
                partner = "Custom"
                custom_desc = item.get('custom_description', 'Custom line item')
                partner_in_hands = "N/A"
                qty = item['quantity']
                cost_per_unit = item.get('total_per_unit', 0)
                cost_total = item.get('product_total', 0)
                cost_verified = "N/A"

                # Custom items use same description for both invoice and PO
                invoice_line_items.append({
                    'PARTNER': partner,
                    'DESCRIPTION (Invoice)': custom_desc,
                    'DESCRIPTION (PO)': custom_desc,
                    'QTY': qty,
                    'IN-HANDS from Partner': partner_in_hands,
                    'COST/UNIT': f"${cost_per_unit:.2f}",
                    'TOTAL COST': f"${cost_total:.2f}",
                    'COST VERIFIED?': cost_verified,
                    'SELL PRICE/UNIT': f"${cost_per_unit:.2f}",
                    'TOTAL SELL PRICE': f"${cost_total:.2f}"
                })
            else:
                # Regular product
                partner = item['partner']

                # Display product name with variant (if applicable)
                from src.helpers import format_product_with_variant
                product_name = format_product_with_variant(
                    item['product_name'],
                    item.get('selected_variant')
                )

                # Get product data from catalog for description lookup
                product_ref = item.get('product_ref', item['product_name'])
                try:
                    if hasattr(st.session_state, 'df_template') and st.session_state.df_template is not None:
                        product_match = st.session_state.df_template[
                            (st.session_state.df_template['Product/Service'] == product_ref) &
                            (st.session_state.df_template['Partner'] == partner)
                        ]
                        if not product_match.empty:
                            product_data = product_match.iloc[0]
                        else:
                            product_data = None
                    else:
                        product_data = None
                except Exception:
                    product_data = None

                # Get descriptions using helper functions
                if product_data is not None:
                    invoice_desc = get_description_for_invoice(product_data, product_name)
                    po_desc = get_description_for_po(product_data, product_name)
                else:
                    # Fallback if product data not found
                    invoice_desc = product_name
                    po_desc = product_name

                # Use edited description if available (overrides spreadsheet descriptions for both invoice and PO)
                edited_desc = item.get('edited_description', '')
                if edited_desc:
                    invoice_desc = edited_desc
                    po_desc = edited_desc

                qty = item['quantity']

                partner_in_hands = item.get('partner_in_hands_date', 'TBD')
                if partner_in_hands and partner_in_hands != 'TBD':
                    partner_in_hands = str(partner_in_hands)

                # Calculate pricing using new system if product data available
                if product_data is not None:
                    try:
                        # Use calculate_pbp_msrp for consistent pricing
                        user_markup_override = item.get('markup_percent') if item.get('manual_override') else None
                        pricing_result = calculate_pbp_msrp(
                            product_data.to_dict(),
                            quantity=qty,
                            user_markup_override=user_markup_override
                        )
                        partner_cost_per_unit = pricing_result['calculation_details']['base_cost']
                        sell_price_per_unit = pricing_result['pbp_msrp']
                    except Exception:
                        # Fallback to stored values if calculation fails
                        partner_cost_per_unit = item.get('partner_cost_per_unit', item.get('base_price', 0))
                        product_subtotal = item.get('product_subtotal', 0)
                        markup_amount = item.get('markup_amount', 0)
                        sell_price_total = product_subtotal + markup_amount
                        sell_price_per_unit = sell_price_total / qty if qty > 0 else 0
                else:
                    # Use stored values if product data not available
                    partner_cost_per_unit = item.get('partner_cost_per_unit', item.get('base_price', 0))
                    product_subtotal = item.get('product_subtotal', 0)
                    markup_amount = item.get('markup_amount', 0)
                    sell_price_total = product_subtotal + markup_amount
                    sell_price_per_unit = sell_price_total / qty if qty > 0 else 0

                partner_cost_total = partner_cost_per_unit * qty
                sell_price_total = sell_price_per_unit * qty

                cost_verified = item.get('cost_verified', 'Pending')

                # Add per-product kitting if included
                if item.get('include_kitting', False):
                    kitting_pbp = item.get('kitting_pbp_cost', 0.0)
                    kitting_client = item.get('kitting_client_price', 0.0)
                    partner_cost_total += kitting_pbp
                    sell_price_total += kitting_client

                    # Append kitting note to invoice description only (client-facing)
                    kitting_desc = item.get('kitting_description', 'Kitting')
                    invoice_desc += f" | {kitting_desc}: +${kitting_client:.2f}"

                # Add base product line
                invoice_line_items.append({
                    'PARTNER': partner,
                    'DESCRIPTION (Invoice)': invoice_desc,
                    'DESCRIPTION (PO)': po_desc,
                    'QTY': qty,
                    'IN-HANDS from Partner': partner_in_hands,
                    'COST/UNIT': f"${partner_cost_per_unit:.2f}",
                    'TOTAL COST': f"${partner_cost_total:.2f}",
                    'COST VERIFIED?': cost_verified,
                    'SELL PRICE/UNIT': f"${sell_price_per_unit:.2f}",
                    'TOTAL SELL PRICE': f"${sell_price_total:.2f}"
                })

                # Add customization line items if present
                if item.get('include_customization', False):
                    customization_desc = item.get('customization_description', 'Custom work')

                    # Client-facing prices
                    customization_setup = item.get('customization_setup_total', 0)
                    customization_unit_total = item.get('customization_unit_total', 0)
                    customization_per_unit = item.get('customization_per_unit', 0)

                    # Partner costs
                    partner_customization_setup = item.get('partner_customization_setup_total', 0)
                    partner_customization_unit_total = item.get('partner_customization_unit_total', 0)
                    partner_customization_per_unit = item.get('partner_customization_per_unit', 0)

                    # Setup fee line item (show partner cost vs. client price)
                    if customization_setup > 0 or partner_customization_setup > 0:
                        invoice_customization_desc = f"  └ Setup Fee: {customization_desc}"
                        po_customization_desc = f"  └ Setup: {customization_desc}"

                        invoice_line_items.append({
                            'PARTNER': partner,
                            'DESCRIPTION (Invoice)': invoice_customization_desc,
                            'DESCRIPTION (PO)': po_customization_desc,
                            'QTY': 1,
                            'IN-HANDS from Partner': partner_in_hands,
                            'COST/UNIT': f"${partner_customization_setup:.2f}",
                            'TOTAL COST': f"${partner_customization_setup:.2f}",
                            'COST VERIFIED?': cost_verified,
                            'SELL PRICE/UNIT': f"${customization_setup:.2f}",
                            'TOTAL SELL PRICE': f"${customization_setup:.2f}"
                        })

                    # Per-unit customization line item (show partner cost vs. client price)
                    if customization_unit_total > 0 or partner_customization_unit_total > 0:
                        invoice_per_unit_desc = f"  └ Customization (per unit): {customization_desc}"
                        po_per_unit_desc = f"  └ Per Unit Customization: {customization_desc}"

                        invoice_line_items.append({
                            'PARTNER': partner,
                            'DESCRIPTION (Invoice)': invoice_per_unit_desc,
                            'DESCRIPTION (PO)': po_per_unit_desc,
                            'QTY': qty,
                            'IN-HANDS from Partner': partner_in_hands,
                            'COST/UNIT': f"${partner_customization_per_unit:.2f}",
                            'TOTAL COST': f"${partner_customization_unit_total:.2f}",
                            'COST VERIFIED?': cost_verified,
                            'SELL PRICE/UNIT': f"${customization_per_unit:.2f}",
                            'TOTAL SELL PRICE': f"${customization_unit_total:.2f}"
                        })

                    # Add customization add-on line items
                    for addon in item.get('customization_addons', []):
                        if addon.get('name'):  # Only show add-ons with names
                            # Use new field names with backward compatibility
                            addon_partner_setup = addon.get('partner_setup_fee', addon.get('setup_fee', 0.0))
                            addon_client_setup = addon.get('client_setup_fee', addon.get('setup_fee', 0.0))
                            addon_partner_perunit = addon.get('partner_per_unit_cost', addon.get('per_unit_cost', 0.0))
                            addon_client_perunit = addon.get('client_per_unit_cost', addon.get('per_unit_cost', 0.0))

                            # Add-on setup fee line item (if applicable)
                            if addon_partner_setup > 0 or addon_client_setup > 0:
                                invoice_addon_setup_desc = f"  └ Add-On Setup: {addon['name']}"
                                po_addon_setup_desc = f"  └ Add-On Setup: {addon['name']}"

                                invoice_line_items.append({
                                    'PARTNER': partner,
                                    'DESCRIPTION (Invoice)': invoice_addon_setup_desc,
                                    'DESCRIPTION (PO)': po_addon_setup_desc,
                                    'QTY': 1,
                                    'IN-HANDS from Partner': partner_in_hands,
                                    'COST/UNIT': f"${addon_partner_setup:.2f}",
                                    'TOTAL COST': f"${addon_partner_setup:.2f}",
                                    'COST VERIFIED?': cost_verified,
                                    'SELL PRICE/UNIT': f"${addon_client_setup:.2f}",
                                    'TOTAL SELL PRICE': f"${addon_client_setup:.2f}"
                                })

                            # Add-on per-unit line item (if applicable)
                            if addon_partner_perunit > 0 or addon_client_perunit > 0:
                                addon_partner_total = addon_partner_perunit * qty
                                addon_client_total = addon_client_perunit * qty
                                invoice_addon_desc = f"  └ Add-On (per unit): {addon['name']}"
                                po_addon_desc = f"  └ Add-On: {addon['name']}"

                                invoice_line_items.append({
                                    'PARTNER': partner,
                                    'DESCRIPTION (Invoice)': invoice_addon_desc,
                                    'DESCRIPTION (PO)': po_addon_desc,
                                    'QTY': qty,
                                    'IN-HANDS from Partner': partner_in_hands,
                                    'COST/UNIT': f"${addon_partner_perunit:.2f}",
                                    'TOTAL COST': f"${addon_partner_total:.2f}",
                                    'COST VERIFIED?': cost_verified,
                                    'SELL PRICE/UNIT': f"${addon_client_perunit:.2f}",
                                    'TOTAL SELL PRICE': f"${addon_client_total:.2f}"
                                })

                # Add tariff line item if applicable
                if item.get('tariff_amount', 0) > 0:
                    tariff_amount_total = item.get('tariff_amount', 0)
                    tariff_rate = item.get('tariff_rate_percent', 0)
                    tariff_per_unit = tariff_amount_total / qty if qty > 0 else 0
                    country = item.get('country_of_origin', 'Unknown')

                    invoice_tariff_desc = f"  └ Tariff ({country}, {tariff_rate:.1f}%)"
                    po_tariff_desc = f"  └ Tariff ({tariff_rate:.1f}%)"

                    invoice_line_items.append({
                        'PARTNER': partner,
                        'DESCRIPTION (Invoice)': invoice_tariff_desc,
                        'DESCRIPTION (PO)': po_tariff_desc,
                        'QTY': qty,
                        'IN-HANDS from Partner': "N/A",
                        'COST/UNIT': f"${tariff_per_unit:.2f}",
                        'TOTAL COST': f"${tariff_amount_total:.2f}",
                        'COST VERIFIED?': "Yes",
                        'SELL PRICE/UNIT': f"${tariff_per_unit:.2f}",
                        'TOTAL SELL PRICE': f"${tariff_amount_total:.2f}"
                    })

        # Add shipping line item (show partner cost vs. client price)
        partner_shipping_cost = st.session_state.partner_shipping
        client_shipping_price = shipping
        if partner_shipping_cost > 0 or client_shipping_price > 0:
            invoice_line_items.append({
                'PARTNER': 'Shipping',
                'DESCRIPTION (Invoice)': 'Shipping to Client',
                'DESCRIPTION (PO)': 'Shipping from Partner',
                'QTY': 1,
                'IN-HANDS from Partner': 'N/A',
                'COST/UNIT': f"${partner_shipping_cost:.2f}",
                'TOTAL COST': f"${partner_shipping_cost:.2f}",
                'COST VERIFIED?': 'Yes',
                'SELL PRICE/UNIT': f"${client_shipping_price:.2f}",
                'TOTAL SELL PRICE': f"${client_shipping_price:.2f}"
            })

        # Add kitting line item for sale-wide kitting only (per-product kitting is merged into product rows)
        kitting_pbp_cost = st.session_state.get('kitting_pbp_cost', 0)
        kitting_client_price = st.session_state.get('kitting_client_price', 0)
        if kitting_pbp_cost > 0 or kitting_client_price > 0:
            invoice_line_items.append({
                'PARTNER': 'Kitting (Sale-wide)',
                'DESCRIPTION (Invoice)': 'Gift Set Assembly & Packaging',
                'DESCRIPTION (PO)': 'Kitting & Assembly Services',
                'QTY': 1,
                'IN-HANDS from Partner': 'N/A',
                'COST/UNIT': f"${kitting_pbp_cost:.2f}",
                'TOTAL COST': f"${kitting_pbp_cost:.2f}",
                'COST VERIFIED?': 'Yes',
                'SELL PRICE/UNIT': f"${kitting_client_price:.2f}",
                'TOTAL SELL PRICE': f"${kitting_client_price:.2f}"
            })

        # Display line items table with better column sizing
        invoice_df = pd.DataFrame(invoice_line_items)
        st.dataframe(
            invoice_df,
            use_container_width=True,
            hide_index=True,
            column_config={
                "PARTNER": st.column_config.TextColumn("PARTNER", width="small"),
                "DESCRIPTION (Invoice)": st.column_config.TextColumn("DESCRIPTION (Invoice)", width="medium"),
                "DESCRIPTION (PO)": st.column_config.TextColumn("DESCRIPTION (PO)", width="medium"),
                "QTY": st.column_config.NumberColumn("QTY", width="small"),
                "IN-HANDS from Partner": st.column_config.TextColumn("IN-HANDS from Partner", width="small"),
                "COST/UNIT": st.column_config.TextColumn("COST/UNIT", width="small"),
                "TOTAL COST": st.column_config.TextColumn("TOTAL COST", width="small"),
                "COST VERIFIED?": st.column_config.TextColumn("COST VERIFIED?", width="small"),
                "SELL PRICE/UNIT": st.column_config.TextColumn("SELL PRICE/UNIT", width="small"),
                "TOTAL SELL PRICE": st.column_config.TextColumn("TOTAL SELL PRICE", width="small")
            }
        )

        # Display totals section
        st.write("")  # Spacing
        st.markdown("**Summary Totals:**")
        totals_data = [
            ["Subtotal (Products + Customization)", f"${products_subtotal:.2f}"]
        ]

        # Add discount line if applicable (discount applies to products only, not customization)
        if discount_percent > 0:
            totals_data.append([f"Discount ({discount_description}) - on products only", f"-${discount_amount:.2f}"])

        totals_data.append(["Shipping", f"${shipping:.2f}"])

        # Tariff is already in line items, so we show it separately
        if tariff > 0:
            totals_data.append(["Tariff", f"${tariff:.2f}"])

        # Add credit card fee if applicable
        if st.session_state.apply_cc_fee and cc_fee_amount > 0:
            totals_data.append([f"Credit Card Fee ({st.session_state.cc_fee_percent}%)", f"${cc_fee_amount:.2f}"])

        totals_data.append(["**TOTAL**", f"**${total_quote:.2f}**"])

        totals_df = pd.DataFrame(totals_data, columns=["Item", "Amount"])
        st.table(totals_df)

        st.divider()

        # === NOTES SECTION ===
        st.markdown("#### NOTES")
        st.caption("""
        Enter any specific details, kitting specs, client requests, add-on samples for
        Peace by Piece to be added to purchase orders. Remember to attach titled artwork
        that matches your purchase order request and any additional spec sheets for our partners.
        """)

        order_notes = st.session_state.order_notes

        # Display all 4 categories with clear section headers
        notes_content = []
        if order_notes.get('internal_pbp_team'):
            notes_content.append(f"**Internal Notes (For PBP Team):**\n{order_notes['internal_pbp_team']}")
        if order_notes.get('internal_bookkeeping'):
            notes_content.append(f"**Internal Notes (For Bookkeeping):**\n{order_notes['internal_bookkeeping']}")
        if order_notes.get('external_partners'):
            notes_content.append(f"**External Notes (For Partners/POs):**\n{order_notes['external_partners']}")
        if order_notes.get('external_clients'):
            notes_content.append(f"**External Notes (For Clients/Invoices):**\n{order_notes['external_clients']}")

        if notes_content:
            for note in notes_content:
                st.markdown(note)
        else:
            st.caption("No notes added")

        st.divider()

        # === DOWNLOAD SECTION ===
        st.markdown("#### Download Invoice & Purchase Order")

        # Add download button for complete invoice/PO
        # Combine line items and totals into one downloadable file
        invoice_complete = invoice_df.copy()

        # Add blank row
        blank_row = pd.DataFrame([{col: "" for col in invoice_df.columns}])
        invoice_complete = pd.concat([invoice_complete, blank_row], ignore_index=True)

        # Add totals section (map to new column names)
        for total_item in totals_data:
            total_row = pd.DataFrame([{
                'PARTNER': '',
                'DESCRIPTION (Invoice)': total_item[0],
                'DESCRIPTION (PO)': '',
                'QTY': '',
                'IN-HANDS from Partner': '',
                'COST/UNIT': '',
                'TOTAL COST': '',
                'COST VERIFIED?': '',
                'SELL PRICE/UNIT': '',
                'TOTAL SELL PRICE': total_item[1]
            }])
            invoice_complete = pd.concat([invoice_complete, total_row], ignore_index=True)

        # Add notes section
        notes_row = pd.DataFrame([{
            'PARTNER': '',
            'DESCRIPTION (Invoice)': '',
            'DESCRIPTION (PO)': '',
            'QTY': '',
            'IN-HANDS from Partner': '',
            'COST/UNIT': '',
            'TOTAL COST': '',
            'COST VERIFIED?': '',
            'SELL PRICE/UNIT': '',
            'TOTAL SELL PRICE': ''
        }])
        invoice_complete = pd.concat([invoice_complete, notes_row], ignore_index=True)

        # Add notes content
        if notes_content:
            for note in notes_content:
                notes_row = pd.DataFrame([{
                    'PARTNER': '',
                    'DESCRIPTION (Invoice)': note.replace('**', '').replace('\n', ' '),
                    'DESCRIPTION (PO)': '',
                    'QTY': '',
                    'IN-HANDS from Partner': '',
                    'COST/UNIT': '',
                    'TOTAL COST': '',
                    'COST VERIFIED?': '',
                    'SELL PRICE/UNIT': '',
                    'TOTAL SELL PRICE': ''
                }])
                invoice_complete = pd.concat([invoice_complete, notes_row], ignore_index=True)

        # Generate HTML Invoice/PO Form
        html_invoice = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>
        body {{ font-family: Arial, sans-serif; max-width: 1200px; margin: 20px auto; padding: 20px; background-color: #ffffff; }}
        h2 {{ color: #2c3e50; background-color: #ffffff; border-bottom: 3px solid #3498db; padding-bottom: 10px; }}
        h3 {{ color: #34495e; margin-top: 25px; margin-bottom: 10px; }}
        table {{ border-collapse: collapse; width: 100%; margin: 15px 0; background-color: #ffffff; }}
        th {{ background-color: #3498db !important; color: #ffffff !important; padding: 12px; text-align: left; font-weight: bold; border: 1px solid #2980b9; }}
        td {{ border: 1px solid #ddd; padding: 10px; background-color: #ffffff; color: #000000; }}
        td:first-child {{ background-color: #f8f9fa !important; color: #000000 !important; font-weight: 500; width: 30%; }}
        .section-header {{ background-color: #2c3e50 !important; color: #ffffff !important; font-weight: bold; padding: 12px; text-align: center; }}
        .notes-section {{ background-color: #e8f4f8; border: 1px solid #3498db; padding: 15px; margin: 15px 0; }}
        .notes-header {{ font-weight: bold; color: #2c3e50; margin-bottom: 10px; }}
        .company-header {{ background-color: #34495e; color: #ffffff; padding: 15px; text-align: center; margin-bottom: 20px; }}
    </style>
</head>
<body>
    <h2>INVOICE AND PURCHASE ORDER REQUEST FORM</h2>

    <h3>1. Client/Company Information</h3>
    <table>
        <tr>
            <td>Company/Client Name</td>
            <td>{client_info.get('company_name', 'Not specified')}</td>
        </tr>
        <tr>
            <td>New Client?</td>
            <td>{'[X] Yes  [ ] No' if client_info.get('is_new_client', False) else '[ ] Yes  [X] No'}</td>
        </tr>
        <tr>
            <td>Contact + Email</td>
            <td>{contact_display.replace('<', '&lt;').replace('>', '&gt;')}</td>
        </tr>
        <tr>
            <td>Company Billing Address + Email</td>
            <td>{client_info.get('billing_address', 'Not specified').replace(chr(10), ', ')} | {billing_email}</td>
        </tr>
        <tr>
            <td>Company Shipping Address</td>
            <td>{client_info.get('shipping_address', 'Not specified').replace(chr(10), ', ')} ({client_info.get('shipping_type', 'Not specified')})</td>
        </tr>
        <tr>
            <td>Client PO #</td>
            <td>{client_info.get('client_po', 'N/A')}</td>
        </tr>
    </table>

    <h3>2. Partners + Point of Contacts</h3>
    <table>
        <tr>
            <th>Partner</th>
            <th>Point of Contact (POC)</th>
        </tr>"""

        # Add partner rows
        if partners_in_order and hasattr(st.session_state, 'partner_contacts'):
            for partner_name in partners_in_order:
                partner_contact = st.session_state.partner_contacts.get(partner_name, {})
                poc_name = partner_contact.get('poc_name', 'Not specified')
                poc_email = partner_contact.get('poc_email', 'Not specified')
                html_invoice += f"""
        <tr>
            <td style="width: 40%;">{partner_name}</td>
            <td>{poc_name} &lt;{poc_email}&gt;</td>
        </tr>"""
        else:
            html_invoice += """
        <tr>
            <td colspan="2" style="text-align: center; color: #7f8c8d;">No partners in order</td>
        </tr>"""

        html_invoice += f"""
    </table>

    <h3>3. Order Details</h3>
    <table>
        <tr>
            <td>Client In-Hands Date</td>
            <td>{format_date_display(client_in_hands)}</td>
        </tr>
        <tr>
            <td>Ship Method</td>
            <td>{ship_method}</td>
        </tr>
        <tr>
            <td>Payment Terms</td>
            <td>{payment_terms}</td>
        </tr>
        <tr>
            <td>Payment Method</td>
            <td>{payment_method}</td>
        </tr>
        <tr>
            <td>Order Submitted By</td>
            <td>{order_submitted_by} (Date: {format_date_display(order_submitted_date)})</td>
        </tr>
        <tr>
            <td>Cost Submitted By</td>
            <td>{cost_submitted_by} (Date: {format_date_display(cost_submitted_date) if cost_submitted_date else 'Not specified'})</td>
        </tr>
    </table>

    <h3>4. Invoice and Purchase Order Item Details</h3>
    <p style="color: #7f8c8d; font-size: 0.9em; margin: 10px 0;">
        This cost-to-sell segment outlines our partners' cost, our sell price to client,
        and our partners' requested in-hands date.
    </p>
    <table>
        <tr>
            <th>Partner</th>
            <th>Description (Invoice)</th>
            <th>Description (PO)</th>
            <th>Qty</th>
            <th>In-Hands from Partner</th>
            <th>Cost/Unit</th>
            <th>Total Cost</th>
            <th>Cost Verified?</th>
            <th>Sell Price/Unit</th>
            <th>Total Sell Price</th>
        </tr>"""

        # Add all line items
        for line_item in invoice_line_items:
            html_invoice += f"""
        <tr>
            <td>{line_item['PARTNER']}</td>
            <td>{line_item['DESCRIPTION (Invoice)'].replace(chr(10), '<br>')}</td>
            <td>{line_item['DESCRIPTION (PO)'].replace(chr(10), '<br>')}</td>
            <td style="text-align: center;">{line_item['QTY']}</td>
            <td>{line_item['IN-HANDS from Partner']}</td>
            <td style="text-align: right;">{line_item['COST/UNIT']}</td>
            <td style="text-align: right;">{line_item['TOTAL COST']}</td>
            <td style="text-align: center;">{line_item['COST VERIFIED?']}</td>
            <td style="text-align: right;">{line_item['SELL PRICE/UNIT']}</td>
            <td style="text-align: right;">{line_item['TOTAL SELL PRICE']}</td>
        </tr>"""

        html_invoice += """
    </table>

    <h3>Summary Totals</h3>
    <table>"""

        # Add totals
        for total_item in totals_data:
            is_total_row = total_item[0].startswith('**')
            style = ' style="background-color: #f8f9fa; font-weight: bold;"' if is_total_row else ''
            html_invoice += f"""
        <tr{style}>
            <td>{total_item[0].replace('**', '')}</td>
            <td style="text-align: right;">{total_item[1].replace('**', '')}</td>
        </tr>"""

        html_invoice += """
    </table>"""

        # Add dropshipping notes section (appears in Invoice only, not PO)
        if st.session_state.dropshipping_notes:
            html_invoice += f"""
    <div class="notes-section">
        <div class="notes-header">DROPSHIPPING INSTRUCTIONS</div>
        <p>{st.session_state.dropshipping_notes.replace(chr(10), '<br/>')}</p>
    </div>"""

        # Add notes if any
        if notes_content:
            html_invoice += """
    <div class="notes-section">
        <div class="notes-header">NOTES</div>"""
            for note in notes_content:
                html_invoice += f"""
        <p><strong>{note.split(':')[0] if ':' in note else 'Note'}:</strong> {note.split(':', 1)[1] if ':' in note else note}</p>"""
            html_invoice += """
    </div>"""

        html_invoice += """
</body>
</html>"""

        # Download buttons side by side
        col1, col2 = st.columns(2)

        with col1:
            invoice_csv = invoice_complete.to_csv(index=False)
            st.download_button(
                label="Download as CSV",
                data=invoice_csv,
                file_name=f"invoice_po_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                mime="text/csv",
                key="download_invoice_po_csv",
                use_container_width=True
            )

        with col2:
            st.download_button(
                label="Download as HTML",
                data=html_invoice,
                file_name=f"invoice_po_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html",
                mime="text/html",
                key="download_invoice_po_html",
                use_container_width=True
            )

        st.caption("CSV for spreadsheet import | HTML for email-ready professional format")

        st.divider()

        # ============================================================
        # SECTION 4: ACCOUNTING EXPORT (FUTURE)
        # ============================================================
        st.subheader("3. Export for Accounting")
        st.caption("Future: QuickBooks export, accounting reports, etc.")
        st.info("Accounting export features will be added in Phase 4")

# ============================================================
# TAB 5: EXECUTIVE PRICING TOOL
# ============================================================

# Helper function to preserve tab state
def preserve_tab_5():
    """Set query params to keep Tab 5 active on rerun"""
    try:
        st.query_params.tab = "4"  # Tab 5 is index 4
    except AttributeError:
        try:
            st.experimental_set_query_params(tab="4")
        except:
            pass  # Fail silently if query params not supported

with tab5:
    st.header("Executive Pricing Tool")
    st.caption("Build and analyze pricing scenarios with detailed cost breakdowns")
    st.divider()

    # Check data is loaded
    if 'df_template' not in st.session_state or st.session_state.df_template is None:
        st.error("Please load pricing data first. Use the sidebar to select a dataset.")
        st.stop()

    df_template = st.session_state.df_template

    # Initialize session state for executive tool
    if 'exec_products' not in st.session_state:
        st.session_state.exec_products = []
    if 'exec_shipping' not in st.session_state:
        st.session_state.exec_shipping = 70.0
    if 'exec_cc_fee' not in st.session_state:
        st.session_state.exec_cc_fee = False
    if 'exec_cc_fee_percent' not in st.session_state:
        st.session_state.exec_cc_fee_percent = 3.0

    # ============================================================
    # SECTION 1: PARTNER & PRODUCT SELECTION
    # ============================================================
    st.subheader("1. Add Products")

    # Auto-collapse after first product is added
    expand_selector = len(st.session_state.exec_products) == 0

    with st.expander("**Product Selector**", expanded=expand_selector):
        col1, col2, col3 = st.columns([2, 3, 1])

        with col1:
            # Partner selector
            all_partners = sorted(df_template['Partner'].unique().tolist())
            selected_partner = st.selectbox(
                "Partner",
                options=all_partners,
                key="exec_partner_select",
                help="Select a partner to view their products"
            )

        # Filter products by selected partner
        partner_products = df_template[df_template['Partner'] == selected_partner].copy()

        with col2:
            # Product selector
            product_options = partner_products['Product/Service'].tolist()
            if product_options:
                selected_product = st.selectbox(
                    "Product",
                    options=product_options,
                    key="exec_product_select",
                    help="Select a product to add"
                )
            else:
                selected_product = None
                st.warning("No products available for this partner")

        with col3:
            # Add button
            if st.button("Add Product", type="primary", use_container_width=True, disabled=not selected_product):
                if selected_product:
                    # Get product data
                    product_row = partner_products[partner_products['Product/Service'] == selected_product].iloc[0]

                    # Check if already added
                    existing = any(p['product_name'] == selected_product for p in st.session_state.exec_products)
                    if not existing:
                        # Get country of origin to determine if tariffs should be included
                        country_of_origin = product_row.get('Country of Origin (Ships From)', 'Unknown')
                        # Auto-check tariffs for non-USA products
                        include_tariffs_default = country_of_origin.upper() not in ['USA', 'UNITED STATES', 'US', 'U.S.', 'AMERICA', 'UNITED STATES OF AMERICA']

                        # Create product entry
                        exec_product = {
                            'product_name': selected_product,
                            'partner': selected_partner,
                            'product_data': product_row.to_dict(),
                            'quantity': 100,
                            'markup_percent': 100.0,
                            'include_customization': False,
                            'custom_setup_fee': 0.0,
                            'custom_per_unit': 0.0,
                            'pbp_setup_fee': 0.0,  # Initialize PBP cost fields
                            'pbp_per_unit_cost': 0.0,  # Initialize PBP cost fields
                            'include_tariffs': include_tariffs_default  # Auto-set based on country
                        }
                        st.session_state.exec_products.append(exec_product)
                        st.toast(f"Added {selected_product}")
                        # Preserve Tab 5 as active tab
                        preserve_tab_5()
                        st.rerun()
                    else:
                        st.warning("Product already added")

    st.divider()

    # ============================================================
    # SECTION 2: PRODUCT DETAILS
    # ============================================================
    if len(st.session_state.exec_products) > 0:
        st.subheader("2. Product Configuration")
        st.caption("Configure quantities, markups, and customization for each product")

        # Display each product as an expandable card
        for idx, product in enumerate(st.session_state.exec_products):
            with st.expander(f"**{product['product_name']}** - {product['partner']}", expanded=True):
                # Get pricing info first
                row = product['product_data']
                base_cost, tier_range, tier_col = get_unit_price_new_system(row, product['quantity'])

                # Top row: Quantity, Markup, Price, Remove
                col1, col2, col3, col4 = st.columns([2, 1.5, 1.5, 1])

                with col1:
                    # Quantity input
                    new_qty = st.number_input(
                        "Quantity",
                        min_value=1,
                        max_value=10000,
                        value=product['quantity'],
                        step=10,
                        key=f"exec_qty_{idx}"
                    )
                    product['quantity'] = new_qty

                    # Show tier info and base cost
                    if base_cost and base_cost > 0:
                        st.caption(f"Tier: {tier_range} | PBP Cost: ${base_cost:.2f}/unit")
                    else:
                        st.caption(f"Tier: {tier_range}")

                with col2:
                    # Calculate client price
                    if base_cost and base_cost > 0:
                        client_price = base_cost * (1 + product['markup_percent'] / 100)

                        # Pricing edit mode toggle
                        if 'edit_mode' not in product:
                            product['edit_mode'] = 'markup'  # Default to markup editing

                        edit_mode = st.radio(
                            "Edit Mode",
                            options=['markup', 'price'],
                            format_func=lambda x: 'Markup %' if x == 'markup' else 'Price Direct',
                            index=0 if product.get('edit_mode') == 'markup' else 1,
                            key=f"edit_mode_{idx}",
                            horizontal=True,
                            help="Choose how to set pricing"
                        )
                        product['edit_mode'] = edit_mode

                        if edit_mode == 'markup':
                            # Edit markup
                            new_markup = st.number_input(
                                "Markup %",
                                min_value=-50,
                                max_value=500,
                                value=int(product['markup_percent']),
                                step=5,
                                key=f"exec_markup_{idx}",
                                help="Profit margin. 100% = 2x the cost"
                            )
                            # Update if changed
                            if new_markup != int(product['markup_percent']):
                                product['markup_percent'] = float(new_markup)
                                preserve_tab_5()
                                st.rerun()
                        else:
                            # Display markup
                            st.markdown("**Markup %**")
                            st.success(f"{product['markup_percent']:.1f}%")
                    else:
                        st.markdown("**Pricing**")
                        st.markdown("—")

                with col3:
                    if base_cost and base_cost > 0:
                        client_price = base_cost * (1 + product['markup_percent'] / 100)

                        if product.get('edit_mode', 'markup') == 'price':
                            # Edit price directly
                            new_client_price = st.number_input(
                                "Client Price/Unit",
                                min_value=0.01,
                                value=client_price,
                                step=1.0,
                                format="%.2f",
                                key=f"exec_price_{idx}",
                                help="Price per unit charged to client"
                            )

                            # Update markup if price changed
                            if abs(new_client_price - client_price) > 0.01:
                                from src.helpers import calculate_markup_from_price
                                new_markup_calc = calculate_markup_from_price(base_cost, new_client_price)
                                product['markup_percent'] = new_markup_calc
                                preserve_tab_5()
                                st.rerun()
                        else:
                            # Display price
                            st.markdown("**Client Price/Unit**")
                            st.success(f"${client_price:.2f}")
                            st.caption(f"Profit: ${client_price - base_cost:.2f}/unit")
                    else:
                        st.markdown("**Client Price/Unit**")
                        st.markdown("—")

                with col4:
                    # Remove button
                    st.write("")  # Add spacing
                    if st.button("Remove", key=f"exec_remove_{idx}", type="secondary"):
                        st.session_state.exec_products.pop(idx)
                        preserve_tab_5()
                        st.rerun()

                # Customization section
                st.markdown("##### Customization Options")
                product['include_customization'] = st.checkbox(
                    "Include Customization",
                    value=product.get('include_customization', False),
                    key=f"exec_custom_{idx}"
                )

                if product['include_customization']:
                    col1, col2 = st.columns(2)
                    with col1:
                        st.markdown("**Setup Fee**")
                        # Get default values from spreadsheet
                        default_setup = clean_price(get_column_value(
                            row, 'Client Price: Customization Setup Fee', 'Customization Setup Fee', 0
                        ))
                        product['custom_setup_fee'] = st.number_input(
                            "Client Price",
                            min_value=0.0,
                            value=float(default_setup) if default_setup else float(product.get('custom_setup_fee', 0)),
                            step=10.0,
                            key=f"exec_setup_{idx}"
                        )

                        # Get PBP cost for setup fee
                        pbp_setup_cost = clean_price(get_column_value(
                            row, 'PBP Cost: Customization Setup Fee', 'Customization Setup Fee', 0
                        ))

                        if pbp_setup_cost and pbp_setup_cost > 0:
                            st.caption(f"PBP Cost: ${pbp_setup_cost:.2f} (from spreadsheet)")
                            # Store in product for calculations
                            product['pbp_setup_fee'] = pbp_setup_cost
                        else:
                            # Allow user to input PBP cost
                            if 'pbp_setup_fee' not in product:
                                product['pbp_setup_fee'] = 0.0

                            product['pbp_setup_fee'] = st.number_input(
                                "PBP Cost (not in spreadsheet - enter manually)",
                                min_value=0.0,
                                value=float(product.get('pbp_setup_fee', 0)),
                                step=10.0,
                                key=f"exec_pbp_setup_{idx}",
                                help="What PBP pays the partner for setup"
                            )

                    with col2:
                        st.markdown("**Per Unit Cost**")
                        default_per_unit = clean_price(get_column_value(
                            row, 'Client Price: Customization Cost per Unit', 'Customization Cost per Unit', 0
                        ))
                        product['custom_per_unit'] = st.number_input(
                            "Client Price",
                            min_value=0.0,
                            value=float(default_per_unit) if default_per_unit else float(product.get('custom_per_unit', 0)),
                            step=0.50,
                            key=f"exec_per_unit_{idx}"
                        )

                        # Get PBP cost per unit
                        pbp_per_unit_cost = clean_price(get_column_value(
                            row, 'PBP Cost: Customization Cost per Unit', 'Customization Cost per Unit', 0
                        ))

                        if pbp_per_unit_cost and pbp_per_unit_cost > 0:
                            st.caption(f"PBP Cost: ${pbp_per_unit_cost:.2f} (from spreadsheet)")
                            # Store in product for calculations
                            product['pbp_per_unit_cost'] = pbp_per_unit_cost
                        else:
                            # Allow user to input PBP cost
                            if 'pbp_per_unit_cost' not in product:
                                product['pbp_per_unit_cost'] = 0.0

                            product['pbp_per_unit_cost'] = st.number_input(
                                "PBP Cost (not in spreadsheet - enter manually)",
                                min_value=0.0,
                                value=float(product.get('pbp_per_unit_cost', 0)),
                                step=0.50,
                                key=f"exec_pbp_per_unit_{idx}",
                                help="What PBP pays the partner per unit"
                            )

                # Tariff section
                st.markdown("##### Tariff Options")

                # Get country of origin first
                product_country = row.get('Country of Origin (Ships From)', 'Unknown')

                # Show country and tariff checkbox
                product['include_tariffs'] = st.checkbox(
                    f"Include Tariffs (Ships From: {product_country})",
                    value=product.get('include_tariffs', False),
                    key=f"exec_tariff_{idx}",
                    help="Apply tariff costs for this product (auto-checked for non-USA products)"
                )

                if product['include_tariffs']:
                    # Calculate client price for tariff calculation (commercial/declared value)
                    client_price_for_tariff = base_cost * (1 + product['markup_percent'] / 100)

                    # Get tariff data from spreadsheet
                    tariff_dollar = clean_price(get_column_value(row, 'Tariff Estimate ($)', 'Tariff Estimate ($)', 0))
                    tariff_percent = clean_price(get_column_value(row, 'Tariff Estimate (%)', 'Tariff Estimate (%)', 0))

                    st.caption("Tariffs are calculated on commercial value (client price) and passed through without markup")

                    # Input method selector
                    if 'tariff_input_method' not in product:
                        # Determine default based on spreadsheet data
                        if tariff_dollar and tariff_dollar > 0:
                            product['tariff_input_method'] = 'dollar'
                        else:
                            product['tariff_input_method'] = 'percentage'

                    tariff_method = st.radio(
                        "Tariff Input Method",
                        options=['percentage', 'dollar'],
                        format_func=lambda x: 'Percentage (% of commercial value)' if x == 'percentage' else 'Dollar Amount ($ per unit)',
                        index=0 if product.get('tariff_input_method') == 'percentage' else 1,
                        key=f"exec_tariff_method_{idx}",
                        horizontal=True
                    )
                    product['tariff_input_method'] = tariff_method

                    col1, col2 = st.columns(2)
                    with col1:
                        if tariff_method == 'percentage':
                            # Calculate default rate
                            if tariff_dollar and tariff_dollar > 0:
                                # Convert dollar to percentage of client price (commercial value)
                                default_rate = (tariff_dollar / (client_price_for_tariff * 100)) * 100 if client_price_for_tariff > 0 else 0
                                st.caption(f"Default from spreadsheet: ${tariff_dollar:.2f} for 100 units ({default_rate:.1f}% of commercial value)")
                            elif tariff_percent and tariff_percent > 0:
                                default_rate = tariff_percent
                                st.caption(f"Default from spreadsheet: {tariff_percent:.1f}%")
                            else:
                                default_rate = 0
                                st.caption("No default tariff specified")

                            product['tariff_rate'] = st.number_input(
                                "Tariff Rate (% of commercial value)",
                                min_value=0.0,
                                max_value=100.0,
                                value=float(product.get('tariff_rate', default_rate)),
                                step=0.5,
                                format="%.1f",
                                key=f"exec_tariff_rate_{idx}",
                                help="Percentage of commercial value (client price)"
                            )
                            # Calculate dollar amount from percentage OF CLIENT PRICE
                            tariff_per_unit = client_price_for_tariff * (product['tariff_rate'] / 100)
                            product['tariff_dollar'] = tariff_per_unit

                        else:  # dollar input
                            # Default dollar amount
                            if tariff_dollar and tariff_dollar > 0:
                                default_dollar = tariff_dollar / 100  # Convert from 100 units to per unit
                                st.caption(f"Default from spreadsheet: ${tariff_dollar:.2f} for 100 units (${default_dollar:.2f} per unit)")
                            elif tariff_percent and tariff_percent > 0:
                                default_dollar = client_price_for_tariff * (tariff_percent / 100)
                                st.caption(f"Default from spreadsheet: {tariff_percent:.1f}% = ${default_dollar:.2f} per unit")
                            else:
                                default_dollar = 0
                                st.caption("No default tariff specified")

                            product['tariff_dollar'] = st.number_input(
                                "Tariff Amount ($ per unit)",
                                min_value=0.0,
                                value=float(product.get('tariff_dollar', default_dollar)),
                                step=1.0,
                                format="%.2f",
                                key=f"exec_tariff_dollar_{idx}",
                                help="Fixed dollar amount per unit"
                            )
                            # Calculate percentage from dollar amount (as % of client price)
                            product['tariff_rate'] = (product['tariff_dollar'] / client_price_for_tariff * 100) if client_price_for_tariff > 0 else 0
                            tariff_per_unit = product['tariff_dollar']

                    with col2:
                        # Show tariff amounts (PBP pays tariff, passes cost to client)
                        tariff_total = tariff_per_unit * product['quantity']
                        if tariff_per_unit > 0:
                            st.markdown("**Tariff Costs (Pass-through):**")
                            st.caption(f"PBP pays: ${tariff_per_unit:.2f} per unit")
                            st.caption(f"Client pays: ${tariff_per_unit:.2f} per unit (no markup)")
                            if tariff_method == 'percentage':
                                st.caption(f"({product['tariff_rate']:.1f}% of commercial value ${client_price_for_tariff:.2f})")
                            st.caption(f"Total ({product['quantity']} units): ${tariff_total:.2f}")
                        else:
                            st.markdown("**Tariff Costs:**")
                            st.caption("No tariffs applied")

                # Calculate and display pricing breakdown
                st.markdown("##### Pricing Breakdown")

                # Calculate all costs and prices
                if base_cost and base_cost > 0:
                    client_price = base_cost * (1 + product['markup_percent'] / 100)

                    # Show margin info prominently
                    margin_per_unit = client_price - base_cost
                    margin_total = margin_per_unit * product['quantity']
                    margin_percent = (margin_per_unit / client_price * 100) if client_price > 0 else 0

                    col_m1, col_m2, col_m3 = st.columns(3)
                    with col_m1:
                        st.metric("Markup %", f"{product['markup_percent']:.0f}%")
                    with col_m2:
                        st.metric("Margin/Unit", f"${margin_per_unit:.2f}")
                    with col_m3:
                        st.metric("Total Margin", f"${margin_total:.2f}")

                    st.caption(f"Margin represents {margin_percent:.1f}% of client price")

                    # PBP costs for customization (what PBP pays)
                    # Use stored values (either from spreadsheet or user input)
                    if product['include_customization']:
                        pbp_setup = product.get('pbp_setup_fee', 0)
                        pbp_per_unit = product.get('pbp_per_unit_cost', 0)
                    else:
                        pbp_setup = 0
                        pbp_per_unit = 0

                    # Build the breakdown table data
                    breakdown_data = []

                    # Base product
                    breakdown_data.append({
                        'Item': f"Base Product: {product['product_name']}",
                        'Qty': product['quantity'],
                        'PBP Cost (Per Unit)': f"${base_cost:.2f}",
                        'PBP Cost': f"${base_cost * product['quantity']:.2f}",
                        'Client Price (Per Unit)': f"${client_price:.2f}",
                        'Client Price': f"${client_price * product['quantity']:.2f}"
                    })

                    # Customization if included
                    if product['include_customization']:
                        if product['custom_setup_fee'] > 0:
                            breakdown_data.append({
                                'Item': f"{product['product_name']} - Setup",
                                'Qty': "one-time",
                                'PBP Cost (Per Unit)': f"${pbp_setup:.2f}",
                                'PBP Cost': f"${pbp_setup:.2f}",
                                'Client Price (Per Unit)': f"${product['custom_setup_fee']:.2f}",
                                'Client Price': f"${product['custom_setup_fee']:.2f}"
                            })

                        if product['custom_per_unit'] > 0:
                            pbp_total = pbp_per_unit * product['quantity'] if pbp_per_unit else 0
                            client_total = product['custom_per_unit'] * product['quantity']
                            breakdown_data.append({
                                'Item': f"{product['product_name']} - Per Unit",
                                'Qty': product['quantity'],
                                'PBP Cost (Per Unit)': f"${pbp_per_unit:.2f}" if pbp_per_unit else "$0.00",
                                'PBP Cost': f"${pbp_total:.2f}",
                                'Client Price (Per Unit)': f"${product['custom_per_unit']:.2f}",
                                'Client Price': f"${client_total:.2f}"
                            })

                    # Tariffs if included (pass-through cost)
                    if product.get('include_tariffs', False):
                        # Use the stored tariff_dollar which is calculated based on PBP cost
                        tariff_per_unit = product.get('tariff_dollar', 0)
                        if tariff_per_unit > 0:
                            tariff_total = tariff_per_unit * product['quantity']
                            # Get country for display
                            product_country = row.get('Country of Origin (Ships From)', 'Unknown')
                            breakdown_data.append({
                                'Item': f"{product['product_name']} - Tariff ({product.get('tariff_rate', 0):.1f}% - {product_country})",
                                'Qty': product['quantity'],
                                'PBP Cost (Per Unit)': f"${tariff_per_unit:.2f}",
                                'PBP Cost': f"${tariff_total:.2f}",
                                'Client Price (Per Unit)': f"${tariff_per_unit:.2f}",
                                'Client Price': f"${tariff_total:.2f}"
                            })

                    # Display the breakdown table
                    df_breakdown = pd.DataFrame(breakdown_data)
                    st.dataframe(df_breakdown, use_container_width=True, hide_index=True)

                    # Calculate subtotals
                    pbp_subtotal = base_cost * product['quantity']
                    client_subtotal = client_price * product['quantity']

                    if product['include_customization']:
                        pbp_subtotal += pbp_setup + (pbp_per_unit * product['quantity'] if pbp_per_unit else 0)
                        client_subtotal += product['custom_setup_fee'] + (product['custom_per_unit'] * product['quantity'])

                    # Add tariffs to BOTH PBP and client subtotals (pass-through cost)
                    if product.get('include_tariffs', False):
                        tariff_per_unit = product.get('tariff_dollar', 0)
                        if tariff_per_unit > 0:
                            tariff_total = tariff_per_unit * product['quantity']
                            pbp_subtotal += tariff_total  # PBP pays the tariff
                            client_subtotal += tariff_total  # Client reimburses PBP

                    # Calculate margin (excludes tariffs since they're pass-through)
                    # Margin is only on the base product and customization markup
                    product_margin = (client_price * product['quantity']) - (base_cost * product['quantity'])
                    if product['include_customization']:
                        # Add customization margin (client price - PBP cost)
                        product_margin += (product['custom_setup_fee'] - pbp_setup)
                        product_margin += ((product['custom_per_unit'] - pbp_per_unit) * product['quantity'] if pbp_per_unit else product['custom_per_unit'] * product['quantity'])

                    # Calculate margin percentage on revenue-generating items only (excluding tariffs)
                    revenue_base = (client_price * product['quantity']) + (product.get('custom_setup_fee', 0) if product.get('include_customization') else 0) + (product.get('custom_per_unit', 0) * product['quantity'] if product.get('include_customization') else 0)
                    product_margin_pct = (product_margin / revenue_base * 100) if revenue_base > 0 else 0

                    st.info(f"**Product Subtotal:** PBP Cost: ${pbp_subtotal:.2f} | Client Price: ${client_subtotal:.2f} | **Margin: ${product_margin:.2f} ({product_margin_pct:.1f}%)**")

        st.divider()

        # Quick Add Bar - convenient way to add more products
        with st.container():
            st.markdown("#### Quick Add Products")
            st.caption("Add more products without scrolling back up")

            col1, col2, col3 = st.columns([2, 3, 1])

            with col1:
                # Partner selector for quick add
                quick_partner = st.selectbox(
                    "Partner",
                    options=sorted(df_template['Partner'].unique().tolist()),
                    key="exec_quick_partner",
                    label_visibility="collapsed"
                )

            with col2:
                # Product selector for quick add
                quick_partner_products = df_template[df_template['Partner'] == quick_partner]
                quick_product_options = quick_partner_products['Product/Service'].tolist()

                if quick_product_options:
                    quick_product = st.selectbox(
                        "Product",
                        options=quick_product_options,
                        key="exec_quick_product",
                        label_visibility="collapsed"
                    )
                else:
                    quick_product = None
                    st.warning("No products for this partner")

            with col3:
                # Quick add button
                if st.button("➕ Add", type="primary", use_container_width=True,
                           disabled=not quick_product, key="exec_quick_add"):
                    if quick_product:
                        # Check if already added
                        existing = any(p['product_name'] == quick_product for p in st.session_state.exec_products)
                        if not existing:
                            # Get product data
                            product_row = quick_partner_products[quick_partner_products['Product/Service'] == quick_product].iloc[0]

                            # Get country of origin for tariff auto-check
                            country_of_origin = product_row.get('Country of Origin (Ships From)', 'Unknown')
                            include_tariffs_default = country_of_origin.upper() not in ['USA', 'UNITED STATES', 'US', 'U.S.', 'AMERICA', 'UNITED STATES OF AMERICA']

                            # Create product entry
                            exec_product = {
                                'product_name': quick_product,
                                'partner': quick_partner,
                                'product_data': product_row.to_dict(),
                                'quantity': 100,
                                'markup_percent': 100.0,
                                'include_customization': False,
                                'custom_setup_fee': 0.0,
                                'custom_per_unit': 0.0,
                                'pbp_setup_fee': 0.0,
                                'pbp_per_unit_cost': 0.0,
                                'include_tariffs': include_tariffs_default
                            }
                            st.session_state.exec_products.append(exec_product)
                            st.toast(f"Added {quick_product}")
                            preserve_tab_5()
                            st.rerun()
                        else:
                            st.warning("Product already added")

        st.divider()
    else:
        st.info("No products added yet. Use the selector above to add products.")
        st.divider()

    # ============================================================
    # SECTION 3: ORDER-LEVEL SETTINGS
    # ============================================================
    st.subheader("3. Order Settings")

    col1, col2 = st.columns(2)

    with col1:
        st.session_state.exec_shipping = st.number_input(
            "Shipping Cost ($)",
            min_value=0.0,
            value=st.session_state.exec_shipping,
            step=10.0,
            format="%.2f",
            help="Total shipping cost for the order"
        )

    with col2:
        st.session_state.exec_cc_fee = st.checkbox(
            "Credit Card Fee",
            value=st.session_state.exec_cc_fee,
            help="Add credit card processing fee"
        )

        if st.session_state.exec_cc_fee:
            st.session_state.exec_cc_fee_percent = st.number_input(
                "CC Fee %",
                min_value=0.0,
                max_value=10.0,
                value=st.session_state.exec_cc_fee_percent,
                step=0.1,
                format="%.1f"
            )

    st.divider()

    # ============================================================
    # SECTION 4: ORDER SUMMARY
    # ============================================================
    st.subheader("4. Order Summary")
    st.caption("Complete pricing breakdown matching Tab 3 format")

    if len(st.session_state.exec_products) == 0:
        st.info("Add products above to see the order summary")
    else:
        # Build summary data
        summary_items = []
        total_units = 0
        total_pbp_cost = 0
        total_client_price = 0
        products_pbp = 0
        products_client = 0
        custom_pbp = 0
        custom_client = 0

        # SECTION 1: Products
        for product in st.session_state.exec_products:
            row = product['product_data']
            qty = product['quantity']

            # Get pricing for this quantity
            base_cost, tier_range, _ = get_unit_price_new_system(row, qty)
            if not base_cost or base_cost <= 0:
                continue

            client_price = base_cost * (1 + product['markup_percent'] / 100)

            # Add to totals
            total_units += qty
            products_pbp += base_cost * qty
            products_client += client_price * qty

            # Add base product row
            summary_items.append([
                f"Base Product: {product['product_name']}",
                str(qty),  # Convert to string for consistent column type
                f"${base_cost:.2f}",
                f"${base_cost * qty:.2f}",
                f"${client_price:.2f}",
                f"${client_price * qty:.2f}"
            ])

            # Add customization if included
            if product.get('include_customization'):
                # Get PBP costs (use stored values - either from spreadsheet or user input)
                pbp_setup = product.get('pbp_setup_fee', 0)
                pbp_per_unit = product.get('pbp_per_unit_cost', 0)

                # Client prices from product settings
                client_setup = product.get('custom_setup_fee', 0)
                client_per_unit = product.get('custom_per_unit', 0)

                if client_setup > 0:
                    summary_items.append([
                        f"{product['product_name']} - Setup",
                        "one-time",
                        f"${pbp_setup:.2f}",
                        f"${pbp_setup:.2f}",
                        f"${client_setup:.2f}",
                        f"${client_setup:.2f}"
                    ])
                    custom_pbp += pbp_setup
                    custom_client += client_setup

                if client_per_unit > 0:
                    pbp_total = pbp_per_unit * qty
                    client_total = client_per_unit * qty
                    summary_items.append([
                        f"{product['product_name']} - Per Unit",
                        str(qty),  # Convert to string for consistent column type
                        f"${pbp_per_unit:.2f}",
                        f"${pbp_total:.2f}",
                        f"${client_per_unit:.2f}",
                        f"${client_total:.2f}"
                    ])
                    custom_pbp += pbp_total
                    custom_client += client_total

        # Products subtotal
        summary_items.append([
            "**Products Subtotal**",
            "",
            "",
            f"**${products_pbp:.2f}**",
            "",
            f"**${products_client:.2f}**"
        ])

        # Customization subtotal if any
        if custom_client > 0:
            summary_items.append(["", "", "", "", "", ""])  # Empty row
            summary_items.append([
                "**Customization Subtotal**",
                "",
                "",
                f"**${custom_pbp:.2f}**",
                "",
                f"**${custom_client:.2f}**"
            ])

        # Additional costs
        summary_items.append(["", "", "", "", "", ""])  # Empty row

        # Shipping
        if st.session_state.exec_shipping > 0:
            summary_items.append([
                "Shipping",
                "",
                "",
                f"${st.session_state.exec_shipping:.2f}",
                "",
                f"${st.session_state.exec_shipping:.2f}"
            ])

        # Tariffs (pass-through costs)
        total_tariff = 0
        for product in st.session_state.exec_products:
            # Check if this product has tariffs enabled
            if product.get('include_tariffs', False):
                tariff_per_unit = product.get('tariff_dollar', 0)
                if tariff_per_unit > 0:
                    row = product['product_data']
                    qty = product['quantity']
                    tariff_amount = tariff_per_unit * qty

                    country = row.get('Country of Origin (Ships From)', 'Unknown')
                    summary_items.append([
                        f"Tariff: {product['product_name']} ({product.get('tariff_rate', 0):.1f}% - {country})",
                        str(qty),  # Convert to string for consistent column type
                        f"${tariff_per_unit:.2f}",
                        f"${tariff_amount:.2f}",
                        f"${tariff_per_unit:.2f}",
                        f"${tariff_amount:.2f}"
                    ])
                    total_tariff += tariff_amount

        # Calculate totals
        total_pbp_cost = products_pbp + custom_pbp + st.session_state.exec_shipping + total_tariff
        total_before_cc = products_client + custom_client + st.session_state.exec_shipping + total_tariff

        # Credit card fee
        cc_fee_amount = 0
        if st.session_state.exec_cc_fee:
            cc_fee_amount = total_before_cc * (st.session_state.exec_cc_fee_percent / 100)
            summary_items.append([
                f"Credit Card Fee ({st.session_state.exec_cc_fee_percent}%)",
                "",
                "",
                "",
                "",
                f"${cc_fee_amount:.2f}"
            ])

        total_client_price = total_before_cc + cc_fee_amount

        # Total row
        summary_items.append([
            "**TOTAL**",
            f"**{total_units} units**",
            "",
            f"**${total_pbp_cost:.2f}**",
            "",
            f"**${total_client_price:.2f}**"
        ])

        # Create and display the summary table
        summary_df = pd.DataFrame(summary_items, columns=[
            "Item", "Qty", "PBP Cost (Per Unit)", "PBP Cost",
            "Client Price (Per Unit)", "Client Price"
        ])

        # Ensure Qty column is string type to avoid Arrow serialization issues
        summary_df['Qty'] = summary_df['Qty'].astype(str)

        st.table(summary_df)

        # Display total and margin
        if total_units > 0:
            avg_per_unit = total_client_price / total_units

            # Calculate true margin (excluding pass-through costs like tariffs and shipping)
            # Margin is only on products and customization
            true_margin = (products_client - products_pbp) + (custom_client - custom_pbp)
            # Calculate margin percentage based on revenue generating items only (excluding pass-throughs)
            revenue_base = products_client + custom_client
            margin_pct = (true_margin / revenue_base * 100) if revenue_base > 0 else 0

            col1, col2, col3 = st.columns(3)
            with col1:
                st.success(f"**Total Quote:** ${total_client_price:.2f}")
            with col2:
                st.info(f"**Avg/Unit:** ${avg_per_unit:.2f}")
            with col3:
                st.warning(f"**True Margin:** ${true_margin:.2f} ({margin_pct:.1f}%)")

            # Show breakdown of pass-through costs
            if total_tariff > 0 or st.session_state.exec_shipping > 0:
                st.caption(f"Pass-through costs: Shipping ${st.session_state.exec_shipping:.2f}, Tariffs ${total_tariff:.2f} (no markup)")

        st.divider()

        # ============================================================
        # SECTION 5: ACTIONS
        # ============================================================
        st.subheader("5. Import & Export")

        col1, col2, col3 = st.columns(3)

        with col1:
            if st.button("Import to Proposal", type="primary", use_container_width=True):
                if 'proposal_products' not in st.session_state:
                    st.session_state.proposal_products = []

                added = 0
                for product in st.session_state.exec_products:
                    # Check for duplicates
                    if not any(p['product_data']['Product/Service'] == product['product_name']
                              for p in st.session_state.proposal_products):
                        st.session_state.proposal_products.append({
                            'product_data': product['product_data'],
                            'markup_percent': product['markup_percent']
                        })
                        added += 1

                if added > 0:
                    st.toast(f"Added {added} products to proposal")
                    time.sleep(0.5)
                    preserve_tab_5()
                    st.rerun()
                else:
                    st.warning("All products already in proposal")

        with col2:
            if st.button("Import to Order", use_container_width=True):
                if 'order_items' not in st.session_state:
                    st.session_state.order_items = []

                added = 0
                for product in st.session_state.exec_products:
                    # Convert to order item format
                    if not any(item['product_name'] == product['product_name']
                              for item in st.session_state.order_items if not item.get('is_custom')):
                        order_item = {
                            'product_name': product['product_name'],
                            'partner': product['partner'],
                            'quantity': product['quantity'],
                            'markup_percent': product['markup_percent'],
                            'product_data': product['product_data'],
                            'customization_setup': product.get('custom_setup_fee', 0),
                            'customization_per_unit': product.get('custom_per_unit', 0),
                            'include_customization': product.get('include_customization', False),
                            'is_custom': False
                        }
                        st.session_state.order_items.append(order_item)
                        added += 1

                if added > 0:
                    st.toast(f"Added {added} products to order")
                    time.sleep(0.5)
                    preserve_tab_5()
                    st.rerun()
                else:
                    st.warning("All products already in order")

        with col3:
            # CSV export
            csv = summary_df.to_csv(index=False)
            st.download_button(
                label="Download CSV",
                data=csv,
                file_name=f"exec_pricing_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                mime="text/csv",
                use_container_width=True
            )
